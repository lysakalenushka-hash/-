#!/usr/bin/env python3
"""Build full Tema 6 welding presentation from 6.docx + source pptx style (6de2)."""

from copy import deepcopy
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________6de2.pptx"
OUT = "/workspace/Презентация_Сварочные_работы.pptx"
PHOTO_DIR = "/workspace/assets/welding_tema6"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

PHOTO_PATHS = []


def extract_photos_from_src():
    """Extract large photos from source deck for visual variety."""
    import os

    os.makedirs(PHOTO_DIR, exist_ok=True)
    paths = []
    src = Presentation(SRC)
    n = 0
    for si, slide in enumerate(src.slides):
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            blob = sh.image.blob
            if len(blob) < 50_000:
                continue
            n += 1
            path = f"{PHOTO_DIR}/src_s{si+1}_{n}.{sh.image.ext}"
            with open(path, "wb") as f:
                f.write(blob)
            paths.append(path)
    return paths


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def content_shapes(slide):
    shapes = [s for s in slide.shapes if s.has_text_frame]
    return [
        s
        for s in shapes
        if not (
            s.top is not None
            and s.left is not None
            and s.top > Emu(6200000)
            and s.left > Emu(8000000)
        )
    ]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 10, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 10, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10637215), Emu(6455664), Emu(914400), Emu(274320)
            )
            set_text(box, f"{i:02d} / {total:02d}", 10, True, RED)


def fill_hero(slide, section, title, body, note=None):
    shapes = content_shapes(slide)
    left = sorted(
        [s for s in shapes if (s.left or 0) < Emu(7000000)],
        key=lambda s: s.top or 0,
    )
    for sh in left:
        set_text(sh, "", 12, False, DARK)
    set_text(left[0], section, 13, True, RED)
    set_text(left[1], title, 28, True, DARK)
    set_text(left[2], body, 14, False, DARK)
    if len(left) > 3:
        set_text(left[3], note or "", 13, True, AMBER if note else GRAY)


def fill_bullet_body(slide, section, title, intro, items, note=None):
    """Section + title + body bullets (source slides 1–4 style)."""
    shapes = content_shapes(slide)
    left = sorted(
        [s for s in shapes if (s.left or 0) < Emu(7500000)],
        key=lambda s: (s.top or 0, s.left or 0),
    )
    for sh in left:
        set_text(sh, "", 12, False, DARK)

    set_text(left[0], section, 13, True, RED)
    set_text(left[1], title, 24, True, DARK)

    # Combine intro + bullets into the tall body box
    body = max(left[2:], key=lambda s: s.height or 0)
    for sh in left[2:]:
        if sh is not body:
            set_text(sh, "", 12, False, DARK)

    size = 14 if len(items) <= 4 else 13
    tf = body.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = intro
    set_run(r0, 14, False, GRAY)

    for item in items:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        set_run(r, size, False, DARK)

    if note:
        # Prefer unused lower-left text box; else append amber line
        note_shape = None
        for sh in left[2:]:
            if sh is body:
                continue
            if (sh.top or 0) > Emu(4500000):
                note_shape = sh
        if note_shape:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = note
            set_run(r, 12, True, AMBER)


def fill_two_col(slide, section, title, subtitle, left_h, left_items, right_h, right_items):
    shapes = content_shapes(slide)
    texts = sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))
    for sh in texts:
        set_text(sh, "", 12, False, DARK)

    # Expected layout from source slide 5: section, title, subtitle, Lh, Lb, Rh, Rb
    roles = []
    for sh in texts:
        if (sh.top or 0) > Emu(6200000):
            continue
        roles.append(sh)
    if len(roles) < 7:
        return
    set_text(roles[0], section, 13, True, RED)
    set_text(roles[1], title, 24, True, DARK)
    set_text(roles[2], subtitle, 14, False, GRAY)
    set_text(roles[3], left_h, 14, True, RED)
    tf = roles[4].text_frame
    tf.clear()
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        set_run(r, 13, False, DARK)
    set_text(roles[5], right_h, 14, True, RED)
    tf = roles[6].text_frame
    tf.clear()
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        set_run(r, 13, False, DARK)


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(
        prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    )
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def chunks(items, n=4):
    for i in range(0, len(items), n):
        yield items[i : i + n]


def replace_main_picture(slide, _unused_path, photo_idx):
    """Replace the large right-side picture to vary visuals."""
    if not PHOTO_PATHS or PHOTO_PATHS == [SRC]:
        return photo_idx
    pics = [
        s
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and (s.height or 0) > Emu(2000000)
    ]
    if not pics:
        return photo_idx
    pic = max(pics, key=lambda s: (s.width or 0) * (s.height or 0))
    path = PHOTO_PATHS[photo_idx % len(PHOTO_PATHS)]
    photo_idx += 1
    try:
        with open(path, "rb") as f:
            blob = f.read()
        left, top, width, height = pic.left, pic.top, pic.width, pic.height
        pic._element.getparent().remove(pic._element)
        slide.shapes.add_picture(BytesIO(blob), left, top, width, height)
    except Exception:
        pass
    return photo_idx


def main():
    global PHOTO_PATHS
    PHOTO_PATHS = extract_photos_from_src()
    if not PHOTO_PATHS:
        PHOTO_PATHS = [SRC]  # unused fallback; replace_main_picture no-ops safely

    prs = Presentation(SRC)

    # Templates before deletion
    hero_elems = [deepcopy(s.element) for s in prs.slides[0].shapes]
    bullet_elems = [deepcopy(s.element) for s in prs.slides[1].shapes]
    twocol_elems = [deepcopy(s.element) for s in prs.slides[5].shapes]

    # Drop all source slides; rebuild cleanly
    delete_slides(prs, list(range(len(prs.slides))))

    def hero():
        return clone_slide(prs, hero_elems)

    def bullets():
        return clone_slide(prs, bullet_elems)

    def twocol():
        return clone_slide(prs, twocol_elems)

    photo_idx = 0

    # —— 1. Hero ——
    s = hero()
    fill_hero(
        s,
        "РАЗДЕЛ 1. СВАРОЧНЫЕ РАБОТЫ НА ОПО",
        "Общие требования к производству сварочных работ",
        "Производство сварочных работ — деятельность персонала сварочного производства "
        "с применением сварочных и родственных процессов, сварочных материалов и оборудования "
        "с соблюдением норм, правил, методик и условий для получения сварных соединений "
        "с качеством, соответствующим нормативным требованиям.",
        note="На ОПО сварку выполняют по аттестованным технологиям и ПТД.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— 2. Personnel duties ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Персонал",
        "Обязанности персонала сварочного производства",
        "Сварщики, операторы, специалисты и контролёры обеспечивают:",
        [
            "техническую и технологическую подготовку и выполнение сварочных работ по ФНП и НД",
            "безопасную эксплуатацию, обслуживание и ремонт сварочного оборудования",
            "соблюдение технологий сварки",
            "контроль качества сварных соединений",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— 3. Certification ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Аттестация персонала",
        "Квалификация и аттестация сварщиков и специалистов",
        "Допуск к сварочным работам на ОПО",
        [
            "квалификация соответствует видам работ и применяемым технологиям сварки",
            "сварщики и специалисты аттестованы для способов, конструкций, положений и материалов",
            "допуск — по положительным результатам аттестационных испытаний",
            "аттестацию проводят независимые аттестационные центры",
        ],
        note="Без положительных результатов аттестационных испытаний к работам не допускают.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— 4. Technology readiness ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Аттестация технологий",
        "Проверка готовности к применению аттестованных технологий",
        "Организации и индивидуальные предприниматели",
        [
            "работы выполняют организации (ИП), прошедшие проверку готовности технологий на ОПО",
            "проверка подтверждает технические, организационные и квалификационные возможности",
            "контрольные сварные соединения выполняют на месте производства работ",
            "положительные результаты оформляют документом с подтверждёнными характеристиками",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— 5. Organization ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Организация работ",
        "Кто обеспечивает организацию и выполнение",
        "Ответственность руководства",
        [
            "организацию обеспечивает руководитель организации (ИП) или уполномоченное лицо",
            "аттестационные процедуры — руководитель независимого аттестационного центра",
            "работы выполняют по ПТД, разработанной специалистом сварочного производства",
            "ПТД утверждает руководитель (технический руководитель) или ИП",
        ],
        note="ПТД разрабатывают на основании проектной документации и НПА/НД РФ.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— 6–10. PTD content (18 items) ——
    ptd_items = [
        "способы сварки",
        "требования к квалификации, аттестации и допускным испытаниям сварщиков",
        "требования к сборке, прихваткам и приварке временных технологических креплений",
        "конструкция нестандартизированных сварных соединений (в т. ч. разной толщины)",
        "требования к хранению и подготовке сварочных материалов",
        "сочетания марок основных и сварочных материалов",
        "типоразмеры сварочных материалов (электрод, проволока, лента)",
        "используемое сварочное оборудование",
        "род и полярность сварочного тока",
        "типы выполняемых сварных соединений",
        "режимы сварки для конкретных сварных соединений",
        "необходимость, методы и режимы предварительного и сопутствующего подогрева",
        "пространственные положения при сварке",
        "порядок и последовательность выполнения шва (наплавки)",
        "способы защиты зоны сварки",
        "порядок и способы маркировки сварных соединений",
        "методы и объёмы НК и механических испытаний (при наличии требований НД)",
        "требования к исправлению дефектов и контролю после исправления",
    ]
    parts = list(chunks(ptd_items, 4))
    for idx, part in enumerate(parts, 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"РАЗДЕЛ 1 · ПТД ({idx}/{len(parts)})",
            "Производственно-технологическая документация по сварке",
            "В технологических инструкциях и картах сварки устанавливают:",
            part,
            note=(
                "Конструктивные элементы, режимы и контроль указывают в технологических картах."
                if idx == len(parts)
                else None
            ),
        )
        photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Assembly ——
    asm = [
        "способы подготовки поверхностей деталей под сварку",
        "приспособления и оборудование для сборки",
        "порядок и последовательность сборки; способы крепления деталей",
        "способы сварки, материалы и режимы прихваток и временных креплений",
        "размеры, количество и расположение прихваток",
        "методы контроля качества сборки",
    ]
    for idx, part in enumerate(chunks(asm, 4), 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"РАЗДЕЛ 1 · Сборка ({idx}/2)",
            "Требования по сборке деталей под сварку",
            "В ПТД по сборке должны быть приведены:",
            part,
            note=(
                "Стыковые кольцевые швы собирают с соосным позиционированием и фиксацией."
                if idx == 2
                else None
            ),
        )
        photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Equipment ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Оборудование и материалы",
        "Сварочное оборудование и сварочные материалы",
        "Соответствие аттестованным технологиям",
        [
            "оборудование и материалы соответствуют применяемым аттестованным технологиям",
            "сварочно-технологические характеристики обеспечивают качество соединений",
            "соответствие технологиям и нормам качества подтверждают независимые аттестационные центры",
            "оборудование содержат в исправном состоянии по указаниям производителя",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Supervisor ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Руководство работами",
        "Обязанности руководителя перед началом работ",
        "Лицо, осуществляющее руководство сварочными работами, обязано:",
        [
            "проверить состав и квалификацию персонала, оборудование, материалы и технологию по ПТД",
            "ознакомить сварщиков с технологическими картами и изменениями — под подпись",
            "организовать проведение операционного контроля",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Admission / workplace ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Допуск и рабочее место",
        "Допускные испытания и комплектация места сварки",
        "Перед допуском к сварке на объекте",
        [
            "при первом допуске или после перерыва дольше нормы НД — допускные соединения",
            "конструкцию, методы и объём контроля допускных соединений определяет руководитель",
            "место работ комплектуют исправным оборудованием, оснасткой и инструментом по ПТД",
            "место сварки защищают от осадков, влаги, сквозняков и иных воздействий",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Control overview (two-col) ——
    s = twocol()
    fill_two_col(
        s,
        "РАЗДЕЛ 1 · Контроль качества",
        "Место контроля в производственном процессе",
        "Виды контроля при подготовке и выполнении сварочных работ",
        "Этапы контроля",
        [
            "Входной — материалы и документация до начала работ",
            "Операционный — подготовка, сборка, прихватка, сварка, обработка",
            "Приёмочный — все выполненные сварные соединения по ПТД",
        ],
        "Кто обеспечивает контроль",
        [
            "персонал сварочного производства (сварщики, специалисты, контролёры)",
            "организация / ИП определяет состав персонала и обеспечивает требования",
            "руководитель сварочных работ организует операционный контроль",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Control types short ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Виды контроля",
        "Входной, операционный и приёмочный контроль",
        "Что подлежит контролю",
        [
            "входной — все партии свариваемых и сварочных материалов до применения",
            "операционный — все операции подготовки кромок, сборки, прихватки, сварки и обработки",
            "приёмочный — все выполненные сварные соединения; методы и объём — по ПТД",
        ],
        note="При трещинах или недопустимых дефектах работы останавливают до устранения причин.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Incoming ——
    incoming = [
        "документы о качестве продукции (идентификация, материалы, результаты испытаний)",
        "маркировка на каждом упаковочном месте: марка, сортамент, номер партии",
        "отсутствие повреждений упаковки и самих материалов",
        "покрытые электроды — номинальные размеры и состояние покрытия по сертификату",
        "проволока и лента — размеры, вид и состояние поверхности, маркировка",
        "флюс — цвет, однородность и гранулометрический состав",
    ]
    for idx, part in enumerate(chunks(incoming, 4), 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"РАЗДЕЛ 1 · Входной контроль ({idx}/2)",
            "Входной контроль материалов",
            "Что проверяют до применения партий материалов",
            part,
        )
        photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Operational ——
    operational = [
        "наличие маркировки",
        "размеры деталей и форму разделки кромок",
        "качество подготовленных под сварку поверхностей (если регламентировано ПТД)",
        "марки и типоразмеры сварочных материалов для прихваток",
        "надёжность фиксации и расположение деталей в сборочных приспособлениях",
        "чистоту и отсутствие повреждений кромок и прилегающих поверхностей",
        "размеры и расположение прихваток и швов приварки временных креплений",
        "зазоры, смещение кромок, перелом осей или плоскостей после прихваток",
        "размеры собранной под сварку конструкции",
    ]
    parts = list(chunks(operational, 4))
    for idx, part in enumerate(parts, 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"РАЗДЕЛ 1 · Операционный контроль ({idx}/{len(parts)})",
            "Операционный контроль сборки и подготовки",
            "Что проверяют на операциях подготовки и сборки",
            part,
        )
        photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— During welding ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Контроль в процессе сварки",
        "Минимальный контроль во время сварки",
        "В процессе сварки проводят, как минимум:",
        [
            "соответствие параметров режима сварки и технологических приёмов",
            "очерёдность выполнения сварных швов и участков наплавки",
            "отсутствие видимых дефектов",
            "иные параметры по технологическим (операционным) картам сварки",
        ],
        note="Устранение дефектов — по ПТД; после исправления — повторный контроль.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Marking / repair ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Маркировка и ремонт",
        "Маркировка швов и исправление дефектов",
        "После выполнения и при выявлении дефектов",
        [
            "сварные соединения маркируют по требованиям ПТД",
            "маркировка — шифр клейма сварщика для однозначной идентификации",
            "недопустимые дефекты исправляют по ПТД и карте на ремонт",
            "число исправлений одного участка — не выше указанного в ПТД",
        ],
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Documentation duties ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Документация",
        "Что обязан обеспечить руководитель сварочных работ",
        "При производстве сварочных работ обеспечивают:",
        [
            "идентификацию применяемых сварочных материалов и оборудования",
            "выполнение соединений по технологическим (операционным) картам сварки",
            "регистрацию сведений о сварщиках, выполняющих соединения",
            "идентификацию мест швов в конструкции и мест исправления дефектов",
        ],
        note="Также регистрируют результаты контроля качества, включая контроль после ремонта.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # —— Executive docs ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Итоги темы",
        "Исполнительная документация и ключевые требования",
        "В процессе сварочных работ оформляют:",
        [
            "журналы сварочных работ",
            "паспорта, акты и заключения по неразрушающему контролю",
            "протоколы испытаний сварных соединений",
            "иную исполнительную и эксплуатационную документацию по ПТД и НД",
        ],
        note="Сварка на ОПО — аттестованным персоналом, по ПТД и с полным циклом контроля.",
    )
    photo_idx = replace_main_picture(s, PHOTO_PATHS[0], photo_idx)

    # Save intermediate to avoid zip issues, then fix pages
    prs.save(OUT)
    prs = Presentation(OUT)
    update_pages(prs)
    prs.save(OUT)

    print(f"Saved {OUT} with {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        titles = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip().split("\n")[0]
                if t and " / " not in t:
                    titles.append(t[:60])
        print(f"  {i:02d}: {' | '.join(titles[:3])}")


if __name__ == "__main__":
    main()
