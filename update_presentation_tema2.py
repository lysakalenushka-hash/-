#!/usr/bin/env python3
"""Update Tema 2 presentation safely: recolor, prune, reload, then append."""

from copy import deepcopy
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________d3fc.pptx"
TMP = "/tmp/tema2_step1.pptx"
OUT = "/workspace/Презентация_Тема_2_Допуск.pptx"
IMG = "/workspace/assets/admission_images"

NAVY = "1B3A5B"
RED_HEX = "E30613"
RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x5B, 0x68, 0x78)
AMBER = RGBColor(0xA3, 0x72, 0x00)

# 0-based: quiz 12-16, dup 22, dup 26, quiz 27
DELETE_INDICES = {11, 12, 13, 14, 15, 21, 25, 26}


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def recolor_tree(element):
    for node in element.iter():
        tag = node.tag.split("}")[-1] if "}" in node.tag else node.tag
        if tag == "srgbClr":
            val = (node.get("val") or "").upper()
            if val in {NAVY, "FF0000"}:
                node.set("val", RED_HEX)


def recolor_presentation(prs):
    for slide in prs.slides:
        recolor_tree(slide._element)
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        c = r.font.color.rgb
                    except Exception:
                        continue
                    if c is not None and str(c).upper() in {NAVY, "FF0000"}:
                        r.font.color.rgb = RED


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    break


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 26, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 26, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    note_shape = None
    for sh in content[4:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    if note and note_shape:
        set_text(note_shape, note, 12, True, AMBER)


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3800000), width=Emu(10700000), height=Emu(2300000))
    else:
        slide.shapes.add_picture(path, Emu(6600000), Emu(1200000), width=Emu(5000000), height=Emu(4500000))


def find_templates(prs):
    text_tpl = list_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if text_tpl is None and 5 <= n <= 7:
            text_tpl = slide
        if list_tpl is None and n >= 10:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(t in {"01", "02", "1", "2"} for t in texts):
                list_tpl = slide
        if text_tpl and list_tpl:
            break
    return text_tpl, list_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    recolor_tree(blank._element)
    return blank


def append_missing(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    # Missing 2.1
    s = nl()
    fill_list_slide(
        s, "Тема 2.1 · Отказ и сроки",
        "Основания для отказа в выдаче разрешения на допуск",
        "Орган энергонадзора отказывает при наличии следующих оснований:",
        "Основания отказа",
        [
            "несоответствие содержания документов требованиям НПА",
            "несоответствие допускаемого объекта требованиям НПА",
            "несоответствие условий эксплуатации требованиям НПА",
            "незавершённость монтажа, наладки и испытаний (без заключения о соответствии)",
            "отсутствие разрешения на допуск энергопринимающих устройств объекта теплоснабжения",
        ],
        note="Разрешение на допуск объекта теплоснабжения — при наличии разрешения на допуск его ЭПУ.",
    )

    s = nl()
    fill_list_slide(
        s, "Тема 2.1 · Сроки действия",
        "Срок действия разрешения и продолжительность осмотра",
        "Ключевые сроки по ПП РФ № 85",
        "Сроки",
        [
            "разрешение на допуск (кроме временного) действует бессрочно",
            "временное разрешение действует 180 суток",
            "общая продолжительность осмотра — не более 7 рабочих дней",
            "при наличии ранее выданного временного разрешения — осмотр не более 4 рабочих дней",
            "при более поздней дате осмотра сроки рассмотрения увеличиваются",
        ],
    )
    add_image(s, f"{IMG}/kipia_instruments.png", bottom=True)

    s = nt()
    fill_text_slide(
        s, "Тема 2.1 · Аварийные работы",
        "Допуск при аварийно-восстановительных работах",
        "Особый порядок без предварительного разрешения",
        "Оформление разрешения на допуск для включения энергоустановок в работу "
        "при аварийно-восстановительных работах не требуется.\n\n"
        "Собственник письменно уведомляет орган энергонадзора в течение 10 рабочих дней "
        "после ввода в работу. Получение разрешения (временного) — в течение 90 календарных "
        "дней после ввода.",
        note="Исключение действует только для аварийно-восстановительных работ.",
    )
    add_image(s, f"{IMG}/pipe_burst_geyser.png")

    # Topic 2.2
    s = nt()
    fill_text_slide(
        s, "Тема 2.2 · ПТЭ потребителей",
        "Правила технической эксплуатации электроустановок потребителей",
        "Приказ Минэнерго России от 12 августа 2022 г. № 811",
        "Правила устанавливают требования к организации и осуществлению технической "
        "эксплуатации электроустановок потребителей.\n\n"
        "Распространяются на юрлиц, ИП и физлиц, владеющих электроустановками. "
        "Не распространяются на физлиц с ЭУ ниже 1000 В для личных или бытовых нужд.",
    )
    add_image(s, f"{IMG}/switchgear_room.png")

    s = nl()
    fill_list_slide(
        s, "Тема 2.2 · ПТЭ потребителей",
        "Что включает техническая эксплуатация электроустановок",
        "Согласно Приказу Минэнерго № 811 эксплуатация должна включать:",
        "Состав эксплуатации",
        [
            "ввод в работу новых и реконструированных электроустановок и оборудования",
            "использование по функциональному назначению и ведение документации",
            "оперативно-технологическое управление",
            "ремонт и техническое обслуживание",
            "консервацию, реконструкцию, модернизацию",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "Тема 2.2 · ПТЭ потребителей",
        "Обязанности потребителя при эксплуатации",
        "Потребитель должен обеспечить:",
        "Ключевые обязанности",
        [
            "содержание ЭУ в исправном состоянии и безопасную эксплуатацию",
            "ТО и ремонт, оперативно-технологическое управление",
            "подготовку и подтверждение готовности персонала",
            "исправность устройств РЗА, учёт и расследование аварий",
            "наличие документации, средств защиты, испытания оборудования",
            "эксплуатацию молниезащиты, приборов и средств учёта",
        ],
    )
    add_image(s, f"{IMG}/switchgear_room.png", bottom=True)

    s = nl()
    fill_list_slide(
        s, "Тема 2.2 · Ответственный за электрохозяйство",
        "Ответственный за электрохозяйство и группы по ЭБ",
        "Назначается из числа административно-технического персонала",
        "Назначение и группы",
        [
            "руководитель назначает ответственного и его заместителя",
            "для ИП обязанность возлагается непосредственно на ИП",
            "группа V — выше 1000 В; группа IV — до 1000 В",
            "при простых ЭУ (до 0,4 кВ / до 150 кВт) ответственность может быть на ЕИО",
            "списки работников с правами переключений — у ответственного и на рабочих местах",
            "ежегодно до 1 января списки передаются в сетевую организацию",
        ],
        note="Назначение — после проверки знаний и присвоения группы по электробезопасности.",
    )

    s = nl()
    fill_list_slide(
        s, "Тема 2.2 · Документация и персонал",
        "Техническая документация и категории персонала",
        "Ключевые требования Приказа № 811",
        "Документация и персонал",
        [
            "перечни документации и схемы — пересмотр/сверка не реже 1 раза в 3 года",
            "производственные инструкции — пересмотр не реже 1 раза в 3 года",
            "персонал: АТП, оперативный, оперативно-ремонтный, ремонтный, вспомогательный",
            "очередная проверка знаний: раз в 12 мес. / раз в 3 года (АТП и ОТ)",
            "комплексное опробование: оборудование 72 ч, ЛЭП — 24 ч",
            "измерение сопротивления изоляции электросварочных установок — не реже 1 раза в 6 мес.",
        ],
    )

    # Topic 2.3
    s = nt()
    fill_text_slide(
        s, "Тема 2.3 · Переключения",
        "Правила переключений в электроустановках",
        "Приказ Минэнерго России от 13 сентября 2018 г. № 757",
        "Правила устанавливают требования к инструкциям, персоналу, командам "
        "(разрешениям, подтверждениям), программам и бланкам переключений.\n\n"
        "Переключения выполняются по программам/типовым программам (ДЦ, ЦУС) "
        "или по бланкам/типовым бланкам (оперативный персонал объектов), "
        "за исключением установленных случаев.",
    )
    add_image(s, f"{IMG}/switchgear_room.png")

    s = nl()
    fill_list_slide(
        s, "Тема 2.3 · Переключения",
        "Кто выполняет и контролирует переключения",
        "Полномочия персонала по Приказу № 757",
        "Персонал",
        [
            "диспетчерский персонал ДЦ; оперативный персонал ЦУС, ВЭС/СЭС, НСО",
            "оперативный (оперативно-ремонтный) персонал объектов электроэнергетики",
            "персонал РЗА — операции с устройствами РЗА",
            "контролирующее лицо обязательно, кроме единичных/несложных/аварийных/дистанционных",
            "не допускается менять обязанности и отвлекать персонал во время переключений",
            "не допускается плановые переключения в ОРУ при грозе или ветре > 20 м/с",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "Тема 2.3 · Переключения",
        "Программы, бланки и сложные переключения",
        "Оформление переключений",
        "Документы переключений",
        [
            "сложные переключения ДЦ/ЦУС — по программам или типовым программам",
            "сложные и при неисправной блокировке — по бланкам или типовым бланкам",
            "без программ/бланков — несложные и для ликвидации нарушений нормального режима",
            "к сложным: вывод/ввод ЛЭП 110 кВ+, АТ/Т 110 кВ+, ШСВ/ОВ и др.",
            "использованные бланки хранятся не менее 20 суток (или до ввода в работу)",
            "электронные программы — не менее 12 месяцев",
        ],
        note="Перед операцией — проверка надписи на аппарате; без прочтения надписи запрещено.",
    )
    add_image(s, f"{IMG}/kipia_instruments.png", bottom=True)

    # Topic 2.4
    s = nt()
    fill_text_slide(
        s, "Тема 2.4 · Молниезащита",
        "Устройство молниезащиты зданий, сооружений и промышленных коммуникаций",
        "Приказ Минэнерго России от 30.06.2003 № 280 (СО 153-34.21.122-2003)",
        "Инструкция распространяется на все виды зданий, сооружений и промышленные "
        "коммуникации.\n\n"
        "Устройство молниезащиты не предотвращает развитие молнии, но существенно "
        "снижает риск ущерба. Тип и размещение устройств выбирают на стадии проектирования.",
    )
    add_image(s, f"{IMG}/lightning_protection.png")

    s = nl()
    fill_list_slide(
        s, "Тема 2.4 · Молниезащита",
        "Основные элементы и классификация объектов",
        "Термины и виды объектов",
        "Ключевые понятия",
        [
            "молниеотвод: молниеприемник + токоотвод + заземлитель",
            "зоны защиты: 0, 0Е, 1 и прочие",
            "обычные объекты — жилые, административные, высота не более 60 м",
            "специальные — опасные для окружения/среды, высота > 60 м и др.",
            "естественные молниеприемники — металлические кровли, фермы, трубы",
            "заземлитель молниезащиты совмещают с заземлителями ЭУ и связи",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "Тема 2.4 · Молниезащита",
        "Приёмка и эксплуатация устройств молниезащиты",
        "Порядок ввода и проверки",
        "Требования",
        [
            "приёмка рабочей комиссией до монтажа технологического оборудования",
            "в комиссию: ответственный за электрохозяйство, подрядчик, пожарная охрана",
            "после приёмки — паспорт молниезащитных устройств и паспорт заземлителей",
            "ежегодно перед грозовым сезоном — проверка и осмотр всех устройств",
            "проверки также после изменений в системе и повреждений объекта",
            "заземлитель-контур — глубина ≥ 0,5 м, расстояние от стен ≥ 1 м",
        ],
        note="Задача эксплуатации — поддержание устройств в исправном и надёжном состоянии.",
    )
    add_image(s, f"{IMG}/lightning_protection.png", bottom=True)

    s = nl()
    fill_list_slide(
        s, "Итоги темы 2",
        "Ключевые требования темы 2",
        "Что должен знать персонал",
        "Главные выводы",
        [
            "допуск в эксплуатацию — по ПП РФ № 85; запись в реестре = разрешение",
            "временное разрешение — на испытания/ПНР, действует 180 суток",
            "ПТЭ потребителей — Приказ № 811; ответственный за электрохозяйство",
            "переключения — Приказ № 757; программы/бланки, контролирующее лицо",
            "молниезащита — Приказ № 280; проверка ежегодно перед грозовым сезоном",
            "незнание НПА не освобождает от ответственности",
        ],
    )


def main():
    # Step 1: recolor + delete + save
    prs = Presentation(SRC)
    recolor_presentation(prs)
    delete_slides(prs, DELETE_INDICES)
    # images on early slides
    try:
        add_image(prs.slides[0], f"{IMG}/electrical.png")
        add_image(prs.slides[1], f"{IMG}/switchgear_room.png")
    except Exception:
        pass
    prs.save(TMP)

    # Step 2: reload clean package, append missing, renumber
    prs = Presentation(TMP)
    append_missing(prs)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
