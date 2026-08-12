#!/usr/bin/env python3
"""Обновить ВСЕ презентации темы (Яндекс.Диск) под № 805 + пособие, сохранив стиль."""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

SRC = Path("/workspace/tema3_yadisk_fixed")
OUT = Path("/workspace/tema3_yadisk_updated")
ZIP = Path("/workspace/Тема3_все_презентации_обновление_805.zip")

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

EYEBROW_THEME = "Тема 3. Оказание первой помощи при отсутствии сознания, остановке дыхания и кровообращения"
EYEBROW_31 = "Тема 3 · п. 3.1  Определение признаков жизни  ·  ред. № 805 с 01.09.2026"
EYEBROW_33 = "Тема 3 · п. 3.3  Техника проведения СЛР  ·  ред. № 805 с 01.09.2026"
EYEBROW_35 = "Тема 3 · п. 3.5  Инородное тело в дыхательных путях  ·  ред. № 805 с 01.09.2026"
EYEBROW_EXCL = "Доп. материал · исключено из Прил. № 2 с 01.09.2026 (ПП РФ № 805)"
EYEBROW_OV = "Тема 3 · Прил. № 2 к № 2464 (ред. № 805 с 01.09.2026)"


def first_run_style(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text:
                continue
            color = None
            try:
                if run.font.color is not None and run.font.color.rgb is not None:
                    color = run.font.color.rgb
            except Exception:
                pass
            return run.font.name, run.font.size, run.font.bold, color
    return "Open Sans", Pt(28), False, RGBColor(0x40, 0x40, 0x40)


def set_shape_plain(shape, text: str):
    name, size, bold, color = first_run_style(shape)
    tf = shape.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for t_el in p0._p.findall(f".//{A_NS}t"):
        t_el.text = ""
    run = p0.add_run()
    run.text = text
    if name:
        run.font.name = name
    if size:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def set_shape_multiline(shape, lines: list[tuple[str, bool]], *, size=None):
    name0, size0, _, color0 = first_run_style(shape)
    font = name0 or "Open Sans"
    size = size or size0 or Pt(28)
    color = color0 or RGBColor(0x53, 0x53, 0x53)
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for t_el in p0._p.findall(f".//{A_NS}t"):
        t_el.text = ""
    for i, (text, bold) in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        if i > 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = size
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def replace_in_slide(slide, replacements: list[tuple[str, str | list]]):
    """replacements: (substring_to_find, new_text_str OR list[(text,bold)])."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        for needle, new in replacements:
            if needle in t:
                if isinstance(new, list):
                    set_shape_multiline(sh, new)
                else:
                    set_shape_plain(sh, new)
                break


def replace_eyebrows(prs, new_eyebrow: str):
    old = "Оказание первой помощи при отсутствии сознания"
    for i in range(len(list(prs.slides._sldIdLst))):
        for sh in prs.slides[i].shapes:
            if not sh.has_text_frame:
                continue
            t = "".join(p.text for p in sh.text_frame.paragraphs)
            if old in t and len(t) < 200:
                set_shape_plain(sh, new_eyebrow)


# ---------- individual updates ----------

def update_overview() -> Path:
    src = next(SRC.glob("2 Оказание*.pptx"))
    out = OUT / "3_Тема3_обзор_сознание_дыхание_кровообращение.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    # Title stays valid for Theme 3
    replace_eyebrows(prs, EYEBROW_OV)

    # S2 signs of life
    s = prs.slides[1]
    replace_in_slide(
        s,
        [
            ("ОСНОВНЫЕ ПРИЗНАКИ ЖИЗНИ", "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ"),
            (
                "К основным признакам жизни",
                [
                    ("П. 3.1 (№ 805): «Определение признаков жизни».", True),
                    ("Для решения о СЛР ориентируйтесь на отсутствие сознания и нормального дыхания.", False),
                    ("Проверка пульса для решения о СЛР не рекомендуется.", False),
                    ("Внезапная смерть — заболевания или внешнее воздействие; алгоритм СЛР единый.", False),
                ],
            ),
        ],
    )

    # Algorithm slides titles
    for i in [2, 3, 4, 5]:
        replace_in_slide(
            prs.slides[i],
            [("АЛГОРИТМ ПРОВЕДЕНИЯ СЕРДЕЧНО-ЛЕГОЧНОЙ РЕАНИМАЦИИ", "ТЕХНИКА ПРОВЕДЕНИЯ СЛР (п. 3.3)")],
        )

    # Errors slide — mark excluded
    replace_in_slide(
        prs.slides[6],
        [
            ("ОШИБКИ И ОСЛОЖНЕНИЯ", "ОШИБКИ И ОСЛОЖНЕНИЯ СЛР (ДОП. · ИСКЛЮЧЕНО ИЗ ПРИЛ. № 2)"),
            (
                "К основным ошибкам",
                [
                    ("Исключено как отдельный пункт с 01.09.2026 (№ 805). По пособию:", True),
                    ("• нарушение последовательности; неверная техника давления/вдохов;", False),
                    ("• соотношение не 30:2; частота <100 или >120; паузы >10 с.", False),
                    ("Осложнение: перелом рёбер (точка давления, избыточная сила, хрупкость костей).", False),
                ],
            ),
        ],
    )
    replace_in_slide(prs.slides[7], [("ПОКАЗАНИЯ К ПРЕКРАЩЕНИЮ СЛР", "ПРЕКРАЩЕНИЕ СЛР (ДОП. · ПО ПОСОБИЮ)")])
    for i in [8, 9]:
        replace_in_slide(
            prs.slides[i],
            [("МЕРОПРИЯТИЯ, ВЫПОЛНЯЕМЫЕ ПОСЛЕ", "ЕСЛИ ПОЯВИЛИСЬ ПРИЗНАКИ ЖИЗНИ → п. 3.2 (боковое положение)")],
        )

    replace_in_slide(
        prs.slides[10],
        [("ОСОБЕННОСТИ СЛР У ДЕТЕЙ", "ОСОБЕННОСТИ СЛР У ДЕТЕЙ (ДОП. · ИСКЛЮЧЕНО ИЗ ПРИЛ. № 2)")],
    )

    for i in [11, 12]:
        replace_in_slide(
            prs.slides[i],
            [("ПЕРВАЯ ПОМОЩЬ ПРИ НАРУШЕНИИ ПРОХОДИМОСТИ", "П. 3.5  ИНОРОДНОЕ ТЕЛО В ДЫХАТЕЛЬНЫХ ПУТЯХ")],
        )

    # Add note slide before thanks? Insert by editing thanks is hard; add text on last content.
    # Soft-add a small note on slide 1 title area via replacing nothing — instead add shape on slide 0 after title
    s0 = prs.slides[0]
    tb = s0.shapes.add_textbox(Emu(1780000), Emu(9000000), Emu(18000000), Emu(1200000))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Структура теории с 01.09.2026 (№ 805): 3.1 признаки жизни · 3.2 проходимость ДП · 3.3 техника СЛР · 3.4 АНД · 3.5 инородное тело · 3.6 иные нарушения дыхания. Ошибки СЛР и особенности у детей — исключены как отдельные пункты."
    run.font.name = "Open Sans"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x53, 0x53, 0x53)

    prs.save(out)
    print("saved", out.name)
    return out


def update_21() -> Path:
    src = SRC / "2.1_Основные признаки жизни у пострадавшего.pptx"
    out = OUT / "3.1_Определение_признаков_жизни.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    replace_in_slide(prs.slides[0], [("ОСНОВНЫЕ ПРИЗНАКИ ЖИЗНИ У ПОСТРАДАВШЕГО", "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ")])
    replace_eyebrows(prs, EYEBROW_31)

    replace_in_slide(
        prs.slides[1],
        [
            ("ОСНОВНЫЕ ПРИЗНАКИ ЖИЗНИ У ПОСТРАДАВШЕГО", "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ"),
            (
                "К основным признакам жизни",
                [
                    ("По № 805 теория: «Определение признаков жизни».", True),
                    ("Для СЛР: нет сознания + нет нормального дыхания.", False),
                    ("Пульс на артериях для решения о СЛР не рекомендуется.", False),
                    ("Причины: заболевания или внешнее воздействие; алгоритм СЛР единый.", False),
                ],
            ),
        ],
    )
    replace_in_slide(prs.slides[2], [("ПРИЧИНЫ НАРУШЕНИЯ", "КОНТЕКСТ: ПРИЧИНЫ ОСТАНОВКИ ДЫХАНИЯ И КРОВООБРАЩЕНИЯ")])
    replace_in_slide(
        prs.slides[3],
        [
            ("ОСНОВНЫЕ ПРИЗНАКИ ЖИЗНИ У ПОСТРАДАВШЕГО", "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ"),
            (
                "Внезапная смерть",
                [
                    ("Заболевания: инфаркт, нарушения ритма и др.", False),
                    ("Внешние воздействия: травма, электротравма, утопление и др.", False),
                    ("Алгоритм СЛР единый вне зависимости от причины (пособие Минздрава).", True),
                ],
            ),
        ],
    )
    replace_in_slide(
        prs.slides[4],
        [("СПОСОБЫ ПРОВЕРКИ СОЗНАНИЯ, ДЫХАНИЯ, КРОВООБРАЩЕНИЯ", "СПОСОБЫ ПРОВЕРКИ СОЗНАНИЯ И ДЫХАНИЯ")],
    )
    replace_in_slide(
        prs.slides[5],
        [
            ("АЛГОРИТМ ПРОВЕДЕНИЯ СЕРДЕЧНО-ЛЕГОЧНОЙ РЕАНИМАЦИИ", "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ: ПОРЯДОК ПРОВЕРКИ"),
            ("Шаг 1 Наличие сознания", "Шаг 1. Сознание: тормошение за плечи + «Что с вами? Нужна ли вам помощь?»"),
            ("Шаг 2 Открытие дыхательных путей", "Шаг 2. Открытие дыхательных путей: запрокинуть голову, поднять подбородок"),
            ("Шаг 3 Проверка дыхания", "Шаг 3. Дыхание ≤ 10 с. Нет / агональное → СМП + СЛР"),
        ],
    )
    prs.save(out)
    print("saved", out.name)
    return out


def update_22() -> Path:
    src = SRC / "2.2_Современный алгоритм проведения.pptx"
    out = OUT / "3.3_Техника_проведения_СЛР.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    # Old title was long App2 old wording; new is shorter
    replace_in_slide(
        prs.slides[0],
        [
            (
                "СОВРЕМЕННЫЙ АЛГОРИТМ",
                "ТЕХНИКА ПРОВЕДЕНИЯ СЕРДЕЧНО-ЛЕГОЧНОЙ РЕАНИМАЦИИ",
            )
        ],
    )
    replace_eyebrows(prs, EYEBROW_33)
    for i in range(1, 5):
        replace_in_slide(
            prs.slides[i],
            [("АЛГОРИТМ ПРОВЕДЕНИЯ СЕРДЕЧНО-ЛЕГОЧНОЙ РЕАНИМАЦИИ", "ТЕХНИКА ПРОВЕДЕНИЯ СЛР")],
        )

    # Clarify step labels per methodicheka / new practice order (compressions before breaths)
    replace_in_slide(
        prs.slides[1],
        [
            ("Шаг 1 Наличие сознания", "Шаг 1. Сознание"),
            ("Шаг 2 Открытие дыхательных путей", "Шаг 2. Открытие дыхательных путей"),
            ("Шаг 3 Проверка дыхания", "Шаг 3. Проверка дыхания (≤ 10 с)"),
        ],
    )
    replace_in_slide(prs.slides[2], [("Шаг 4 Вызов", "Шаг 4. Вызов СМП (112 / 103)")])
    replace_in_slide(prs.slides[3], [("Шаг 5 Давление", "Шаг 5. Давление на грудину (30; 100–120/мин; ~5–6 см)")])
    replace_in_slide(
        prs.slides[4],
        [("Шаг 6 Искусственное", "Шаг 6. Искусственное дыхание (2 вдоха; цикл 30:2)")],
    )

    # Note on title slide about ANД being separate 3.4
    s0 = prs.slides[0]
    tb = s0.shapes.add_textbox(Emu(1780000), Emu(9500000), Emu(18000000), Emu(900000))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "П. 3.3 по № 805. АНД — отдельный п. 3.4 (при наличии). Порядок: компрессии → искусственное дыхание."
    run.font.name = "Open Sans"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x53, 0x53, 0x53)

    prs.save(out)
    print("saved", out.name)
    return out


def update_23() -> Path:
    src = SRC / "2.3_Ошибки и осложнения возникающие при проведении.pptx"
    out = OUT / "доп_исключено_Ошибки_осложнения_и_прекращение_СЛР.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    replace_in_slide(prs.slides[0], [("ОШИБКИ И ОСЛОЖНЕНИЯ", "ОШИБКИ И ОСЛОЖНЕНИЯ СЛР (ДОП. МАТЕРИАЛ)")])
    replace_eyebrows(prs, EYEBROW_EXCL)

    replace_in_slide(
        prs.slides[1],
        [
            ("ОШИБКИ И ОСЛОЖНЕНИЯ", "ОШИБКИ И ОСЛОЖНЕНИЯ ПРИ СЛР"),
            (
                "К основным ошибкам",
                [
                    ("С 01.09.2026 отдельный пункт исключён (№ 805). По пособию Минздрава:", True),
                    ("• нарушение последовательности СЛР;", False),
                    ("• неверная техника давления на грудину (точка, глубина, частота, расправление);", False),
                    ("• неверная техника искусственного дыхания; соотношение не 30:2;", False),
                    ("• частота < 100 или > 120 в минуту; паузы между циклами > 10 с.", False),
                    ("Осложнение: перелом рёбер — неверная точка, избыточная сила, хрупкость костей.", False),
                ],
                ),
        ],
    )
    # fix multiline call - I passed size wrongly inside tuple. Fix by not using size kw in tuple.
    # Actually set_shape_multiline got list - good. The trailing comma in tuple with size= was wrong - I used nested wrong.
    # Looking at my code - the second element of replace pair for errors is a list - good.
    # Wait I had `], ),`  with size inside - let me check the file... I wrote:
    # [
    #   ("ОШИБКИ...
    #   ( "К основным...", [ ... ], ),  # WRONG - third element size
    # ]
    # That would make replacements item a 3-tuple and break. Need to fix and re-run.

    replace_in_slide(prs.slides[2], [("ПОКАЗАНИЯ К ПРЕКРАЩЕНИЮ РЕАНИМАЦИИ", "ПРЕКРАЩЕНИЕ СЛР (ПО ПОСОБИЮ)")])
    replace_in_slide(
        prs.slides[3],
        [
            ("ПОКАЗАНИЯ К ПРЕКРАЩЕНИЮ СЛР", "КОГДА ПРЕКРАТИТЬ / НЕ НАЧИНАТЬ СЛР"),
            (
                "Реанимационные мероприятия продолжаются",
                [
                    ("До прибытия СМП/спецслужб и их распоряжения;", False),
                    ("до появления признаков жизни; можно прекратить при угрозе себе.", False),
                ],
            ),
            (
                "В случае длительного проведения",
                [
                    ("При усталости — привлечь помощника (смена ~ каждые 2 минуты).", False),
                    ("Можно не начинать при явных признаках нежизнеспособности или терминальной стадии неизлечимого заболевания.", False),
                ],
            ),
        ],
    )
    replace_in_slide(prs.slides[4], [("МЕРОПРИЯТИЯ, ВЫПОЛНЯЕМЫЕ ПОСЛЕ ПРЕКРАЩЕНИЯ", "ЕСЛИ ПОЯВИЛИСЬ ПРИЗНАКИ ЖИЗНИ")])
    for i in [5, 6]:
        replace_in_slide(
            prs.slides[i],
            [
                ("МЕРОПРИЯТИЯ, ВЫПОЛНЯЕМЫЕ ПОСЛЕ СЕРДЕЧНО-ЛЕГОЧНОЙ", "ПОДДЕРЖАНИЕ ПРОХОДИМОСТИ ПОСЛЕ СЛР (п. 3.2)"),
                ("Шаг 1", "Шаг 1. Ближняя рука под 90°"),
                ("Шаг 2", "Шаг 2. Дальняя рука к щеке"),
                ("Шаг 3", "Шаг 3. Поворот набок"),
                ("Шаг 4", "Шаг 4. Нога к животу, проверить дыхание"),
                ("Шаг 5", "Шаг 5. Наблюдение; каждые 30 мин — другой бок"),
            ],
        )

    prs.save(out)
    print("saved", out.name)
    return out


def update_24() -> Path:
    src = SRC / "2.4_Особенности реанимации у детей.pptx"
    out = OUT / "доп_исключено_Особенности_СЛР_у_детей.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    replace_in_slide(prs.slides[0], [("ОСОБЕННОСТИ РЕАНИМАЦИИ У ДЕТЕЙ", "ОСОБЕННОСТИ СЛР У ДЕТЕЙ (ДОП. МАТЕРИАЛ)")])
    replace_eyebrows(prs, EYEBROW_EXCL)
    replace_in_slide(prs.slides[1], [("ОСОБЕННОСТИ СЛР У ДЕТЕЙ", "ОСОБЕННОСТИ СЛР У ДЕТЕЙ")])

    s = prs.slides[1]
    # move pictures down
    for sh in s.shapes:
        if sh.shape_type == 13 and sh.top < Emu(5500000):
            sh.top = Emu(5300000)

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(1800000), Emu(3450000), Emu(20800000), Emu(1650000))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF8)
    box.line.fill.background()
    try:
        box.adjustments[0] = 0.04
    except Exception:
        pass
    tb = s.shapes.add_textbox(Emu(2000000), Emu(3520000), Emu(20400000), Emu(1500000))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        ("С 01.09.2026 отдельный пункт исключён из Прил. № 2 (№ 805). По пособию Минздрава:", True),
        ("• Та же последовательность и 30 : 2, что у взрослых.", False),
        ("• Глубина ≈ ⅓ размера груди (~4 см до 1 года; ~5 см старше). До 1 года — двумя пальцами; старше — одной/двумя руками.", False),
        ("• После отсутствия признаков жизни у ребёнка эффективнее 5 вдохов, затем 30:2. При утоплении — так же.", False),
    ]
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.name = "Open Sans"
        run.font.size = Pt(20)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    prs.save(out)
    print("saved", out.name)
    return out


def update_25() -> Path:
    src = SRC / "2.5_Порядок оказания первой помощи.pptx"
    out = OUT / "3.5_Инородное_тело_в_дыхательных_путях.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    # New App2 short title; keep methodicheka details inside
    replace_in_slide(
        prs.slides[0],
        [
            (
                "ПОРЯДОК ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ ПРИ ЧАСТИЧНОМ",
                "ПЕРВАЯ ПОМОЩЬ ПРИ НАРУШЕНИИ ПРОХОДИМОСТИ ДЫХАТЕЛЬНЫХ ПУТЕЙ ИНОРОДНЫМ ТЕЛОМ",
            )
        ],
    )
    replace_eyebrows(prs, EYEBROW_35)

    for i in [1, 2, 4]:
        replace_in_slide(
            prs.slides[i],
            [("ПЕРВАЯ ПОМОЩЬ ПРИ НАРУШЕНИИ ПРОХОДИМОСТИ", "П. 3.5  ИНОРОДНОЕ ТЕЛО В ДЫХАТЕЛЬНЫХ ПУТЯХ")],
        )

    # Special cases divider — still in methodicheka but generalized out of App2 wording
    replace_in_slide(
        prs.slides[3],
        [
            (
                "ОСОБЕННОСТИ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ ТУЧНОМУ",
                "ОСОБЕННОСТИ (ПО ПОСОБИЮ): ТУЧНЫЕ, БЕРЕМЕННЫЕ, ДЕТИ",
            )
        ],
    )

    s0 = prs.slides[0]
    tb = s0.shapes.add_textbox(Emu(1780000), Emu(9500000), Emu(18000000), Emu(900000))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "П. 3.5 по № 805 сформулирован кратко. Детализация частичное/полное, тучные/беременные/дети — по пособию Минздрава."
    run.font.name = "Open Sans"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x53, 0x53, 0x53)

    prs.save(out)
    print("saved", out.name)
    return out


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    paths = [
        update_overview(),
        update_21(),
        update_22(),
        update_23(),
        update_24(),
        update_25(),
    ]
    readme = OUT / "README.md"
    readme.write_text(
        """# Тема 3 — обновление всех презентаций с Яндекс.Диска

Источник: https://disk.360.yandex.ru/d/f1iR1Cn4F9c6Gw  
Учтены **ПП РФ № 805** (Прил. № 2 с 01.09.2026) и пособие Минздрава (2025).  
Стиль исходников сохранён (шрифты, картинки, размер слайдов).

| Было (на диске) | Стало | Статус № 805 |
|-----------------|-------|--------------|
| 2 Оказание первой помощи… | `3_Тема3_обзор_…pptx` | обзор темы 3 |
| 2.1 Признаки жизни | `3.1_Определение_признаков_жизни.pptx` | **п. 3.1** |
| 2.2 Современный алгоритм СЛР | `3.3_Техника_проведения_СЛР.pptx` | **п. 3.3** |
| 2.3 Ошибки и осложнения | `доп_исключено_Ошибки_…pptx` | исключено |
| 2.4 Особенности у детей | `доп_исключено_Особенности_СЛР_у_детей.pptx` | исключено |
| 2.5 Порядок при инородном теле | `3.5_Инородное_тело_в_дыхательных_путях.pptx` | **п. 3.5** |

**Пока нет отдельных исходников на диске для новых пунктов:** 3.2 (проходимость ДП), 3.4 (АНД), 3.6 (иные нарушения дыхания) — в обзоре указана полная структура теории.
""",
        encoding="utf-8",
    )
    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in paths:
            zf.write(p, arcname=p.name)
        zf.write(readme, arcname="README.md")
    print("ZIP", ZIP, ZIP.stat().st_size)


if __name__ == "__main__":
    main()
