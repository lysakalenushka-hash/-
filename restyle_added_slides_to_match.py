#!/usr/bin/env python3
"""Replace Arial-styled added slides with clones of original deck slides (same fonts/layout)."""

from copy import deepcopy
from pathlib import Path
import re
import shutil
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

OUT = Path("/workspace/first_aid_tema1_numbered")
ZIP = Path("/workspace/Тема1_нумерация_по_списку_Прил2_2464.zip")

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

TITLE_FONT = "PFDinDisplayPro-Regular"
TITLE_SIZE = Pt(46)  # 584200 EMU
TITLE_COLOR = RGBColor(0x40, 0x40, 0x40)
EYEBROW_FONT = "Open Sans"
EYEBROW_SIZE = Pt(28)
EYEBROW_COLOR = RGBColor(0x53, 0x53, 0x53)
BODY_SIZE = Pt(28)
BODY_COLOR = RGBColor(0x40, 0x40, 0x40)


# Content for supplements (same as before, slightly tighter for body box)
SUPPLEMENTS = {
    "1_Организация_оказания_первой_помощи_в_РФ_Нормативно-правовая_база.pptx": [
        {
            "after_title": True,
            "title": "ПЯТЬ КОМПОНЕНТОВ СИСТЕМЫ ПЕРВОЙ ПОМОЩИ",
            "eyebrow": "Тема 1 · п. 1  Организация и нормативно-правовая база",
            "body": [
                "По учебному пособию Минздрава (2025) система оказания первой помощи в РФ состоит из пяти основных компонентов:",
                "",
                "1. Организация и нормативно-правовое обеспечение.",
                "2. Обучение участников оказания первой помощи правилам и навыкам ее оказания.",
                "3. Оснащение участников средствами для ее оказания (аптечками, укладками, наборами, комплектами).",
                "4. Мотивирование на обучение и оказание первой помощи.",
                "5. Учет и анализ эффективности оказания первой помощи.",
            ],
        },
        {
            "before_thanks": True,
            "title": "ПОРЯДОК № 220н: СОГЛАСИЕ, ПРИОРИТЕТ, АНД",
            "eyebrow": "Тема 1 · п. 1  Организация и нормативно-правовая база",
            "body": [
                "По приказу Минздрава России от 03.05.2024 № 220н (с 01.09.2024):",
                "",
                "• Оказание первой помощи допускается, если отсутствует выраженный до начала оказания первой помощи отказ гражданина или его законного представителя.",
                "• Первая помощь оказывается при условии отсутствия угрожающих жизни и здоровью оказывающего ее лица факторов.",
                "• При двух и более пострадавших первоочередность — по тяжести состояния; приоритет отдаётся детям (несовершеннолетним).",
                "• Разрешено применение автоматических наружных дефибрилляторов (АНД) при наличии.",
                "• Допускается помощь пострадавшему в принятии лекарственных препаратов, назначенных врачом; могут использоваться подручные средства.",
            ],
        },
    ],
    "2_Современные_аптечки_укладки_комплекты_и_наборы.pptx": [
        {
            "after_title": True,
            "title": "ПРИКАЗЫ № 260н, 261н, 262н И СОСТАВ АПТЕЧЕК",
            "eyebrow": "Тема 1 · п. 2  Аптечки, укладки, комплекты и наборы",
            "body": [
                "Требования к комплектации (приказы Минздрава России от 24.05.2024):",
                "",
                "• № 260н — аптечка для оказания первой помощи пострадавшим в ДТП (автомобильная);",
                "• № 262н — аптечка для оказания работниками первой помощи пострадавшим;",
                "• № 261н — аптечка для организаций, осуществляющих образовательную деятельность.",
                "",
                "Типовые компоненты: жгут; бинты; салфетки стерильные; лейкопластыри; устройство для ИВЛ «Рот–Устройство–Рот»; ножницы; перчатки медицинские; маска медицинская; покрывало спасательное изотермическое.",
                "",
                "Пополнять аптечку — по мере расхода и/или истечения срока годности.",
            ],
        },
    ],
    "3_Порядок_и_приоритетность_оказания_первой_помощи.pptx": [
        {
            "after_title": True,
            "title": "ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            "eyebrow": "Тема 1 · п. 3  Порядок и приоритетность; извлечение и перемещение",
            "body": [
                "Если пострадавших несколько, а участников оказания помощи недостаточно — определяют приоритетность:",
                "",
                "• в первую очередь — наиболее тяжело пострадавшим и несовершеннолетним детям;",
                "• для взрослых приоритетность определяется последовательностью мероприятий Порядка оказания первой помощи (приказ № 220н);",
                "• в ряде случаев допустима самопомощь (например, прямое давление на рану), пока оказывается помощь другому пострадавшему;",
                "• более опытный участник может координировать действия остальных, направляя их к наиболее тяжелым пострадавшим.",
            ],
        },
    ],
    "4_Перечень_состояний_и_мероприятий_первой_помощи.pptx": [
        {
            "after_title": True,
            "title": "СТРУКТУРА ПОРЯДКА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ (№ 220н)",
            "eyebrow": "Тема 1 · п. 4  Состояния и мероприятия (приказ Минздрава № 220н)",
            "body": [
                "Порядок оказания первой помощи (приказ Минздрава № 220н) включает:",
                "",
                "• общие организационные положения;",
                "• перечень из 9 состояний (приложение № 1), в т.ч. острые психологические реакции на стресс; судорожный приступ с потерей сознания; укусы и ужаливания;",
                "• перечень из 9 мероприятий и последовательность их проведения (приложение № 2).",
                "",
                "Важные изменения относительно прежнего приказа № 477н:",
                "• не проверяют пульс для оценки кровообращения;",
                "• из перечня мероприятий убраны пальцевое прижатие артерии и максимальное сгибание конечности в суставе как обязательные техники для широкого обучения.",
            ],
        },
    ],
    "5_Обеспечение_безопасных_условий_и_профилактика_инфекций.pptx": [
        {
            "before_thanks": True,
            "title": "МАСКА И ПЕРЧАТКИ ИЗ АПТЕЧКИ",
            "eyebrow": "Тема 1 · п. 5  Безопасные условия и профилактика инфекций",
            "body": [
                "По учебному пособию Минздрава (2025):",
                "",
                "• В аптечке для оказания первой помощи работниками есть медицинские маски — для снижения риска инфицирования оказывающего помощь.",
                "• Эти маски не используются для проведения искусственного дыхания (для ИВЛ — отдельное устройство «Рот–Устройство–Рот»).",
                "• Перчатки медицинские — защита от контакта с кровью и другими биологическими жидкостями пострадавшего.",
            ],
        },
    ],
    "6_Основные_правила_вызова_скорой_медицинской_помощи.pptx": [
        {
            "after_title": True,
            "title": "НОМЕР 112 И ПОВОДЫ К ВЫЗОВУ СМП",
            "eyebrow": "Тема 1 · п. 6  Вызов СМП и специальных служб",
            "body": [
                "Единый номер экстренных служб — 112 (также 101, 102, 103 и региональные номера).",
                "",
                "Поводы к вызову скорой медицинской помощи (не все входят в перечень состояний первой помощи, но требуют вызова СМП): нарушения сознания, дыхания, кровообращения; психические расстройства с опасностью для себя/окружающих; болевой синдром; травмы, отравления, ранения; термические и химические ожоги; кровотечения; роды, угроза прерывания беременности.",
                "",
                "Сообщить диспетчеру: место и суть происшествия; число пострадавших, повреждения, тяжесть; какая помощь оказывается.",
                "Трубку отключать после сообщения диспетчера о том, что вызов принят.",
            ],
        },
    ],
}


def is_added_slide(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.name == "Arial":
                    return True
                try:
                    if r.font.color and r.font.color.rgb and str(r.font.color.rgb) == "E30613":
                        return True
                except Exception:
                    pass
    # also detect red left bar fill
    for sh in slide.shapes:
        try:
            if sh.fill.type is not None and sh.fill.fore_color.rgb and str(sh.fill.fore_color.rgb) == "E30613":
                if (sh.width or 0) < Emu(400000):
                    return True
        except Exception:
            pass
    return False


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(R_NS + "id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def duplicate_slide(prs, index):
    """Duplicate slide keeping image relationships."""
    source = prs.slides[index]
    # use blank-ish layout
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    dest = prs.slides.add_slide(blank_layout)

    # clear default shapes
    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)

    # copy relationships from source slide part to dest (images etc.)
    # Map old rId -> new rId
    rid_map = {}
    for rel in source.part.rels.values():
        if "image" in rel.reltype or rel.reltype.endswith("/image"):
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        elif rel.reltype.endswith("/notesSlide"):
            continue
        elif rel.reltype.endswith("/slideLayout"):
            continue
        else:
            # chart, hyperlink, etc.
            try:
                new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
                rid_map[rel.rId] = new_rid
            except Exception:
                pass

    # deepcopy shapes and rewrite rIds for blips
    for shape in source.shapes:
        el = deepcopy(shape.element)
        # rewrite r:embed / r:link
        for blip in el.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"):
            embed = blip.get(R_NS + "embed")
            if embed and embed in rid_map:
                blip.set(R_NS + "embed", rid_map[embed])
            link = blip.get(R_NS + "link")
            if link and link in rid_map:
                blip.set(R_NS + "link", rid_map[link])
        dest.shapes._spTree.insert_element_before(el, "p:extLst")

    return dest


def find_roles(slide):
    title = eyebrow = body = page = None
    pictures = []
    underline = None
    for sh in slide.shapes:
        if sh.shape_type is not None and sh.shape_type == 13:  # PICTURE
            pictures.append(sh)
            continue
        if not sh.has_text_frame:
            # gradient underline often has no/short text
            if (sh.height or 0) < Emu(150000) and (sh.top or 0) > Emu(2400000) and (sh.top or 0) < Emu(2900000):
                underline = sh
            continue
        t = sh.text_frame.text.strip()
        top = sh.top or 0
        left = sh.left or 0
        w = sh.width or 0
        h = sh.height or 0
        # page number
        if left < Emu(900000) and w < Emu(1200000) and top > Emu(5000000):
            page = sh
            continue
        # title: upper wide
        if top < Emu(2300000) and w > Emu(8000000) and h < Emu(1200000) and t:
            title = sh
            continue
        # eyebrow
        if ("Тема 1" in t or "Аспекты оказания" in t) and h < Emu(900000):
            eyebrow = sh
            continue
        # body: large text area
        if t and h > Emu(2000000) and w > Emu(5000000):
            if body is None or (h * w) > ((body.height or 0) * (body.width or 0)):
                body = sh
    return {"title": title, "eyebrow": eyebrow, "body": body, "page": page, "pictures": pictures, "underline": underline}


def clear_text_keep_style(shape, paragraphs_text):
    """Replace text in shape: reuse first paragraph run style, add paragraphs as needed."""
    tf = shape.text_frame
    # Clear all a:t
    for t_el in shape.element.findall(f".//{A_NS}t"):
        t_el.text = ""

    # Collect paragraph elements
    txBody = shape.text_frame._txBody
    p_elems = txBody.findall(f"{A_NS}p")
    if not p_elems:
        return

    # Use first paragraph as style template
    template_p = p_elems[0]
    # Find a run with rPr in template
    template_r = template_p.find(f".//{A_NS}r")
    template_rPr = None
    if template_r is not None:
        template_rPr = template_r.find(f"{A_NS}rPr")

    # Remove all paragraphs except we'll rebuild
    for p in list(p_elems):
        txBody.remove(p)

    for i, line in enumerate(paragraphs_text):
        p = etree.SubElement(txBody, f"{A_NS}p")
        # copy pPr from template if any
        pPr_src = template_p.find(f"{A_NS}pPr")
        if pPr_src is not None:
            p.append(deepcopy(pPr_src))
        r = etree.SubElement(p, f"{A_NS}r")
        if template_rPr is not None:
            rPr = deepcopy(template_rPr)
            # ensure readable size/color for body
            r.append(rPr)
        else:
            rPr = etree.SubElement(r, f"{A_NS}rPr")
            rPr.set("sz", "2800")
            rPr.set("dirty", "0")
            solid = etree.SubElement(rPr, f"{A_NS}solidFill")
            srgb = etree.SubElement(solid, f"{A_NS}srgbClr")
            srgb.set("val", "404040")
            latin = etree.SubElement(rPr, f"{A_NS}latin")
            latin.set("typeface", "Open Sans")
        t = etree.SubElement(r, f"{A_NS}t")
        # preserve spaces
        if line.startswith(" ") or line.endswith(" "):
            t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        t.text = line if line else " "
        # end para marker
        etree.SubElement(p, f"{A_NS}endParaRPr")


def set_simple_text(shape, text, font_name=None, size=None, bold=None, color=None):
    for t_el in shape.element.findall(f".//{A_NS}t"):
        t_el.text = ""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if not p0.runs:
        run = p0.add_run()
    else:
        run = p0.runs[0]
        # clear other runs
        for r in p0.runs[1:]:
            r.text = ""
    run.text = text
    if font_name:
        run.font.name = font_name
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def remove_pictures(slide, pictures):
    for sh in pictures:
        sh.element.getparent().remove(sh.element)


def widen_body(body_shape):
    """Move body to a full-width content area under eyebrow."""
    body_shape.left = Emu(2061252)
    body_shape.top = Emu(3600000)
    body_shape.width = Emu(20000000)
    body_shape.height = Emu(9000000)


def find_template_index(prs):
    """Pick an original content slide with title+eyebrow+body."""
    for i, slide in enumerate(prs.slides):
        if is_added_slide(slide):
            continue
        roles = find_roles(slide)
        if roles["title"] and roles["body"] and roles["eyebrow"]:
            return i
    for i, slide in enumerate(prs.slides):
        if is_added_slide(slide):
            continue
        roles = find_roles(slide)
        if roles["title"] and roles["body"]:
            return i
    # fallback: first non-title non-thanks
    for i, slide in enumerate(prs.slides):
        texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if "БЛАГОДАР" in texts.upper():
            continue
        if i == 0:
            continue
        if not is_added_slide(slide):
            return i
    return 1 if len(prs.slides) > 1 else 0


def fill_cloned_slide(slide, title, eyebrow, body_lines):
    roles = find_roles(slide)
    remove_pictures(slide, roles["pictures"])
    roles = find_roles(slide)  # refresh

    if roles["title"]:
        set_simple_text(roles["title"], title, TITLE_FONT, TITLE_SIZE, False, TITLE_COLOR)
    if roles["eyebrow"]:
        set_simple_text(roles["eyebrow"], eyebrow, EYEBROW_FONT, EYEBROW_SIZE, False, EYEBROW_COLOR)
    if roles["body"]:
        widen_body(roles["body"])
        clear_text_keep_style(roles["body"], body_lines)
        # force Open Sans / size on runs
        for p in roles["body"].text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Open Sans"
                r.font.size = BODY_SIZE
                try:
                    r.font.color.rgb = BODY_COLOR
                except Exception:
                    pass
    else:
        # create body textbox if missing - rare
        from pptx.enum.shapes import MSO_SHAPE
        box = slide.shapes.add_textbox(Emu(2061252), Emu(3600000), Emu(20000000), Emu(9000000))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.name = "Open Sans"
            run.font.size = BODY_SIZE
            run.font.color.rgb = BODY_COLOR


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    new_items = [items[i] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in new_items:
        sldIdLst.append(el)


def renumber_pages(prs):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if (sh.left or 0) < Emu(1200000) and (sh.width or 0) < Emu(1500000) and (sh.top or 0) > Emu(5000000):
                for t_el in sh.element.findall(f".//{A_NS}t"):
                    t_el.text = ""
                p0 = sh.text_frame.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = str(i)
                else:
                    p0.add_run().text = str(i)


def process_deck(path: Path, specs: list):
    prs = Presentation(str(path))
    # 1) delete added slides
    added = [i for i, s in enumerate(prs.slides) if is_added_slide(s)]
    print(f"  remove added indices: {added}")
    delete_slides(prs, added)
    # save/reload after delete
    tmp = Path("/tmp") / ("restyle_" + path.name)
    prs.save(str(tmp))
    prs = Presentation(str(tmp))

    template_idx = find_template_index(prs)
    print(f"  template slide index: {template_idx}")

    # Detect thanks index
    thanks_idx = None
    for i, s in enumerate(prs.slides):
        texts = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        if "БЛАГОДАР" in texts.upper():
            thanks_idx = i

    created_indices = []
    for spec in specs:
        dest = duplicate_slide(prs, template_idx)
        fill_cloned_slide(dest, spec["title"], spec["eyebrow"], spec["body"])
        created_indices.append(len(prs.slides) - 1)

    # Reorder: after_title ones right after 0; before_thanks before thanks
    n = len(prs.slides)
    # Build order of original slides 0..n-1-len(created) then place created
    after = [created_indices[i] for i, sp in enumerate(specs) if sp.get("after_title")]
    before = [created_indices[i] for i, sp in enumerate(specs) if sp.get("before_thanks")]
    other = [created_indices[i] for i, sp in enumerate(specs) if not sp.get("after_title") and not sp.get("before_thanks")]

    orig = [i for i in range(n) if i not in created_indices]
    order = []
    # title
    if orig:
        order.append(orig[0])
        rest_orig = orig[1:]
    else:
        rest_orig = []
    order.extend(after)
    # if thanks at end of rest_orig
    if thanks_idx is not None and thanks_idx in rest_orig:
        # rest without thanks
        mid = [i for i in rest_orig if i != thanks_idx]
        order.extend(mid)
        order.extend(other)
        order.extend(before)
        order.append(thanks_idx)
    else:
        order.extend(rest_orig)
        order.extend(other)
        order.extend(before)

    # sanity
    if sorted(order) != list(range(n)):
        # fallback keep append order
        print(f"  WARN reorder mismatch {order}, keeping natural order")
    else:
        reorder_slides(prs, order)

    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    renumber_pages(prs)
    prs.save(str(path))
    print(f"  saved {path.name} ({len(prs.slides)} slides)")


def main():
    for name, specs in SUPPLEMENTS.items():
        path = OUT / name
        print(f"\n=== {name} ===")
        process_deck(path, specs)

    # rebuild zip
    if ZIP.exists():
        ZIP.unlink()
    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in sorted(OUT.glob("*.pptx")):
            zf.write(p, p.name)
        for extra in OUT.glob("*.xlsx"):
            zf.write(extra, extra.name)
        readme = OUT / "README.md"
        if readme.exists():
            zf.write(readme, readme.name)
    print(f"\nZIP updated: {ZIP}")

    # verify no Arial left in content decks
    print("\n=== VERIFY ===")
    for p in sorted(OUT.glob("*.pptx")):
        prs = Presentation(str(p))
        for i, s in enumerate(prs.slides):
            fonts = set()
            for sh in s.shapes:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.name:
                                fonts.add(r.font.name)
            flag = "BAD" if "Arial" in fonts else "ok"
            title = ""
            for sh in s.shapes:
                if sh.has_text_frame:
                    t = " ".join(sh.text_frame.text.split())
                    if len(t) > 20:
                        title = t[:70]
                        break
            print(f"{p.name[:20]:20} {i:02d} {flag} {fonts} | {title}")


if __name__ == "__main__":
    main()
