#!/usr/bin/env python3
"""Add missing facts from 3.docx into presentation.pptx preserving style."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________5348.pptx"
OUT = "/workspace/Презентация_ПС_обновленная.pptx"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_shape_text(shape, text, size_pt=11, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return dest


def insert_slide_copy(prs, template_index, insert_after_index):
    new_slide = duplicate_slide(prs, template_index)
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(insert_after_index + 1, new_id)
    return prs.slides[insert_after_index + 1]


def fill_list_slide(slide, section, title, intro, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    # layout: section, title, intro, heading, (num, text)*, note, page
    set_shape_text(shapes[0], section, 11, True, RED)
    set_shape_text(shapes[1], title, 30, True, DARK)
    set_shape_text(shapes[2], intro, 11, False, GRAY)

    pairs = []
    for sh in shapes[4:]:
        t = sh.text_frame.text.strip()
        if t in {"01", "02", "03", "04", "05", "06", "07", "08"}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("При полном") or t.startswith("Частичное") or t.startswith("Эксплуатация"):
            note_shape = sh
        elif " / " in t and t.strip().endswith("12"):
            page_shape = sh

    for i, item in enumerate(items):
        if i < len(pairs):
            set_shape_text(pairs[i][0], f"{i+1:02d}", 9, True, RED)
            set_shape_text(pairs[i][1], item, 11, False, GRAY)

    if note:
        for sh in shapes:
            t = sh.text_frame.text
            if t.startswith("При полном") or t.startswith("Частичное") or t.startswith("Эксплуатация"):
                set_shape_text(sh, note, 11, True, AMBER)
                break


def fill_simple_slide(slide, section, title, subtitle, body, footer=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_shape_text(shapes[0], section, 11, True, RED)
    set_shape_text(shapes[1], title, 30, True, DARK)
    set_shape_text(shapes[2], subtitle, 11, False, GRAY)
    set_shape_text(shapes[3], body, 11, False, GRAY)
    if footer and len(shapes) > 4:
        set_shape_text(shapes[4], footer, 11, False, GRAY)


def find_by_prefix(slide, prefix):
    for sh in text_shapes(slide):
        if sh.text_frame.text.startswith(prefix):
            return sh
    return None


def update_page_numbers(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_shape_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def main():
    prs = Presentation(SRC)

    # Slide 3
    s3 = prs.slides[2]
    set_shape_text(find_by_prefix(s3, "Виды и периодичность"), "Требования к организации работ с применением ПС", 30, True, DARK)
    set_shape_text(find_by_prefix(s3, "Пункты ФНП"), "ППР разрабатывается эксплуатирующей или специализированной организацией по пп. 155–163 ФНП. ППР и ТК утверждаются организацией, эксплуатирующей ПС. Изменения вносятся только разработчиком.", 11, False, GRAY)

    # Slide 4
    s4 = prs.slides[3]
    hidden = [s for s in s4.shapes if s.has_text_frame and s.text_frame.text.strip() == "… ещё 1"]
    if hidden:
        set_shape_text(hidden[0], "Погрузочно-разгрузочные работы на базах и складах выполняются по ТК (п. 155–163 ФНП), если иное не предусмотрено п. 98 ФНП.", 9, False, GRAY)
    if len(hidden) > 1:
        set_shape_text(hidden[1], "ППР и ТК утверждает организация, эксплуатирующая ПС. Эксплуатация с отступлениями не допускается.", 9, False, GRAY)

    # Slide 5 -> installation summary
    s5 = prs.slides[4]
    set_shape_text(find_by_prefix(s5, "Ключевые выводы"), "Установка ПС и производство работ", 30, True, DARK)
    set_shape_text(find_by_prefix(s5, "Что важно помнить"), "Ключевые требования при эксплуатации ПС", 11, False, GRAY)
    install_items = [
        "Работы в стеснённых условиях выполняются только по утверждённому ППР и (или) ТК.",
        "ПС устанавливаются по руководству по эксплуатации; груз поднимают минимум на 0,5 м над препятствиями.",
        "Работа вблизи ВЛ — только по наряду-допуску и под руководством ответственного ИТР.",
        "Подъём начинают с 0,2–0,3 м для проверки строповки и тормоза; под грузом людей быть не должно.",
        "Кантовка — только на кантовальных площадках или по ППР; при массе > 75 % грузоподъёмности — под руководством ИТР.",
        "Работы прекращают при ветре выше допустимого, неблагоприятной погоде и неисправных ограничителях.",
    ]
    bodies = []
    for sh in text_shapes(s5):
        t = sh.text_frame.text
        if t.startswith(("Техническое освидетельствование —", "Различайте виды", "В объём", "Результаты освидетельствования", "Специалист, ответственный")):
            bodies.append(sh)
    for sh, item in zip(bodies, install_items):
        set_shape_text(sh, item, 11, False, GRAY)
    note = find_by_prefix(s5, "Эксплуатация подъёмного сооружения без актуального")
    if note:
        set_shape_text(note, "Эксплуатация ПС с отступлениями от ППР, ТК и требований ФНП запрещена.", 11, True, AMBER)

    new_specs = [
        ("list", 6, {
            "section": "Установка ПС и производство работ",
            "title": "Перемещение грузов и основные запреты",
            "intro": "При перемещении груза ПС и в процессе работ необходимо соблюдать требования промышленной безопасности.",
            "items": [
                "Горизонтальное перемещение — на 0,5 м выше встречающихся предметов; опускать груз только на предназначенное место.",
                "Запрещён подъём груза неизвестной массы, подтаскивание по земле, оттягивание груза при подъёме и перемещении.",
                "Запрещены работа при неисправных ограничителях и тормозах, нахождение людей под стрелой, перемещение людей грузовыми подъемниками.",
                "Запрещено выравнивать груз руками и менять стропы на подвешенном грузе; разворот вручную — только до 1 м.",
            ],
            "note": "По окончании работ на грузозахватном органе не должно оставаться подвешенного груза.",
        }),
        ("list", 6, {
            "section": "Установка ПС и производство работ",
            "title": "Установка ПС: расстояния и опасные зоны",
            "intro": "Установка кранов и других ПС должна исключать задевание конструкций и обеспечивать безопасные зазоры.",
            "items": [
                "До потолка и стропил — не менее 0,1 м; до пола в зоне возможного нахождения людей — не менее 2 м.",
                "Стреловые краны и манипуляторы — на подготовленной площадке; зазор до препятствий не менее 1 м.",
                "Установка на выносные опоры — по руководству по эксплуатации; при отсутствии требований — на все опоры.",
                "Работа ближе 30 м от провода ВЛ > 50 В — только по наряду-допуску; в охранной зоне — с разрешением владельца линии.",
            ],
            "note": "Краны без координатной защиты в стеснённых условиях применять запрещается.",
        }),
        ("simple", 5, {
            "section": "Система сигнализации при выполнении работ",
            "title": "Обмен сигналами между крановщиками и стропальщиками",
            "subtitle": "Знаковая сигнализация и порядок связи",
            "body": (
                "Эксплуатирующая организация устанавливает порядок обмена сигналами между стропальщиками и крановщиками. "
                "Знаковая сигнализация и система сигналов при радиосвязи вносятся в производственные инструкции. "
                "При смене участка работы крановщики и стропальщики инструктируются под подпись о применяемой сигнализации.\n\n"
                "Для подъемника связь с машинистом: до 10 м — голосом; свыше 10 м — знаковая сигнализация; свыше 22 м — радио- или телефонная связь."
            ),
            "footer": "При работе подъемника связь между люлькой и машинистом должна поддерживаться непрерывно.",
        }),
        ("list", 6, {
            "section": "Пуск ПС в работу и постановка на учёт",
            "title": "Основания и документы для пуска ПС в работу",
            "intro": "Решение о пуске выдаёт ответственный за производственный контроль; для ряда ПС создаётся комиссия.",
            "items": [
                "Перед пуском: после постановки на учёт, монтажа на новом месте, реконструкции, замены расчётных элементов со сваркой.",
                "Комплект документов: паспорт ПС, сертификаты, руководство по эксплуатации, акт монтажа, ППР/ТК, акт рельсового пути.",
                "Для башенных, мостовых и портальных кранов — акт готовности и комиссия; уведомление за 10 рабочих дней.",
                "ОПО с ПС регистрируются в реестре; часть ПС (до 10 т, отдельные типы) учёту не подлежит.",
            ],
            "note": "Решение руководителя о пуске оформляется внутренним распорядительным документом.",
        }),
        ("simple", 5, {
            "section": "Виды технического освидетельствования",
            "title": "Периодическое техническое освидетельствование ПС",
            "subtitle": "Периодичность и объём проверок",
            "body": (
                "ПС подвергаются техническому освидетельствованию до пуска и в процессе эксплуатации. "
                "Периодическое ТО включает: частичное — не реже 1 раза в 12 месяцев; полное — не реже 1 раза в 3 года "
                "(для ПС, используемых только при ремонте оборудования, — 1 раз в 5 лет).\n\n"
                "При полном ТО проводятся осмотр, статические и динамические испытания, испытания на устойчивость (если предусмотрено паспортом). "
                "При частичном ТО статические и динамические испытания не выполняются."
            ),
            "footer": "Результаты фиксируются в паспорте ПС с указанием срока следующего освидетельствования.",
        }),
        ("list", 6, {
            "section": "Действия в аварийных ситуациях",
            "title": "Инструкции для работников ОПО, эксплуатирующих ПС",
            "intro": "В организации должны быть разработаны и доведены под подпись инструкции о действиях в аварийных ситуациях.",
            "items": [
                "Оперативные действия по предотвращению и локализации аварий; способы и методы ликвидации.",
                "Схемы эвакуации при взрыве, пожаре, выбросе токсичных веществ, если ситуацию нельзя локализовать.",
                "Порядок приведения ПС в безопасное положение и эвакуации крановщика (оператора) из кабины.",
                "Места отключения электропитания ПС, расположение аптечек, порядок оповещения работников об авариях.",
            ],
            "note": "Инструкции должны содержать порядок оказания первой помощи и использования системы пожаротушения.",
        }),
    ]

    insert_pos = 3
    for kind, template_idx, spec in new_specs:
        slide = insert_slide_copy(prs, template_idx, insert_pos)
        insert_pos += 1
        if kind == "list":
            fill_list_slide(slide, spec["section"], spec["title"], spec["intro"], spec["items"], spec.get("note"))
        else:
            fill_simple_slide(slide, spec["section"], spec["title"], spec["subtitle"], spec["body"], spec.get("footer"))

    for slide in prs.slides:
        if find_by_prefix(slide, "1. Статическое испытание"):
            sh = find_by_prefix(slide, "Нагрузка 125")
            if sh:
                set_shape_text(
                    sh,
                    "125 % — для большинства ПС; 140 % — краны-трубоукладчики; 200 % — грузопассажирские и фасадные подъемники; 150 % — иные подъемники (вышки). Выдержка не менее 10 минут.",
                    11,
                    False,
                    GRAY,
                )
            extra = find_by_prefix(slide, "… ещё 2")
            if extra:
                set_shape_text(
                    extra,
                    "4. Испытания на устойчивость — при первичном ТО стреловых кранов, если в паспорте нет ссылок на ранее проведённые испытания.",
                    11,
                    False,
                    GRAY,
                )
        if find_by_prefix(slide, "Длительный простой"):
            set_shape_text(find_by_prefix(slide, "Длительный простой"), "Замена грузозахватного органа", 11, True, DARK)
            set_shape_text(find_by_prefix(slide, "Неиспользование ПС более 12 месяцев"), "Проводятся только статические испытания; динамические испытания не выполняются", 11, False, GRAY)
            set_shape_text(find_by_prefix(slide, "Авария или инцидент"), "Монтаж на новом месте", 11, True, DARK)
            set_shape_text(find_by_prefix(slide, "Подъёмное сооружение подверглось воздействию"), "Кроме подъемников, вышек, стреловых и быстромонтируемых башенных кранов", 11, False, GRAY)

    update_page_numbers(prs)
    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
