#!/usr/bin/env python3
"""Build individual Theme 1 presentations for new App.2 PP2464 structure (1.1–1.7)."""

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.util import Emu

SRC = Path(
    "/workspace/first_aid_tema1/"
    "1 Организационно-правовые аспекты оказания первой помощи.pptx"
)
OUT_DIR = Path("/workspace/first_aid_tema1_new")
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Combined source indices (0-based):
# 0 title, 1-6 org/NPA, 7 concept, 8-9 states/measures,
# 10-11 kits, 12-15 sequence, 16 safety, 17-20 extract,
# 21 infection, 22 call, 23 thanks

DECKS = [
    {
        "file": "1.1_НПА_и_организация_оказания_первой_помощи.pptx",
        "title": "1.1. НОРМАТИВНО-ПРАВОВАЯ БАЗА И ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 1.1  НПА и организация оказания первой помощи",
        "keep": [0, 1, 2, 3, 4, 5, 6, 23],
        "content_titles": {
            1: "ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ В РОССИЙСКОЙ ФЕДЕРАЦИИ",
            2: "ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ В РОССИЙСКОЙ ФЕДЕРАЦИИ",
            3: "НОРМАТИВНО-ПРАВОВАЯ БАЗА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            4: "НОРМАТИВНО-ПРАВОВАЯ БАЗА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            5: "НОРМАТИВНО-ПРАВОВАЯ БАЗА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            6: "НОРМАТИВНО-ПРАВОВАЯ БАЗА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        },
        "note": "Из сводной: организация + НПА. Старого отдельного файла 1.1 с понятием ПП здесь нет.",
    },
    {
        "file": "1.2_Укладки_наборы_комплекты_и_аптечки.pptx",
        "title": "1.2. УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 1.2  Укладки, наборы, комплекты и аптечки",
        "keep": [0, 10, 11, 23],
        "content_titles": {
            10: "УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            11: "УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        },
        "note": "Было: 1.2 Современные наборы… Заголовки обновлены под новую терминологию.",
    },
    {
        "file": "1.3_Порядок_и_приоритетность_оказания_первой_помощи.pptx",
        "title": "1.3. ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 1.3  Порядок и приоритетность оказания первой помощи",
        "keep": [0, 12, 13, 14, 15, 23],
        "content_titles": {
            12: "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            13: "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            14: "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            15: "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        },
        "note": "Было частью 1.3 (общая последовательность). Извлечение вынесено в 1.6.",
    },
    {
        "file": "1.4_Перечень_состояний_и_мероприятий.pptx",
        "title": "1.4. ПЕРЕЧЕНЬ СОСТОЯНИЙ И МЕРОПРИЯТИЙ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 1.4  Состояния и мероприятия первой помощи",
        "keep": [0, 7, 8, 9, 23],
        "content_titles": {
            7: "ПЕРЕЧЕНЬ СОСТОЯНИЙ И МЕРОПРИЯТИЙ ПЕРВОЙ ПОМОЩИ",
            8: "ПЕРЕЧЕНЬ СОСТОЯНИЙ, ПРИ КОТОРЫХ ОКАЗЫВАЕТСЯ ПЕРВАЯ ПОМОЩЬ",
            9: "ПЕРЕЧЕНЬ МЕРОПРИЯТИЙ И ПОСЛЕДОВАТЕЛЬНОСТЬ ИХ ВЫПОЛНЕНИЯ",
        },
        "note": "Было: 1.1 Понятие первая помощь. Понятие не отдельный элемент новой программы — слайд с определением оставлен как вводный.",
    },
    {
        "file": "1.5_Безопасные_условия_и_профилактика_инфекций.pptx",
        "title": "1.5. БЕЗОПАСНЫЕ УСЛОВИЯ И ПРОФИЛАКТИКА ИНФЕКЦИЙ",
        "eyebrow": "Тема 1 · п. 1.5  Безопасные условия и профилактика инфекций",
        "keep": [0, 16, 21, 23],
        "content_titles": {
            16: "ОБЕСПЕЧЕНИЕ БЕЗОПАСНЫХ УСЛОВИЙ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            21: "ПРОФИЛАКТИКА ИНФЕКЦИОННЫХ ЗАБОЛЕВАНИЙ ПРИ ОКАЗАНИИ ПЕРВОЙ ПОМОЩИ",
        },
        "note": "Было: 1.4 Соблюдение правил личной безопасности.",
    },
    {
        "file": "1.6_Извлечение_и_перемещение_пострадавших.pptx",
        "title": "1.6. ИЗВЛЕЧЕНИЕ И ПЕРЕМЕЩЕНИЕ ПОСТРАДАВШИХ",
        "eyebrow": "Тема 1 · п. 1.6  Извлечение и перемещение пострадавших",
        "keep": [0, 17, 18, 19, 20, 23],
        "content_titles": {
            17: "ИЗВЛЕЧЕНИЕ ПОСТРАДАВШИХ ИЗ ТРУДНОДОСТУПНЫХ МЕСТ И ПЕРЕМЕЩЕНИЕ",
            18: "ИЗВЛЕЧЕНИЕ ПОСТРАДАВШИХ ИЗ ТРУДНОДОСТУПНЫХ МЕСТ И ПЕРЕМЕЩЕНИЕ",
            19: "ИЗВЛЕЧЕНИЕ ПОСТРАДАВШИХ ИЗ ТРУДНОДОСТУПНЫХ МЕСТ И ПЕРЕМЕЩЕНИЕ",
            20: "ИЗВЛЕЧЕНИЕ ПОСТРАДАВШИХ ИЗ ТРУДНОДОСТУПНЫХ МЕСТ И ПЕРЕМЕЩЕНИЕ",
        },
        "note": "НОВЫЙ теоретический пункт. Материал был внутри старого 1.3 (последовательность).",
    },
    {
        "file": "1.7_Правила_вызова_скорой_медицинской_помощи.pptx",
        "title": "1.7. ПРАВИЛА ВЫЗОВА СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 1.7  Вызов скорой медицинской помощи и спецслужб",
        "keep": [0, 22, 23],
        "content_titles": {
            22: "ПРАВИЛА ВЫЗОВА СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ И ДРУГИХ СПЕЦИАЛЬНЫХ СЛУЖБ",
        },
        "note": "Было: 1.5 Основные правила вызова…",
    },
]


def replace_shape_text(shape, text):
    for t_el in shape.element.findall(f".//{NS}t"):
        t_el.text = ""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
    else:
        p0.add_run().text = text


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def set_title_like(slide, text):
    """Set main title: widest text box in upper area, or any large title-like box."""
    cands = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if (sh.top or 0) < Emu(3500000) and (sh.width or 0) > Emu(8000000):
            cands.append((-(sh.width or 0), sh))
    if not cands:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                if (sh.width or 0) > Emu(5000000):
                    cands.append((-(sh.width or 0), sh))
    if cands:
        cands.sort()
        replace_shape_text(cands[0][1], text)
        return True
    return False


def set_eyebrow(slide, text):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        top = sh.top or 0
        h = sh.height or 0
        w = sh.width or 0
        if "Аспекты оказания" in t or "Тема 1 ·" in t:
            replace_shape_text(sh, text)
            return True
        if Emu(2500000) < top < Emu(3600000) and h < Emu(900000) and w < Emu(9000000) and t:
            replace_shape_text(sh, text)
            return True
    return False


def update_pages(prs):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if (
                (sh.left or 0) < Emu(900000)
                and (sh.width or 0) < Emu(1200000)
                and (sh.top or 0) > Emu(5000000)
            ):
                replace_shape_text(sh, str(i))


def build_deck(spec):
    tmp = Path(f"/tmp/fa_deck_{spec['file']}")
    shutil.copy(SRC, tmp)
    prs = Presentation(str(tmp))

    # Update texts while original indices still valid
    set_title_like(prs.slides[0], spec["title"])
    for idx, title in spec["content_titles"].items():
        set_title_like(prs.slides[idx], title)
        set_eyebrow(prs.slides[idx], spec["eyebrow"])

    # Delete slides not kept
    drop = [i for i in range(len(prs.slides)) if i not in spec["keep"]]
    delete_slides(prs, drop)

    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    update_pages(prs)
    out = OUT_DIR / spec["file"]
    prs.save(str(out))
    return out, len(prs.slides)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    lines = [
        "# Тема 1 — отдельные презентации под новую структуру (Прил. 2 № 2464)",
        "",
        "Исходники: папка на Яндекс.Диске (старые 1.1–1.5 + сводная).",
        "Собрано из сводной презентации с переименованием под пункты 1.1–1.7.",
        "",
        "| Новый файл | Было | Примечание |",
        "|---|---|---|",
    ]
    print(f"Building into {OUT_DIR}")
    for spec in DECKS:
        out, n = build_deck(spec)
        print(f"  OK {out.name} ({n} slides)")
        # verify titles
        prs = Presentation(str(out))
        for i, s in enumerate(prs.slides, 1):
            title = ""
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    if (sh.width or 0) > Emu(5000000):
                        title = sh.text_frame.text.strip().replace("\n", " ")[:70]
                        break
            print(f"     {i:02d}| {title}")
        old = {
            "1.1_НПА_и_организация_оказания_первой_помощи.pptx": "сводная (организация/НПА); старый 1.1 был про понятие ПП → ушёл в 1.4",
            "1.2_Укладки_наборы_комплекты_и_аптечки.pptx": "1.2_Современные наборы…",
            "1.3_Порядок_и_приоритетность_оказания_первой_помощи.pptx": "часть 1.3_Общая последовательность…",
            "1.4_Перечень_состояний_и_мероприятий.pptx": "1.1_Понятие первая помощь",
            "1.5_Безопасные_условия_и_профилактика_инфекций.pptx": "1.4_Соблюдение правил личной безопасности",
            "1.6_Извлечение_и_перемещение_пострадавших.pptx": "часть 1.3 (извлечение) — новый пункт",
            "1.7_Правила_вызова_скорой_медицинской_помощи.pptx": "1.5_Основные правила вызова…",
        }[spec["file"]]
        lines.append(f"| `{spec['file']}` | {old} | {spec['note']} |")

    lines.append("")
    lines.append("## Что проверить методисту")
    lines.append("- 1.3: при необходимости добавить явный слайд про приоритетность.")
    lines.append("- 1.2: визуально всё ещё авто/работникам — нужны ли другие укладки/комплекты.")
    lines.append("- 1.4: определение «первая помощь» оставлено как вводный слайд.")
    readme = OUT_DIR / "README.md"
    readme.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"Wrote {readme}")


if __name__ == "__main__":
    main()
