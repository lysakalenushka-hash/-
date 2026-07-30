#!/usr/bin/env python3
"""Tema 1 lifts/towers: prune quizzes, add missing from 1.docx."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________3d83.pptx"
# Single-column list template (source pptx only has dual-column checks slide)
LIST_TMPL_SRC = "/workspace/Презентация_Тема_5_Расследование_аварий.pptx"
TMP = "/tmp/tema1_lifts_step1.pptx"
OUT = "/workspace/Презентация_Тема_1_Подъемники.pptx"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep title, definitions, construction classification, nodes, specs, checks; drop quizzes 3-5
KEEP = {0, 1, 2, 6, 7, 8}


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10271455), Emu(6492240), Emu(1200000), Emu(228600)
            )
            set_text(box, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    # Prefer left-column content on dual-column layouts
    leftish = [s for s in content if (s.left or 0) < Emu(5500000)]
    hdr = leftish if len(leftish) >= 4 else content
    set_text(hdr[0], section, 11, True, RED)
    set_text(hdr[1], title, 24, True, DARK)
    set_text(hdr[2], intro, 12, False, GRAY)
    set_text(hdr[3], list_title, 11, True, RED)

    pairs, more_shape, note_shape = [], None, None
    for sh in content[4:]:
        if (sh.left or 0) >= Emu(5500000):
            # Clear right column leftovers from dual-column template
            t = sh.text_frame.text.strip()
            if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)} or len(t) > 0:
                if sh.top < Emu(6200000):
                    set_text(sh, "", 12, False, DARK)
            continue
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    shown = items[:4]
    for i, item in enumerate(shown):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i + 1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(shown), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    if note_shape is None:
        for sh in content[4:]:
            if (
                sh.top > Emu(5600000)
                and sh.left < Emu(5500000)
                and sh not in [p[0] for p in pairs] + [p[1] for p in pairs if p[1]]
            ):
                note_shape = sh
                break
    if note_shape:
        if note:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            extra = items[4:]
            set_text(
                note_shape,
                ("· " + " · ".join(extra)) if extra else "",
                11,
                True,
                AMBER if extra else GRAY,
            )
    elif note:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, note, 12, True, AMBER)
    elif len(items) > 4:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, "· " + " · ".join(items[4:]), 11, True, AMBER)


def find_templates(prs):
    text_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if text_tpl is None and 5 <= n <= 8:
            titles = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(len(t) > 40 for t in titles):
                text_tpl = slide
                break
    if text_tpl is None:
        text_tpl = prs.slides[min(1, len(prs.slides) - 1)]

    list_tpl = None
    # Prefer external single-column list template
    try:
        donor = Presentation(LIST_TMPL_SRC)
        for slide in donor.slides:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            ones = sum(1 for t in texts if t == "01")
            if ones == 1 and any(t == "02" for t in texts) and len(texts) >= 10:
                list_tpl = slide
                break
    except Exception:
        list_tpl = None
    if list_tpl is None:
        for slide in prs.slides:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(t in {"01", "02"} for t in texts):
                list_tpl = slide
                break
    return text_tpl, list_tpl or text_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def enrich_kept(prs):
    # After prune: 0 title, 1 defs, 2 class construction, 3 nodes, 4 specs, 5 checks
    shapes = sorted(
        [s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: (s.top, s.left),
    )
    if len(shapes) >= 5:
        set_text(shapes[0], "Общее устройство", 36, True, DARK)
        set_text(shapes[1], "и классификация", 36, True, DARK)
        set_text(shapes[2], "подъёмников (вышек)", 36, True, RED)
        set_text(
            shapes[3],
            "Тема 1 · Виды, конструкция, характеристики, документация и допуск к работе",
            13,
            False,
            GRAY,
        )
        set_text(
            shapes[4],
            "Подъёмники · Вышки · Классификация · Паспорт · Безопасность",
            12,
            False,
            GRAY,
        )

    fill_text_slide(
        prs.slides[1],
        "1 · Определения",
        "Подъёмник и вышка",
        "Базовые понятия",
        "Подъёмник — грузоподъёмная машина прерывного действия для перемещения людей "
        "с инструментом и материалами и проведения работ в пределах зоны обслуживания. "
        "Вышка — грузоподъёмная машина прерывного действия для перемещения людей "
        "с инструментом и материалами и проведения работ в вертикальном направлении (вверх, вниз). "
        "Подъёмники и вышки обеспечивают механизацию строительно-ремонтных, электромонтажных, "
        "вентиляционных и других работ.",
        note="Машинист должен знать конструкцию, безопасно управлять и соблюдать меры безопасности.",
    )

    fill_text_slide(
        prs.slides[3],
        "1 · Конструкция",
        "Основные узлы и элементы подъёмника (вышки)",
        "Состав конструкции",
        "Основные узлы: ходовая рама (шасси); выносные опоры; поворотная часть; "
        "рабочее оборудование (колена, стрела, мачта); люлька (рабочая платформа); "
        "механизмы подъёма, поворота и выдвижения; гидросистема и/или электропривод; "
        "пульты управления (нижний и в люльке); приборы и устройства безопасности. "
        "Конкретный состав зависит от типа: ножничный, телескопический, шарнирный, мачтовый.",
        note="Знание конструкции нужно для эксплуатации, ТО и оценки технического состояния.",
    )

    # Section labels on classification / specs / checks
    for idx, label in ((2, "1 · По конструкции и приводу"), (4, "1 · Характеристики"), (5, "1 · Допуск к работе")):
        try:
            shapes = sorted(text_shapes(prs.slides[idx]), key=lambda s: (s.top, s.left))
            content = [
                s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))
            ]
            set_text(content[0], label, 11, True, RED)
        except Exception:
            pass


def append_missing(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    enrich_kept(prs)

    s = nl()
    fill_list_slide(
        s,
        "1 · Области применения",
        "Где применяются подъёмники и вышки",
        "Комплексная механизация работ на высоте",
        "Типовые задачи",
        [
            "ремонтно-строительные работы на фасадах и внутри зданий",
            "электромонтажные и вентиляционные работы",
            "складирование товаров на высоте",
            "замена освещения и обслуживание инженерных сетей",
        ],
        note="Эффективность зависит от умелого использования и высокой подготовки персонала.",
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Самоходные",
        "Подъёмники самоходные",
        "Оборудованы механизмом передвижения по площадке и/или по дорогам",
        "Типы шасси",
        [
            "автомобильный — на автомобильном шасси",
            "на спецшасси автомобильного типа",
            "пневмоколёсный — на пневмоколёсном шасси",
            "гусеничный — на гусеничном шасси",
            "железнодорожный — на дрезине, движение по ж/д пути",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Прицепные и передвижные",
        "Классификация по способу транспортирования",
        "Прицепные и передвижные подъёмники",
        "Виды",
        [
            "прицепной пневмоколёсный — буксируется механизированным ТС",
            "прицепной железнодорожный — на ж/д платформе",
            "передвижной самоходный — механизм на площадке + перевозка на ТС по дорогам",
            "передвижной несамоходный — вручную по площадке, перевозка на ТС",
        ],
        note="Способ передвижения определяет требования к площадке, транспорту и организации работ.",
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Документация",
        "Комплект документации подъёмника",
        "После приёмо-сдаточных испытаний и приёмки ОТК",
        "В комплект входят",
        [
            "паспорт подъёмника",
            "руководство по эксплуатации",
            "документация по стандарту / ТУ на изготовление",
            "копия сертификата соответствия и копия разрешения на применение (изготовление)",
        ],
        note="Металлоконструкции сторонних изготовителей снабжают документом, удостоверяющим качество.",
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Паспорт",
        "Основные разделы паспорта подъёмника",
        "Что фиксируется в паспорте",
        "Разделы",
        [
            "общие сведения: изготовитель, тип, заводской №, год, назначение, привод, среда, ветер",
            "технические данные: грузоподъёмность, высота, вылет, база, колея, масса, габариты",
            "данные сборочных единиц: электродвигатели, гидронасос/мотор, цилиндры, канаты, цепи",
            "указатели, ограничители и регистраторы безопасности; данные о металле; приёмка; учёт; ремонт; освидетельствование",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Устройства безопасности",
        "Указатели, ограничители и системы безопасности",
        "Из раздела паспорта о приборах безопасности",
        "Обязательные устройства",
        [
            "ограничитель предельного груза; ориентация пола люльки; ограничение зоны обслуживания",
            "блокировка подъёма/поворота при невыставленных опорах; система опускания при отказе привода",
            "устройство эвакуации из люльки; защита опор от самопроизвольного выдвижения",
            "указатель угла наклона; аварийный останов двигателя из люльки и с нижнего пульта; анемометр / переговорное устройство",
        ],
        note="Работа без исправных устройств безопасности запрещена.",
    )

    s = nl()
    fill_list_slide(
        s,
        "1 · Руководство по эксплуатации",
        "Что должно быть в руководстве по эксплуатации",
        "Обязательные сведения для машиниста и службы",
        "Включает",
        [
            "периодичность осмотра и смазки; способы осмотра металлоконструкций",
            "быстроизнашивающиеся детали и допуски на износ; нормы браковки элементов",
            "типовые повреждения металлоконструкций и способы устранения",
            "требования безопасности; порядок техосвидетельствования; срок службы; эвакуация из люльки при аварийном останове",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "Итоги темы 1",
        "Устройство и классификация подъёмников (вышек)",
        "Что должен знать машинист",
        "Главные выводы",
        [
            "различать подъёмник и вышку; знать назначение и зоны применения",
            "классификация: по конструкции/приводу и по способу передвижения (самоходные, прицепные, передвижные)",
            "знать узлы конструкции, характеристики и комплект документации (паспорт, РЭ, сертификат)",
            "перед работой — полный осмотр и функциональный контроль; неисправности — запрет пуска",
        ],
        note="Безопасность определяется соблюдением паспорта, СИЗ и требований эксплуатационной документации.",
    )


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order length {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for item in new_items:
        sldIdLst.append(item)


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    prs.save(TMP)

    prs = Presentation(TMP)
    append_missing(prs)
    # After prune+append indices:
    # 0 title, 1 defs, 2 class constr, 3 nodes, 4 specs, 5 checks,
    # 6 applications, 7 self-propelled, 8 trailer/mobile, 9 docs, 10 passport,
    # 11 safety, 12 manual, 13 summary
    order = [
        0,  # title
        1,  # definitions
        6,  # applications
        2,  # construction/drive class
        7,  # self-propelled
        8,  # trailer/mobile
        3,  # nodes
        4,  # specs
        9,  # documentation
        10,  # passport
        11,  # safety devices
        12,  # manual
        5,  # pre-work checks
        13,  # summary
    ]
    reorder_slides(prs, order)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
