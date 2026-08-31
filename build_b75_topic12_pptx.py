#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Добавление в презу темы 1.2 пунктов 4–6: таблицы, карточки, формат ОШИБКА/ПРАВИЛЬНО."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

SOURCE = Path("/home/ubuntu/.cursor/projects/workspace/uploads/_____1.2__6d53.pptx")
OUTPUT = Path("Тема_1.2_Контрольно-надзорная_деятельность.pptx")
INSERT_AFTER = 2  # после слайда 3 (индекс 2)

RED = RGBColor(0xE3, 0x06, 0x13)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
FONT = "Inter"

ML = Emu(548640)
CW = Emu(11094415)
SW = Emu(12192000)
SH = Emu(6858000)


def set_font(run, *, size=13, bold=False, color=GRAY):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_tb(slide, left, top, width, height, text="", *, size=13, bold=False, color=GRAY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return box


def slide_shell(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    el = bg._element
    slide.shapes._spTree.remove(el)
    slide.shapes._spTree.insert(2, el)


def footer_bar(slide, note: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Emu(5852160), CW, Emu(457200))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CREAM
    bar.line.fill.background()
    add_tb(slide, Emu(731520), Emu(5897880), Emu(10607040), Emu(365760), note, size=11, color=GRAY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Emu(6355080), CW, Emu(91440))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()


def header(slide, tag: str, title: str, subtitle: str):
    add_tb(slide, ML, Emu(731520), Emu(4572000), Emu(228600), tag, size=12, bold=True, color=RED)
    add_tb(slide, ML, Emu(1005840), CW, Emu(640080), title, size=24, bold=True, color=DARK)
    add_tb(slide, ML, Emu(1737360), CW, Emu(457200), subtitle, size=14, color=GRAY)


def set_page_num(slide, num: int, total: int):
    add_tb(
        slide,
        Emu(10271455),
        Emu(6492240),
        Emu(1188720),
        Emu(228600),
        f"{num:02d} / {total:02d}",
        size=11,
        color=RED,
        align=PP_ALIGN.RIGHT,
    )


def insert_slide(prs: Presentation, after_idx: int) -> object:
    layout = prs.slides[0].slide_layout
    prs.slides.add_slide(layout)
    sid = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(sid)
    prs.slides._sldIdLst.insert(after_idx + 1, sid)
    return prs.slides[after_idx + 1]


def add_table(slide, left, top, width, rows, col_widths, row_h=320000):
    nrows = len(rows)
    ncols = len(rows[0])
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, Emu(row_h * nrows))
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci else PP_ALIGN.CENTER
                for r in p.runs:
                    if ri == 0:
                        set_font(r, size=10, bold=True, color=WHITE)
                    elif ci == 0:
                        set_font(r, size=10, bold=True, color=RED)
                    else:
                        set_font(r, size=10, color=GRAY)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RED
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def build_divider_slide(slide):
    slide_shell(slide)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9753600), Emu(-457200), Emu(2743200), Emu(2743200))
    oval.fill.solid()
    oval.fill.fore_color.rgb = CREAM
    oval.line.fill.background()
    add_tb(slide, ML, Emu(731520), Emu(3200400), Emu(365760), "ТЕМА 1.2", size=14, bold=True, color=RED)
    add_tb(
        slide,
        ML,
        Emu(1188720),
        CW,
        Emu(914400),
        "Контрольно-надзорная и разрешительная деятельность Ростехнадзора",
        size=26,
        bold=True,
        color=DARK,
    )
    add_tb(
        slide,
        ML,
        Emu(2194560),
        CW,
        Emu(640080),
        "Оценка соответствия, материалы идентификации, разрешение на строительство",
        size=16,
        color=GRAY,
    )
    add_tb(
        slide,
        ML,
        Emu(2926080),
        CW,
        Emu(914400),
        "Государственный контроль (надзор) при проектировании, строительстве, эксплуатации, "
        "реконструкции, капремонте, монтаже, консервации и ликвидации сетей газораспределения "
        "и газопотребления осуществляет Ростехнадзор.",
        size=13,
        color=GRAY,
    )
    tagline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Emu(4114800), CW, Emu(457200))
    tagline.fill.solid()
    tagline.fill.fore_color.rgb = CREAM
    tagline.line.fill.background()
    add_tb(
        slide,
        Emu(731520),
        Emu(4206240),
        Emu(10607040),
        Emu(365760),
        "Надзор · Экспертиза · Идентификация · Разрешение на строительство",
        size=12,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )
    footer_bar(slide, "ПП РФ № 870 · Технический регламент (§ 12–13) · Область аттестации Б.7.5")


def build_point4_slide(slide):
    slide_shell(slide)
    header(
        slide,
        "ТЕМА 1.2 · ПУНКТ 4",
        "Материалы идентификации объектов",
        "Перечень документов для идентификации сетей газораспределения и газопотребления (технический регламент)",
    )

    rows = [
        ("№", "Материал", "Назначение"),
        ("1", "Проектная документация", "Идентификация объекта и принятых технических решений"),
        ("2", "Заключение экспертизы проектной документации", "Подтверждение соответствия при проектировании"),
        (
            "3",
            "Заключение экспертизы промышленной безопасности",
            "Для проектов на консервацию и ликвидацию сетей",
        ),
        ("4", "Разрешение на строительство", "Правовое основание начала строительно-монтажных работ"),
        ("5", "Исполнительная документация", "Подтверждение фактически выполненных работ"),
        ("6", "Акт приёмки сетей", "Подтверждение соответствия после строительства или реконструкции"),
        ("7", "Разрешение на ввод в эксплуатацию", "Основание для начала эксплуатации объекта"),
    ]
    add_table(
        slide,
        ML,
        Emu(2011680),
        CW,
        rows,
        [Emu(640080), Emu(4572000), Emu(5882235)],
    )
    footer_bar(
        slide,
        "К материалам идентификации относятся только документы из перечня технического регламента.",
    )


def build_point5_slide(slide):
    slide_shell(slide)
    header(
        slide,
        "ТЕМА 1.2 · ПУНКТ 5",
        "Запрет на иные материалы идентификации",
        "Использование иных материалов в качестве материалов для идентификации не допускается",
    )

    err_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, Emu(2011680), CW, Emu(1463040))
    err_box.fill.solid()
    err_box.fill.fore_color.rgb = CREAM
    err_box.line.color.rgb = BORDER
    add_tb(slide, Emu(731520), Emu(2103120), Emu(2000000), Emu(274320), "ОШИБКА", size=13, bold=True, color=RED)
    add_tb(
        slide,
        Emu(731520),
        Emu(2420000),
        Emu(10400000),
        Emu(960000),
        "Для идентификации сети газораспределения можно использовать акт технического осмотра, "
        "договор подряда, паспорт организации или справку об отсутствии замечаний.",
        size=12,
        color=GRAY,
    )

    ok_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, Emu(3657600), CW, Emu(1550000))
    ok_box.fill.solid()
    ok_box.fill.fore_color.rgb = GREEN_BG
    ok_box.line.color.rgb = BORDER
    add_tb(slide, Emu(731520), Emu(3749040), Emu(2000000), Emu(274320), "ПРАВИЛЬНО", size=13, bold=True, color=RED)
    add_tb(
        slide,
        Emu(731520),
        Emu(4057120),
        Emu(10400000),
        Emu(1050000),
        "Идентификация объекта технического регулирования выполняется только по документам, "
        "прямо перечисленным в техническом регламенте. Любые иные документы использовать нельзя.",
        size=12,
        color=GRAY,
    )
    footer_bar(slide, "На экзамене: перечислите 7 материалов и укажите запрет на любые «другие» документы.")


def build_point6_slide(slide):
    slide_shell(slide)
    header(
        slide,
        "ТЕМА 1.2 · ПУНКТ 6",
        "Связь экспертизы и разрешения на строительство",
        "Порядок перехода от проекта к началу работ на объекте",
    )

    cards = [
        (
            "01",
            "Экспертиза проекта",
            "При проектировании оценка соответствия осуществляется в форме экспертизы "
            "проектной документации и результатов инженерных изысканий.",
        ),
        (
            "02",
            "Заключение экспертизы",
            "Положительное заключение экспертизы включается в состав доказательственных "
            "материалов при получении разрешения на строительство.",
        ),
        (
            "03",
            "Разрешение на строительство",
            "Без заключения экспертизы и комплекта идентификационных материалов "
            "разрешение на строительство сети не выдаётся.",
        ),
        (
            "04",
            "Приёмка и ввод",
            "После завершения строительства или реконструкции — приёмка комиссией, "
            "акт приёмки и разрешение на ввод в эксплуатацию.",
        ),
    ]

    x_positions = [ML, Emu(6324447)]
    y_start = Emu(1920240)
    y_step = Emu(2100000)
    for i, (num, title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = x_positions[col]
        y = y_start + y_step * row
        w = Emu(5300000)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Emu(1950000))
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM if row == 0 else WHITE
        card.line.color.rgb = BORDER
        add_tb(slide, x + Emu(182880), y + Emu(137160), Emu(640080), Emu(365760), num, size=14, bold=True, color=RED)
        add_tb(slide, x + Emu(960000), y + Emu(137160), w - Emu(1140000), Emu(365760), title, size=13, bold=True, color=DARK)
        add_tb(slide, x + Emu(182880), y + Emu(640080), w - Emu(365760), Emu(1200000), body, size=11, color=GRAY)

    for y in (Emu(2700000), Emu(4800000)):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(5780000), y, Emu(400000), Emu(200000))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RED
        arrow.line.fill.background()

    footer_bar(slide, "Цепочка: проект → экспертиза → разрешение на строительство → СМР → приёмка → ввод.")


def renumber_all(prs: Presentation):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        found = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip()
            if " / " in txt:
                parts = txt.split()
                if len(parts) >= 3 and parts[-1].isdigit():
                    sh.text_frame.paragraphs[0].clear()
                    run = sh.text_frame.paragraphs[0].add_run()
                    run.text = f"{i:02d} / {total:02d}"
                    set_font(run, size=11, color=RED)
                    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    found = True
                    break
        if not found:
            set_page_num(slide, i, total)


def build():
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    builders = [
        build_divider_slide,
        build_point4_slide,
        build_point5_slide,
        build_point6_slide,
    ]

    for offset, builder in enumerate(builders):
        slide = insert_slide(prs, INSERT_AFTER + offset)
        builder(slide)

    renumber_all(prs)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
