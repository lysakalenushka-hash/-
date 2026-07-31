#!/usr/bin/env python3
"""Build pipeline presentation using uploaded template styling."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_Трубопроводы.pptx"
IMG = "/workspace/assets/pipeline_images/generated"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

IMG_LEFT = Emu(6600000)
IMG_W = Emu(5200000)
TEXT_W = Emu(5800000)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def hide_shape(shape):
    shape.left = Emu(0)
    shape.top = Emu(0)
    shape.width = Emu(0)
    shape.height = Emu(0)


def fill_text_slide(slide, section, title, intro, body, footer=None, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], body, 14, False, DARK)
    if footer and len(shapes) > 4:
        set_text(shapes[4], footer, 11, False, GRAY)
    elif len(shapes) > 4:
        set_text(shapes[4], "", 11, False, GRAY)
    if note:
        for sh in shapes:
            t = sh.text_frame.text
            if "Эксплуатация" in t or "Каждое утверждение" in t or sh.top > Emu(5800000):
                if sh.top > Emu(5800000) and sh.left < Emu(5000000):
                    set_text(sh, note, 14, True, AMBER)
                    break


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    for sh in shapes[5:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5800000) and sh.left > Emu(500000):
            pass
        elif sh.top > Emu(5800000) and sh.left < Emu(2000000) and not t:
            pass

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            set_text(pairs[i][1], item, 14, False, DARK)

    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        set_text(pairs[i][1], "", 14, False, DARK)

    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)

    if note:
        for sh in shapes:
            if sh.top > Emu(5900000) and sh.left > Emu(700000):
                set_text(sh, note, 14, True, AMBER)
                break


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3600000), width=Emu(10700000), height=Emu(2500000))
    else:
        slide.shapes.add_picture(path, IMG_LEFT, Emu(1200000), width=IMG_W, height=Emu(4800000))


def add_highlight(slide, parts):
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(548640), Emu(2400000), TEXT_W, Emu(650000))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    box.line.color.rgb = RED
    box.line.width = Pt(1.5)
    tf = slide.shapes.add_textbox(Emu(640000), Emu(2520000), TEXT_W - Emu(100000), Emu(500000)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        set_run(r, 14, bold, color)


def main():
    tpl = Presentation(TEMPLATE)
    text_tpl = [deepcopy(s.element) for s in tpl.slides[0].shapes]
    list_tpl = [deepcopy(s.element) for s in tpl.slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Промышленная безопасность",
        "Устройство и назначение трубопроводов пара и горячей воды",
        "ФНП при использовании оборудования, работающего под избыточным давлением",
        "Модуль охватывает нормативную базу, общие положения ФНП и устройство "
        "трубопроводов пара и горячей воды: состав, арматуру, КИПиА, опоры, "
        "дренаж, компенсаторы, заглушки, теплоизоляцию и опознавательную окраску.",
        note="Все требования основаны на ФНП при использовании оборудования, работающего под избыточным давлением.",
    )

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Нормативная база",
        "Федеральные нормы и правила",
        "Приказ Ростехнадзора от 15.12.2020 № 536",
        "Об утверждении федеральных норм и правил в области промышленной "
        "безопасности «Правила промышленной безопасности при использовании "
        "оборудования, работающего под избыточным давлением» (далее — ФНП).",
    )

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Общие положения",
        "Назначение ФНП",
        "Требования промышленной безопасности при эксплуатации оборудования под давлением",
        "ФНП устанавливают требования промышленной безопасности при разработке, "
        "размещении, монтаже, наладке, эксплуатации, ремонте, реконструкции, "
        "техническом освидетельствовании и диагностировании оборудования, "
        "работающего под избыточным давлением.",
    )
    add_image(s, f"{IMG}/industrial_steam_sidebar.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Общие положения",
        "Область применения ФНП",
        "Требования обязательны при следующих видах работ",
        "Обязательные требования ФНП",
        [
            "при разработке технологических процессов",
            "при техническом перевооружении опасного производственного объекта (ОПО)",
            "при размещении, монтаже и ремонте",
            "при реконструкции (модернизации) оборудования",
            "при наладке и эксплуатации",
            "при техническом освидетельствовании",
            "при техническом диагностировании и экспертизе промышленной безопасности",
        ],
    )
    add_image(s, f"{IMG}/pipeline_system_intro.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Общие положения",
        "К какому оборудованию применяются ФНП",
        "Оборудование, работающее под избыточным давлением",
        "Объекты применения ФНП",
        [
            "паровые и водогрейные котлы",
            "сосуды, работающие под давлением",
            "трубопроводы пара и горячей воды",
            "газовые баллоны и резервуары для сжатых газов",
            "арматуру и предохранительные устройства",
        ],
        note="ФНП обязательны для всех стадий жизненного цикла оборудования под давлением.",
    )
    add_image(s, f"{IMG}/industrial_steam_sidebar.png")

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Устройство трубопроводов",
        "Конструкция трубопроводов пара и горячей воды",
        "Назначение и общие требования к устройству",
        "Трубопроводы пара и горячей воды предназначены для транспортировки "
        "теплоносителя от источника тепла к потребителям. Конструкция должна "
        "обеспечивать прочность, герметичность, компенсацию температурных "
        "деформаций и безопасную эксплуатацию.",
    )
    add_image(s, f"{IMG}/pipeline_system_intro.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Устройство трубопроводов",
        "Состав трубопровода (1/2)",
        "Трубопровод состоит из основных элементов",
        "Элементы трубопровода",
        [
            "плотно соединённых между собой прямых участков труб",
            "фасонных деталей — отводы, переходники, тройники",
            "крепёжных элементов — фланцы, болты, шпильки",
            "арматуры — краны, вентили, задвижки, регулирующие клапаны",
            "редуцирующих и предохранительных клапанов",
        ],
    )
    add_image(s, f"{IMG}/pipeline_assembly_3d.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Устройство трубопроводов",
        "Приборы КИПиА",
        "Контрольно-измерительные приборы и автоматика",
        "Приборы на трубопроводе",
        [
            "манометры — контроль давления в трубопроводе",
            "термометры — контроль температуры теплоносителя",
            "расходомеры — учёт и контроль расхода среды",
            "диафрагмы — измерение расхода по перепаду давления",
        ],
        note="Приборы КИПиА обеспечивают контроль параметров и безопасную эксплуатацию.",
    )
    add_image(s, f"{IMG}/kipia_instruments.png", bottom=True)

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Устройство трубопроводов",
        "Опорно-подвесная система",
        "Крепление и поддержание трубопровода",
        "Типы опор и подвесок",
        [
            "неподвижные опоры — фиксируют положение трубопровода",
            "подвижные опоры — допускают перемещение при температурных деформациях",
            "скользящие, катковые и подвесные опоры",
            "пружинные и жёсткие подвески",
            "направляющие и ограничители перемещения",
        ],
    )
    add_image(s, f"{IMG}/pipe_supports_3d.png")

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Устройство трубопроводов",
        "Конденсатоотводчики, дренажи и патрубки",
        "Устройства для отвода конденсата и воздуха",
        "В нижних точках каждого отключаемого задвижками участка трубопровода "
        "должны предусматриваться спускные штуцера с запорной арматурой для "
        "опорожнения. В верхних точках устанавливаются воздушники. Нижние "
        "концевые точки паропроводов и изгибов снабжаются устройствами для "
        "продувки. Непрерывный отвод конденсата обязателен для паропроводов "
        "насыщенного пара и тупиковых участков паропроводов перегретого пара.",
    )
    add_image(s, f"{IMG}/condensate_drainage.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Устройство трубопроводов",
        "Компенсаторы температурных деформаций",
        "Для компенсации удлинения трубопровода при нагреве",
        "Типы компенсаторов",
        [
            "резиновые компенсаторы с фланцами",
            "П-образные компенсаторы (Z-образные участки)",
            "U-образные (омегаобразные) компенсаторы",
            "Г-образные компенсаторы",
            "сильфонные компенсаторы",
            "линзовые компенсаторы",
            "сальниковые компенсаторы",
        ],
    )
    add_highlight(s, [
        ("При нагреве трубопровода на ", False, DARK),
        ("100 °C", True, RED),
        (" удлинение стальной трубы составляет около ", False, DARK),
        ("1,2 мм", True, RED),
        (" на 1 м длины.", False, DARK),
    ])
    add_image(s, f"{IMG}/compensators_7types.png")

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Устройство трубопроводов",
        "Заглушки",
        "Элементы для глухого запечатывания выходных отверстий",
        "Заглушка — элемент трубопровода, предназначенный для глухого "
        "запечатывания его выходных отверстий; как правило используется при "
        "проведении пневмо- и гидроиспытаний. Типы: плоские сварные, фланцевые, "
        "эллиптические, сферические, быстросъёмные, межфланцевые.",
        note="Металлические концевые заглушки различаются по виду исполнения и типу крепления.",
    )
    add_image(s, f"{IMG}/pipe_caps_4types.png")

    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Устройство трубопроводов",
        "Теплоизоляция",
        "Защита от потерь тепла и ожогов",
        "Теплоизоляция трубопроводов предназначена для снижения теплопотерь, "
        "предотвращения конденсации влаги на поверхности, защиты персонала от "
        "термических ожогов и обеспечения стабильности технологических параметров.",
    )
    add_image(s, f"{IMG}/thermal_insulation_cutaway.png")

    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Устройство трубопроводов",
        "Опознавательная окраска",
        "Маркировка трубопроводов по ГОСТ 14202-69",
        "Цветовая маркировка",
        [
            "коммуникации делятся на 10 основных групп по транспортируемым веществам",
            "зелёный цвет — 1 группа, транспортирует воду",
            "красный цвет — 2 группа, транспортирует пар",
            "предупреждающие кольца: красные — легковоспламеняющиеся и взрывоопасные",
            "жёлтые — токсичные и радиоактивные вещества",
            "зелёные с белой каймой — внутренняя безопасность",
        ],
        note="Количество предупреждающих колец зависит от давления и температуры среды.",
    )
    add_image(s, f"{IMG}/pipe_identification_gost.png")

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
