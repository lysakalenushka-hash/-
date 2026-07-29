#!/usr/bin/env python3
"""Fill missing facts from 5.docx into presentation.pptx preserving style."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SRC = '/home/ubuntu/.cursor/projects/workspace/uploads/____________ec19.pptx'
OUT = '/workspace/Презентация_5_обновленная.pptx'

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def set_text(shape, text, size=14, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = 'Inter'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def fill_labels(slide, replacements):
    """Replace 'label — desc' placeholders in top order with provided texts."""
    labels = [s for s in text_shapes(slide) if s.text_frame.text.strip() == 'label — desc']
    labels.sort(key=lambda s: (s.top, s.left))
    for sh, text in zip(labels, replacements):
        set_text(sh, text, 14, False, DARK)


def find_exact(slide, txt):
    for sh in text_shapes(slide):
        if sh.text_frame.text.strip() == txt:
            return sh
    return None


def find_prefix(slide, prefix):
    for sh in text_shapes(slide):
        if sh.text_frame.text.startswith(prefix):
            return sh
    return None


prs = Presentation(SRC)

# --- Slide 4: appointment placeholders ---
s4 = prs.slides[3]
# Order by top: right, left, right, left, right
fill_labels(s4, [
    'Назначение: приказом руководителя Компании',
    'Категория должности: «Рабочие»',
    'Освобождение от должности: приказом руководителя Компании',
    'Разряд: машинист подъёмника 5-го разряда',
    'Порядок: в соответствии с действующим трудовым законодательством',
])

# --- Slide 15: missing knowledge item ---
s15 = prs.slides[14]
h = find_exact(s15, '… ещё 1')
if h:
    set_text(h, 'Ремонт подъёмных лебёдок: способы ремонта двигателя, трансмиссии и ходовой части.', 9, False, GRAY)

# --- Slide 32: overtime placeholders ---
s32 = prs.slides[31]
fill_labels(s32, [
    'Сверхурочная работа — по служебной необходимости',
    'Основание: служебная необходимость',
    'Привлечение допускается только в порядке, предусмотренном законодательством',
    'Порядок: в соответствии с действующим законодательством РФ',
    'Норма: должностная инструкция машиниста подъёмника 5-го разряда',
])
note = find_prefix(s32, '⚠ Точный порядок')
if note:
    set_text(note, 'В случае служебной необходимости машинист может привлекаться к сверхурочной работе только в порядке, предусмотренном законодательством.', 14, True, AMBER)

# --- Slide 33: replace generic duties with facts from DI ---
s33 = prs.slides[32]
body = None
for sh in text_shapes(s33):
    if sh.text_frame.text.strip().startswith('1. Проверять'):
        body = sh
        break
if body:
    set_text(
        body,
        '1. Обслуживать подъёмник (агрегат) при капитальном, текущем ремонте и опробовании скважин.\n'
        '2. Готовить подъёмник к работе, участвовать в монтаже/демонтаже и оснастке талевой системы.\n'
        '3. Управлять лебёдкой при всех спускоподъёмных операциях и силовым электрогенератором.\n'
        '4. Наблюдать за исправностью регистратора и механизмов, вести журнал учёта работы.\n'
        '5. Управлять автомобилем/трактором, выполнять текущий ремонт механизмов и обслуживать электростанции до 100 кВт.\n'
        '6. Соблюдать трудовую дисциплину, требования охраны труда и пожарной безопасности на рабочем месте.',
        14, False, DARK,
    )

# --- Slide 37: rights placeholders ---
s37 = prs.slides[36]
# left column then right alternating by top, but sorted by top,left:
# left, right, left, right, left, right
fill_labels(s37, [
    'Знакомиться с проектами решений руководства, касающимися его деятельности',
    'Сообщать руководителю о выявленных недостатках и предлагать их устранение',
    'Запрашивать информацию и документы, необходимые для выполнения обязанностей',
    'Привлекать специалистов подразделений к решению возложенных задач',
    'Вносить предложения по совершенствованию работы по должностной инструкции',
    'Требовать от руководства содействия в исполнении обязанностей и прав',
])

# --- Slide 39: missing responsibility ground ---
s39 = prs.slides[38]
h = find_exact(s39, '… ещё 1')
if h:
    set_text(
        h,
        'Неправомерное использование служебных полномочий (в т. ч. в личных целях); невыполнение указаний руководителя и трудовых функций.',
        9, False, GRAY,
    )

# --- Slide 29: ensure prep-final works mentioned if incomplete ---
s29 = prs.slides[28]
prep = find_prefix(s29, 'Подготовка подъёмника к работе')
# find description under it
for sh in text_shapes(s29):
    if sh.text_frame.text.strip().startswith('Подготовительно-заключительные'):
        set_text(
            sh,
            'Подготовительно-заключительные работы, монтаж и демонтаж, оснастка талевой системы, монтаж и обслуживание вспомогательных механизмов',
            11, False, GRAY,
        )
        break

prs.save(OUT)
print('saved', OUT, 'slides', len(prs.slides))

# verify placeholders gone
left = 0
for i,s in enumerate(prs.slides,1):
    for sh in text_shapes(s):
        t = sh.text_frame.text.strip()
        if t == 'label — desc' or t.startswith('… ещё'):
            print('REMAINING', i, t)
            left += 1
print('remaining placeholders:', left)
