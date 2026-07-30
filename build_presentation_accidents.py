#!/usr/bin/env python3
"""Build accidents & traumatism presentation using uploaded template styling."""

from copy import deepcopy

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_Аварии_травматизм.pptx"
IMG = "/workspace/assets/accidents_images"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

IMG_LEFT = Emu(6600000)
IMG_W = Emu(5200000)
TEXT_W = Emu(5800000)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, footer=None, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], body, 13, False, DARK)
    if footer and len(shapes) > 4:
        set_text(shapes[4], footer, 11, False, GRAY)
    elif len(shapes) > 4:
        set_text(shapes[4], "", 11, False, GRAY)
    if note:
        for sh in shapes:
            if sh.top > Emu(5800000) and sh.left < Emu(5000000):
                set_text(sh, note, 13, True, AMBER)
                break


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    for sh in shapes[5:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            set_text(pairs[i][1], item, 13, False, DARK)

    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        set_text(pairs[i][1], "", 14, False, DARK)

    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)

    if note:
        for sh in shapes:
            if sh.top > Emu(5900000) and sh.left > Emu(700000):
                set_text(sh, note, 13, True, AMBER)
                break


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3600000), width=Emu(10700000), height=Emu(2500000))
    else:
        slide.shapes.add_picture(path, IMG_LEFT, Emu(1200000), width=IMG_W, height=Emu(4800000))


def add_definition_block(slide, items, start_top=Emu(2100000)):
    top = start_top
    for term, definition in items:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(548640), top, TEXT_W, Emu(1050000))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        box.line.width = Pt(0.75)
        tf = slide.shapes.add_textbox(
            Emu(640000), top + Emu(70000), TEXT_W - Emu(100000), Emu(900000)
        ).text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{term} — "
        set_run(r1, 12, True, RED)
        r2 = p.add_run()
        r2.text = definition
        set_run(r2, 12, False, GRAY)
        top += Emu(1150000)


def add_cause_cards(slide, items):
    """Four cause cards in two columns."""
    positions = [
        (Emu(548640), Emu(2100000)),
        (Emu(6200000), Emu(2100000)),
        (Emu(548640), Emu(4200000)),
        (Emu(6200000), Emu(4200000)),
    ]
    card_w = Emu(5400000)
    card_h = Emu(1800000)
    for (left, top), (title, body) in zip(positions, items):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RED
        box.line.width = Pt(1.25)
        tf = slide.shapes.add_textbox(
            left + Emu(120000), top + Emu(120000), card_w - Emu(240000), card_h - Emu(240000)
        ).text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = title
        set_run(r1, 14, True, RED)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        set_run(r2, 12, False, GRAY)


def add_pie_chart(slide):
    chart_data = CategoryChartData()
    chart_data.categories = [
        "Нарушение трудовой и производственной дисциплины",
        "Дефекты вследствие эксплуатации",
        "Низкое качество обслуживания, освидетельствования и ЭПБ",
        "Низкое качество ремонта оборудования",
    ]
    chart_data.add_series("Доля", (39, 23, 23, 15))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Emu(548640),
        Emu(2000000),
        Emu(11000000),
        Emu(4200000),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = False

    # Color series points
    colors = [
        RGBColor(0xE3, 0x06, 0x13),
        RGBColor(0x6B, 0x72, 0x80),
        RGBColor(0x0D, 0x94, 0x88),
        RGBColor(0xF5, 0x9E, 0x0B),
    ]
    series = chart.series[0]
    for i, color in enumerate(colors):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color


def main():
    tpl = Presentation(TEMPLATE)
    text_tpl = [deepcopy(s.element) for s in tpl.slides[0].shapes]
    list_tpl = [deepcopy(s.element) for s in tpl.slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # 1. Title
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Промышленная безопасность",
        "Примеры и причины аварий и несчастных случаев при эксплуатации трубопроводов пара и горячей воды",
        "Обучение персонала, обслуживающего трубопроводы пара и горячей воды",
        "Модуль охватывает нагрузки на элементы трубопроводов, причины травматизма и аварий, "
        "характерные аварии тепловых сетей, последствия аварий, производственный и "
        "непроизводственный травматизм, а также классификацию производственных травм.",
        note="Разбор реальных случаев — основа профилактики аварий и травм.",
    )
    add_image(s, f"{IMG}/pipe_burst_geyser.png")

    # 2. Pipeline loads
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Трубопроводы ПиГВ",
        "Трубопроводы пара и горячей воды (ПиГВ)",
        "Их элементы находятся под воздействием:",
        "Воздействующие факторы",
        [
            "внутреннего давления рабочей среды",
            "веса труб, арматуры и тепловой изоляции",
            "напряжений самокомпенсации, возникающих в результате теплового расширения",
        ],
        note="Внутреннее давление вызывает напряжение растяжения, а нагрузки — напряжение изгиба.",
    )

    # 3. Causes of traumatism
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Травматизм",
        "Причины травматизма",
        "Основные группы причин травматизма на производстве:",
        "Группы причин",
        [
            "технические и технологические — неисправности, недостатки конструкций и процессов",
            "организационные — нарушения правил эксплуатации, режима труда и отдыха",
            "личностные (психофизиологические) — дисциплина, переутомление, человеческий фактор",
        ],
    )

    # 4. Pie chart causes
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Аварии и НС",
        "Основные причины возникновения аварий и несчастных случаев",
        "Статистическое распределение причин",
        "",
    )
    # Clear body area and put chart
    shapes = sorted(text_shapes(s), key=lambda sh: (sh.top, sh.left))
    if len(shapes) > 3:
        set_text(shapes[3], "", 13, False, DARK)
    add_pie_chart(s)

    # 5. Typical heating network accidents
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Аварии",
        "Характерные аварии тепловых сетей",
        "Типичные повреждения при эксплуатации тепловых сетей:",
        "Виды аварий",
        [
            "разрывы или повреждения стыков труб",
            "нарушения герметичности фланцевых соединений",
            "образования течей в местах установки регулирующей арматуры и сальниковых компенсаторов",
        ],
        note="Места аварий обнаруживаются по выходу пара или выбиванию горячей воды.",
    )
    add_image(s, f"{IMG}/pipe_burst_geyser.png")

    # 6. Causes of pipeline accidents
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Причины аварий",
        "Причины аварий трубопроводов ПиГВ",
        "Причинами аварий трубопроводов пара и горячей воды являются:",
        "Основные причины",
        [
            "дефекты их изготовления и монтажа",
            "гидравлические удары",
            "нарушение нормального режима эксплуатации",
            "отсутствие или неисправность предохранительных устройств",
            "физический износ стенок трубопроводов",
            "коррозия металла",
            "некачественный или несвоевременный ремонт",
        ],
    )

    # 7. Consequences
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Последствия аварий",
        "Последствия аварий тепловых сетей",
        "Аварии тепловых сетей приводят к следующим последствиям:",
        "Последствия",
        [
            "разрушение трубопровода и оборудования тепловых сетей",
            "разрыв отопительных приборов",
            "ожоговый травматизм, в том числе гражданских лиц",
            "затопление жилища или порча имущества физических лиц",
            "нарушение нормальной работы городской инфраструктуры жизнеобеспечения",
        ],
    )
    add_image(s, f"{IMG}/pipe_burst_geyser.png")

    # 8. Non-production: Caramel hotel
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Непроизводственный травматизм",
        "Гостиница «Карамель», Пермь (2020)",
        "Реальный случай гибели граждан при аварии теплосети",
        "Поздней ночью 20 января 2020 года номера пермской гостиницы «Карамель», "
        "работавшей в подвале дома, заполнил кипяток. Вода прибывала слишком быстро "
        "и лилась прямо из стены. Все девять номеров затопило теплоносителем, "
        "перегретым до 150 °C. Запасной выход был закрыт.\n\n"
        "В ту ночь погибли четверо взрослых и четырёхлетняя девочка. Прорвало "
        "магистральный трубопровод в 100 м от подвала — вода затекла по бетонному "
        "желобу через негерметизированную гильзу трубы ХВС.",
    )
    add_image(s, f"{IMG}/hotel_flood_scene.png")

    # 9. Non-production: Kupino
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Непроизводственный травматизм",
        "Котлован в г. Купино, Новосибирская область (2019)",
        "Гибель гражданского лица из-за незакрытого котлована",
        "В сентябре 2019 г. в г. Купино для устранения прорыва трубопровода "
        "работниками МУП «Теплосети» разрабатывался котлован. Прорыв устранили, "
        "но котлован временно не засыпали для последующего ремонта. В нарушение "
        "норм ограждения были из подручных средств, яма не закрыта.\n\n"
        "В первых числах ноября в котлован, заполненный водой, упала 48-летняя "
        "женщина. Она не смогла самостоятельно выбраться и утонула.",
        note="Незакрытый котлован — прямая угроза жизни граждан.",
    )
    add_image(s, f"{IMG}/flooded_pit_winter.png")

    # 10. Non-production: homeless in heating chambers
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Непроизводственный травматизм",
        "Жертвы на теплосетях — лица без определённого места жительства",
        "Опасность тепловых камер для граждан",
        "Часто жертвами на теплосетях становятся лица без определённого места "
        "жительства, которые живут в тепловых камерах и пользуются горячей водой.\n\n"
        "Нередко из-за неисправности дренажной запорной арматуры происходит "
        "быстрое запаривание тепловой камеры, ухудшая видимость, а заполнение "
        "горячей водой не оставляет никаких шансов на спасение.",
    )
    add_image(s, f"{IMG}/manhole_steam.png", bottom=True)

    # 11. Non-production: motorists
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Непроизводственный травматизм",
        "Автомобилисты в опасной зоне",
        "Гражданский травматизм при авариях тепловых сетей",
        "Половина пострадавших из числа гражданского населения приходится на "
        "автомобилистов, которые в силу личной невнимательности или "
        "неудовлетворительной организации аварийно-восстановительных работ "
        "попадают в опасную зону.",
        note="Организация ограждения и освещения места аварии обязательна.",
    )
    add_image(s, f"{IMG}/overturned_car_steam.png", bottom=True)

    # 12. Production traumatism definitions
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Производственный травматизм",
        "Производственный травматизм и несчастный случай",
        "Ключевые определения",
        "",
    )
    add_definition_block(
        s,
        [
            (
                "Производственный травматизм",
                "совокупность травм, полученных работающими на производстве и вызванных "
                "несоблюдением требований безопасности труда. Производственная травма "
                "всегда является результатом несчастного случая.",
            ),
            (
                "Несчастный случай на производстве",
                "случай воздействия на работающих опасного производственного фактора "
                "при выполнении трудовых обязанностей или заданий руководителя работ. "
                "Также — событие, в результате которого застрахованный получил увечье "
                "при исполнении обязанностей по трудовому договору (ст. 3 ФЗ № 125-ФЗ, "
                "ст. 227 ТК РФ).",
            ),
        ],
        start_top=Emu(2000000),
    )

    # 13. Case: Perm Network Company
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Производственный травматизм",
        "Случай в тепловой камере, Пермь (21.03.2018)",
        "Тяжёлый ожог слесаря при ремонте трубопровода",
        "21 марта 2018 г. ремонтная бригада «Пермской сетевой компании» устраняла "
        "дефект подающего трубопровода. В 4:30 слесарь спустился в тепловую камеру "
        "для открытия дренажной задвижки. Пол камеры был покрыт горячей водой "
        "(около 80 °C) глубиной около 200 мм.\n\n"
        "При спуске слесарь по неизвестной причине упал в горячую воду. Самостоятельно "
        "выбрался и был госпитализирован с ожогами III степени обеих ног и одной руки "
        "(51% поверхности тела). Травма квалифицирована как тяжёлая.",
    )
    add_image(s, f"{IMG}/heating_chamber_manhole.png")

    # 14. Injury types intro
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Классификация травматизма",
        "Основные виды производственных травм",
        "В зависимости от травмирующего объекта принято выделять:",
        "Виды травм",
        [
            "порезы",
            "ожоги",
            "переломы",
            "ушибы",
            "ампутации",
            "вывихи",
        ],
        note="Классификация травм шире: по фактору, числу пострадавших и тяжести.",
    )

    # 15. Full classification
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Классификация травматизма",
        "Классификация производственного травматизма",
        "Производственные травмы классифицируются по нескольким признакам:",
        "Оси классификации",
        [
            "виды травм: порезы, ожоги, переломы, ушибы, ампутации, вывихи",
            "вид травмирующего фактора: механические, термические, химические, электрические, комбинированные",
            "количество пострадавших: одиночные, групповые",
            "тяжесть повреждения: лёгкие, тяжёлые, со смертельным исходом",
        ],
    )

    # 16. Causes classification
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Классификация травматизма",
        "Классификация причин производственного травматизма",
        "На законодательном уровне единой классификации нет — приведён общий перечень",
        "Приведённая классификация позволяет выделять не только виды травм, но и их причины.",
    )
    add_cause_cards(
        s,
        [
            (
                "Технические причины",
                "Конструкторские недостатки, неисправности машин и механизмов, "
                "несовершенство технологического процесса, низкий уровень механизации.",
            ),
            (
                "Санитарно-гигиенические",
                "Нарушение санитарных норм, отсутствие санитарно-бытовых помещений, "
                "плохая организация рабочего места.",
            ),
            (
                "Организационные факторы",
                "Нарушение правил эксплуатации транспорта и оборудования, низкий уровень "
                "организации погрузочно-разгрузочных работ, нарушение режима труда и отдыха.",
            ),
            (
                "Психофизиологические",
                "Нарушение дисциплины труда, употребление алкоголя на рабочих местах, "
                "умышленное нанесение травм, переутомление.",
            ),
        ],
    )

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
