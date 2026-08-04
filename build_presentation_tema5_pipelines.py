#!/usr/bin/env python3
"""Build full Tema 5 pipeline operation presentation from 5.docx + source pptx style."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________cdd8.pptx"
TMP = "/tmp/tema5_pipe_base.pptx"
OUT = "/workspace/Презентация_Эксплуатация_трубопроводов.pptx"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep useful content slides; drop premature итог (4) and final итог (9) — rebuild summaries
KEEP = {0, 1, 2, 3, 5, 6, 7, 8}


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def content_shapes(slide):
    shapes = []
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        if s.top is not None and s.left is not None and s.top > Emu(6200000) and s.left > Emu(8000000):
            continue
        if s.text_frame.text.strip() or True:
            shapes.append(s)
    # Prefer shapes that already had text / typical content positions
    shapes = [s for s in slide.shapes if s.has_text_frame]
    shapes = [
        s
        for s in shapes
        if not (
            s.top is not None
            and s.left is not None
            and s.top > Emu(6200000)
            and s.left > Emu(8000000)
        )
    ]
    return sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 10, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 10, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10637215), Emu(6455664), Emu(914400), Emu(274320)
            )
            set_text(box, f"{i:02d} / {total:02d}", 10, True, RED)


def fill_hero(slide, section, title, body, note=None):
    """Title/hero style: section, title, body, note (left column)."""
    shapes = content_shapes(slide)
    # Left column texts only
    left = [s for s in shapes if (s.left or 0) < Emu(7000000)]
    left = sorted(left, key=lambda s: s.top or 0)
    for sh in left:
        set_text(sh, "", 12, False, DARK)
    set_text(left[0], section, 13, True, RED)
    set_text(left[1], title, 28, True, DARK)
    set_text(left[2], body, 14, False, DARK)
    if len(left) > 3:
        set_text(left[3], note or "", 13, True, AMBER if note else GRAY)


def fill_bullet_body(slide, section, title, intro, items, note=None):
    """Bullet body slide (like source slide 2)."""
    shapes = content_shapes(slide)
    left = [s for s in shapes if (s.left or 0) < Emu(7500000)]
    left = sorted(left, key=lambda s: (s.top or 0, s.left or 0))

    # Identify: section, title, intro, body (tall), note
    for sh in left:
        set_text(sh, "", 12, False, DARK)

    set_text(left[0], section, 13, True, RED)
    set_text(left[1], title, 24, True, DARK)
    set_text(left[2], intro, 14, False, GRAY)

    body = max(left[3:], key=lambda s: s.height or 0)
    note_shape = None
    for sh in left[3:]:
        if sh is body:
            continue
        if (sh.top or 0) > Emu(5400000):
            note_shape = sh
        else:
            set_text(sh, "", 12, False, DARK)

    size = 14 if len(items) <= 5 else 13
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size, False, DARK)

    if note:
        if note_shape:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            # amber box text if auto-shape note exists nearby
            for sh in slide.shapes:
                if sh.has_text_frame and (sh.top or 0) > Emu(5600000) and (sh.left or 0) < Emu(7000000):
                    set_text(sh, note, 12, True, AMBER)
                    break


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def chunks(items, n=5):
    for i in range(0, len(items), n):
        yield items[i : i + n]


def enrich_kept(prs):
    """Refresh kept slides from docx; indices after KEEP prune: 0..7."""
    # 0 hero schemes
    fill_hero(
        prs.slides[0],
        "РАЗДЕЛ 1. ЭКСПЛУАТАЦИЯ ТРУБОПРОВОДОВ",
        "Исполнительные и эксплуатационные схемы",
        "На рабочих местах персонала, обслуживающего трубопровод, эксплуатирующая организация "
        "должна обеспечить наличие в доступной для постоянного использования форме комплекта "
        "исполнительных схем трубопроводов или разработанных на их основе эксплуатационных "
        "(технологических) схем, а также производственных инструкций.",
        note="Схемы нужны для безопасного пуска, отключения, ремонта и испытаний.",
    )

    # 1 two-column docs slide — rewrite both columns cleanly
    s1 = prs.slides[1]
    shapes = content_shapes(s1)
    # Map by position: section, title, intro, left header, left body, right header, right body
    texts = sorted(
        [s for s in shapes if s.has_text_frame],
        key=lambda s: (s.top or 0, s.left or 0),
    )
    # Clear all text first
    for sh in texts:
        if not (sh.top and sh.top > Emu(6200000) and sh.left and sh.left > Emu(8000000)):
            set_text(sh, "", 12, False, DARK)
    # Assign by approximate roles
    section_sh = title_sh = intro_sh = None
    left_h = left_b = right_h = right_b = None
    for sh in texts:
        if sh.top and sh.top > Emu(6200000):
            continue
        if (sh.top or 0) < Emu(700000) and section_sh is None:
            section_sh = sh
        elif (sh.top or 0) < Emu(1600000) and title_sh is None:
            title_sh = sh
        elif (sh.top or 0) < Emu(2200000) and intro_sh is None:
            intro_sh = sh
        elif (sh.height or 0) < Emu(500000) and (sh.top or 0) < Emu(2600000):
            if (sh.left or 0) < Emu(3500000) and left_h is None:
                left_h = sh
            elif (sh.left or 0) >= Emu(3500000) and right_h is None:
                right_h = sh
        elif (sh.height or 0) > Emu(2000000):
            if (sh.left or 0) < Emu(3500000) and left_b is None:
                left_b = sh
            elif (sh.left or 0) >= Emu(3500000) and right_b is None:
                right_b = sh
    if section_sh:
        set_text(section_sh, "РАЗДЕЛ 1. ЭКСПЛУАТАЦИЯ ТРУБОПРОВОДОВ", 13, True, RED)
    if title_sh:
        set_text(title_sh, "Документы на рабочем месте персонала", 22, True, DARK)
    if intro_sh:
        set_text(
            intro_sh,
            "Обязательные документы для персонала, обслуживающего трубопровод",
            14,
            False,
            GRAY,
        )
    if left_h:
        set_text(left_h, "Исполнительные и эксплуатационные схемы", 14, True, RED)
    if right_h:
        set_text(right_h, "Производственная инструкция", 14, True, RED)
    if left_b:
        tf = left_b.text_frame
        tf.clear()
        for i, line in enumerate(
            [
                "•  Обеспечивают безопасный пуск, отключение, ремонт и испытания",
                "•  Дополняют указания производственных инструкций",
                "•  Должны быть доступны для постоянного использования",
                "•  Персонал обязан уметь применять схемы на практике",
            ]
        ):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ""
            r = p.add_run()
            r.text = line
            set_run(r, 13, False, DARK)
    if right_b:
        tf = right_b.text_frame
        tf.clear()
        for i, line in enumerate(
            [
                "•  Состав схемы и назначение трубопровода",
                "•  Обязанности персонала в смену",
                "•  Проверка КИП, арматуры и предохранителей",
                "•  Пуск, остановка, ремонт, авария, журнал смены",
            ]
        ):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ""
            r = p.add_run()
            r.text = line
            set_run(r, 13, False, DARK)

    # 2 instruction part 1
    fill_bullet_body(
        prs.slides[2],
        "РАЗДЕЛ 1 · Производственная инструкция",
        "Что регламентирует производственная инструкция (1/2)",
        "Обязательные разделы инструкции по эксплуатации трубопровода",
        [
            "трубопровод (система) и входящее оборудование; назначение и состав схемы",
            "обязанности персонала во время дежурства по наблюдению и контролю",
            "порядок, сроки и способы проверки КИП, арматуры, предохранительных устройств, автоматики",
            "порядок подготовки к пуску (заполнение, прогрев), пуска и остановки трубопровода",
            "меры безопасности при выводе в ремонт и сливе рабочей среды",
        ],
        note="Схемы и инструкция должны быть доступны на рабочем месте постоянно.",
    )

    # 3 instruction part 2
    fill_bullet_body(
        prs.slides[3],
        "РАЗДЕЛ 1 · Производственная инструкция",
        "Что регламентирует производственная инструкция (2/2)",
        "Аварийные случаи и документация смены",
        [
            "случаи немедленной остановки трубопровода и связанного оборудования (по ФНП и специфике схемы)",
            "порядок действий персонала при аварии или инциденте",
            "порядок ведения сменного (оперативного) журнала и форм приёма/сдачи смены",
            "описание схемы и порядка пуска/остановки, при которых воздействуют на арматуру и аппаратуру",
        ],
        note="Каждый работник обязан знать порядок аварийной остановки.",
    )

    # 4 акт термин (was 5)
    fill_hero(
        prs.slides[4],
        "ВВОД В ЭКСПЛУАТАЦИЮ · АКТ ГОТОВНОСТИ",
        "Акт готовности оборудования",
        "Результаты проверки готовности оборудования к пуску и организации надзора оформляют "
        "актом готовности оборудования под давлением к вводу в эксплуатацию. "
        "Акт прилагают к паспорту и передают руководителю эксплуатирующей организации.",
        note="Без решения о вводе, оформленного по акту, пуск оборудования недопустим.",
    )

    # 5-7 act structure — light label refresh
    for idx, label in (
        (5, "АКТ ГОТОВНОСТИ"),
        (6, "РАБОТА КОМИССИИ"),
        (7, "АКТ ГОТОВНОСТИ ОБОРУДОВАНИЯ"),
    ):
        shapes = content_shapes(prs.slides[idx])
        left = sorted([s for s in shapes if (s.left or 0) < Emu(7000000)], key=lambda s: s.top or 0)
        if left:
            set_text(left[0], label, 13, True, RED)


def append_missing(prs):
    # Templates from current deck
    hero_elems = [deepcopy(s.element) for s in prs.slides[0].shapes]
    bullet_elems = [deepcopy(s.element) for s in prs.slides[2].shapes]

    def hero():
        return clone_slide(prs, hero_elems)

    def bullets():
        return clone_slide(prs, bullet_elems)

    # —— Monitoring ——
    mon = [
        "величины тепловых перемещений и соответствие расчётным значениям по индикаторам (реперам)",
        "отсутствие защемлений и повышенной вибрации трубопроводов",
        "плотность предохранительных устройств, арматуры и фланцевых соединений",
        "температурный режим металла при пусках и остановах",
        "степень затяжки пружин подвесок и опор в рабочем и холодном состоянии — не реже 1 раза в 2 года",
        "герметичность сальниковых уплотнений арматуры",
        "соответствие указателей положения регулирующей арматуры на щитах её фактическому положению",
        "наличие смазки подшипников, узлов приводов, винтовых пар, редукторов — по руководству",
    ]
    parts = list(chunks(mon, 4))
    for idx, part in enumerate(parts, 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"РАЗДЕЛ 1 · Контроль при эксплуатации ({idx}/{len(parts)})",
            "Что контролируют при эксплуатации трубопроводов и арматуры",
            "По производственным инструкциям, руководству и проекту",
            part,
            note=(
                "При заполнении неостывших паропроводов контролируют разность температур стенки и среды."
                if idx == len(parts)
                else None
            ),
        )

    # —— Valve marking ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Арматура",
        "Маркировка арматуры и указатели положения",
        "Требования к обозначению и индикации",
        [
            "на арматуре или бирке — названия и номера по технологическим схемам",
            "указатели направления вращения штурвала (маховика)",
            "регулирующие клапаны — указатели степени открытия регулирующего органа",
            "запорная арматура — указатели положения запорного органа (открыто / закрыто)",
        ],
    )

    # —— Manometer frequency ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Манометры и ПК",
        "Сроки проверки исправности манометров и предохранительных клапанов",
        "О результатах — запись в сменном (оперативном) журнале",
        [
            "давление до 1,4 МПа включительно — не реже одного раза в смену",
            "свыше 1,4 до 4,0 МПа включительно — не реже одного раза в сутки",
            "свыше 4 МПа и все трубопроводы на ТЭС — в сроки по утверждённой инструкции",
            "исправность ПК — кратковременным подрывом или на испытательных стендах",
        ],
    )

    # —— Manometer class ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Манометры",
        "Класс точности, красная черта и обслуживание",
        "Требования к выбору и установке манометров",
        [
            "до 2,5 МПа — класс точности не ниже 2,5; свыше 2,5 до 14 МПа — не ниже 1,5; свыше 14 МПа — не ниже 1",
            "на шкале — красная черта разрешённого рабочего давления (или отдельный указатель max)",
            "перед манометром — трёхходовой кран или аналог для продувки и отключения",
            "поверка не реже одного раза в 12 месяцев (если иное не указано в документации) с клеймом/пломбой",
        ],
    )

    # —— Safety valves / reducing ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Предохранительные и редуцирующие устройства",
        "Отвод от ПК и редуцирование давления",
        "Когда давление источника выше разрешённого для трубопровода",
        [
            "ПК должны иметь отводящие трубопроводы, защищающие персонал от ожогов",
            "на отводящих трубопроводах запорные устройства не устанавливают",
            "при меньшем разрешённом давлении трубопровода, чем у источника — редуцирующее устройство",
            "редуцирующие устройства — авторегулирование давления; РОУ — также температуры",
        ],
    )

    # —— Repair ——
    s = bullets()
    fill_bullet_body(
        s,
        "РАЗДЕЛ 1 · Ремонт",
        "Ремонтный журнал и подготовка к ремонту",
        "Организация ремонтных работ",
        [
            "ведётся ремонтный журнал (бумажный или электронный с резервированием)",
            "фиксируют работы, вызывающие внеочередное освидетельствование, материалы и документы качества",
            "до ремонта трубопровод отделяют заглушками или отсоединяют от действующего оборудования",
            "ремонт арматуры, ДУ, установка/снятие заглушек — по наряду-допуску",
        ],
    )

    # —— Commissioning checks ——
    ready_docs = [
        "документация изготовителя и соответствие техрегламентам и правилам ПБ",
        "документация качества монтажа (ремонта/реконструкции) и приёмки",
        "положительные результаты технического освидетельствования",
        "результаты пусконаладки и комплексного опробования (если требуются проектом)",
        "документы о приёмке после ПНР и комплексного опробования (при необходимости)",
        "соответствие требованиям техрегулирования и ст. 7 Федерального закона о ПБ",
        "наличие, соответствие проекту и исправность арматуры, КИП, приборов безопасности и защит",
        "правильность установки, размещения и обвязки оборудования",
    ]
    for idx, part in enumerate(chunks(ready_docs, 4), 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"ВВОД В ЭКСПЛУАТАЦИЮ · Проверка готовности ({idx}/2)",
            "Что контролируют при проверке готовности к пуску",
            "Решение о вводе принимает руководитель эксплуатирующей организации",
            part,
            note=(
                "Проверку проводят ответственный за производственный контроль совместно с ответственным "
                "за исправное состояние либо комиссия по распорядительному документу."
                if idx == 2
                else None
            ),
        )

    s = bullets()
    fill_bullet_body(
        s,
        "ВВОД В ЭКСПЛУАТАЦИЮ · Надзор",
        "Проверка организации надзора за эксплуатацией",
        "Что контролируют дополнительно",
        [
            "наличие обслуживающего персонала, обученного и допущенного по ФНП и распорядительным документам",
            "наличие должностных инструкций для ответственных лиц и специалистов",
            "наличие производственных инструкций и эксплуатационной документации, их соответствие ФНП",
            "решение о вводе оформляют распорядительным документом; сведения записывают в паспорт",
        ],
        note="Член комиссии вправе изложить особое мнение письменно — оно прилагается к акту.",
    )

    # —— Nameplate / coloring ——
    s = bullets()
    fill_bullet_body(
        s,
        "ВВОД В ЭКСПЛУАТАЦИЮ · Идентификация",
        "Табличка перед пуском и опознавательная окраска",
        "Перед пуском на каждой единице оборудования",
        [
            "номер оборудования по системе нумерации эксплуатирующей организации",
            "учётный номер, присвоенный территориальным органом Ростехнадзора",
            "разрешённые параметры (давление, температура рабочей среды)",
            "даты следующих НВО/ГИ (котлы, сосуды) или НО (трубопроводы) и срок службы",
        ],
        note="Трубопроводы окрашивают и маркируют по назначению и параметрам среды.",
    )

    # —— Routing ——
    s = bullets()
    fill_bullet_body(
        s,
        "ПРОКЛАДКА ТРУБОПРОВОДОВ",
        "Уклоны, каналы и тоннели",
        "Требования к размещению трубопроводов пара и горячей воды",
        [
            "горизонтальные участки пара и ГВ — уклон не менее 0,004; тепловых сетей — не менее 0,002",
            "полупроходные каналы — высота в свету ≥ 1,5 м; ширина прохода между изолированными трубами — по нормам",
            "проходные тоннели (коллекторы) — высота ≥ 2 м; ширина прохода — по проекту/нормам",
            "надземная открытая прокладка допускается совместно с технологическими трубопроводами (с исключениями)",
        ],
    )

    s = bullets()
    fill_bullet_body(
        s,
        "ПРОКЛАДКА ТРУБОПРОВОДОВ",
        "Люки, камеры, доступ к арматуре",
        "Обслуживание подземных и надземных участков",
        [
            "проходные каналы — входные люки с лестницей/скобами; расстояние между люками ≤ 300 м",
            "на всех трубопроводах тепловых сетей — антикоррозионная, тепловая и гидроизоляционная защита",
            "камеры обслуживания подземных трубопроводов — не менее двух люков с лестницами/скобами",
            "арматура — в местах, доступных для безопасного обслуживания и ремонта; чугунная — защита от изгиба",
        ],
    )

    s = bullets()
    fill_bullet_body(
        s,
        "ПРОКЛАДКА ТРУБОПРОВОДОВ",
        "Тепловые перемещения, ползучесть, секционирование",
        "Проектные требования к оснащению",
        [
            "паропроводы Øвн > 150 мм при t ≥ 300 °C — указатели тепловых перемещений (по проекту)",
            "при температуре, вызывающей ползучесть, — устройства контроля роста остаточных деформаций",
            "запорная арматура на всех выводах тепловых сетей от источников теплоты",
            "на водяных сетях Øвн ≥ 100 мм — секционирующие задвижки не реже чем через 1000 м с перемычкой",
        ],
        note="Также запорная арматура на ответвлениях Øвн ≥ 100 мм и на конденсатопроводах к сборному баку.",
    )

    # —— Drainage ——
    drain = [
        "дренажи для слива после гидравлического испытания и воздушники в верхних точках",
        "паропроводы — дренажные устройства в местах возможного скопления конденсата при пуске и работе",
        "в нижних точках водяных сетей и конденсатопроводов — штуцера с арматурой для спуска воды",
        "из паропроводов тепловых сетей в нижних точках и перед вертикальными подъёмами — непрерывный отвод конденсата через конденсатоотводчики",
        "пусковой дренаж: через 400–500 м при попутном и 200–300 м при встречном уклоне",
        "отключаемые участки паропроводов — штуцер с запорным устройством в концевых точках для прогрева и продувки",
        "контроль работы дренажей при прогреве; нижние концевые точки и изгибы — устройства продувки",
    ]
    for idx, part in enumerate(chunks(drain, 4), 1):
        s = bullets()
        fill_bullet_body(
            s,
            f"ПРОКЛАДКА · Дренаж и воздушники ({idx}/2)",
            "Дренажные и воздухоудаляющие устройства",
            "Обязательное оснащение трубопроводов",
            part,
        )

    # —— Electric drives / insulation ——
    s = bullets()
    fill_bullet_body(
        s,
        "ПРОКЛАДКА · Приводы и изоляция",
        "Электроприводы арматуры и тепловая изоляция",
        "Дополнительные требования",
        [
            "электроприводы на водяных сетях Øвн ≥ 500 мм при P ≥ 1,6 МПа и Øвн ≥ 300 мм при P ≥ 2,5 МПа "
            "(а также на паровых — по нормам проекта)",
            "материалы изоляции — по параметрам и условиям эксплуатации трубопровода",
            "изоляция фланцев, арматуры и контрольных участков (швы, бобышки ползучести) — съёмная",
            "на открытом воздухе и у масло-/мазутопроводов — металлическое или иное защитное покрытие изоляции",
        ],
    )

    # —— Summary ——
    s = bullets()
    fill_bullet_body(
        s,
        "ИТОГИ",
        "Ключевые требования эксплуатации трубопроводов пара и горячей воды",
        "Что должен знать обслуживающий персонал",
        [
            "на рабочем месте — схемы, инструкция; контроль перемещений, плотности, КИП и ПК",
            "манометры и ПК проверяют по срокам; класс точности и красная черта обязательны",
            "ремонт — с отглушением и по наряду-допуску; ввод — по акту готовности и распоряжению",
            "прокладка, дренажи, изоляция и арматура — строго по проекту и ФНП",
        ],
        note="Нарушение порядка пуска, ремонта или контроля — прямая угроза аварии на ОПО.",
    )


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for item in new_items:
        sldIdLst.append(item)


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    prs.save(TMP)

    prs = Presentation(TMP)
    enrich_kept(prs)
    n_before = len(prs.slides)
    append_missing(prs)
    n_after = len(prs.slides)
    # kept 0-7; appended start at 8
    # Logical order: schemes/instruction → monitoring → valves/manometers → repair
    # → commissioning/act kept → routing/drain → summary
    kept = list(range(n_before))  # 0..7
    added = list(range(n_before, n_after))
    # added: mon(2), valve, man_freq, man_class, pk_redu, repair, ready(2), nadzor,
    #        nameplate, route3, drain2, drives, summary  = count them
    # Place act block after commissioning checks
    # kept: 0 schemes, 1 docs, 2 instr1, 3 instr2, 4 act term, 5 act struct, 6 opinion, 7 remarks
    order = [
        0, 1, 2, 3,  # schemes + instruction
        *added[:6],  # monitoring..repair (2 mon + valve + manf + manc + pk = 6? )
    ]
    # Recalculate added groups explicitly by rebuild count
    # Safer: don't reorder complex — append is already mostly chronological if we reorder kept act later
    # Rebuild order manually based on known append sequence:
    # append order:
    # 8-9 mon, 10 valve, 11 man_freq, 12 man_class, 13 pk, 14 repair,
    # 15-16 ready, 17 nadzor, 18 nameplate, 19-21 route, 22-23 drain, 24 drives, 25 summary
    if n_after >= 26:
        order = [
            0, 1, 2, 3,
            8, 9, 10, 11, 12, 13, 14,
            15, 16, 17,
            4, 5, 6, 7,  # act block
            18, 19, 20, 21, 22, 23, 24,
            25,
        ]
        if len(order) == n_after:
            reorder_slides(prs, order)

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides), before_append={n_before}")


if __name__ == "__main__":
    main()
