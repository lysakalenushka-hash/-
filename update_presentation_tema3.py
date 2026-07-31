#!/usr/bin/env python3
"""Compact Tema 3 presentation to ~27 slides covering 3.docx."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________e977.pptx"
TMP = "/tmp/tema3_compact.pptx"
OUT = "/workspace/Презентация_Тема_3_Персонал.pptx"
IMG = "/workspace/assets/personnel_images"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep strong base slides (0-based). Quizzes were 10,11,12,24 — exclude them.
# 0 title, 1–2 rules/scope, 4–5 p.2.4, 9 categories, 14–15 groups,
# 17 leader, 19–20 knowledge, 26–27 compat/nadzor, 30–31 drills
KEEP = {0, 1, 2, 4, 5, 9, 14, 15, 17, 19, 20, 26, 27, 30, 31}


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    break


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 26, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 26, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], list_title, 11, True, RED)

    pairs, more_shape, note_shape = [], None, None
    for sh in content[4:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    if note and note_shape:
        set_text(note_shape, note, 12, True, AMBER)


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3800000), width=Emu(10700000), height=Emu(2300000))
    else:
        slide.shapes.add_picture(path, Emu(6600000), Emu(1200000), width=Emu(5000000), height=Emu(4500000))


def find_templates(prs):
    text_tpl = list_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if text_tpl is None and 5 <= n <= 7:
            text_tpl = slide
        if list_tpl is None and n >= 10:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(t in {"01", "02", "1", "2"} for t in texts):
                list_tpl = slide
        if text_tpl and list_tpl:
            break
    return text_tpl, list_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def append_compact(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    s = nl()
    fill_list_slide(
        s, "3.2 · Работа с персоналом",
        "Обязательные формы работы с персоналом",
        "Приказ Минэнерго № 796 — по категориям персонала",
        "Формы работы",
        [
            "АТП: предэкзаменационная подготовка, проверка знаний, производственный инструктаж",
            "диспетчерский / оперативный / ОРп: стажировка, новая должность, проверка, дублирование, тренировки, спецподготовка, инструктаж",
            "ремонтный: стажировка, новая должность, проверка знаний, производственный инструктаж",
            "вспомогательный: предэкзаменационная подготовка и проверка знаний",
            "персонал РЗА — дополнительно подготовка и допуск к ТО устройств РЗА",
            "спецподготовка: 5–20% рабочего времени; тренировки, теория, плановый инструктаж",
        ],
        note="Требования не распространяются на потребителей — физических лиц.",
    )

    s = nl()
    fill_list_slide(
        s, "3.3 · Новая должность",
        "Подготовка по новой должности (рабочему месту)",
        "Для диспетчерского, оперативного, ОРп и ремонтного персонала",
        "Этапы",
        [
            "теоретическая подготовка → стажировка → проверка знаний",
            "дублирование (диспетчерский, оперативный, ОРп)",
            "противоаварийные и противопожарные тренировки",
            "допуск к самостоятельной работе",
            "программа индивидуальная; аттестация по безопасности — до дублирования",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "3.4 · Стажировка",
        "Стажировка в организациях",
        "Для диспетчерского, оперативного, ОРп и ремонтного персонала",
        "Ключевые требования",
        [
            "под руководством назначенного ответственного за стажировку",
            "минимум 2 рабочих дня (смены) на каждом рабочем месте",
            "максимум — не более 14 рабочих дней (смен)",
            "допуск к стажировке оформляется ОРД",
            "ознакомление диспетчеров с объектами — не менее 1 рабочего дня",
        ],
    )
    add_image(s, f"{IMG}/switchgear_room.png", bottom=True)

    s = nl()
    fill_list_slide(
        s, "3.5 · Группы I–V",
        "Группы по электробезопасности",
        "Пять классов допуска по опыту и правам работы",
        "Группы",
        [
            "I — неэлектротехнический персонал (запись в журнале)",
            "II — начинающие, переносной инструмент, водители в ЭУ",
            "III — единоличная работа до 1000 В; в бригаде — и выше",
            "IV — обслуживание > 1000 В; старший смены / ответственный до 1000 В",
            "V — ответственность за электроустановки > 1000 В",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "3.5 · Проверка знаний",
        "Внеочередная проверка и состав комиссии",
        "По Приказу Минэнерго № 796",
        "Основания и комиссия",
        [
            "внеочередная: новые НПА, перевод, новое оборудование, нарушения, аварии",
            "внеочередная: предписание надзора; перерыв в работе > 6 месяцев",
            "комиссия ≥ 5 человек; при проверке ≥ 3 членов, включая председателя",
            "при присвоении группы ЭБ — ≥ 3 с группой, один не ниже присваиваемой",
            "протокол, журнал, удостоверение; «неуд.» — повтор не более 1 месяца",
            "ознакомление с графиком — в течение месяца, не позднее чем за 14 дней",
        ],
        note="Допуск к самостоятельной работе без очередной проверки знаний запрещён.",
    )
    add_image(s, f"{IMG}/kipia_instruments.png", bottom=True)

    s = nl()
    fill_list_slide(
        s, "3.6–3.7 · Дублирование и допуск",
        "Дублирование и допуск к самостоятельной работе",
        "Для диспетчерского, оперативного и ОРп персонала",
        "Ключевые правила",
        [
            "дублирование при новой должности — ≥ 12 рабочих смен после проверки знаний",
            "после перерыва > 30 дней — ≥ 1 смена; 60 дней–6 мес. — обязательно",
            "в дублировании — контрольные противоаварийная и противопожарная тренировки",
            "допуск действует до очередной проверки знаний",
            "отзыв: «неуд.» по знаниям/тренировке, нарушения НПА, акты НС/аварий, предписания",
        ],
        note="Ответственность несут дублёр и руководящий работник в равной мере.",
    )

    s = nl()
    fill_list_slide(
        s, "3.8 · Инструктаж",
        "Производственный инструктаж",
        "Плановый и внеплановый для диспетчерского, оперативного, ОРп и ремонтного персонала",
        "Периодичность",
        [
            "плановый для диспетчерского / оперативного / ОРп — ежемесячно",
            "плановый для ремонтного — не реже одного раза в квартал",
            "программа актуализируется ежегодно",
            "внеплановый — при новых документах, схемах, нарушениях, пропуске планового",
            "проверка усвоения опросом; регистрация в журнале",
        ],
    )
    add_image(s, f"{IMG}/ppe.png", bottom=True)

    s = nl()
    fill_list_slide(
        s, "3.9 · Тренировки",
        "Противоаварийные и противопожарные тренировки",
        "Учебные и контрольные тренировки персонала",
        "Периодичность",
        [
            "учебная противопожарная — 1 раз в 3 календарных месяца",
            "контрольная противопожарная — 1 раз в 6 календарных месяцев",
            "допускается совмещать контрольные противоаварийные и противопожарные",
            "проводятся в свободное от дежурства время (или в смену по решению)",
            "при «неуд.» по контрольной ППТ — повтор не позднее 1 месяца",
        ],
        note="Время тренировок включается в рабочее время.",
    )

    s = nl()
    fill_list_slide(
        s, "3.10–3.11 · ДПО и обходы",
        "Повышение квалификации и обходы рабочих мест",
        "Дополнительное профессиональное образование и контроль готовности",
        "Требования",
        [
            "ДПО для АТП, диспетчерского, оперативного, ОРп и ремонтного — не реже 1 раза в 5 лет",
            "с отрывом от основных обязанностей в лицензированной организации",
            "обходы — для диспетчерского, оперативного и ОРп персонала",
            "проверка: режим, документация, дефекты, наряды-допуски, СИЗ и инструмент",
            "результаты фиксируются; при нарушениях — меры по устранению",
        ],
    )

    s = nl()
    fill_list_slide(
        s, "Итоги темы 3",
        "Ключевые требования к персоналу и его подготовке",
        "Что должен знать каждый специалист",
        "Главные выводы",
        [
            "пять категорий персонала + обязательные формы работы (Приказ № 796)",
            "группа по ЭБ — по результатам проверки знаний (кроме I группы)",
            "новая должность: теория → стажировка → проверка → дублирование → допуск",
            "проверка знаний: 12 месяцев / 3 года; комиссия ≥ 5 человек",
            "инструктажи, тренировки, ДПО раз в 5 лет, обходы рабочих мест",
            "допуск без подготовки и проверки знаний запрещён",
        ],
    )


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    try:
        add_image(prs.slides[0], f"{IMG}/ppe.png")
    except Exception:
        pass
    prs.save(TMP)

    prs = Presentation(TMP)
    append_compact(prs)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
