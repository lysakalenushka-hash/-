#!/usr/bin/env python3
"""
Restyle «Организационно-правовые аспекты…» to match «Общие сведения о СУОТ» style.
Keeps text meaning and pictures; changes slide size, fonts, colors, title chrome.
"""

from __future__ import annotations

import re
import shutil
import zipfile
from copy import deepcopy
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

SRC_ORG = Path("/workspace/restyle_org_suot/src_org.pptx")
SRC_SUOT = Path("/workspace/restyle_org_suot/src_suot.pptx")
OUT = Path(
    "/workspace/1_Организационно-правовые_аспекты_оказания_первой_помощи_стиль_СУОТ.pptx"
)

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

SCALE = 0.5  # ORG is 2× SUOT canvas

# SUOT palette / typography
TEAL = "1ABD9E"
TEAL_LIGHT = "8DDBCD"
YELLOW_OLD = "FFC000"
TITLE_COLOR = "404040"
EYEBROW_COLOR = "535353"
BODY_COLOR = "535353"


def scale_attr(el, attr, factor):
    if el.get(attr) is None:
        return
    try:
        el.set(attr, str(int(round(int(el.get(attr)) * factor))))
    except ValueError:
        pass


def scale_tree_geometry(root, factor):
    """Scale positions/sizes in slide XML."""
    for el in root.iter():
        tag = etree.QName(el).localname
        if tag in ("off", "ext", "chOff", "chExt"):
            scale_attr(el, "x", factor)
            scale_attr(el, "y", factor)
            scale_attr(el, "cx", factor)
            scale_attr(el, "cy", factor)
        # EMU in some other attrs
        if tag == "sndSz":
            continue
        # font size sz is in 100ths of a point — also ~2x in ORG vs SUOT for titles
        if "sz" in el.attrib and tag in ("rPr", "defRPr", "endParaRPr"):
            try:
                sz = int(el.get("sz"))
                # Only scale large title-like sizes; keep already-small
                if sz >= 30000:  # >=30pt hundredths? actually sz=584200 is invalid — pptx uses 4600 for 46pt
                    pass
            except Exception:
                pass
        # python-pptx / OOXML font sz is like "4600" for 46pt (hundredths of a point)
        if tag in ("rPr", "defRPr", "endParaRPr") and el.get("sz"):
            try:
                sz = int(el.get("sz"))
                # ORG titles ~4600, body ~2800; SUOT titles ~2300, body ~1600, eyebrow ~1100
                new_sz = max(1100, int(round(sz * factor)))
                el.set("sz", str(new_sz))
            except Exception:
                pass


def replace_fonts_in_tree(root):
    for el in root.iter():
        tag = etree.QName(el).localname
        if tag in ("latin", "ea", "cs"):
            tf = el.get("typeface")
            if not tf:
                continue
            low = tf.lower()
            if "open sans" in low:
                el.set("typeface", "PFDinDisplayPro-Medium")
            elif tf == "Montserrat Semi Bold":
                el.set("typeface", "Montserrat SemiBold")


def replace_colors_in_tree(root):
    for el in root.iter():
        if etree.QName(el).localname == "srgbClr":
            val = (el.get("val") or "").upper()
            if val == YELLOW_OLD:
                el.set("val", TEAL)
            # keep blues/grays of shared system


def process_slide_xml(xml_bytes: bytes) -> bytes:
    root = etree.fromstring(xml_bytes)
    scale_tree_geometry(root, SCALE)
    replace_fonts_in_tree(root)
    replace_colors_in_tree(root)
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def copy_theme_from_suot(dst_zip: Path, suot: Path):
    """Replace theme1.xml in dst with SUOT theme."""
    tmp = dst_zip.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(suot, "r") as zs, zipfile.ZipFile(dst_zip, "r") as zd, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as zo:
        theme = zs.read("ppt/theme/theme1.xml")
        for item in zd.infolist():
            data = zd.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                data = theme
            zo.writestr(item, data)
    tmp.replace(dst_zip)


def restyle_package(src: Path, out: Path):
    """Scale slides + fonts/colors in the pptx package, set sldSz."""
    tmp = out.with_suffix(".build.pptx")
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            name = info.filename
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                data = process_slide_xml(data)
            elif name == "ppt/presentation.xml":
                root = etree.fromstring(data)
                sldSz = root.find(f"{P}sldSz")
                if sldSz is not None:
                    sldSz.set("cx", "12192000")
                    sldSz.set("cy", "6858000")
                notesSz = root.find(f"{P}notesSz")
                # leave notes
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            elif name.startswith("ppt/slideLayouts/") and name.endswith(".xml"):
                root = etree.fromstring(data)
                scale_tree_geometry(root, SCALE)
                replace_fonts_in_tree(root)
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            elif name.startswith("ppt/slideMasters/") and name.endswith(".xml"):
                root = etree.fromstring(data)
                scale_tree_geometry(root, SCALE)
                replace_fonts_in_tree(root)
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(info, data)
    tmp.replace(out)


def get_slides(prs):
    slides = []
    for sldId in prs.slides._sldIdLst:
        rId = sldId.get(qn("r:id"))
        slides.append(prs.part.related_part(rId).slide)
    return slides


def set_run_font(run, name=None, size=None, bold=None, color=None):
    if name:
        run.font.name = name
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def polish_with_pptx(path: Path):
    """Final pass: title slide chrome like SUOT + role-based typography."""
    prs = Presentation(str(path))
    slides = get_slides(prs)
    assert prs.slide_width == 12192000

    # --- Title slide: rebuild chrome like SUOT (dark bar + Montserrat SemiBold) ---
    title_slide = slides[0]
    # Capture title text
    title_text = ""
    for sh in title_slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if len(t) > 20:
                title_text = t
                break
    if not title_text:
        title_text = "ОРГАНИЗАЦИОННО-ПРАВОВЫЕ АСПЕКТЫ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ"

    # Remove existing non-picture shapes on title (keep pictures if any)
    for sh in list(title_slide.shapes):
        if sh.shape_type == 13:  # picture — keep
            continue
        sh.element.getparent().remove(sh.element)

    # Soft white/gray base already may exist via bg; add light overlay + dark bar
    # Full-slide light veil (like SUOT bg1 alpha) — skip if no photo; use solid soft gray panel style
    # Right gray panel (SUOT often has gray freeform) — use soft full background accent
    bg_rect = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    bg_rect.line.fill.background()
    # send to back by reordering xml
    spTree = title_slide.shapes._spTree
    spTree.remove(bg_rect.element)
    # insert after nvGrpSpPr
    spTree.insert(2, bg_rect.element)

    # Dark translucent bar
    bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(3429000), prs.slide_width, Emu(1314547)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    try:
        # approx 54% opacity via solid — pptx alpha:
        from pptx.oxml.ns import nsmap

        spPr = bar.element.find(qn("p:spPr"))
        solid = spPr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            # might be scheme
            pass
        else:
            # replace with srgb + alpha
            srgb.set("val", "000000")
            for child in list(srgb):
                srgb.remove(child)
            alpha = etree.SubElement(srgb, f"{A}alpha")
            alpha.set("val", "54000")
    except Exception:
        pass
    bar.line.fill.background()

    # Teal accent strip on bar top
    strip = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(3429000), prs.slide_width, Emu(36000)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x1A, 0xBD, 0x9E)
    strip.line.fill.background()

    # Title text
    box = title_slide.shapes.add_textbox(Emu(263352), Emu(3699646), Emu(9000000), Emu(800000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_run_font(run, "Montserrat SemiBold", Pt(20), True, RGBColor(0xFF, 0xFF, 0xFF))

    # --- Content slides: normalize title / eyebrow / body fonts ---
    for idx, slide in enumerate(slides[1:], start=2):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            top = sh.top or 0
            left = sh.left or 0
            w = sh.width or 0
            h = sh.height or 0
            # page number — leave
            if left < Emu(500000) and w < Emu(600000) and top > Emu(2500000):
                continue
            # main title
            if top < Emu(1200000) and w > Emu(4000000) and h < Emu(600000) and t and t == t.upper():
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        set_run_font(
                            r,
                            "PFDinDisplayPro-Regular",
                            Pt(23),
                            False,
                            RGBColor(0x40, 0x40, 0x40),
                        )
                continue
            # eyebrow
            if ("Аспекты оказания" in t or "Тема 1" in t) and h < Emu(500000):
                # unify eyebrow text style like SUOT subtitle
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        set_run_font(
                            r,
                            "PFDinDisplayPro-Medium",
                            Pt(11),
                            False,
                            RGBColor(0x53, 0x53, 0x53),
                        )
                continue
            # body-ish
            if t and h > Emu(800000):
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        # don't force if empty
                        if not r.text:
                            continue
                        name = r.font.name
                        if name in (None, "Open Sans", "Calibri"):
                            set_run_font(r, "PFDinDisplayPro-Medium", None, None, None)
                        try:
                            if r.font.size and r.font.size > Pt(20):
                                r.font.size = Pt(16)
                        except Exception:
                            pass
                        try:
                            r.font.color.rgb = RGBColor(0x53, 0x53, 0x53)
                        except Exception:
                            pass

    # Fix yellow fills remaining via XML on each slide part
    for slide in slides:
        for el in slide._element.iter():
            if etree.QName(el).localname == "srgbClr" and (el.get("val") or "").upper() == YELLOW_OLD:
                el.set("val", TEAL)

    prs.save(str(path))


def verify(path: Path):
    prs = Presentation(str(path))
    slides = get_slides(prs)
    print("size", prs.slide_width, prs.slide_height, "slides", len(slides))
    fonts = set()
    colors = []
    for i, s in enumerate(slides):
        pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = " ".join(sh.text_frame.text.split())
                if t:
                    texts.append(t[:70])
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name:
                            fonts.add(r.font.name)
        xml = etree.tostring(s._element).decode()
        if "FFC000" in xml.upper():
            colors.append(f"slide{i+1}:has-yellow")
        if "1ABD9E" in xml.upper():
            colors.append(f"slide{i+1}:has-teal")
        print(f"{i+1:02d} pics={pics} | {' | '.join(texts)[:150]}")
    print("fonts", fonts)
    print("color markers", colors)


def main():
    shutil.copy(SRC_ORG, OUT)
    print("1) Scaling package + fonts/colors…")
    restyle_package(SRC_ORG, OUT)
    print("2) Copy SUOT theme…")
    copy_theme_from_suot(OUT, SRC_SUOT)
    print("3) Polish title + typography…")
    polish_with_pptx(OUT)
    print("4) Verify…")
    verify(OUT)
    print("OUT", OUT, OUT.stat().st_size)


if __name__ == "__main__":
    main()
