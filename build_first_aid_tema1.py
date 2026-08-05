#!/usr/bin/env python3
"""Reassemble Theme 1 first-aid presentation under new App.2 PP2464 structure (from 01.09.2026)."""

from copy import deepcopy
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = Path(
    "/workspace/first_aid_tema1/"
    "1 Организационно-правовые аспекты оказания первой помощи.pptx"
)
OUT = Path(
    "/workspace/Презентация_Тема_1_Первая_помощь_новая_структура.pptx"
)
TMP = Path("/tmp/first_aid_tema1_build.pptx")

DARK = RGBColor(0x1A, 0x1A, 0x1A)
# Keep existing brand colors where possible; only rewrite run text/size when needed


def set_run_font(run, size_pt=None, bold=None, color=None):
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def replace_shape_text(shape, text, keep_first_run_style=True):
    """Replace shape text robustly (clear all <a:t> nodes, then set first run)."""
    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for t_el in shape.element.findall(f".//{NS}t"):
        t_el.text = ""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
    else:
        r = p0.add_run()
        r.text = text


def find_shapes_by_role(slide):
    """Heuristic roles for this deck's content slides."""
    texts = [sh for sh in slide.shapes if sh.has_text_frame]
    page = title = eyebrow = None
    bodies = []
    for sh in texts:
        t = sh.text_frame.text.strip()
        top = sh.top or 0
        left = sh.left or 0
        w = sh.width or 0
        h = sh.height or 0
        if left < Emu(900000) and top > Emu(5000000) and w < Emu(1200000) and t.isdigit():
            page = sh
        elif top < Emu(2300000) and w > Emu(10000000) and h > Emu(500000):
            title = sh
        elif "Аспекты оказания" in t or (top > Emu(2500000) and top < Emu(3600000) and h < Emu(900000) and w < Emu(9000000)):
            if eyebrow is None and t:
                eyebrow = sh
        elif t and sh is not page:
            bodies.append(sh)
    return {"page": page, "title": title, "eyebrow": eyebrow, "bodies": bodies}


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def clone_slide(prs, source_slide):
    """Clone a slide inside the same presentation (keeps image relationships)."""
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    # remove default shapes
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in source_slide.shapes:
        el = shape.element
        newel = deepcopy(el)
        slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    # copy slide-level relationships that pictures need: already in package via rIds
    # deepcopy of blip rIds still points to same slide's relationships — WRONG for new slide.
    # Safer approach for dividers: only clone title slide (no pictures).
    return slide


def reorder_slides(prs, order):
    """Reorder slides by list of current indices. order length == n slides."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order len {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in new_items:
        sldIdLst.append(el)


def update_pages(prs):
    from pptx.util import Emu as _Emu
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            # page index box: far left, mid-height, narrow
            if (sh.left or 0) < _Emu(900000) and (sh.width or 0) < _Emu(1200000) and (sh.top or 0) > _Emu(5000000):
                # clear regardless of current text (may be corrupted digits)
                replace_shape_text(sh, str(i))


def set_eyebrow(slide, text):
    roles = find_shapes_by_role(slide)
    if roles["eyebrow"] is not None:
        replace_shape_text(roles["eyebrow"], text)


def set_title(slide, text):
    roles = find_shapes_by_role(slide)
    if roles["title"] is not None:
        replace_shape_text(roles["title"], text)
        return True
    # fallback: largest top text
    cands = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if (sh.top or 0) < Emu(3500000) and (sh.width or 0) > Emu(8000000):
            cands.append(sh)
    if cands:
        replace_shape_text(cands[0], text)
        return True
    return False


def make_divider(prs, template_slide, heading, sub):
    """Create a section divider from title slide template (no pictures)."""
    slide = clone_slide(prs, template_slide)
    # title slide: find the big title text box
    titled = False
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "ОРГАНИЗАЦИОННО" in t or "ПЕРВОЙ ПОМОЩИ" in t or len(t) > 20:
            replace_shape_text(sh, heading)
            titled = True
            break
    if not titled:
        for sh in slide.shapes:
            if sh.has_text_frame and (sh.width or 0) > Emu(8000000):
                replace_shape_text(sh, heading)
                break
    # optional subtitle line under title
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if (sh.top or 0) > Emu(7000000) and (sh.width or 0) > Emu(5000000):
            replace_shape_text(sh, sub)
            break
    return slide


def main():
    shutil.copy(SRC, TMP)
    prs = Presentation(str(TMP))
    n = len(prs.slides)
    print(f"Loaded {n} slides from source")

    # Source index map (0-based):
    # 0 title
    # 1-6 = org/NPA (old) -> 1.1
    # 7 concept -> keep under 1.4 as brief
    # 8-9 states/measures -> 1.4
    # 10-11 kits -> 1.2
    # 12-15 sequence -> 1.3
    # 16 safety -> 1.5
    # 17-20 extraction -> 1.6
    # 21 infection -> 1.5
    # 22 call EMS -> 1.7
    # 23 thanks

    title_slide = prs.slides[0]

    # --- Update title ---
    set_title(
        title_slide,
        "ТЕМА 1. ОРГАНИЗАЦИОННО-ПРАВОВЫЕ АСПЕКТЫ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
    )

    # --- Update content titles / eyebrows before reorder ---
    # 1.1
    for i in range(1, 7):
        set_eyebrow(prs.slides[i], "Тема 1 · п. 1.1  НПА и организация оказания первой помощи")
    set_title(prs.slides[1], "ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ В РОССИЙСКОЙ ФЕДЕРАЦИИ")
    set_title(prs.slides[2], "ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ В РОССИЙСКОЙ ФЕДЕРАЦИИ")
    for i in range(3, 7):
        set_title(prs.slides[i], "НОРМАТИВНО-ПРАВОВАЯ БАЗА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ")

    # 1.4 pieces currently at 7-9
    set_eyebrow(prs.slides[7], "Тема 1 · п. 1.4  Состояния и мероприятия первой помощи")
    set_title(prs.slides[7], "ПЕРЕЧЕНЬ СОСТОЯНИЙ И МЕРОПРИЯТИЙ ПЕРВОЙ ПОМОЩИ")
    set_eyebrow(prs.slides[8], "Тема 1 · п. 1.4  Состояния и мероприятия первой помощи")
    set_title(prs.slides[8], "ПЕРЕЧЕНЬ СОСТОЯНИЙ, ПРИ КОТОРЫХ ОКАЗЫВАЕТСЯ ПЕРВАЯ ПОМОЩЬ")
    set_eyebrow(prs.slides[9], "Тема 1 · п. 1.4  Состояния и мероприятия первой помощи")
    set_title(
        prs.slides[9],
        "ПЕРЕЧЕНЬ МЕРОПРИЯТИЙ И ПОСЛЕДОВАТЕЛЬНОСТЬ ИХ ВЫПОЛНЕНИЯ",
    )

    # 1.2 kits at 10-11
    for i in (10, 11):
        set_eyebrow(
            prs.slides[i],
            "Тема 1 · п. 1.2  Укладки, наборы, комплекты и аптечки",
        )
        set_title(
            prs.slides[i],
            "УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        )

    # 1.3 sequence at 12-15
    for i in range(12, 16):
        set_eyebrow(
            prs.slides[i],
            "Тема 1 · п. 1.3  Порядок и приоритетность оказания первой помощи",
        )
        set_title(
            prs.slides[i],
            "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        )

    # 1.5 safety 16, infection 21
    set_eyebrow(prs.slides[16], "Тема 1 · п. 1.5  Безопасные условия и профилактика инфекций")
    set_title(prs.slides[16], "ОБЕСПЕЧЕНИЕ БЕЗОПАСНЫХ УСЛОВИЙ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ")
    set_eyebrow(prs.slides[21], "Тема 1 · п. 1.5  Безопасные условия и профилактика инфекций")
    set_title(prs.slides[21], "ПРОФИЛАКТИКА ИНФЕКЦИОННЫХ ЗАБОЛЕВАНИЙ ПРИ ОКАЗАНИИ ПЕРВОЙ ПОМОЩИ")

    # 1.6 extraction 17-20
    for i in range(17, 21):
        set_eyebrow(
            prs.slides[i],
            "Тема 1 · п. 1.6  Извлечение и перемещение пострадавших",
        )
        set_title(
            prs.slides[i],
            "ИЗВЛЕЧЕНИЕ ПОСТРАДАВШИХ ИЗ ТРУДНОДОСТУПНЫХ МЕСТ И ПЕРЕМЕЩЕНИЕ",
        )

    # 1.7 call 22
    set_eyebrow(prs.slides[22], "Тема 1 · п. 1.7  Вызов скорой медицинской помощи и спецслужб")
    set_title(
        prs.slides[22],
        "ПРАВИЛА ВЫЗОВА СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ И ДРУГИХ СПЕЦИАЛЬНЫХ СЛУЖБ",
    )

    # Save text updates, then add dividers at end and reorder
    prs.save(str(TMP))
    prs = Presentation(str(TMP))

    # Create section divider slides by cloning title (index 0)
    dividers_spec = [
        ("1.1. НОРМАТИВНО-ПРАВОВАЯ БАЗА И ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ", "Теоретическое занятие · Тема 1"),
        ("1.2. УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ", "Основные компоненты и их назначение"),
        ("1.3. ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ", "Теоретическое занятие · Тема 1"),
        ("1.4. ПЕРЕЧЕНЬ СОСТОЯНИЙ И МЕРОПРИЯТИЙ ПЕРВОЙ ПОМОЩИ", "Последовательность выполнения мероприятий"),
        ("1.5. БЕЗОПАСНЫЕ УСЛОВИЯ И ПРОФИЛАКТИКА ИНФЕКЦИЙ", "Теоретическое занятие · Тема 1"),
        ("1.6. ИЗВЛЕЧЕНИЕ И ПЕРЕМЕЩЕНИЕ ПОСТРАДАВШИХ", "В безопасное место"),
        ("1.7. ПРАВИЛА ВЫЗОВА СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ", "И других специальных служб"),
    ]

    # After reload, indices same 0..23
    title_slide = prs.slides[0]
    div_indices = []
    for heading, sub in dividers_spec:
        make_divider(prs, title_slide, heading, sub)
        div_indices.append(len(prs.slides) - 1)

    # Current indices after adding 7 dividers at end:
    # content 0..23, dividers 24..30
    d1_1, d1_2, d1_3, d1_4, d1_5, d1_6, d1_7 = div_indices

    # Target order:
    # title,
    # d1.1, slides 1-6,
    # d1.2, slides 10-11,
    # d1.3, slides 12-15,
    # d1.4, slides 7-9,
    # d1.5, slides 16, 21,
    # d1.6, slides 17-20,
    # d1.7, slide 22,
    # thanks 23
    order = (
        [0]
        + [d1_1]
        + list(range(1, 7))
        + [d1_2]
        + [10, 11]
        + [d1_3]
        + list(range(12, 16))
        + [d1_4]
        + [7, 8, 9]
        + [d1_5]
        + [16, 21]
        + [d1_6]
        + list(range(17, 21))
        + [d1_7]
        + [22]
        + [23]
    )

    assert len(order) == len(prs.slides), (len(order), len(prs.slides))
    reorder_slides(prs, order)

    # Intermediate save before page numbers (avoids zip issues)
    prs.save(str(TMP))
    prs = Presentation(str(TMP))
    update_pages(prs)
    prs.save(str(OUT))

    # Verification printout
    prs = Presentation(str(OUT))
    print(f"\nSaved {OUT} ({len(prs.slides)} slides)\n")
    for i, slide in enumerate(prs.slides, 1):
        roles = find_shapes_by_role(slide)
        title = ""
        if roles["title"] is not None:
            title = roles["title"].text_frame.text.strip().replace("\n", " ")[:90]
        else:
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    title = sh.text_frame.text.strip().replace("\n", " ")[:90]
                    break
        eye = ""
        if roles["eyebrow"] is not None:
            eye = roles["eyebrow"].text_frame.text.strip()[:60]
        print(f"{i:02d}| {title}")
        if eye:
            print(f"    ({eye})")

    # Gap note file for Maria/Alena
    notes = OUT.with_suffix(".md")
    notes.write_text(
        """# Тема 1 — пересборка под Прил. 2 к Порядку № 2464 (с 01.09.2026)

## Файл
`Презентация_Тема_1_Первая_помощь_новая_структура.pptx`

## Структура
1. Титул
2. **1.1** НПА и организация оказания первой помощи (из прежних слайдов организации/НПА)
3. **1.2** Укладки, наборы, комплекты и аптечки (переименованы «современные наборы…»)
4. **1.3** Порядок и приоритетность оказания первой помощи (из «общей последовательности действий»)
5. **1.4** Перечень состояний и мероприятий + последовательность (включая краткое определение ПП)
6. **1.5** Безопасные условия + профилактика инфекций
7. **1.6** Извлечение и перемещение (вынесено в отдельный теоретический пункт)
8. **1.7** Вызов СМП и спецслужб
9. Заключение

## Пробелы / на проверку методисту
- **Приоритетность** (п. 1.3): отдельного исходного слайда не было — блок собран из алгоритма последовательности; при необходимости добавить явный слайд «что делать в первую очередь».
- **Терминология аптечек** (п. 1.2): заголовки обновлены на «укладки, наборы, комплекты и аптечки»; визуальный контент по-прежнему про автомобильную аптечку и аптечку работникам — уточнить, нужны ли другие типы укладок/комплектов.
- **Понятие «первая помощь»**: в новой программе не выделено отдельным элементом; слайд с определением оставлен внутри п. 1.4 как вводный.
""",
        encoding="utf-8",
    )
    print(f"Notes: {notes}")


if __name__ == "__main__":
    main()
