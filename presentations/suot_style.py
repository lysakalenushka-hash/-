#!/usr/bin/env python3
"""Общий стиль презентаций по образцу «1.1. Общие сведения о СУОТ»."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path("/workspace/presentations")
ASSETS = ROOT / "assets"
OUT = ROOT / "out"
OUT.mkdir(parents=True, exist_ok=True)

# 16:9 (как эталон 960×540)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF8, 0xF8)
BG_PANEL = RGBColor(0xF0, 0xF0, 0xF0)
BG_BAR = RGBColor(0xDD, 0xDD, 0xDD)
LINE = RGBColor(0x99, 0x99, 0x99)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x77, 0x77, 0x77)
NUM_BG = RGBColor(0x60, 0x60, 0x60)
# нейтральные серые акценты (без зелёного/бирюзы)
TEAL = RGBColor(0xDD, 0xDD, 0xDD)       # заливка облачков / карточек
TEAL_DARK = RGBColor(0x60, 0x60, 0x60)  # обводка / средняя заливка
TEAL_TEXT = RGBColor(0x33, 0x33, 0x33)  # текст на акцентах
ACCENT_RED = RGBColor(0xED, 0x1C, 0x24)
BANNER = RGBColor(65, 87, 98)
TABLE_HDR = RGBColor(0x55, 0x55, 0x55)
ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "Calibri"
SUBTITLE = "Оказание первой помощи при наружных кровотечениях"
COURSE = "Тема 2 · Оказание первой помощи при наружных кровотечениях"


def asset(name: str) -> Path:
    return ASSETS / name


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", FONT)


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    return sh


def oval(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def tbox(slide, l, t, w, h, text, *, size=18, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def rich_tbox(slide, l, t, w, h, parts, *, size=16, color=TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold in parts:
        r = p.add_run()
        r.text = text
        font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, *, size=15, marker="•", alert=None):
    alert = alert or set()
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        font(r, size, bold=(i in alert), color=ACCENT_RED if i in alert else TEXT)
    return box


def pic_fit(slide, name, l, t, max_w, max_h):
    from PIL import Image as PILImage
    path = asset(name)
    if not path.exists():
        rect(slide, l, t, max_w, max_h, WHITE, line=LINE)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def pic_cover(slide, name, l=0, t=0, w=None, h=None):
    """Фон на весь слайд (или заданную область)."""
    w = w or SLIDE_W
    h = h or SLIDE_H
    path = asset(name)
    if path.exists():
        return slide.shapes.add_picture(str(path), l, t, width=w, height=h)
    rect(slide, l, t, w, h, BG)
    return None


def slide_number(slide, n: int):
    size = Emu(400000)
    top = Emu(2700000)
    rect(slide, Emu(0), top, size, size, NUM_BG)
    tbox(slide, Emu(0), top, size, size, str(n), size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, title: str, num: int | None = None, subtitle: str = SUBTITLE):
    """Заголовок контентного слайда в стиле СУОТ."""
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    # лёгкая серая «шапка»
    rect(slide, 0, 0, SLIDE_W, Emu(1300000), BG)
    margin = Emu(700000)
    tbox(slide, margin, Emu(220000), Emu(11200000), Emu(550000),
         title, size=20, bold=True, color=TEXT)
    rect(slide, margin, Emu(820000), Emu(8500000), Emu(28000), NUM_BG)
    tbox(slide, margin, Emu(880000), Emu(11000000), Emu(320000),
         subtitle, size=11, color=MUTED)
    if num is not None:
        slide_number(slide, num)


def teal_bubble(slide, l, t, w, h, text, *, size=13):
    """Серое «облачко»-выноска (нейтральный акцент)."""
    sh = round_rect(slide, l, t, w, h, TEAL, line=TEAL_DARK)
    tbox(slide, l + Emu(200000), t + Emu(180000), w - Emu(400000), h - Emu(360000),
         text, size=size, color=TEAL_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    return sh


def panel_right(slide, fill=BG_PANEL):
    """Светло-серая правая панель."""
    return rect(slide, Emu(7200000), Emu(1400000), Emu(5200000), Emu(5200000), fill)


def slide_title_suot(prs, title_lines: str):
    """Титул: фото-фон + тёмная лента + белый CAPS."""
    slide = blank(prs)
    pic_cover(slide, "suot_title_bg.png")
    # затемнение
    overlay = rect(slide, 0, Emu(2600000), SLIDE_W, Emu(2000000), BANNER)
    try:
        from pptx.oxml.ns import qn as _qn
        spPr = overlay._element.spPr
        solid = spPr.find(_qn("a:solidFill"))
        if solid is not None:
            srgb = solid.find(_qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, _qn("a:alpha"))
                alpha.set("val", "85000")  # ~85%
    except Exception:
        pass
    tbox(slide, Emu(800000), Emu(2900000), Emu(10500000), Emu(1400000),
         title_lines, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_thanks_suot(prs):
    """Финал в стиле СУОТ: серый сплит."""
    slide = blank(prs)
    rect(slide, 0, 0, Emu(3600000), SLIDE_H, RGBColor(0x8A, 0x8A, 0x8A))
    rect(slide, Emu(3600000), 0, Emu(8600000), SLIDE_H, BG)
    tbox(slide, Emu(4200000), Emu(3000000), Emu(7200000), Emu(1000000),
         "БЛАГОДАРИМ ЗА ВНИМАНИЕ", size=30, bold=True, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_toc_suot(prs, items, num=2):
    slide = blank(prs)
    content_header(slide, "СОДЕРЖАНИЕ", num)
    y = Emu(1550000)
    for i, item in enumerate(items, 1):
        oval(slide, Emu(700000), y, Emu(520000), Emu(520000), TEAL, line=TEAL_DARK)
        tbox(slide, Emu(700000), y, Emu(520000), Emu(520000), str(i),
             size=15, bold=True, color=TEAL_TEXT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        round_rect(slide, Emu(1400000), y, Emu(10000000), Emu(520000), BG_PANEL)
        tbox(slide, Emu(1650000), y, Emu(9500000), Emu(520000),
             item, size=15, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(850000)
    return slide


def verify(path: Path):
    prs = Presentation(str(path))
    full = pics = texts = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
                # фон титула — ок, но не считаем ошибкой
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts += 1
    print(f"OK: {path.name} · slides={len(prs.slides)} · pics={pics} · texts={texts} · "
          f"{path.stat().st_size // 1024} KB")
    return prs
