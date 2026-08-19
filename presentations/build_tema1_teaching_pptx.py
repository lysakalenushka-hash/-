#!/usr/bin/env python3
"""Build teacher PPTX for Topic 1 (organizational/legal aspects).

Clones the sample deck style (1.1. Проведение занятий.pptx):
title → recommendation slides → thanks.
Content focuses on teaching emphasis points and common mistakes —
not student study text.
"""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parent
SAMPLE = ROOT / "assets" / "sample_style" / "template_1.1_provedenie_zanyatiy.pptx"
OUT = ROOT / "1.1_Organizatsionno-pravovye_aspekty_rekomendatsii.pptx"

# Content slides: (title, attention bullets, error bullets, errors_label)
SLIDES = [
    (
        "ТЕМА 1. ОРГАНИЗАЦИОННО-ПРАВОВЫЕ АСПЕКТЫ",
        [
            "Цель занятия: понятия, НПА, важность ПП, алгоритм и оснащение;",
            "Формат: интерактивная лекция (2 часа), диалог, а не монолог;",
            "Особый акцент — актуальное российское законодательство;",
            "Развеять миф о наказании за оказание первой помощи;",
            "Итог темы: значимость ПП и её законодательно установленный объём.",
        ],
        [
            "Свести занятие к «чтению закона» без разбора страхов очевидцев;",
            "Не связать правовые нормы с действиями на месте происшествия;",
            "Не зафиксировать понимание объёма ПП после разбора темы.",
        ],
        "Возможные ошибки при обучении:",
    ),
    (
        "ВАЖНОСТЬ ПЕРВОЙ ПОМОЩИ",
        [
            "ПП — отдельный вид помощи до медицинской помощи;",
            "Около 25% погибших в ДТП могли выжить при своевременной ПП;",
            "Критическое окно: остановка сердца ~10 мин; сильное кровотечение 1,5–2 мин;",
            "Главные барьеры очевидцев: «сделаю хуже» и страх ответственности;",
            "ПП полезна пострадавшему и юридически защищает оказывающего.",
        ],
        [
            "Надеяться, что «скорая успеет»;",
            "Отказываться помогать из страха юридической ответственности;",
            "Путать первую помощь со скорой медицинской помощью.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ОРГАНИЗАЦИЯ СИСТЕМЫ ПЕРВОЙ ПОМОЩИ В РФ",
        [
            "Пять компонентов системы: НПА, обучение, оснащение, мотивация, учёт;",
            "Три категории участников: обязанные; само-/взаимопомощь; добровольные;",
            "Роль очевидцев до прибытия специалистов и передачи бригаде СМП;",
            "Даже простые действия без аптечки могут спасти жизнь.",
        ],
        [
            "Путать право и обязанность оказывать первую помощь;",
            "Считать, что без аптечки «нельзя ничего делать»;",
            "Забывать про передачу пострадавшего прибывшей бригаде СМП.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "НОРМАТИВНО-ПРАВОВАЯ БАЗА (СТ. 31 323-ФЗ, ПРИКАЗ 220н)",
        [
            "ПП оказывается только по порядкам Минздрава России;",
            "Не по правилам ОТ, памяткам, программам и «личному опыту»;",
            "Приказ № 220н: 9 состояний и 9 мероприятий + последовательность;",
            "С 01.09.2024 — новая редакция ст. 31 (ФЗ № 135-ФЗ);",
            "Показывать НПА и короткие цитаты на слайдах.",
        ],
        [
            "Опираться на устаревший приказ № 477н;",
            "Учить объём ПП «по интернету» вне Порядка;",
            "Подменять Порядок локальными инструкциями и памятками.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ПРАВА, ОБЯЗАННОСТИ И ОТВЕТСТВЕННОСТЬ",
        [
            "«Вправе» — право выбора, а не обязанность для добровольцев;",
            "Обязанные: полиция, пожарные, спасатели и др. по НПА;",
            "Водитель — участник ДТП с пострадавшими обязан принять меры;",
            "Медработник во внерабочее время: право, не обязанность по ПП;",
            "Крайняя необходимость; ПП — смягчающее обстоятельство;",
            "Риск ответственности — при выходе за рамки Порядка.",
        ],
        [
            "Утверждать, что «за любую ошибку при ПП посадят»;",
            "Путать неоказание помощи обязанным лицом и добровольным очевидцем;",
            "Считать медработников всегда обязанными оказывать именно ПП.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "АПТЕЧКИ, УКЛАДКИ И ПОДРУЧНЫЕ СРЕДСТВА",
        [
            "Ключевые аптечки: автомобильная (260н), работников (262н), образования (261н);",
            "Состав обязательный — замена компонентов не допускается;",
            "С 2024 г. в комплектацию по приказам могут входить лекарственные средства;",
            "Подручные средства допустимы;",
            "Пополнение — по расходу и истечении срока годности.",
        ],
        [
            "Заменять компоненты аптечки «на своё усмотрение»;",
            "Использовать медицинскую маску для искусственного дыхания;",
            "Не контролировать доступность и готовность аптечки.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ПОРЯДОК 220н И УНИВЕРСАЛЬНЫЙ АЛГОРИТМ",
        [
            "Соблюдать последовательность мероприятий Приложения № 2;",
            "Приоритет: безопасность → кровотечение → признаки жизни → СЛР;",
            "Признаки жизни: сознание и дыхание (пульс не проверяется);",
            "Нет назначения лекарств; допустимо содействие в приёме своих препаратов;",
            "Устное информирование о начале ПП; учёт выраженного отказа.",
        ],
        [
            "Начинать «лечение» без оценки обстановки и безопасности;",
            "Назначать/давать чужие лекарственные препараты;",
            "Делать проверку пульса обязательным шагом алгоритма.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ОБЕСПЕЧЕНИЕ БЕЗОПАСНЫХ УСЛОВИЙ",
        [
            "Собственная безопасность — условие оказания ПП (п. 4 Порядка);",
            "Типичные угрозы: ток, движение, пожар/взрыв, токсины, агрессия, обрушение;",
            "Сначала устранить/снизить опасность, затем помогать;",
            "Обезопасить пострадавшего и прекратить действие повреждающих факторов;",
            "При сохраняющейся угрозе — вызов спецслужб, а не героизм.",
        ],
        [
            "Забыть устранить опасный фактор;",
            "Бросаться к пострадавшему и становиться вторым пострадавшим;",
            "Игнорировать информирование окружающих о начале оказания ПП.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ПРОФИЛАКТИКА ИНФЕКЦИОННЫХ ЗАБОЛЕВАНИЙ",
        [
            "Обязательно: перчатки и устройство «Рот–Устройство–Рот»;",
            "Маска из аптечки — для защиты оказывающего, не для ИВЛ;",
            "При контакте с биожидкостями — сразу промыть кожу/руки;",
            "После ИВЛ рекомендуется прополоскать рот.",
        ],
        [
            "Пренебрегать средствами индивидуальной защиты;",
            "Путать назначение маски и устройства для искусственного дыхания;",
            "Игнорировать контакт с кровью и другими биожидкостями.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ИЗВЛЕЧЕНИЕ И ПЕРЕМЕЩЕНИЕ ПОСТРАДАВШИХ",
        [
            "Экстренное извлечение — только при угрозе жизни/невозможности помочь на месте;",
            "Перед извлечением из автомобиля: стояночный тормоз, выключить зажигание;",
            "Приём Раутека; при подозрении на травму ШОП — фиксация головы и шеи;",
            "Способ перемещения выбирать по состоянию, травмам и числу участников;",
            "Приоритет при нескольких пострадавших: тяжёлые и дети.",
        ],
        [
            "Извлекать «на всякий случай» без реальной угрозы;",
            "Не фиксировать голову/шею при подозрении на травму позвоночника;",
            "Выбирать травмоопасный способ переноски «для скорости».",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ВЫЗОВ СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ",
        [
            "Единый номер 112 (также 103 и региональные номера);",
            "Сообщить: место и что произошло; число пострадавших; повреждения; какая помощь оказывается;",
            "Вызов — после оценки состояния, чтобы точнее передать информацию;",
            "Не класть трубку первым: дождаться подтверждения принятия вызова;",
            "Диспетчер может давать команды по оказанию ПП.",
        ],
        [
            "Забыть вызвать скорую медицинскую помощь;",
            "Не сообщить число пострадавших и характер повреждений;",
            "Положить трубку до подтверждения, что вызов принят.",
        ],
        "Возможные ошибки слушателей:",
    ),
    (
        "ТИПИЧНЫЕ ОШИБКИ ПРИ ОБУЧЕНИИ ТЕМЕ",
        [
            "Держать фокус на актуальных НПА и коротких цитатах;",
            "Интерактивно разбирать страхи («сделаю хуже», «посадят»);",
            "Связывать право → алгоритм → оснащение → вызов СМП;",
            "Закреплять контрольными вопросами из учебного пособия;",
            "Подчеркнуть итог: ПП проста, законна и спасает жизни.",
        ],
        [
            "Уходить в «юридический ликбез» без связи с практикой;",
            "Пугать ответственностью вместо развенчания мифа;",
            "Давать устаревшие перечни состояний и мероприятий;",
            "Пропускать безопасность, инфекцию и вызов СМП как «очевидное».",
        ],
        "Возможные ошибки при обучении:",
    ),
]


def _set_run_text(run, text: str) -> None:
    run.text = text


def _clear_paragraphs_keep_first(tf):
    """Remove extra paragraphs from a text frame, keep first empty."""
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r.text = ""
    # remove additional paragraphs via XML
    p_elements = tf._txBody.findall(qn("a:p"))
    for pe in p_elements[1:]:
        tf._txBody.remove(pe)


def _add_paragraph(tf, text: str, *, bold: bool = False, size_pt: float = 18, space_before_pt: float = 0):
    p = tf.add_paragraph()
    p.level = 0
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.name = "Open Sans Light"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    # Force latin/ea fonts via XML for consistent look
    rPr = run._r.get_or_add_rPr()
    for tag, typeface in (
        ("a:latin", "Open Sans Light"),
        ("a:ea", "Open Sans Light"),
        ("a:cs", "Open Sans Light"),
    ):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", typeface)
    return p


def _fill_content_textbox(
    shape,
    title: str,
    attention: list[str],
    errors: list[str] | None,
    errors_label: str = "Возможные ошибки слушателей:",
):
    tf = shape.text_frame
    tf.word_wrap = True
    _clear_paragraphs_keep_first(tf)

    # Use first paragraph for title
    p0 = tf.paragraphs[0]
    if not p0.runs:
        run = p0.add_run()
    else:
        run = p0.runs[0]
        for extra in p0.runs[1:]:
            extra.text = ""
    run.text = title
    run.font.name = "Open Sans Light"
    run.font.size = Pt(18)
    run.font.bold = True
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", "Open Sans Light")

    _add_paragraph(tf, "", bold=False, size_pt=18)
    _add_paragraph(tf, "Важно заострить внимание:", bold=True, size_pt=18)
    for item in attention:
        _add_paragraph(tf, item, bold=False, size_pt=18, space_before_pt=2.8)

    if errors:
        _add_paragraph(tf, "", bold=False, size_pt=18)
        _add_paragraph(tf, errors_label, bold=True, size_pt=18)
        for item in errors:
            _add_paragraph(tf, item, bold=False, size_pt=18, space_before_pt=2.8)


def _find_content_textbox(slide):
    """Largest text box that is not the header."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text or ""
        if "РЕКОМЕНДАЦИИ" in text:
            continue
        if shape.name and "Номер" in shape.name:
            continue
        candidates.append(shape)
    if not candidates:
        return None
    # Prefer name object 7 / largest area
    for shape in candidates:
        if shape.name == "object 7":
            return shape
    return max(candidates, key=lambda s: s.width * s.height)


def _set_notes(slide, text: str):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def _delete_slide(prs: Presentation, index: int):
    """Delete slide at index (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def build() -> Path:
    if not SAMPLE.exists():
        raise FileNotFoundError(f"Sample PPTX not found: {SAMPLE}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SAMPLE, OUT)
    prs = Presentation(str(OUT))

    # Slide 1 title — keep visual; update notes
    _set_notes(
        prs.slides[0],
        "Тема 1. Организационно-правовые аспекты оказания первой помощи\n"
        "Презентация для преподавателей: на что обратить внимание и типичные ошибки.",
    )

    # Content slides in sample: indices 1..15 (slides 2–16), thanks is last (16)
    # We need: title + len(SLIDES) content + thanks
    content_count = len(SLIDES)
    # Sample has 15 content-ish slides (2-16) then thanks (17) = indices 1..15 content, 16 thanks
    # We'll rewrite first content_count slides (1..content_count) and delete extras before thanks

    thanks_index = len(prs.slides) - 1

    # Ensure we have enough content slots between title and thanks
    available = thanks_index - 1  # slides strictly between title and thanks
    if content_count > available:
        raise RuntimeError(
            f"Need {content_count} content slides, sample has only {available}"
        )

    for i, (title, attention, errors, errors_label) in enumerate(SLIDES):
        slide = prs.slides[1 + i]
        box = _find_content_textbox(slide)
        if box is None:
            raise RuntimeError(f"No content textbox on slide {i + 2}")
        _fill_content_textbox(box, title, attention, errors, errors_label)
        _set_notes(
            slide,
            f"{title}\n\nВажно:\n- "
            + "\n- ".join(attention)
            + (
                (f"\n\n{errors_label}\n- " + "\n- ".join(errors)) if errors else ""
            ),
        )

    # Delete unused content slides between last used content and thanks (from the end)
    # After rewrite, unused indices are: 1+content_count .. thanks_index-1
    for idx in range(thanks_index - 1, content_count, -1):
        _delete_slide(prs, idx)

    # Update thanks notes
    thanks = prs.slides[len(prs.slides) - 1]
    _set_notes(thanks, "Благодарим за внимание")

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
