#!/usr/bin/env python3
"""Build teacher PPTX for Topic 2 (external bleeding).

Style/layout from sample «1.1. Проведение занятий.pptx».
Images from Topic 2.1–2.5 decks. Content = teaching focus points
and typical mistakes — not student study text.
"""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.oxml.ns import qn
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent
SAMPLE = ROOT / "assets" / "sample_style" / "template_1.1_provedenie_zanyatiy.pptx"
ASSETS = ROOT / "assets" / "tema2_teaching"
OUT = ROOT / "2_Naruzhnye_krovotecheniya_rekomendatsii.pptx"

# (title, attention, errors, errors_label, image_paths)
# image_paths: 1–2 files to replace picture fills on the slide (primary, optional secondary)
SLIDES = [
    (
        "ТЕМА 2. НАРУЖНЫЕ КРОВОТЕЧЕНИЯ",
        [
            "Цель: виды кровотечений и навыки временной остановки;",
            "Формат: 4 часа (2 теория + 2 практика), демонстрация и задачи;",
            "Учесть непереносимость вида крови у части обучающихся;",
            "Донести: сильное кровотечение — одна из главных причин гибели на месте;",
            "Приоритет способов: прямое давление → повязка; жгут — крайняя мера.",
        ],
        [
            "Игнорировать реакцию на вид крови и не адаптировать практику;",
            "Сразу учить жгут, минуя прямое давление и давящую повязку;",
            "Не связать теорию с отработкой навыка на практике.",
        ],
        "Возможные ошибки при обучении:",
        ["hero_bleed.png", "overview_diagram.png"],
    ),
    (
        "ОПАСНОСТЬ КРОВОПОТЕРИ",
        [
            "Развести понятия «кровотечение» и «острая кровопотеря»;",
            "Признаки кровопотери возможны и без видимой крови;",
            "Интенсивное кровотечение даёт минуты, а не «время до скорой»;",
            "Ориентир на месте — интенсивность, а не точный «вид сосуда».",
        ],
        [
            "Ждать «классический фонтан», игнорируя промокшую одежду и лужу крови;",
            "Считать, что без явной струи кровь «неопасна»;",
            "Путать внутреннюю кровопотерю с «отсутствием проблемы».",
        ],
        "Возможные ошибки слушателей:",
        ["blood_loss_signs.jpeg", "blood_loss2.jpeg"],
    ),
    (
        "ПРИЗНАКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ",
        [
            "Артериальное / венозное / капиллярное — кратко, по признакам;",
            "На практике вид определить сложно — останавливать по интенсивности;",
            "Смешанные кровотечения (в т.ч. при отрыве) особенно опасны;",
            "Не тратить время на «точный диагноз вида» вместо остановки.",
        ],
        [
            "Зацикливаться на определении вида вместо действий;",
            "Недооценивать венозное кровотечение («ручьём»);",
            "Считать капиллярное всегда «несерьёзным» без осмотра.",
        ],
        "Возможные ошибки слушателей:",
        ["arterial.png", "venous.png"],
    ),
    (
        "ОБЗОРНЫЙ ОСМОТР ПОСТРАДАВШЕГО",
        [
            "Цель осмотра — быстро найти угрожающее жизни кровотечение;",
            "Признаки: промокшая одежда, лужа крови, видимая струя/рана;",
            "Осмотр — за несколько секунд, с головы до ног;",
            "Обнаружили — сразу к остановке всеми доступными способами;",
            "Сначала безопасность места происшествия.",
        ],
        [
            "Долго «диагностировать» вместо быстрого осмотра;",
            "Пропустить кровотечение под одеждой / сзади;",
            "Начать помощь, не оценив безопасность.",
        ],
        "Возможные ошибки слушателей:",
        ["overview_diagram.png", "overview_exam.jpeg"],
    ),
    (
        "ПРЯМОЕ ДАВЛЕНИЕ НА РАНУ",
        [
            "Самый простой и приоритетный способ при интенсивном кровотечении;",
            "Давление — через салфетки/бинт/ткань, в перчатках;",
            "Сила достаточна, чтобы остановить кровь — не «слегка прижать»;",
            "Если сработало — далее давящая повязка;",
            "Если нет повязки/жгута — продолжать давление до СМП.",
        ],
        [
            "Слабое давление «для вида», кровь продолжает течь;",
            "Сразу хвататься за жгут при доступном давлении на рану;",
            "Забыть перчатки / защиту от контакта с кровью.",
        ],
        "Возможные ошибки слушателей:",
        ["direct_pressure_hero.png", "direct_pressure.jpeg"],
    ),
    (
        "ДАВЯЩАЯ ПОВЯЗКА",
        [
            "Задача повязки — остановить кровь: накладывать с усилием;",
            "На рану — салфетки/бинт/ткань, затем тугие туры;",
            "Слабо промокает — ещё одна давящая повязка сверху;",
            "Быстро промокает / неэффективна — переход к жгуту;",
            "Закрепить свободный конец; контролировать эффект.",
        ],
        [
            "Накладывать «слабо», без давления на рану;",
            "Считать любую повязку достаточной без контроля промокания;",
            "Снимать промокшую повязку вместо усиления / смены тактики.",
        ],
        "Возможные ошибки слушателей:",
        ["pressure_bandage.jpg", "direct_pressure.jpeg"],
    ),
    (
        "ИНОРОДНОЕ ТЕЛО И ОТКРЫТЫЙ ПЕРЕЛОМ",
        [
            "Инородный предмет из раны не извлекать;",
            "Фиксировать предмет салфетками/бинттами и накладывать повязку;",
            "Прямое давление на предмет / отломки — опасно или неэффективно;",
            "В этих случаях — повязка с фиксацией и/или жгут по показаниям.",
        ],
        [
            "Пытаться вытащить осколок / предмет из раны;",
            "Давить прямо на инородное тело;",
            "Тратить время на «идеальное» давление там, где оно противопоказано.",
        ],
        "Возможные ошибки слушателей:",
        ["foreign_body_bandage.jpg", "tourniquet_when.png"],
    ),
    (
        "КРОВООСТАНАВЛИВАЮЩИЙ ЖГУТ",
        [
            "Жгут — когда давление/повязка невозможны, неэффективны или при отрыве;",
            "Только на конечность; выше раны (5–7 см) или максимально к туловищу;",
            "Обычно на прокладку/одежду; турникет — по инструкции (иногда на голое тело);",
            "Остановка — первым растянутым туром; жгут на виду; время указано;",
            "Срок до 2 часов; снятие >2 ч вне медорганизации не рекомендуется.",
        ],
        [
            "Накладывать жгут «на всякий случай» при останавливаемом давлении;",
            "Прятать жгут под повязкой/одеждой; не указывать время;",
            "Использовать тонкую проволоку/шнур как импровизированный жгут;",
            "Снимать жгут самостоятельно после долгого срока наложения.",
        ],
        "Возможные ошибки слушателей:",
        ["tourniquet1.png", "tourniquet2.png"],
    ),
    (
        "ПОСЛЕДОВАТЕЛЬНОСТЬ ОСТАНОВКИ КРОВОТЕЧЕНИЯ",
        [
            "Безопасность → обзорный осмотр → выбор способа;",
            "Интенсивное кровотечение → сначала прямое давление;",
            "Давление нельзя/опасно → повязка и/или жгут сразу;",
            "Отрыв / разрушение конечности → жгут немедленно;",
            "Давление помогло → повязка; не помогло → жгут выше раны.",
        ],
        [
            "Нарушать порядок и начинать с жгута без показаний;",
            "Останавливаться после «слабой» попытки давления;",
            "Забывать контроль: остановилась ли кровь после приёма.",
        ],
        "Возможные ошибки слушателей:",
        ["sequence.png", "sequence_photo.jpeg"],
    ),
    (
        "ГОЛОВА, ШЕЯ, ГРУДЬ, ЖИВОТ",
        [
            "Голова: давление/повязка; при риске перелома черепа — без усиленного давления на кости;",
            "Инородное в ране головы — фиксировать, не извлекать;",
            "Шея: давление → повязка через противоположную подмышку;",
            "Грудь/спина: остановить наружное; внутреннее крупных сосудов — только СМП/хирургия;",
            "Живот/таз: не давить на выпавшие органы давящей повязкой.",
        ],
        [
            "Сильно давить при подозрении на перелом черепа;",
            "Накладывать круговую повязку на шею «как на конечность»;",
            "Пытаться «остановить» внутреннее кровотечение груди подручными средствами;",
            "Давить/бинтовать выпавшие органы живота.",
        ],
        "Возможные ошибки слушателей:",
        ["head.jpg", "neck_bandage.jpg"],
    ),
    (
        "КОНЕЧНОСТИ, ОТРЫВ, СМЕЖНЫЕ ЗОНЫ",
        [
            "Конечности: все способы по общей последовательности;",
            "Выбор: интенсивность, место раны, наличие средств, срок СМП;",
            "Отрыв части конечности — показание к немедленному жгуту;",
            "Смежные зоны (пах, подмышка и т.п.): часто только прямое давление;",
            "Жгут/повязку в смежных зонах наложить сложно — не терять время.",
        ],
        [
            "Пытаться наложить жгут там, где анатомически нельзя;",
            "При отрыве долго пробовать только давление;",
            "Игнорировать сильное кровотечение из паха/подмышки.",
        ],
        "Возможные ошибки слушателей:",
        ["limb.jpg", "zones.jpeg"],
    ),
    (
        "ТИПИЧНЫЕ ОШИБКИ ПРИ ОБУЧЕНИИ ТЕМЕ",
        [
            "Сначала отработать давление и повязку, жгут — осознанно как крайнюю меру;",
            "Обязательно отработать жгут на практике (навык критичен);",
            "Разбирать непереносимость вида крови и технику безопасности;",
            "Закреплять алгоритм последовательности, а не «набор приёмов»;",
            "Контролировать: перчатки, время жгута, видимость жгута, фиксация инородного тела.",
        ],
        [
            "Свести практику к «намотали жгут» без сценариев выбора способа;",
            "Пугать жгутом или, наоборот, подавать его как универсальный приём;",
            "Не проверять указание времени и видимость жгута;",
            "Пропустить особенности областей тела (шея, живот, смежные зоны).",
        ],
        "Возможные ошибки при обучении:",
        ["methods_end.png", "improvised_ok.png"],
    ),
]


def _clear_paragraphs_keep_first(tf):
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r.text = ""
    for pe in tf._txBody.findall(qn("a:p"))[1:]:
        tf._txBody.remove(pe)


def _force_font(run, name: str = "Open Sans Light"):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def _add_paragraph(tf, text: str, *, bold: bool = False, size_pt: float = 17, space_before_pt: float = 0):
    p = tf.add_paragraph()
    p.level = 0
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    _force_font(run)
    return p


def _fill_content_textbox(shape, title: str, attention: list[str], errors: list[str] | None, errors_label: str):
    tf = shape.text_frame
    tf.word_wrap = True
    _clear_paragraphs_keep_first(tf)

    p0 = tf.paragraphs[0]
    run = p0.runs[0] if p0.runs else p0.add_run()
    for extra in p0.runs[1:]:
        extra.text = ""
    run.text = title
    run.font.size = Pt(17)
    run.font.bold = True
    _force_font(run)

    _add_paragraph(tf, "")
    _add_paragraph(tf, "Важно заострить внимание:", bold=True)
    for item in attention:
        _add_paragraph(tf, item, space_before_pt=2.4)

    if errors:
        _add_paragraph(tf, "")
        _add_paragraph(tf, errors_label, bold=True)
        for item in errors:
            _add_paragraph(tf, item, space_before_pt=2.4)


def _find_content_textbox(slide):
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text or ""
        if "РЕКОМЕНДАЦИИ" in text:
            continue
        if shape.name and "Номер" in shape.name:
            continue
        candidates.append(shape)
    for shape in candidates:
        if shape.name == "object 7":
            return shape
    return max(candidates, key=lambda s: s.width * s.height) if candidates else None


def _picture_fill_shapes(slide):
    out = []
    for shape in slide.shapes:
        try:
            if shape.fill.type == MSO_FILL.PICTURE:
                out.append(shape)
        except Exception:
            continue
    # Stable order: larger first
    out.sort(key=lambda s: s.width * s.height, reverse=True)
    return out


def _set_picture_fill(slide, shape, image_path: Path):
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    _part, rId = slide.part.get_or_add_image_part(str(image_path))
    for blip in shape._element.findall(".//" + qn("a:blip")):
        blip.set(qn("r:embed"), rId)


def _set_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def _delete_slide(prs: Presentation, index: int):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def build() -> Path:
    if not SAMPLE.exists():
        raise FileNotFoundError(SAMPLE)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SAMPLE, OUT)
    prs = Presentation(str(OUT))

    _set_notes(
        prs.slides[0],
        "Тема 2. Оказание первой помощи при наружных кровотечениях\n"
        "Презентация для преподавателей: акценты и типичные ошибки.",
    )

    content_count = len(SLIDES)
    thanks_index = len(prs.slides) - 1
    available = thanks_index - 1
    if content_count > available:
        raise RuntimeError(f"Need {content_count} content slides, sample has {available}")

    for i, (title, attention, errors, errors_label, images) in enumerate(SLIDES):
        slide = prs.slides[1 + i]
        box = _find_content_textbox(slide)
        if box is None:
            raise RuntimeError(f"No content textbox on slide {i + 2}")
        _fill_content_textbox(box, title, attention, errors, errors_label)

        pic_shapes = _picture_fill_shapes(slide)
        for j, img_name in enumerate(images):
            if j >= len(pic_shapes):
                break
            _set_picture_fill(slide, pic_shapes[j], ASSETS / img_name)

        # Hide unused second picture shape if we only provided one image
        # (keep layout clean)
        if len(images) == 1 and len(pic_shapes) > 1:
            for extra in pic_shapes[1:]:
                extra._element.getparent().remove(extra._element)

        _set_notes(
            slide,
            f"{title}\n\nВажно:\n- "
            + "\n- ".join(attention)
            + (f"\n\n{errors_label}\n- " + "\n- ".join(errors) if errors else ""),
        )

    for idx in range(thanks_index - 1, content_count, -1):
        _delete_slide(prs, idx)

    _set_notes(prs.slides[len(prs.slides) - 1], "Благодарим за внимание")
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
