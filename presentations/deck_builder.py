#!/usr/bin/env python3
"""Универсальный сборщик презентаций в legal-стиле из структурированного контента."""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from PIL import Image as PILImage
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

from legal_style import (
    ACCENT_RED, ASSETS, LINE, MUTED, NUM_BG, OUT, ROOT, TEXT, WHITE,
    blank, bullets, content_header, cream_note, new_prs, oval, rect,
    set_context, slide_thanks, slide_title, slide_toc, tbox, verify,
)

INFO = ROOT / "infographics"
INFO.mkdir(parents=True, exist_ok=True)
TEMA4_ASSETS = ASSETS / "tema4"


def _resolve_image(name: str) -> Path | None:
    if not name:
        return None
    candidates = [
        ASSETS / name,
        TEMA4_ASSETS / name,
        TEMA4_ASSETS / Path(name).name,
        INFO / name,
        Path("/opt/cursor/artifacts/assets") / name,
    ]
    for p in candidates:
        if p.exists():
            return p
    return None


def _pic(slide, name: str, l, t, max_w, max_h):
    path = _resolve_image(name)
    if path is None:
        rect(slide, l, t, max_w, max_h, WHITE, line=LINE)
        tbox(slide, l, t, max_w, max_h, f"[{name}]", size=11, color=MUTED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def _slide_intro(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    text = s["text"]
    size = 15 if len(text) < 450 else (13 if len(text) < 800 else 12)
    bottom = Emu(5600000) if s.get("note") else Emu(6400000)
    tbox(slide, Emu(700000), Emu(1500000), Emu(11000000), bottom - Emu(1500000),
         text, size=size, color=TEXT)
    if s.get("note"):
        cream_note(slide, Emu(700000), Emu(5600000), Emu(11000000), Emu(1000000),
                   s["note"], size=12)
    return slide


def _slide_bullets(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    items = s["items"]
    alert = set(s.get("alert") or [])
    img = s.get("image")
    if img:
        _pic(slide, img, Emu(500000), Emu(1450000), Emu(5200000), Emu(4000000))
        bullets(slide, Emu(6000000), Emu(1450000), Emu(6200000), Emu(3800000),
                items, size=13, alert=alert)
        if s.get("note"):
            cream_note(slide, Emu(6000000), Emu(5400000), Emu(6200000), Emu(1100000),
                       s["note"], size=12)
    else:
        size = 14 if len(items) <= 6 else 12
        bullets(slide, Emu(700000), Emu(1450000), Emu(11000000), Emu(4000000),
                items, size=size, alert=alert)
        if s.get("note"):
            cream_note(slide, Emu(700000), Emu(5600000), Emu(11000000), Emu(1000000),
                       s["note"], size=12)
    return slide


def _slide_alert_bullets(prs, s, num):
    s = dict(s)
    if "alert" not in s:
        s["alert"] = [0] if s.get("items") else []
    return _slide_bullets(prs, s, num)


def _slide_image_bullets(prs, s, num):
    s = dict(s)
    s["type"] = "bullets"
    return _slide_bullets(prs, s, num)


def _slide_image_text(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    _pic(slide, s.get("image", ""), Emu(500000), Emu(1450000), Emu(5200000), Emu(4000000))
    text = s.get("text", "")
    size = 14 if len(text) < 500 else 12
    tbox(slide, Emu(6000000), Emu(1500000), Emu(6200000), Emu(3500000),
         text, size=size, color=TEXT)
    if s.get("note"):
        cream_note(slide, Emu(6000000), Emu(5200000), Emu(6200000), Emu(1300000),
                   s["note"], size=12)
    return slide


def _slide_cards(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    cards = s["cards"]
    n = len(cards)
    if n <= 3:
        cols, rows = n, 1
    elif n == 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 3
    gap_x, gap_y = Emu(200000), Emu(160000)
    top = Emu(1400000)
    left = Emu(700000)
    avail_w = Emu(11000000)
    avail_h = Emu(5000000)
    cw = int((avail_w - gap_x * (cols - 1)) / cols)
    rh = int((avail_h - gap_y * (rows - 1)) / rows)
    for i, c in enumerate(cards):
        r, col = divmod(i, cols)
        if r >= rows:
            break
        x = left + col * (cw + gap_x)
        y = top + r * (rh + gap_y)
        rect(slide, x, y, cw, rh, WHITE, line=LINE)
        oval(slide, x + Emu(120000), y + Emu(150000), Emu(450000), Emu(450000), NUM_BG)
        tbox(slide, x + Emu(120000), y + Emu(150000), Emu(450000), Emu(450000),
             str(c.get("n", i + 1)), size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, x + Emu(650000), y + Emu(180000), cw - Emu(800000), Emu(400000),
             c["title"], size=12, bold=True, color=TEXT)
        desc_size = 11 if len(c.get("desc", "")) > 120 else 12
        tbox(slide, x + Emu(120000), y + Emu(700000), cw - Emu(240000), rh - Emu(850000),
             c["desc"], size=desc_size, color=MUTED)
    return slide


def _slide_steps(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    steps = s["steps"]
    y = Emu(1400000)
    avail = Emu(4300000) if s.get("note") else Emu(5000000)
    h = min(Emu(700000), int(avail / max(len(steps), 1)))
    step_size = 13 if len(steps) <= 6 else 11
    for i, step in enumerate(steps, 1):
        oval(slide, Emu(700000), y, Emu(500000), Emu(500000),
             ACCENT_RED if i == 1 else NUM_BG)
        tbox(slide, Emu(700000), y, Emu(500000), Emu(500000), str(i),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, Emu(1400000), y, Emu(10550000), max(Emu(500000), h - Emu(80000)),
             WHITE, line=LINE)
        tbox(slide, Emu(1600000), y, Emu(10150000), max(Emu(500000), h - Emu(80000)),
             step, size=step_size, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += h
    if s.get("note"):
        cream_note(slide, Emu(700000), Emu(5800000), Emu(11000000), Emu(900000),
                   s["note"], size=12)
    return slide


def _slide_two_col(prs, s, num):
    slide = blank(prs)
    content_header(slide, s["title"], num)
    tbox(slide, Emu(700000), Emu(1400000), Emu(5400000), Emu(400000),
         s.get("left_title", ""), size=14, bold=True, color=TEXT)
    tbox(slide, Emu(6600000), Emu(1400000), Emu(5400000), Emu(400000),
         s.get("right_title", ""), size=14, bold=True, color=TEXT)
    left_items = s.get("left", [])
    right_items = s.get("right", [])
    size = 12 if max(len(left_items), len(right_items)) > 5 else 13
    bullets(slide, Emu(700000), Emu(1850000), Emu(5400000), Emu(4000000),
            left_items, size=size)
    bullets(slide, Emu(6600000), Emu(1850000), Emu(5400000), Emu(4000000),
            right_items, size=size)
    if s.get("note"):
        cream_note(slide, Emu(700000), Emu(5800000), Emu(11000000), Emu(900000),
                   s["note"], size=12)
    return slide


def _slide_summary(prs, points, num):
    slide = blank(prs)
    content_header(slide, "ГЛАВНОЕ ЗАПОМНИТЬ", num)
    n = len(points)
    h = min(Emu(850000), int(Emu(5000000) / max(n, 1)))
    y = Emu(1450000)
    for pt in points:
        rect(slide, Emu(700000), y, Emu(11000000), h, WHITE, line=LINE)
        rect(slide, Emu(700000), y, Emu(90000), h, ACCENT_RED)
        tbox(slide, Emu(1000000), y + Emu(80000), Emu(10400000), h - Emu(160000),
             pt, size=12 if n > 5 else 13, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += h + Emu(80000)
    return slide


def _slide_memo(prs, memo_name, num):
    slide = blank(prs)
    content_header(slide, "ПАМЯТКА", num)
    path = _resolve_image(memo_name) or (INFO / memo_name if (INFO / memo_name).exists() else None)
    if path is None:
        cream_note(slide, Emu(700000), Emu(2500000), Emu(11000000), Emu(1500000),
                   "Памятка будет добавлена рядом с презентацией.", size=14)
        return slide
    with PILImage.open(path) as im:
        iw, ih = im.size
    max_w, max_h = Emu(11000000), Emu(5200000)
    l, t = Emu(700000), Emu(1400000)
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    return slide


HANDLERS = {
    "intro": _slide_intro,
    "bullets": _slide_bullets,
    "alert_bullets": _slide_alert_bullets,
    "image_bullets": _slide_image_bullets,
    "image_text": _slide_image_text,
    "cards": _slide_cards,
    "steps": _slide_steps,
    "two_col": _slide_two_col,
}


def build_deck(deck: dict) -> Path:
    prs = new_prs()
    slide_title(prs, deck["title"])
    toc = deck.get("toc") or []
    # TOC layout supports up to 5 rows comfortably
    slide_toc(prs, toc[:5], 2)
    num = 3
    for s in deck["slides"]:
        handler = HANDLERS.get(s["type"])
        if handler is None:
            raise ValueError(f"Unknown slide type: {s['type']} in {deck['id']}")
        handler(prs, s, num)
        num += 1
    _slide_summary(prs, deck["memo_points"], num)
    num += 1
    if deck.get("memo_file"):
        _slide_memo(prs, deck["memo_file"], num)
        num += 1
    slide_thanks(prs)

    name = deck["filename"]
    path = OUT / name
    prs.save(path)
    verify(path)
    path2 = ROOT / name
    prs.save(path2)
    print("Saved:", path2)
    return path2


def build_all(decks, *, subtitle: str, course: str):
    set_context(subtitle, course)
    paths = []
    for deck in decks:
        paths.append(build_deck(deck))
    return paths
