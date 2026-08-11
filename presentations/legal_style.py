#!/usr/bin/env python3
"""Стиль «Организационно-правовые аспекты оказания первой помощи»."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path("/workspace/presentations")
ASSETS = ROOT / "assets"
OUT = ROOT / "out"
OUT.mkdir(parents=True, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF8, 0xF8)
BG_BAR = RGBColor(221, 221, 221)  # RGB(221)
LINE = RGBColor(0x99, 0x99, 0x99)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
NUM_BG = RGBColor(0x55, 0x55, 0x55)
CREAM = RGBColor(0xF5, 0xF0, 0xE1)
ACCENT_RED = RGBColor(0xCC, 0x00, 0x00)
TABLE_HDR = RGBColor(0x44, 0x44, 0x44)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)

FONT = "Open Sans"
SUBTITLE = "Оказание первой помощи при наружных кровотечениях"
COURSE = "Тема 2 · Оказание первой помощи при наружных кровотечениях"


def set_context(subtitle: str, course: str) -> None:
    """Переключает подзаголовок и строку курса для серии слайдов."""
    global SUBTITLE, COURSE
    SUBTITLE = subtitle
    COURSE = course


def asset(name: str) -> Path:
    return ASSETS / name


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", FONT)


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def tbox(slide, l, t, w, h, text, *, size=18, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def rich_tbox(slide, l, t, w, h, parts, *, size=16, color=TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold in parts:
        r = p.add_run()
        r.text = text
        font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, *, size=15, marker="•", alert=None):
    alert = alert or set()
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        font(r, size, bold=(i in alert), color=ACCENT_RED if i in alert else TEXT)
    return box


def pic_fit(slide, name, l, t, max_w, max_h):
    from PIL import Image as PILImage
    path = asset(name)
    if not path.exists():
        rect(slide, l, t, max_w, max_h, WHITE, line=LINE)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def slide_number(slide, n: int):
    size = Emu(420000)
    rect(slide, Emu(0), Emu(2800000), size, size, NUM_BG)
    tbox(slide, Emu(0), Emu(2800000), size, size, str(n), size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, title: str, num: int | None = None):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    m = Emu(700000)
    tbox(slide, m, Emu(250000), Emu(11200000), Emu(550000),
         title.upper(), size=22, bold=True, color=TEXT)
    rect(slide, m, Emu(850000), Emu(9000000), Emu(28000), NUM_BG)
    tbox(slide, m, Emu(920000), Emu(11000000), Emu(350000),
         SUBTITLE, size=12, color=MUTED)
    if num is not None:
        slide_number(slide, num)


def slide_title(prs, title: str):
    """Титул: серые панели + CAPS + линия."""
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    rect(slide, Emu(8800000), 0, Emu(3400000), SLIDE_H, BG_BAR)
    rect(slide, 0, Emu(2800000), Emu(2800000), Emu(1800000), BG_BAR)
    tbox(slide, Emu(700000), Emu(3000000), Emu(7800000), Emu(1500000),
         title, size=26, bold=True, color=TEXT)
    rect(slide, Emu(700000), Emu(4700000), Emu(5200000), Emu(30000), TEXT)
    tbox(slide, Emu(700000), Emu(4900000), Emu(7500000), Emu(450000),
         COURSE, size=13, color=MUTED)
    return slide


def slide_thanks(prs):
    slide = blank(prs)
    rect(slide, 0, 0, Emu(3400000), SLIDE_H, BG_BAR)
    rect(slide, Emu(3400000), 0, Emu(8800000), SLIDE_H, BG)
    tbox(slide, Emu(4000000), Emu(3000000), Emu(7500000), Emu(1000000),
         "БЛАГОДАРИМ ЗА ВНИМАНИЕ", size=30, bold=True, color=TEXT,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Emu(4000000), Emu(4100000), Emu(4500000), Emu(30000), TEXT)
    return slide


def slide_toc(prs, items, num=2):
    slide = blank(prs)
    content_header(slide, "СОДЕРЖАНИЕ", num)
    y = Emu(1550000)
    for i, item in enumerate(items, 1):
        oval(slide, Emu(700000), y, Emu(520000), Emu(520000), NUM_BG)
        tbox(slide, Emu(700000), y, Emu(520000), Emu(520000), str(i),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, Emu(1400000), y, Emu(10000000), Emu(520000), WHITE, line=LINE)
        tbox(slide, Emu(1600000), y, Emu(9600000), Emu(520000),
             item, size=15, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(850000)
    return slide


def cream_note(slide, l, t, w, h, text, *, size=14):
    round_rect(slide, l, t, w, h, CREAM)
    tbox(slide, l + Emu(250000), t + Emu(200000), w - Emu(500000), h - Emu(400000),
         text, size=size, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)


def verify(path: Path):
    prs = Presentation(str(path))
    pics = texts = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts += 1
    print(f"OK: {path.name} · slides={len(prs.slides)} · pics={pics} · texts={texts} · "
          f"{path.stat().st_size // 1024} KB")
    return prs
