#!/usr/bin/env python3
"""
Презентация: «Признаки наружного кровотечения и кровопотери»

ДОПОЛНЯЕТ (не дублирует) презентацию
«Кровотечение. Обзорный осмотр пострадавшего»:
- там: понятие, виды «как выглядят», обзорный осмотр, действия;
- здесь: углублённое распознавание признаков кровопотери и интенсивности,
  скрытая кровопотеря, сравнительная таблица, смешанное кровотечение,
  маскирующие ситуации.

Оформление — как в «Организационно-правовые аспекты оказания первой помощи».
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path("/workspace/presentations")
ASSETS = ROOT / "assets"
OUT = ROOT / "out"
OUT.mkdir(parents=True, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF8, 0xF8)
BG_BAR = RGBColor(221, 221, 221)
LINE = RGBColor(0x99, 0x99, 0x99)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED = RGBColor(0xED, 0x1C, 0x24)
CREAM = RGBColor(0xF5, 0xF0, 0xE1)
NUM_BG = RGBColor(0x55, 0x55, 0x55)
TABLE_HDR = RGBColor(0x44, 0x44, 0x44)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE = RGBColor(0x00, 0x55, 0xAA)

FONT = "Open Sans"
SUBTITLE = "Оказание первой помощи при наружных кровотечениях"
COURSE = "Тема 2 · Оказание первой помощи при наружных кровотечениях"


def asset(name: str) -> Path:
    return ASSETS / name


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", FONT)


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def tbox(slide, l, t, w, h, text, *, size=20, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, *, size=16, marker="•", alert=None):
    alert = alert or set()
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        font(r, size, bold=(i in alert), color=ACCENT_RED if i in alert else TEXT)
    return box


def slide_number(slide, n: int):
    size = Emu(420000)
    top = Emu(2800000)
    rect(slide, Emu(0), top, size, size, NUM_BG)
    tbox(slide, Emu(0), top, size, size, str(n), size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, title: str, num: int | None = None):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    margin = Emu(700000)
    tbox(slide, margin, Emu(280000), Emu(11000000), Emu(550000),
         title, size=22, bold=True, color=TEXT)
    rect(slide, margin, Emu(880000), Emu(9000000), Emu(25000), NUM_BG)
    tbox(slide, margin, Emu(950000), Emu(11000000), Emu(350000),
         SUBTITLE, size=12, color=MUTED)
    if num is not None:
        slide_number(slide, num)


def slide_title(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_LIGHT)
    rect(slide, Emu(8800000), 0, Emu(3400000), SLIDE_H, BG_BAR)
    rect(slide, Emu(0), Emu(2800000), Emu(2800000), Emu(1800000), BG_BAR)
    tbox(slide, Emu(700000), Emu(3000000), Emu(7800000), Emu(1400000),
         "ПРИЗНАКИ НАРУЖНОГО\nКРОВОТЕЧЕНИЯ И КРОВОПОТЕРИ",
         size=28, bold=True, color=TEXT)
    rect(slide, Emu(700000), Emu(4600000), Emu(5500000), Emu(30000), TEXT)
    tbox(slide, Emu(700000), Emu(4800000), Emu(7500000), Emu(500000),
         COURSE, size=14, color=MUTED)
    return slide


def slide_thanks(prs):
    slide = blank(prs)
    rect(slide, 0, 0, Emu(3200000), SLIDE_H, BG_BAR)
    rect(slide, Emu(3200000), 0, Emu(9000000), SLIDE_H, BG_LIGHT)
    tbox(slide, Emu(3800000), Emu(3000000), Emu(7500000), Emu(1000000),
         "БЛАГОДАРИМ ЗА ВНИМАНИЕ", size=32, bold=True, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Emu(3800000), Emu(4100000), Emu(4500000), Emu(30000), TEXT)
    return slide


def slide_toc(prs, items, num=2):
    slide = blank(prs)
    content_header(slide, "СОДЕРЖАНИЕ", num)
    y = Emu(1500000)
    for i, item in enumerate(items, 1):
        oval(slide, Emu(700000), y, Emu(500000), Emu(500000), NUM_BG)
        tbox(slide, Emu(700000), y, Emu(500000), Emu(500000), str(i),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, Emu(1400000), y, Emu(10000000), Emu(500000),
             item, size=16, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(780000)
    return slide


def slide_bridge(prs, num=3):
    """Связь с презентацией про обзорный осмотр — без повтора её содержания."""
    slide = blank(prs)
    content_header(slide, "ЧЕМ ЭТА ТЕМА ДОПОЛНЯЕТ ОБЗОРНЫЙ ОСМОТР", num)
    # left — уже изучено
    rect(slide, Emu(700000), Emu(1500000), Emu(5400000), Emu(4800000), WHITE, line=LINE)
    rect(slide, Emu(700000), Emu(1500000), Emu(5400000), Emu(700000), BG_BAR)
    tbox(slide, Emu(900000), Emu(1600000), Emu(5000000), Emu(500000),
         "Уже в теме «Обзорный осмотр»", size=15, bold=True, color=TEXT,
         align=PP_ALIGN.CENTER)
    bullets(slide, Emu(950000), Emu(2450000), Emu(4900000), Emu(3500000), [
        "понятие кровотечения;",
        "как выглядят виды кровотечения;",
        "как провести обзорный осмотр;",
        "что делать при обнаружении крови.",
    ], size=15)
    # right — здесь
    rect(slide, Emu(6500000), Emu(1500000), Emu(5400000), Emu(4800000), WHITE, line=LINE)
    rect(slide, Emu(6500000), Emu(1500000), Emu(5400000), Emu(700000), NUM_BG)
    tbox(slide, Emu(6700000), Emu(1600000), Emu(5000000), Emu(500000),
         "Разбираем сейчас", size=15, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    bullets(slide, Emu(6750000), Emu(2450000), Emu(4900000), Emu(3500000), [
        "как читать признаки кровопотери;",
        "когда крови «не видно», а угроза есть;",
        "как оценить интенсивность;",
        "сравнение видов и смешанное кровотечение;",
        "что маскирует признаки.",
    ], size=15)
    return slide


def slide_signs_groups(prs, num=4):
    """Углублённый разбор признаков кровопотери по группам."""
    slide = blank(prs)
    content_header(slide, "ПРИЗНАКИ КРОВОПОТЕРИ — ПО ГРУППАМ", num)
    groups = [
        ("Самочувствие", [
            "резкая общая слабость;",
            "чувство жажды;",
            "головокружение;",
            "мелькание «мушек» перед глазами;",
            "обморок (чаще при попытке встать).",
        ]),
        ("Кожа", [
            "бледная;",
            "влажная;",
            "холодная на ощупь.",
        ]),
        ("Дыхание и сердце", [
            "учащённое сердцебиение;",
            "частое дыхание.",
        ]),
    ]
    x = Emu(700000)
    widths = [Emu(4000000), Emu(3200000), Emu(3400000)]
    for (title, items), w in zip(groups, widths):
        rect(slide, x, Emu(1500000), w, Emu(4800000), WHITE, line=LINE)
        rect(slide, x, Emu(1500000), w, Emu(700000), NUM_BG)
        tbox(slide, x + Emu(100000), Emu(1600000), w - Emu(200000), Emu(500000),
             title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bullets(slide, x + Emu(150000), Emu(2450000), w - Emu(300000), Emu(3500000),
                items, size=14, marker="☐")
        x += w + Emu(250000)
    return slide


def slide_when_signs(prs, num=5):
    slide = blank(prs)
    content_header(slide, "КОГДА ПОЯВЛЯЮТСЯ ПРИЗНАКИ КРОВОПОТЕРИ", num)
    cards = [
        ("Кровотечение продолжается",
         "Признаки нарастают на фоне видимой крови."),
        ("Кровотечение уже остановлено",
         "Признаки могут сохраняться — потерянная кровь не возвращается сразу."),
        ("Видимой крови нет",
         "Возможно внутреннее (скрытое) кровотечение — оценивайте состояние."),
    ]
    y = Emu(1550000)
    for title, text in cards:
        rect(slide, Emu(700000), y, Emu(11000000), Emu(1300000), WHITE, line=LINE)
        rect(slide, Emu(700000), y, Emu(90000), Emu(1300000), ACCENT_RED)
        tbox(slide, Emu(1000000), y + Emu(150000), Emu(10400000), Emu(400000),
             title, size=16, bold=True, color=TEXT)
        tbox(slide, Emu(1000000), y + Emu(600000), Emu(10400000), Emu(500000),
             text, size=14, color=MUTED)
        y += Emu(1500000)
    return slide


def slide_hidden(prs, num=6):
    slide = blank(prs)
    content_header(slide, "СКРЫТАЯ (ВНУТРЕННЯЯ) КРОВОПОТЕРЯ", num)
    tbox(slide, Emu(700000), Emu(1500000), Emu(11000000), Emu(900000),
         "Снаружи крови может не быть видно, но состояние ухудшается. "
         "Ориентируйтесь на признаки кровопотери и механизм травмы.",
         size=16, color=TEXT)
    bullets(slide, Emu(700000), Emu(2600000), Emu(11000000), Emu(2800000), [
        "Настораживающие ситуации: удар в живот / грудь, падение с высоты, ДТП.",
        "Боль и напряжение живота, нарастающая слабость, холодный пот, жажда.",
        "На месте происшествия внутреннее кровотечение не останавливают — "
        "нужны вызов СМП, покой, положение, контроль состояния.",
    ], size=15)
    round_rect(slide, Emu(700000), Emu(5600000), Emu(11000000), Emu(900000), CREAM)
    tbox(slide, Emu(950000), Emu(5750000), Emu(10500000), Emu(600000),
         "Нет видимой крови ≠ нет угрозы. Смотрите на состояние пострадавшего.",
         size=15, bold=True, color=ACCENT_RED, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_intensity_scale(prs, num=7):
    """Практическая оценка интенсивности — дополнение к обзорному осмотру."""
    slide = blank(prs)
    content_header(slide, "КАК ОЦЕНИТЬ ИНТЕНСИВНОСТЬ КРОВОТЕЧЕНИЯ", num)
    levels = [
        (ACCENT_BLUE, "Слабая",
         "Капиллярное сочение, небольшие ссадины/порезы. "
         "Угрозы жизни обычно нет, но остановить нужно."),
        (NUM_BG, "Средняя",
         "Стойкий «ручей», пропитывание повязки/ткани. "
         "Без остановки возможна значительная кровопотеря."),
        (ACCENT_RED, "Интенсивная",
         "Струя / быстрое пропитывание одежды, лужа крови, "
         "рана с сильным истечением — останавливать немедленно."),
    ]
    y = Emu(1500000)
    for color, title, text in levels:
        oval(slide, Emu(700000), y + Emu(150000), Emu(550000), Emu(550000), color)
        tbox(slide, Emu(700000), y + Emu(150000), Emu(550000), Emu(550000),
             title[0], size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, Emu(1450000), y, Emu(10000000), Emu(400000),
             title, size=16, bold=True, color=TEXT)
        tbox(slide, Emu(1450000), y + Emu(450000), Emu(10000000), Emu(700000),
             text, size=14, color=MUTED)
        y += Emu(1450000)
    return slide


def slide_compare_table(prs, num=8):
    slide = blank(prs)
    content_header(slide, "СРАВНЕНИЕ ПРИЗНАКОВ ПО ВИДАМ", num)
    tbox(slide, Emu(700000), Emu(1400000), Emu(11000000), Emu(400000),
         "Детализация к уже известным видам (артериальное / венозное / капиллярное).",
         size=13, color=MUTED)
    headers = ["Признак", "Артериальное", "Венозное", "Капиллярное"]
    rows = [
        ["Цвет", "Алый, яркий", "Тёмный, вишнёвый", "Красный"],
        ["Характер", "Пульсирующая струя", "Равномерный «ручей»", "Сочится"],
        ["Скорость потери", "Очень быстрая", "Умеренная / быстрая", "Медленная"],
        ["Опасность", "Критическая — минуты", "Высокая", "Обычно низкая"],
        ["Типичная картина", "Лужа алого цвета, одежда быстро мокнет",
         "Стойкое истечение без пульсации", "Ссадина, порез, царапина"],
    ]
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Emu(700000), Emu(1900000), Emu(11000000), Emu(4300000)
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                font(r, 12, True, WHITE)
    for ri, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
            cell.text = val
            alert = any(k in val for k in ("Алый", "Пульсир", "Критическ", "Очень"))
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    font(r, 11, alert, ACCENT_RED if alert else TEXT)
    return slide


def slide_mixed(prs, num=9):
    slide = blank(prs)
    content_header(slide, "СМЕШАННОЕ КРОВОТЕЧЕНИЕ — КОГДА ПОДОЗРЕВАТЬ", num)
    bullets(slide, Emu(700000), Emu(1550000), Emu(11000000), Emu(2800000), [
        "Одновременно признаки артериального, венозного и капиллярного кровотечения.",
        "Частый пример — отрыв (ампутация) конечности или обширное размозжение.",
        "Опасно из‑за артериального компонента: быстрая массивная кровопотеря.",
        "Не тратьте время на «точную классификацию» — останавливайте по интенсивности.",
    ], size=15)
    round_rect(slide, Emu(700000), Emu(4700000), Emu(11000000), Emu(1500000), CREAM)
    rect(slide, Emu(700000), Emu(4700000), Emu(100000), Emu(1500000), ACCENT_RED)
    tbox(slide, Emu(1000000), Emu(4950000), Emu(10400000), Emu(1000000),
         "Правило: есть интенсивное истечение крови — действуйте как при угрозе жизни.",
         size=16, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_masking(prs, num=10):
    slide = blank(prs)
    content_header(slide, "ЧТО МАСКИРУЕТ ПРИЗНАКИ КРОВОТЕЧЕНИЯ", num)
    items = [
        ("Тёмная / плотная одежда",
         "Кровь плохо видна — ощупайте, приподнимите одежду при осмотре."),
        ("Кровь под пострадавшим",
         "Лужа может быть сзади / под телом — осмотрите вокруг."),
        ("Холод, дождь, грязь",
         "Сложно оценить цвет кожи и объём крови — опирайтесь на слабость, пульс, дыхание."),
        ("Несколько пострадавших",
         "Интенсивное кровотечение ищите в первую очередь у каждого."),
    ]
    y = Emu(1450000)
    for title, text in items:
        rect(slide, Emu(700000), y, Emu(11000000), Emu(1100000), WHITE, line=LINE)
        tbox(slide, Emu(950000), y + Emu(120000), Emu(10400000), Emu(350000),
             title, size=15, bold=True, color=TEXT)
        tbox(slide, Emu(950000), y + Emu(500000), Emu(10400000), Emu(450000),
             text, size=13, color=MUTED)
        y += Emu(1200000)
    return slide


def slide_volume(prs, num=11):
    """Ориентиры по объёму — дополнение, которого нет в презентации про осмотр."""
    slide = blank(prs)
    content_header(slide, "ОРИЕНТИРЫ ПО ОБЪЁМУ КРОВОПОТЕРИ", num)
    headers = ["Объём (ориентир)", "Доля ОЦК", "Типичное состояние"]
    rows = [
        ["До ~500 мл", "≈ 10%", "Слабость, жажда"],
        ["500–1000 мл", "≈ 10–20%", "Бледность, учащённый пульс"],
        ["1000–1500 мл", "≈ 20–30%", "Обмороки, холодный пот"],
        [">1500–2000 мл", "> 30%", "Шок, угроза жизни"],
    ]
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Emu(700000), Emu(1550000), Emu(11000000), Emu(3800000)
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                font(r, 13, True, WHITE)
    for ri, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
            cell.text = val
            alert = ">" in val or "Шок" in val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    font(r, 13, alert, ACCENT_RED if alert else TEXT)
    tbox(slide, Emu(700000), Emu(5600000), Emu(11000000), Emu(800000),
         "На месте точный объём измерить нельзя — таблица помогает понять тяжесть "
         "по признакам состояния пострадавшего.",
         size=14, color=MUTED)
    return slide


def slide_summary(prs, num=12):
    slide = blank(prs)
    content_header(slide, "ГЛАВНОЕ ЗАПОМНИТЬ", num)
    points = [
        "Признаки кровопотери смотрите по самочувствию, коже, дыханию и пульсу.",
        "Признаки возможны и без видимой крови — думайте о скрытой кровопотере.",
        "Интенсивность важнее «точного вида сосуда» для решения об остановке.",
        "Тёмная одежда, положение тела и погода могут маскировать кровь — осматривайте внимательно.",
    ]
    y = Emu(1500000)
    for pt in points:
        rect(slide, Emu(700000), y, Emu(11000000), Emu(1000000), WHITE, line=LINE)
        rect(slide, Emu(700000), y, Emu(90000), Emu(1000000), ACCENT_RED)
        tbox(slide, Emu(1000000), y + Emu(150000), Emu(10400000), Emu(700000),
             pt, size=14, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(1150000)
    return slide


def verify(path: Path):
    prs = Presentation(str(path))
    full = pics = texts = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
                if sh.width >= int(SLIDE_W * 0.95) and sh.height >= int(SLIDE_H * 0.95):
                    full += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts += 1
    if full:
        raise SystemExit(f"FAIL: full-slide images = {full}")
    print(f"OK: {path.name} · slides={len(prs.slides)} · pics={pics} · texts={texts} · "
          f"{path.stat().st_size // 1024} KB")


def build():
    prs = new_prs()
    slide_title(prs)
    slide_toc(prs, [
        "Чем тема дополняет обзорный осмотр",
        "Признаки кровопотери по группам",
        "Когда признаки есть без видимой крови",
        "Оценка интенсивности и сравнение видов",
        "Смешанное кровотечение, маскировка, объём",
    ], 2)
    slide_bridge(prs, 3)
    slide_signs_groups(prs, 4)
    slide_when_signs(prs, 5)
    slide_hidden(prs, 6)
    slide_intensity_scale(prs, 7)
    slide_compare_table(prs, 8)
    slide_mixed(prs, 9)
    slide_masking(prs, 10)
    slide_volume(prs, 11)
    slide_summary(prs, 12)
    slide_thanks(prs)

    name = "Признаки_наружного_кровотечения_и_кровопотери.pptx"
    path = OUT / name
    prs.save(path)
    verify(path)
    path2 = ROOT / name
    prs.save(path2)
    print(f"Saved: {path2}")
    return path2


if __name__ == "__main__":
    build()
