#!/usr/bin/env python3
"""
Презентация: «Кровотечение. Обзорный осмотр пострадавшего (пострадавших)»

Оформление — как в «Организационно-правовые аспекты оказания первой помощи»:
- 16:9, белый фон, заголовок CAPS + линия + подзаголовок
- номер слайда в сером квадрате слева
- редактируемый текст + отдельные Picture-объекты

Содержание — по учебному пособию (Тема 2).
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path("/workspace/presentations")
ASSETS = ROOT / "assets"
OUT = ROOT / "out"
OUT.mkdir(parents=True, exist_ok=True)

# 16:9 как в эталонном PDF (1920×1080)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF8, 0xF8)
BG_BAR = RGBColor(221, 221, 221)  # RGB(221) из эталона
LINE = RGBColor(0x99, 0x99, 0x99)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED = RGBColor(0xED, 0x1C, 0x24)
CREAM = RGBColor(0xF5, 0xF0, 0xE1)
NUM_BG = RGBColor(0x55, 0x55, 0x55)

FONT = "Open Sans"
SUBTITLE = "Оказание первой помощи при наружных кровотечениях"
COURSE = "Тема 2 · Оказание первой помощи при наружных кровотечениях"


def asset(name: str) -> Path:
    return ASSETS / name


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", FONT)


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def tbox(slide, l, t, w, h, text, *, size=20, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def rich_tbox(slide, l, t, w, h, parts, *, size=20, color=TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold in parts:
        r = p.add_run()
        r.text = text
        font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, *, size=18, marker="•", alert=None):
    alert = alert or set()
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        font(r, size, bold=(i in alert), color=ACCENT_RED if i in alert else TEXT)
    return box


def pic_fit(slide, name, l, t, max_w, max_h):
    from PIL import Image as PILImage
    path = asset(name)
    if not path.exists():
        rect(slide, l, t, max_w, max_h, WHITE, line=LINE)
        tbox(slide, l, t + max_h // 2 - Emu(200000), max_w, Emu(400000),
             f"[нет файла: {name}]", size=12, color=MUTED, align=PP_ALIGN.CENTER)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def slide_number(slide, n: int):
    """Серый квадрат с номером — как в эталоне."""
    size = Emu(420000)
    top = Emu(2800000)
    left = Emu(0)
    rect(slide, left, top, size, size, NUM_BG)
    tbox(slide, left, top, size, size, str(n), size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, title: str, num: int | None = None):
    """Заголовок контентного слайда в стиле эталона."""
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    margin = Emu(700000)
    tbox(slide, margin, Emu(280000), Emu(11000000), Emu(550000),
         title, size=26, bold=True, color=TEXT)
    # линия под заголовком
    rect(slide, margin, Emu(880000), Emu(9000000), Emu(25000), NUM_BG)
    tbox(slide, margin, Emu(950000), Emu(11000000), Emu(350000),
         SUBTITLE, size=12, color=MUTED)
    if num is not None:
        slide_number(slide, num)


# ───────── слайды ─────────

def slide_title(prs):
    """Титул — как в эталоне: серые панели + CAPS + линия."""
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_LIGHT)
    # правая серая полоса
    rect(slide, Emu(8800000), 0, Emu(3400000), SLIDE_H, BG_BAR)
    # декоративный прямоугольник слева
    rect(slide, Emu(0), Emu(2800000), Emu(2800000), Emu(1800000), BG_BAR)
    tbox(slide, Emu(700000), Emu(3000000), Emu(7800000), Emu(1400000),
         "КРОВОТЕЧЕНИЕ.\nОБЗОРНЫЙ ОСМОТР ПОСТРАДАВШЕГО\n(ПОСТРАДАВШИХ)",
         size=28, bold=True, color=TEXT)
    rect(slide, Emu(700000), Emu(4600000), Emu(5500000), Emu(30000), TEXT)
    tbox(slide, Emu(700000), Emu(4800000), Emu(7500000), Emu(500000),
         COURSE, size=14, color=MUTED)
    return slide


def slide_definition(prs, num=3):
    """Понятие кровотечения — персонаж + облачко + признаки (как стиль эталона)."""
    slide = blank(prs)
    content_header(slide, "ПОНЯТИЕ «КРОВОТЕЧЕНИЕ»", num)
    pic_fit(slide, "def_man.png", Emu(500000), Emu(1500000),
            Emu(3200000), Emu(5200000))
    round_rect(slide, Emu(3400000), Emu(1700000), Emu(4200000), Emu(2600000), CREAM)
    rich_tbox(
        slide, Emu(3650000), Emu(1900000), Emu(3700000), Emu(2200000),
        [
            ("Под ", False),
            ("кровотечением понимают", True),
            (" ситуацию, когда кровь по разным причинам "
             "(чаще всего в результате травмы) покидает сосудистое русло, "
             "что приводит к кровопотере — безвозвратной утрате части крови.", False),
        ],
        size=14, color=TEXT,
    )
    tbox(slide, Emu(7900000), Emu(1600000), Emu(4800000), Emu(450000),
         "Основные признаки острой кровопотери:", size=16, bold=True, color=TEXT)
    bullets(slide, Emu(7900000), Emu(2150000), Emu(4800000), Emu(4500000), [
        "резкая общая слабость;",
        "чувство жажды;",
        "головокружение;",
        "мелькание «мушек» перед глазами;",
        "обморок (чаще при попытке встать);",
        "бледная, влажная и холодная кожа;",
        "учащённое сердцебиение;",
        "частое дыхание.",
    ], size=14, marker="☐")
    return slide


def slide_blood_loss_note(prs, num=4):
    slide = blank(prs)
    content_header(slide, "ОСТРАЯ МАССИВНАЯ КРОВОПОТЕРЯ", num)
    round_rect(slide, Emu(700000), Emu(1600000), Emu(11000000), Emu(1600000), CREAM)
    tbox(slide, Emu(950000), Emu(1800000), Emu(10500000), Emu(1200000),
         "Наиболее опасным является интенсивное кровотечение, приводящее "
         "к быстрой потере большого количества крови, — острая массивная кровопотеря.",
         size=18, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    bullets(slide, Emu(700000), Emu(3600000), Emu(11000000), Emu(3000000), [
        "При повреждении крупных сосудов без остановки кровотечения гибель "
        "может наступить в течение нескольких минут.",
        "Признаки кровопотери возможны и при уже остановленном кровотечении, "
        "и при отсутствии видимой крови (внутреннее кровотечение).",
        "Кровотечения слабой и средней интенсивности тоже нужно останавливать — "
        "чтобы не допустить поздних осложнений.",
    ], size=16)
    return slide


def slide_bleed_types(prs, num=5):
    slide = blank(prs)
    content_header(slide, "ВИДЫ НАРУЖНОГО КРОВОТЕЧЕНИЯ", num)
    tbox(slide, Emu(700000), Emu(1450000), Emu(11000000), Emu(500000),
         "Наружное кровотечение: кровь изливается наружу при повреждении кожи и слизистых.",
         size=15, color=MUTED)
    cols = [
        ("bleed_arterial.png", "Артериальное",
         "Пульсирующая алая струя; быстро растекающаяся лужа; "
         "одежда быстро пропитывается кровью. Наиболее опасно."),
        ("bleed_venous.png", "Венозное",
         "Кровь тёмно-вишнёвая, вытекает «ручьём». "
         "Скорость кровопотери меньше, но остановка обязательна."),
        ("bleed_capillary.png", "Капиллярное",
         "Ссадины, порезы, царапины. "
         "Как правило, непосредственной угрозы жизни не представляет."),
    ]
    x = Emu(700000)
    cw = Emu(3600000)
    gap = Emu(300000)
    for img, title, desc in cols:
        tbox(slide, x, Emu(2050000), cw, Emu(450000), title, size=16, bold=True,
             color=TEXT, align=PP_ALIGN.CENTER)
        pic_fit(slide, img, x + Emu(150000), Emu(2550000), cw - Emu(300000), Emu(2800000))
        tbox(slide, x, Emu(5500000), cw, Emu(1300000), desc, size=12, color=TEXT,
             align=PP_ALIGN.CENTER)
        x += cw + gap
    return slide


def slide_intensity(prs, num=6):
    slide = blank(prs)
    content_header(slide, "ОРИЕНТИР — ИНТЕНСИВНОСТЬ КРОВОТЕЧЕНИЯ", num)
    bullets(slide, Emu(700000), Emu(1600000), Emu(11000000), Emu(2500000), [
        "Смешанные кровотечения (артериальное + венозное + капиллярное) "
        "часто бывают при отрыве конечности и опасны из‑за артериального компонента.",
        "При оказании первой помощи вид кровотечения определить сложно.",
        "Ориентируйтесь на интенсивность и останавливайте кровотечение "
        "любым доступным способом или их комбинацией.",
    ], size=17)
    round_rect(slide, Emu(700000), Emu(4800000), Emu(11000000), Emu(1400000), CREAM)
    rect(slide, Emu(700000), Emu(4800000), Emu(100000), Emu(1400000), ACCENT_RED)
    tbox(slide, Emu(1000000), Emu(5050000), Emu(10400000), Emu(900000),
         "Приоритет: угрожающее жизни кровотечение останавливают немедленно.",
         size=18, bold=True, color=ACCENT_RED, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_overview_goal(prs, num=7):
    slide = blank(prs)
    content_header(slide, "ЦЕЛЬ ОБЗОРНОГО ОСМОТРА", num)
    tbox(slide, Emu(700000), Emu(1600000), Emu(11000000), Emu(1200000),
         "Обзорный осмотр пострадавшего (пострадавших) проводится прежде всего "
         "для определения наличия и расположения ранений с интенсивным наружным "
         "кровотечением, требующим немедленной остановки.",
         size=18, color=TEXT)
    tbox(slide, Emu(700000), Emu(3100000), Emu(11000000), Emu(450000),
         "Признаки интенсивного наружного кровотечения:", size=16, bold=True, color=TEXT)
    bullets(slide, Emu(700000), Emu(3650000), Emu(11000000), Emu(2500000), [
        "одежда, пропитанная кровью;",
        "скопление значительного количества крови на земле возле пострадавшего;",
        "видимые раны с интенсивно вытекающей из них кровью.",
    ], size=17, marker="☐")
    return slide


def slide_overview_how(prs, num=8):
    """Ключевой слайд: картинка + текст «очень быстро… с головы до ног»."""
    slide = blank(prs)
    content_header(slide, "ОБЗОРНЫЙ ОСМОТР ПОСТРАДАВШЕГО", num)
    # слева — иллюстрация (отдельный объект)
    pic_fit(slide, "overview_exam.png", Emu(500000), Emu(1500000),
            Emu(6000000), Emu(5200000))
    # справа — редактируемый текст
    rich_tbox(
        slide, Emu(6800000), Emu(2400000), Emu(5800000), Emu(2200000),
        [
            ("Обзорный осмотр", True),
            (" производится очень быстро, в течение ", False),
            ("нескольких секунд", True),
            (", с головы до ног.", False),
        ],
        size=20, color=TEXT,
    )
    tbox(slide, Emu(6800000), Emu(4800000), Emu(5800000), Emu(1600000),
         "Сразу после обнаружения признаков угрожающего жизни кровотечения "
         "приступают к его остановке всеми доступными способами.",
         size=15, color=MUTED)
    return slide


def slide_actions(prs, num=9):
    slide = blank(prs)
    content_header(slide, "ДЕЙСТВИЯ ПРИ ОБНАРУЖЕНИИ КРОВОТЕЧЕНИЯ", num)
    steps = [
        "Оценить обстановку и безопасность.",
        "Провести обзорный осмотр — найти интенсивное наружное кровотечение.",
        "Немедленно начать временную остановку кровотечения.",
        "Прямое давление на рану → давящая повязка → жгут (по показаниям).",
        "Вызвать скорую медицинскую помощь (если ещё не вызвана).",
        "Контролировать состояние пострадавшего до прибытия СМП.",
    ]
    y = Emu(1550000)
    for i, step in enumerate(steps, 1):
        oval(slide, Emu(700000), y, Emu(500000), Emu(500000),
             ACCENT_RED if i <= 3 else NUM_BG)
        tbox(slide, Emu(700000), y, Emu(500000), Emu(500000), str(i),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, Emu(1400000), y, Emu(10000000), Emu(500000),
             step, size=16, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(750000)
    return slide


def slide_compare(prs, num=10):
    slide = blank(prs)
    content_header(slide, "ОБЗОРНЫЙ И ПОДРОБНЫЙ ОСМОТР", num)
    # left
    rect(slide, Emu(700000), Emu(1550000), Emu(5400000), Emu(4800000), WHITE, line=LINE)
    rect(slide, Emu(700000), Emu(1550000), Emu(5400000), Emu(700000), NUM_BG)
    tbox(slide, Emu(900000), Emu(1650000), Emu(5000000), Emu(500000),
         "Обзорный осмотр", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(slide, Emu(950000), Emu(2500000), Emu(4900000), Emu(3500000), [
        "Цель: найти наружное кровотечение",
        "Время: несколько секунд",
        "Объём: быстро с головы до ног",
        "Когда: сразу после оценки обстановки",
    ], size=15)
    # right
    rect(slide, Emu(6500000), Emu(1550000), Emu(5400000), Emu(4800000), WHITE, line=LINE)
    rect(slide, Emu(6500000), Emu(1550000), Emu(5400000), Emu(700000), BG_BAR)
    tbox(slide, Emu(6700000), Emu(1650000), Emu(5000000), Emu(500000),
         "Подробный осмотр", size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    bullets(slide, Emu(6750000), Emu(2500000), Emu(4900000), Emu(3500000), [
        "Цель: выявить травмы и другие угрозы",
        "Время: несколько минут",
        "Объём: голова → шея → грудь → … → руки",
        "Когда: после остановки кровотечения",
    ], size=15)
    return slide


def slide_summary(prs, num=11):
    slide = blank(prs)
    content_header(slide, "ГЛАВНОЕ ЗАПОМНИТЬ", num)
    points = [
        "Кровотечение — выход крови из сосудов с риском кровопотери.",
        "Обзорный осмотр — быстрый поиск интенсивного наружного кровотечения.",
        "Осмотр — за несколько секунд, с головы до ног.",
        "Сначала останавливаем угрожающее жизни кровотечение, затем — подробный осмотр.",
    ]
    y = Emu(1600000)
    for pt in points:
        rect(slide, Emu(700000), y, Emu(11000000), Emu(1000000), WHITE, line=LINE)
        rect(slide, Emu(700000), y, Emu(90000), Emu(1000000), ACCENT_RED)
        tbox(slide, Emu(1000000), y + Emu(200000), Emu(10400000), Emu(600000),
             pt, size=16, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(1150000)
    return slide


def slide_thanks(prs):
    """Финал — как в эталоне."""
    slide = blank(prs)
    rect(slide, 0, 0, Emu(3200000), SLIDE_H, BG_BAR)
    rect(slide, Emu(3200000), 0, Emu(9000000), SLIDE_H, BG_LIGHT)
    tbox(slide, Emu(3800000), Emu(3000000), Emu(7500000), Emu(1000000),
         "БЛАГОДАРИМ ЗА ВНИМАНИЕ", size=32, bold=True, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Emu(3800000), Emu(4100000), Emu(4500000), Emu(30000), TEXT)
    return slide


def verify(path: Path):
    prs = Presentation(str(path))
    full = pics = texts = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
                if sh.width >= int(SLIDE_W * 0.95) and sh.height >= int(SLIDE_H * 0.95):
                    full += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts += 1
    if full:
        raise SystemExit(f"FAIL: full-slide images = {full}")
    print(f"OK: {path.name} · slides={len(prs.slides)} · pics={pics} · texts={texts} · "
          f"{path.stat().st_size // 1024} KB")


def build():
    prs = new_prs()
    slide_title(prs)

    # TOC without using broken helper
    slide = blank(prs)
    content_header(slide, "СОДЕРЖАНИЕ", 2)
    items = [
        "Понятие «кровотечение» и признаки кровопотери",
        "Виды наружного кровотечения",
        "Цель и признаки для обзорного осмотра",
        "Как проводится обзорный осмотр",
        "Действия при обнаружении кровотечения",
    ]
    y = Emu(1600000)
    for i, item in enumerate(items, 1):
        oval(slide, Emu(700000), y, Emu(550000), Emu(550000), NUM_BG)
        tbox(slide, Emu(700000), y, Emu(550000), Emu(550000), str(i),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, Emu(1450000), y, Emu(10000000), Emu(550000),
             item, size=18, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(900000)

    slide_definition(prs, 3)
    slide_blood_loss_note(prs, 4)
    slide_bleed_types(prs, 5)
    slide_intensity(prs, 6)
    slide_overview_goal(prs, 7)
    slide_overview_how(prs, 8)
    slide_actions(prs, 9)
    slide_compare(prs, 10)
    slide_summary(prs, 11)
    slide_thanks(prs)

    path = OUT / "Кровотечение_Обзорный_осмотр_пострадавшего.pptx"
    prs.save(path)
    verify(path)

    # также в корень presentations
    path2 = ROOT / path.name
    prs.save(path2)
    print(f"Saved: {path}")
    print(f"Saved: {path2}")
    return path


if __name__ == "__main__":
    build()
