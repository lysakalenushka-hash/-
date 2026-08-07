#!/usr/bin/env python3
"""Добавить недостающие слайды темы 1: оглавление, терминология аптечек, приоритетность, резюме."""

from copy import deepcopy
from pathlib import Path
import shutil
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt, Inches
from lxml import etree
import openpyxl

COMBINED = Path(
    "/workspace/first_aid_tema1_new/"
    "1 Организационно-правовые аспекты оказания первой помощи.pptx"
)
OUT_COMBINED_ROOT = Path(
    "/workspace/Презентация_Тема_1_Первая_помощь_новая_структура.pptx"
)
INDIV_DIR = Path("/workspace/first_aid_tema1_new")
ZIP_OUT = Path("/workspace/Тема1_первая_помощь_1.1-1.7.zip")
CHANGE_XLSX = Path("/workspace/Таблица_изменений_Тема1_первая_помощь.xlsx")
NOTES = Path("/workspace/Презентация_Тема_1_Первая_помощь_новая_структура.md")

DARK = RGBColor(0x1A, 0x1A, 0x1A)
RED = RGBColor(0xE3, 0x06, 0x13)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF7, 0xF8)

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def set_run(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_multiline(slide, left, top, width, height, lines, size=16, color=DARK, spacing=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if spacing:
            p.space_after = Pt(8)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=False, color=color)
    return box


def clear_slide_shapes(slide):
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def add_blank_content_slide(prs):
    """Append blank slide using first available layout; clear default placeholders."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    clear_slide_shapes(slide)
    # light background rectangle
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    # red accent bar left
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(180000), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    return slide


def style_content_slide(slide, page_num, title, eyebrow, body_lines, footer_note=None):
    # page number
    add_textbox(slide, Emu(220000), Emu(6540000), Emu(900000), Emu(600000), str(page_num), size=22, bold=True, color=GRAY)
    # title
    add_textbox(slide, Emu(2050000), Emu(1600000), Emu(21000000), Emu(900000), title, size=28, bold=True, color=DARK)
    # red underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(2060000), Emu(2550000), Emu(12000000), Emu(70000))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    # eyebrow
    add_textbox(slide, Emu(2050000), Emu(2700000), Emu(16000000), Emu(500000), eyebrow, size=14, bold=False, color=GRAY)
    # body panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(2050000), Emu(3600000), Emu(20000000), Emu(8500000)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.fill.background()
    add_multiline(slide, Emu(2400000), Emu(3900000), Emu(19000000), Emu(7800000), body_lines, size=17, color=DARK)
    if footer_note:
        add_textbox(slide, Emu(2050000), Emu(12400000), Emu(20000000), Emu(600000), footer_note, size=12, color=GRAY)


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order len {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in new_items:
        sldIdLst.append(el)


def renumber_pages(prs):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            # left page placeholder style
            if (sh.left or 0) < Emu(1200000) and (sh.width or 0) < Emu(1500000) and (sh.top or 0) > Emu(5000000):
                # clear a:t then set
                for t_el in sh.element.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                    t_el.text = ""
                p0 = sh.text_frame.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = str(i)
                else:
                    p0.add_run().text = str(i)


TOC_LINES = [
    "1.1. Нормативно-правовая база и организация оказания первой помощи",
    "1.2. Укладки, наборы, комплекты и аптечки для оказания первой помощи",
    "1.3. Порядок и приоритетность оказания первой помощи",
    "1.4. Перечень состояний и мероприятий первой помощи; последовательность выполнения",
    "1.5. Обеспечение безопасных условий и профилактика инфекций",
    "1.6. Извлечение пострадавших из труднодоступных мест и перемещение",
    "1.7. Правила вызова скорой медицинской помощи и других специальных служб",
]

KITS_LINES = [
    "В новой программе (Прил. 2 к Порядку № 2464) используется терминология:",
    "• аптечки первой помощи — готовые комплекты по назначению (автомобильная; для работников);",
    "• укладки и наборы — оснащение для оказания первой помощи по видам деятельности;",
    "• комплекты медицинских изделий — средства и устройства по приказу Минздрава РФ.",
    "",
    "Типовой состав (ориентир): средства остановки кровотечения, перевязочный материал,",
    "средства индивидуальной защиты оказывающего помощь, устройства для проведения ИВЛ,",
    "покрывала/термоодеяла и иные изделия по утверждённым требованиям к аптечкам.",
    "",
    "Важно: применять только изделия из аптечки/укладки по инструкции; лекарственные",
    "препараты в аптечки первой помощи для работников/автомобильные не входят.",
]

PRIORITY_LINES = [
    "Приоритет — состояния, непосредственно угрожающие жизни пострадавшего.",
    "",
    "1. Оценить обстановку и обеспечить безопасность себе и пострадавшему.",
    "2. Определить признаки жизни; при их отсутствии — начать СЛР, параллельно вызвать СМП.",
    "3. Восстановить проходимость дыхательных путей при их нарушении.",
    "4. Остановить угрожающее жизни наружное кровотечение.",
    "5. Затем выполнить остальные мероприятия первой помощи по состоянию пострадавшего.",
    "",
    "При нескольких пострадавших в первую очередь помогают тем, у кого есть угроза жизни",
    "(отсутствие сознания/дыхания, массивное кровотечение). Не перемещать пострадавшего",
    "без необходимости, кроме угрозы на месте происшествия.",
]

SUMMARY_LINES = [
    "• Оказание первой помощи регулируется НПА (ФЗ-323, ТК РФ, приказ Минздрава № 220н и др.).",
    "• Для оказания помощи используют аптечки, укладки, наборы и комплекты по назначению.",
    "• Действуют в порядке приоритетности: безопасность → угрозы жизни → остальные меры.",
    "• Знают перечень состояний и мероприятий и выполняют их в установленной последовательности.",
    "• Соблюдают безопасные условия и меры профилактики инфекций.",
    "• При необходимости извлекают и перемещают пострадавшего безопасными способами.",
    "• Корректно вызывают СМП и другие специальные службы.",
]


def enhance_combined():
    tmp = Path("/tmp/fa_combined_enhanced.pptx")
    shutil.copy(COMBINED, tmp)
    prs = Presentation(str(tmp))
    n0 = len(prs.slides)
    print(f"Combined before: {n0} slides")

    # Append 4 new slides
    s_toc = add_blank_content_slide(prs)
    style_content_slide(
        s_toc, 0,
        "СОДЕРЖАНИЕ ТЕМЫ 1",
        "Тема 1 · структура по Прил. 2 к Порядку № 2464 (с 01.09.2026)",
        TOC_LINES,
    )

    s_kits = add_blank_content_slide(prs)
    style_content_slide(
        s_kits, 0,
        "ТЕРМИНОЛОГИЯ: УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ",
        "Тема 1 · п. 1.2  Укладки, наборы, комплекты и аптечки",
        KITS_LINES,
        footer_note="Дополнено: в исходнике были только слайды по автомобильной аптечке и аптечке работникам.",
    )

    s_prio = add_blank_content_slide(prs)
    style_content_slide(
        s_prio, 0,
        "ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "Тема 1 · п. 1.3  Порядок и приоритетность оказания первой помощи",
        PRIORITY_LINES,
        footer_note="Дополнено по алгоритмам первой помощи (приказ Минздрава РФ № 220н).",
    )

    s_sum = add_blank_content_slide(prs)
    style_content_slide(
        s_sum, 0,
        "РЕЗЮМЕ ТЕМЫ 1",
        "Тема 1 · ключевые выводы",
        SUMMARY_LINES,
    )

    # Current indices: 0..n0-1 old, n0 toc, n0+1 kits, n0+2 prio, n0+3 summary
    toc_i, kits_i, prio_i, sum_i = n0, n0 + 1, n0 + 2, n0 + 3

    # Old combined order (31 slides):
    # 0 title
    # 1 d1.1, 2-7 content 1.1
    # 8 d1.2, 9-10 kits visuals
    # 11 d1.3, 12-15 sequence
    # 16 d1.4, 17-19 states
    # 20 d1.5, 21-22 safety/infection
    # 23 d1.6, 24-27 extract
    # 28 d1.7, 29 call
    # 30 thanks

    order = (
        [0, toc_i]
        + list(range(1, 8))          # d1.1 + 1.1 content
        + [8, kits_i, 9, 10]         # d1.2 + new kits + visuals
        + [11, prio_i] + list(range(12, 16))  # d1.3 + priority + sequence
        + list(range(16, 20))        # d1.4 + states
        + list(range(20, 23))        # d1.5
        + list(range(23, 28))        # d1.6
        + list(range(28, 30))        # d1.7 + call
        + [sum_i, 30]                # summary + thanks
    )
    assert len(order) == len(prs.slides), (len(order), len(prs.slides))
    reorder_slides(prs, order)
    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    renumber_pages(prs)
    prs.save(str(COMBINED))
    shutil.copy(COMBINED, OUT_COMBINED_ROOT)
    print(f"Combined after: {len(prs.slides)} slides -> {COMBINED}")
    return len(prs.slides)


def enhance_individual():
    """Patch 1.2 and 1.3 individual decks with missing theory slides."""
    # 1.2 kits terminology
    f12 = INDIV_DIR / "1.2_Современные наборы средств и устройств.pptx"
    tmp = Path("/tmp/fa_1_2_enh.pptx")
    shutil.copy(f12, tmp)
    prs = Presentation(str(tmp))
    n0 = len(prs.slides)
    s = add_blank_content_slide(prs)
    style_content_slide(
        s, 0,
        "ТЕРМИНОЛОГИЯ: УКЛАДКИ, НАБОРЫ, КОМПЛЕКТЫ И АПТЕЧКИ",
        "Тема 1 · п. 1.2  Укладки, наборы, комплекты и аптечки",
        KITS_LINES,
        footer_note="Дополнено: в исходнике не было отдельного слайда по новой терминологии.",
    )
    # order: title(0), new, content..., thanks(last)
    kits_i = n0
    thanks = n0 - 1
    order = [0, kits_i] + list(range(1, thanks)) + [thanks]
    reorder_slides(prs, order)
    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    renumber_pages(prs)
    prs.save(str(f12))
    print(f"Updated {f12.name}: {len(prs.slides)} slides")

    # 1.3 priority
    f13 = INDIV_DIR / "1.3_Общая последовательность действий на месте происшествия.pptx"
    tmp = Path("/tmp/fa_1_3_enh.pptx")
    shutil.copy(f13, tmp)
    prs = Presentation(str(tmp))
    n0 = len(prs.slides)
    s = add_blank_content_slide(prs)
    style_content_slide(
        s, 0,
        "ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "Тема 1 · п. 1.3  Порядок и приоритетность оказания первой помощи",
        PRIORITY_LINES,
        footer_note="Дополнено по алгоритмам первой помощи (приказ Минздрава РФ № 220н).",
    )
    prio_i = n0
    thanks = n0 - 1
    # title, priority, sequence slides (1..4 were sequence before extract), keep rest
    order = [0, prio_i] + list(range(1, thanks)) + [thanks]
    reorder_slides(prs, order)
    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    renumber_pages(prs)
    prs.save(str(f13))
    print(f"Updated {f13.name}: {len(prs.slides)} slides")


def update_change_table():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Изменения"
    headers = [
        "№",
        "Название презентации (без изменения)",
        "Было (содержание / акцент)",
        "Стало (под Прил. 2 № 2464 с 01.09.2026)",
        "Пункт новой структуры",
        "Что изменено",
    ]
    ws.append(headers)
    rows = [
        [
            1,
            "1 Организационно-правовые аспекты оказания первой помощи.pptx",
            "Сводная тема 1 в старой компоновке",
            "Сводная тема 1 по пунктам 1.1–1.7 + оглавление, слайд приоритетности, терминология аптечек, резюме",
            "1.1–1.7",
            "Перекомпоновка исходных слайдов; добавлены недостающие структурные и теоретические слайды",
        ],
        [
            2,
            "1.1_Понятие первая помощь.pptx",
            "Понятие ПП; состояния; мероприятия",
            "Содержание под п. 1.4 (состояния и мероприятия + последовательность)",
            "п. 1.4",
            "Подписи/акцент; имя файла сохранено",
        ],
        [
            3,
            "1.2_Современные наборы средств и устройств.pptx",
            "Автомобильная аптечка; аптечка работникам",
            "п. 1.2 + добавлен слайд терминологии «укладки, наборы, комплекты и аптечки»",
            "п. 1.2",
            "Добавлен недостающий теоретический слайд; визуалы аптечек сохранены",
        ],
        [
            4,
            "1.3_Общая последовательность действий на месте происшествия.pptx",
            "Последовательность + извлечение в одном файле",
            "п. 1.3 (порядок/приоритетность) + п. 1.6 (извлечение) + слайд приоритетности",
            "п. 1.3 и п. 1.6",
            "Добавлен слайд приоритетности; извлечение оставлено в том же файле (имя не менялось)",
        ],
        [
            5,
            "1.4_Соблюдение правил личной безопасности.pptx",
            "Безопасность и инфекции",
            "Подпись к п. 1.5 новой структуры",
            "п. 1.5",
            "Подписи; имя файла сохранено",
        ],
        [
            6,
            "1.5_Основные правила вызова скорой медицинской помощи.pptx",
            "Вызов СМП и спецслужб",
            "Подпись к п. 1.7 новой структуры",
            "п. 1.7",
            "Подписи; имя файла сохранено",
        ],
    ]
    for r in rows:
        ws.append(r)
    for col in ws.columns:
        maxlen = max(len(str(c.value or "")) for c in col)
        ws.column_dimensions[col[0].column_letter].width = min(maxlen + 2, 55)
    wb.save(CHANGE_XLSX)
    shutil.copy(CHANGE_XLSX, INDIV_DIR / CHANGE_XLSX.name)
    print(f"Updated {CHANGE_XLSX}")


def rebuild_zip():
    if ZIP_OUT.exists():
        ZIP_OUT.unlink()
    with zipfile.ZipFile(ZIP_OUT, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in sorted(INDIV_DIR.glob("*.pptx")):
            zf.write(p, p.name)
        xlsx = INDIV_DIR / "Таблица_изменений_Тема1_первая_помощь.xlsx"
        if xlsx.exists():
            zf.write(xlsx, xlsx.name)
        readme = INDIV_DIR / "README.md"
        if readme.exists():
            zf.write(readme, readme.name)
    print(f"Zip: {ZIP_OUT} ({ZIP_OUT.stat().st_size} bytes)")


def write_notes(n_slides):
    NOTES.write_text(
        f"""# Тема 1 — пересборка под Прил. 2 к Порядку № 2464 (с 01.09.2026)

## Файлы
- Сводная: `Презентация_Тема_1_Первая_помощь_новая_структура.pptx` / `first_aid_tema1_new/1 Организационно-правовые…pptx` ({n_slides} слайдов)
- Отдельные (имена исходные): `first_aid_tema1_new/1.1_…` … `1.5_…`
- Архив: `Тема1_первая_помощь_1.1-1.7.zip`
- Таблица изменений: `Таблица_изменений_Тема1_первая_помощь.xlsx`

## Структура сводной
1. Титул → **Оглавление**
2. **1.1** НПА и организация
3. **1.2** Укладки/наборы/комплекты/аптечки (+ слайд терминологии) + визуалы аптечек
4. **1.3** Порядок и приоритетность (+ явный слайд приоритетности) + алгоритм последовательности
5. **1.4** Состояния и мероприятия (+ определение ПП как вводный)
6. **1.5** Безопасные условия + профилактика инфекций
7. **1.6** Извлечение и перемещение
8. **1.7** Вызов СМП
9. **Резюме** → благодарность

## Что добавлено (не было в исходнике как отдельные блоки)
- Оглавление темы 1
- Слайд терминологии укладок/наборов/комплектов/аптечек (п. 1.2)
- Слайд приоритетности оказания первой помощи (п. 1.3) — по приказу Минздрава № 220н
- Резюме темы

Исходное содержание слайдов сохранено; имена отдельных файлов не менялись.
""",
        encoding="utf-8",
    )


def main():
    n = enhance_combined()
    enhance_individual()
    update_change_table()
    rebuild_zip()
    write_notes(n)
    # print outline
    prs = Presentation(str(COMBINED))
    print("\n=== FINAL OUTLINE ===")
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = " ".join(sh.text_frame.text.split())
                if len(t) > 15:
                    title = t[:110]
                    break
        print(f"{i:02d}: {title}")


if __name__ == "__main__":
    main()
