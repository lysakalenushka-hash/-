#!/usr/bin/env python3
"""Generate B.7.5 training program (Word) and presentation (PPTX, red style)."""

from __future__ import annotations

import copy
import re
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt as PptxPt

# --- palette (same as themes 1–4) ---
RED = PptxRGB(0xE3, 0x06, 0x13)
GRAY = PptxRGB(0x6B, 0x72, 0x80)
DARK = PptxRGB(0x1A, 0x1A, 0x1A)
WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
CREAM = PptxRGB(0xFE, 0xF3, 0xC7)
AMBER = PptxRGB(0xF5, 0x9E, 0x0B)
LIGHT_BG = PptxRGB(0xFA, 0xFA, 0xFA)
BORDER = PptxRGB(0xE5, 0xE7, 0xEB)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN_L = Emu(640080)
CONTENT_W = Emu(10911535)
FONT = "Inter"

OUT_DOCX = Path("Б.7.5_Учебный_план_24ч.docx")
OUT_PPTX = Path("Б.7.5_Проектирование_строительство_газовых_сетей.pptx")
TEMPLATE = Path("/tmp/tema4_template.pptx")

MODULES = [
    {
        "num": 1,
        "title": "Общие требования промышленной безопасности",
        "hours": 6,
        "topics": [
            ("1.1", "Правовое регулирование в области промышленной безопасности", 2),
            ("1.2", "Контрольно-надзорная и разрешительная деятельность Ростехнадзора", 1),
            ("1.3", "Производственный контроль на ОПО", 1),
            ("1.4", "Требования промышленной безопасности на объектах газораспределения и газопотребления", 2),
        ],
    },
    {
        "num": 2,
        "title": "Проектирование сетей газораспределения и газопотребления",
        "hours": 8,
        "topics": [
            ("2.1", "Общие требования к сетям газораспределения и газопотребления (ФНП № 531, ТР, СП)", 4),
            ("2.2", "Требования к проектированию: трассы, материалы, ГРП/ГРУ, вентиляция, КИП", 4),
        ],
    },
    {
        "num": 3,
        "title": "Технический надзор, строительство, реконструкция и капитальный ремонт",
        "hours": 6,
        "topics": [
            ("3.1", "Строительство и реконструкция газопроводов: прокладка, контроль СМР, охранные зоны", 3),
            ("3.2", "Капитальный ремонт, техническое перевооружение, консервация и расконсервация", 3),
        ],
    },
    {
        "num": 4,
        "title": "Сварочные работы, испытания, пуск газа и приёмка",
        "hours": 4,
        "topics": [
            ("4.1", "Требования к производству и контролю сварочных (наплавочных) работ на ОПО", 2),
            ("4.2", "Испытания на прочность и герметичность; врезка, продувка, пуск газа; приёмка в эксплуатацию", 2),
        ],
    },
]

SLIDES = [
    {
        "type": "title",
        "course": "КУРС",
        "line1": "Требования промышленной безопасности на объектах газораспределения и газопотребления",
        "line2": "Б.7.5. Проектирование, строительство, реконструкция, техническое перевооружение и капитальный ремонт сетей газораспределения и газопотребления",
        "line3": "Программа повышения квалификации · 24 академических часа · 4 модуля",
        "tags": "Газораспределение · Газопотребление · Промышленная безопасность · ФНП № 531",
    },
    {
        "type": "intro",
        "title": "Б.7.5. Область аттестации и цель программы",
        "subtitle": "Подготовка к аттестации в Ростехнадзоре",
        "section": "Назначение программы",
        "body": (
            "Программа предназначена для подготовки руководителей и специалистов организаций, "
            "осуществляющих проектирование, строительство, реконструкцию, техническое перевооружение "
            "и капитальный ремонт сетей газораспределения и газопотребления, к аттестации по области Б.7.5."
        ),
        "note": (
            "Нормативная база: 116-ФЗ; ФНП № 531; ТР о безопасности сетей газораспределения и газопотребления; "
            "СП 62.13330, СП 42-101, СП 42-102, СП 42-103 и иные документы по стандартизации."
        ),
    },
    {
        "type": "module",
        "module": 1,
        "title": "Общие требования промышленной безопасности",
        "subtitle": "Правовые основы и производственный контроль",
        "desc": "116-ФЗ, полномочия Ростехнадзора, производственный контроль, общие требования к объектам газоснабжения",
        "tags": "Промышленная безопасность · 116-ФЗ · Производственный контроль",
    },
    {
        "type": "items",
        "title": "Правовое регулирование промышленной безопасности",
        "items": [
            "116-ФЗ «О промышленной безопасности опасных производственных объектов» — цели, принципы, категории ОПО.",
            "Постановление Правительства РФ № 2168 — категории работников, подлежащих аттестации.",
            "Положение об аттестации (ПП РФ № 978) — порядок подготовки и проверки знаний.",
            "Ответственность за нарушение требований промышленной безопасности (административная и уголовная).",
        ],
    },
    {
        "type": "items",
        "title": "Производственный контроль и надзор на объектах газоснабжения",
        "items": [
            "Производственный контроль: цели, программа, документирование результатов.",
            "Контрольно-надзорная деятельность Ростехнадзора при строительстве и реконструкции сетей.",
            "Регистрация объектов в государственном реестре ОПО; декларирование промышленной безопасности.",
            "Общие требования ФНП № 531 к организациям, выполняющим работы на сетях газораспределения и газопотребления.",
        ],
    },
    {
        "type": "module",
        "module": 2,
        "title": "Проектирование сетей газораспределения и газопотребления",
        "subtitle": "Проектные решения и нормативные ограничения",
        "desc": "Классификация сетей, выбор материалов, прокладка, здания ГРП/ГРУ, вентиляция и контрольно-измерительные приборы",
        "tags": "Проектирование · Газопровод · ГРП · СП 42-102",
    },
    {
        "type": "items",
        "title": "Общие требования к проектированию газовых сетей",
        "items": [
            "Классификация сетей газораспределения и газопотребления по давлению, назначению и способу прокладки.",
            "Выбор материалов труб (сталь, полиэтилен) с учётом давления, среды и условий прокладки.",
            "Расстояния до зданий, сооружений и инженерных коммуникаций; пересечения и параллельное следование.",
            "Охранные зоны газопроводов; требования к маркировке и сигнальным знакам на трассе.",
        ],
    },
    {
        "type": "detail",
        "title": "Проектирование ГРП, ГРУ и внутренних газопроводов",
        "subtitle": "Ключевые проектные решения",
        "section": "Пункты редуцирования газа и установки учёта",
        "body": (
            "Проектом определяются тип ГРП/ГРУ, схема подключения, резервирование, система отопления и вентиляции "
            "помещений, газоопасные сигнализаторы, аварийное отключение и освещение. "
            "Внутренние газопроводы проектируются с учётом давления, категории помещений и требований взрывобезопасности."
        ),
        "note": "После строительства на каждый газопровод, ГРП и технологическую установку оформляется паспорт (ФНП № 531).",
    },
    {
        "type": "items",
        "title": "Документация и экспертиза проектной документации",
        "items": [
            "Состав проектной и рабочей документации на объекты газоснабжения.",
            "Требования к разделам «Технологические решения», «Пожарная безопасность», «Мероприятия по ПБ».",
            "Государственная и негосударственная экспертиза; учёт заключений при строительстве.",
            "Изменения в проектной документации в процессе строительства — порядок согласования.",
        ],
    },
    {
        "type": "module",
        "module": 3,
        "title": "Строительство, реконструкция и технический надзор",
        "subtitle": "Организация СМР и контроль качества",
        "desc": "Прокладка газопроводов, строительный контроль, реконструкция, капремонт, консервация",
        "tags": "Строительство · Реконструкция · Технический надзор · Охранные зоны",
    },
    {
        "type": "items",
        "title": "Строительство и реконструкция газопроводов",
        "items": [
            "Подготовка трассы, траншеи, способы прокладки (подземная, наземная, в коллекторах).",
            "Монтаж стальных и полиэтиленовых газопроводов; требования к стыковым соединениям.",
            "Строительный контроль и авторский надзор; журналы работ и акты скрытых работ.",
            "Реконструкция изношенных газопроводов; методы бестраншейной замены и параллельной прокладки.",
        ],
    },
    {
        "type": "detail",
        "title": "Капитальный ремонт и техническое перевооружение",
        "subtitle": "Работы без нарушения безопасной эксплуатации действующих сетей",
        "section": "Организация ремонтных работ",
        "body": (
            "Капитальный ремонт и техническое перевооружение выполняются по проекту или технической документации "
            "с отключением и продувкой участков, оформлением нарядов-допусков на газоопасные работы. "
            "Консервация и расконсервация газопроводов — по регламенту с испытаниями перед пуском газа."
        ),
        "note": "Запрещается ввод в эксплуатацию объектов без испытаний и оформления исполнительной документации.",
    },
    {
        "type": "items",
        "title": "Газоопасные работы и охрана газопроводов",
        "items": [
            "Перечень газоопасных работ: врезка, пуск газа, снятие заглушек, сварка на действующих газопроводах.",
            "Наряд-допуск и план газоопасных работ; роли ответственного руководителя и исполнителей.",
            "Правила охраны газораспределительных сетей (ПП РФ № 878): охранная зона, ограничения для земляных работ.",
            "Действия при повреждении газопровода третьими лицами; аварийно-диспетчерское обеспечение.",
        ],
    },
    {
        "type": "module",
        "module": 4,
        "title": "Сварка, испытания, пуск газа и приёмка",
        "subtitle": "Контроль качества и ввод в эксплуатацию",
        "desc": "Сварочные работы, неразрушающий контроль, опрессовка, врезка, продувка, приёмка объекта",
        "tags": "Сварка · Испытания · Пуск газа · Приёмка · НК",
    },
    {
        "type": "items",
        "title": "Сварочные работы и контроль качества соединений",
        "items": [
            "Требования к сварщикам и технологии сварки на ОПО; допуск к газоопасным сварочным работам.",
            "Визуальный и неразрушающий контроль сварных швов; оформление результатов контроля.",
            "Запрет эксплуатации газопроводов с дефектными сварными соединениями.",
            "Ремонт дефектов: повторная сварка, выборочная замена участков по результатам контроля.",
        ],
    },
    {
        "type": "items",
        "title": "Испытания газопроводов на прочность и герметичность",
        "items": [
            "После СМР и ремонта — испытания по проекту; при отсутствии проектных требований — по СП.",
            "Испытание на прочность и герметичность после сварочных работ обязательно.",
            "Контрольная опрессовка оборудования и газопроводов ГРП давлением 0,01 МПа.",
            "Оформление актов испытаний; допуск к пуску газа только после положительных результатов.",
        ],
    },
    {
        "type": "detail",
        "title": "Врезка, продувка, пуск газа и приёмка объекта",
        "subtitle": "Завершающий этап строительства и реконструкции",
        "section": "Ввод в эксплуатацию",
        "body": (
            "Пуск газа выполняется по плану после продувки, проверки герметичности запорной арматуры и отсутствия "
            "заглушек. Снятие заглушек — после обхода, ТО и испытаний. "
            "Приёмка законченного строительством объекта — комиссией с оформлением актов и внесением сведений в реестр."
        ),
        "note": "Ввод в эксплуатацию без паспортов, исполнительной документации и актов испытаний не допускается.",
    },
    {
        "type": "summary",
        "part": "1/2",
        "items": [
            ("Правовые основы", "116-ФЗ, аттестация по Б.7.5, производственный контроль и надзор Ростехнадзора."),
            ("Проектирование", "ФНП № 531, СП; выбор трасс, материалов, ГРП/ГРУ; экспертиза проектной документации."),
            ("Строительство", "СМР, реконструкция, капремонт, газоопасные работы, охранные зоны."),
            ("Контроль качества", "Сварка, НК, испытания на прочность и герметичность, пуск газа, приёмка."),
        ],
    },
    {
        "type": "summary",
        "part": "2/2",
        "items": [
            ("Итоговая аттестация", "Проверка знаний в Едином портале тестирования Ростехнадзора по области Б.7.5."),
            ("Документ об обучении", "Удостоверение о повышении квалификации; срок действия аттестации — 5 лет."),
            ("Периодичность обучения", "Не реже одного раза в 5 лет (ч. 1 ст. 14.1 116-ФЗ)."),
        ],
        "closing": (
            "Соблюдение требований ФНП № 531 на этапах проектирования, строительства и ремонта "
            "обеспечивает безопасность сетей газораспределения и газопотребления."
        ),
    },
]


def set_run_font(run, size=14, bold=False, color=DARK, name=FONT):
    run.font.name = name
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text="", size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_bg(slide, color=LIGHT_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_accent_bar(slide, top=Emu(6355080)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, top, CONTENT_W, Emu(91440))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()


def add_page_num(slide, num, total):
    add_textbox(
        slide,
        Emu(10180015),
        Emu(6492240),
        Emu(1188720),
        Emu(228600),
        f"{num:02d} / {total:02d}",
        size=11,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def add_title_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9753600), Emu(-457200), Emu(2743200), Emu(2743200))
    oval.fill.solid()
    oval.fill.fore_color.rgb = CREAM
    oval.line.fill.background()
    add_textbox(slide, MARGIN_L, Emu(731520), Emu(4572000), Emu(228600), data["course"], size=12, bold=True, color=RED)
    add_textbox(slide, MARGIN_L, Emu(1188720), CONTENT_W, Emu(640080), data["line1"], size=28, bold=True, color=DARK)
    add_textbox(slide, MARGIN_L, Emu(1920240), CONTENT_W, Emu(914400), data["line2"], size=18, color=GRAY)
    add_textbox(slide, MARGIN_L, Emu(2926080), CONTENT_W, Emu(457200), data["line3"], size=14, color=GRAY)
    add_accent_bar(slide)
    add_textbox(slide, MARGIN_L, Emu(5486400), CONTENT_W, Emu(365760), data["tags"], size=12, color=GRAY)
    add_page_num(slide, page, total)


def add_intro_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, MARGIN_L, Emu(731520), CONTENT_W, Emu(640080), data["title"], size=24, bold=True, color=DARK)
    add_textbox(slide, MARGIN_L, Emu(1463040), CONTENT_W, Emu(365760), data["subtitle"], size=16, color=RED)
    add_textbox(slide, MARGIN_L, Emu(2103120), Emu(3200400), Emu(365760), data["section"], size=14, bold=True, color=RED)
    add_textbox(slide, MARGIN_L, Emu(2552700), Emu(7772400), Emu(1280160), data["body"], size=13, color=GRAY)
    add_textbox(slide, MARGIN_L, Emu(4206240), CONTENT_W, Emu(914400), data["note"], size=12, color=GRAY)
    add_accent_bar(slide)
    add_page_num(slide, page, total)


def add_module_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9753600), Emu(-457200), Emu(2743200), Emu(2743200))
    oval.fill.solid()
    oval.fill.fore_color.rgb = CREAM
    oval.line.fill.background()
    add_textbox(slide, MARGIN_L, Emu(731520), Emu(3200400), Emu(365760), f"МОДУЛЬ {data['module']}", size=14, bold=True, color=RED)
    add_textbox(slide, MARGIN_L, Emu(1188720), CONTENT_W, Emu(640080), data["title"], size=26, bold=True, color=DARK)
    add_textbox(slide, MARGIN_L, Emu(1920240), CONTENT_W, Emu(457200), data["subtitle"], size=16, color=GRAY)
    add_textbox(slide, MARGIN_L, Emu(2468880), CONTENT_W, Emu(640080), data["desc"], size=14, color=GRAY)
    add_accent_bar(slide)
    add_textbox(slide, MARGIN_L, Emu(5486400), CONTENT_W, Emu(365760), data["tags"], size=12, color=GRAY)
    add_page_num(slide, page, total)


def add_items_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, MARGIN_L, Emu(731520), CONTENT_W, Emu(640080), data["title"], size=24, bold=True, color=DARK)
    y = Emu(1650720)
    for i, item in enumerate(data["items"], 1):
        add_textbox(slide, MARGIN_L, y, Emu(640080), Emu(365760), f"{i:02d}", size=14, bold=True, color=RED)
        add_textbox(slide, Emu(1371600), y, Emu(9539980), Emu(640080), item, size=13, color=GRAY)
        y += Emu(960000)
    add_accent_bar(slide)
    add_page_num(slide, page, total)


def add_detail_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, MARGIN_L, Emu(731520), CONTENT_W, Emu(640080), data["title"], size=24, bold=True, color=DARK)
    add_textbox(slide, MARGIN_L, Emu(1463040), CONTENT_W, Emu(365760), data["subtitle"], size=16, color=GRAY)
    add_textbox(slide, MARGIN_L, Emu(2103120), Emu(3200400), Emu(365760), data["section"], size=14, bold=True, color=RED)
    add_textbox(slide, MARGIN_L, Emu(2552700), CONTENT_W, Emu(1280160), data["body"], size=13, color=GRAY)
    add_textbox(slide, MARGIN_L, Emu(4206240), CONTENT_W, Emu(640080), data["note"], size=12, color=AMBER)
    add_accent_bar(slide)
    add_page_num(slide, page, total)


def add_summary_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, MARGIN_L, Emu(731520), Emu(4572000), Emu(228600), "Итоги курса", size=14, bold=True, color=RED)
    add_textbox(
        slide,
        MARGIN_L,
        Emu(1005840),
        CONTENT_W,
        Emu(914400),
        f"Ключевые положения программы Б.7.5 ({data['part']})",
        size=22,
        bold=True,
        color=DARK,
    )
    y = Emu(2286000)
    for title, body in data["items"]:
        add_textbox(slide, MARGIN_L, y, Emu(3474720), Emu(365760), title, size=13, bold=True, color=RED)
        add_textbox(slide, Emu(4480560), y, Emu(7071055), Emu(640080), body, size=12, color=GRAY)
        y += Emu(1028700)
    if data.get("closing"):
        add_textbox(slide, MARGIN_L, Emu(5221224), CONTENT_W, Emu(548640), data["closing"], size=12, color=GRAY)
    add_accent_bar(slide)
    add_page_num(slide, page, total)


def build_pptx():
    # reuse theme fonts from template if available
    prs = Presentation(TEMPLATE)
    # remove existing slides
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    total = len(SLIDES)
    builders = {
        "title": add_title_slide,
        "intro": add_intro_slide,
        "module": add_module_slide,
        "items": add_items_slide,
        "detail": add_detail_slide,
        "summary": add_summary_slide,
    }
    for i, slide_data in enumerate(SLIDES, 1):
        builders[slide_data["type"]](prs, slide_data, i, total)
    prs.save(OUT_PPTX)


def build_docx():
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(1.5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run(
        "Дополнительная профессиональная программа повышения квалификации\n"
        "«Требования промышленной безопасности на объектах газораспределения и газопотребления»\n"
        "(область аттестации Б.7.5)"
    )
    r.bold = True
    r.font.size = Pt(14)
    r.font.name = "Times New Roman"

    doc.add_paragraph()
    meta = [
        "Срок обучения: 24 академических часа (4 модуля)",
        "Форма обучения: очная, очно-заочная, заочная с применением ДОТ",
        "Категория слушателей: руководители и специалисты организаций, осуществляющих проектирование, "
        "строительство, реконструкцию, техническое перевооружение и капитальный ремонт сетей "
        "газораспределения и газопотребления",
        "Цель: повышение квалификации и подготовка к аттестации в Ростехнадзоре по области Б.7.5",
        "Нормативная база: 116-ФЗ; ФНП № 531 (приказ Ростехнадзора от 15.12.2020); "
        "ТР о безопасности сетей газораспределения и газопотребления; СП 62.13330, СП 42-101, "
        "СП 42-102, СП 42-103; Правила охраны газораспределительных сетей (ПП РФ № 878)",
    ]
    for line in meta:
        p = doc.add_paragraph(line)
        p.runs[0].font.name = "Times New Roman"
        p.runs[0].font.size = Pt(12)

    doc.add_paragraph()
    h = doc.add_paragraph("Учебно-тематический план")
    h.runs[0].bold = True
    h.runs[0].font.size = Pt(13)
    h.runs[0].font.name = "Times New Roman"

    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    headers = ["№", "Наименование раздела (модуля) и темы", "Всего, ч", "Форма занятий"]
    for cell, text in zip(hdr, headers):
        cell.text = text
        for p in cell.paragraphs:
            for run in p.runs:
                run.bold = True
                run.font.name = "Times New Roman"
                run.font.size = Pt(11)

    row_num = 1
    total_hours = 0
    for mod in MODULES:
        row = table.add_row().cells
        row[0].text = str(row_num)
        row[1].text = f"Модуль {mod['num']}. {mod['title']}"
        row[2].text = str(mod["hours"])
        row[3].text = "Лекция / самостоятельная работа (ДОТ)"
        row_num += 1
        total_hours += mod["hours"]
        for code, topic, hours in mod["topics"]:
            row = table.add_row().cells
            row[0].text = ""
            row[1].text = f"  Тема {code}. {topic}"
            row[2].text = str(hours)
            row[3].text = "Лекция / ДОТ"
            row_num += 1

    row = table.add_row().cells
    row[0].text = ""
    row[1].text = "Итоговая аттестация (тестирование по области Б.7.5)"
    row[2].text = "—"
    row[3].text = "Тест"

    row = table.add_row().cells
    row[0].text = ""
    row[1].text = "ИТОГО"
    row[1].paragraphs[0].runs[0].bold = True
    row[2].text = str(total_hours)
    row[2].paragraphs[0].runs[0].bold = True
    row[3].text = ""

    doc.add_paragraph()
    h2 = doc.add_paragraph("Краткое содержание модулей")
    h2.runs[0].bold = True
    h2.runs[0].font.size = Pt(13)

    summaries = {
        1: (
            "Правовые основы промышленной безопасности: 116-ФЗ, аттестация, производственный контроль, "
            "полномочия Ростехнадзора, общие требования к объектам газораспределения и газопотребления."
        ),
        2: (
            "Требования к проектированию сетей: классификация, материалы, трассы, ГРП/ГРУ, "
            "внутренние газопроводы, состав проектной документации и экспертиза."
        ),
        3: (
            "Строительство, реконструкция, капитальный ремонт и технический надзор; "
            "газоопасные работы, наряды-допуски, охранные зоны газопроводов."
        ),
        4: (
            "Сварочные работы и контроль качества; испытания на прочность и герметичность; "
            "врезка, продувка, пуск газа; приёмка и ввод объекта в эксплуатацию."
        ),
    }
    for mod in MODULES:
        p = doc.add_paragraph()
        r1 = p.add_run(f"Модуль {mod['num']}. {mod['title']} ({mod['hours']} ч). ")
        r1.bold = True
        r1.font.name = "Times New Roman"
        r1.font.size = Pt(12)
        r2 = p.add_run(summaries[mod["num"]])
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(12)

    doc.add_paragraph()
    p = doc.add_paragraph(
        "По итогам обучения слушатели сдают итоговую аттестацию (тестирование) по области Б.7.5 "
        "и получают удостоверение о повышении квалификации установленного образца."
    )
    p.runs[0].font.name = "Times New Roman"
    p.runs[0].font.size = Pt(12)

    doc.save(OUT_DOCX)


def main():
    build_docx()
    build_pptx()
    print(f"Created: {OUT_DOCX} ({OUT_DOCX.stat().st_size} bytes)")
    print(f"Created: {OUT_PPTX} ({OUT_PPTX.stat().st_size} bytes, {len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
