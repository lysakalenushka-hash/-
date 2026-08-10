#!/usr/bin/env python3
"""
Нумерация презентаций темы 1 по порядку тем (Прил. № 2 к Порядку № 2464 /
оглавление темы 1 учебного пособия Минздрава 2025) и дополнение пробелов
по методичке.
Первый файл: «Организация оказания первой помощи в РФ. Нормативно-правовая база».
"""

from pathlib import Path
import shutil
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt
from openpyxl import Workbook

SRC = Path("/workspace/first_aid_uploads_src")
OUT = Path("/workspace/first_aid_tema1_numbered")
ZIP = Path("/workspace/Тема1_нумерация_по_списку_Прил2_2464.zip")
XLSX = Path("/workspace/Таблица_нумерация_и_дополнения_Тема1.xlsx")
README = OUT / "README.md"

DARK = RGBColor(0x1A, 0x1A, 0x1A)
RED = RGBColor(0xE3, 0x06, 0x13)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF7, 0xF8)
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Порядок по теме 1 (первый пункт — как указал пользователь)
# Соответствует оглавлению темы 1 пособия Минздрава 2025 и логике Прил. № 2 № 2464
FILES = [
    {
        "num": 1,
        "src": "src_1_org_npa.pptx",
        "out": "1_Организация_оказания_первой_помощи_в_РФ_Нормативно-правовая_база.pptx",
        "title": "ОРГАНИЗАЦИЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ В РФ. НОРМАТИВНО-ПРАВОВАЯ БАЗА",
        "eyebrow": "Тема 1 · п. 1  Организация и нормативно-правовая база",
        "list_item": "Организация оказания первой помощи в Российской Федерации. Нормативно-правовая база, определяющая права, обязанности и ответственность при оказании первой помощи.",
    },
    {
        "num": 2,
        "src": "src_1.2_kits.pptx",
        "out": "2_Современные_аптечки_укладки_комплекты_и_наборы.pptx",
        "title": "СОВРЕМЕННЫЕ АПТЕЧКИ, УКЛАДКИ, КОМПЛЕКТЫ И НАБОРЫ ДЛЯ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 2  Аптечки, укладки, комплекты и наборы",
        "list_item": "Современные аптечки, укладки, комплекты и наборы средств и устройств, использующиеся для оказания первой помощи. Основные компоненты, их назначение.",
    },
    {
        "num": 3,
        "src": "src_1.3_order.pptx",
        "out": "3_Порядок_и_приоритетность_оказания_первой_помощи.pptx",
        "title": "ПОРЯДОК И ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
        "eyebrow": "Тема 1 · п. 3  Порядок и приоритетность; извлечение и перемещение",
        "list_item": "Порядок оказания первой помощи. Приоритетность оказания первой помощи. Способы извлечения пострадавших из труднодоступных мест и их перемещения в безопасное место.",
    },
    {
        "num": 4,
        "src": "src_1.1_states.pptx",
        "out": "4_Перечень_состояний_и_мероприятий_первой_помощи.pptx",
        "title": "ПЕРЕЧЕНЬ СОСТОЯНИЙ И МЕРОПРИЯТИЙ ПЕРВОЙ ПОМОЩИ. ПОСЛЕДОВАТЕЛЬНОСТЬ ВЫПОЛНЕНИЯ",
        "eyebrow": "Тема 1 · п. 4  Состояния и мероприятия (приказ Минздрава № 220н)",
        "list_item": "Перечень состояний, при которых оказывается первая помощь. Перечень мероприятий по оказанию первой помощи и последовательность их выполнения.",
    },
    {
        "num": 5,
        "src": "src_1.4_safety.pptx",
        "out": "5_Обеспечение_безопасных_условий_и_профилактика_инфекций.pptx",
        "title": "ОБЕСПЕЧЕНИЕ БЕЗОПАСНЫХ УСЛОВИЙ И ПРОФИЛАКТИКА ИНФЕКЦИЙ",
        "eyebrow": "Тема 1 · п. 5  Безопасные условия и профилактика инфекций",
        "list_item": "Обеспечение безопасных условий для оказания первой помощи. Простейшие меры профилактики инфекционных заболеваний при оказании первой помощи.",
    },
    {
        "num": 6,
        "src": "src_1.5_call.pptx",
        "out": "6_Основные_правила_вызова_скорой_медицинской_помощи.pptx",
        "title": "ОСНОВНЫЕ ПРАВИЛА ВЫЗОВА СКОРОЙ МЕДИЦИНСКОЙ ПОМОЩИ И ДРУГИХ СПЕЦИАЛЬНЫХ СЛУЖБ",
        "eyebrow": "Тема 1 · п. 6  Вызов СМП и специальных служб",
        "list_item": "Основные правила вызова скорой медицинской помощи, других специальных служб, сотрудники которых обязаны оказывать первую помощь.",
    },
]


def replace_shape_text(shape, text):
    for t_el in shape.element.findall(f".//{NS}t"):
        t_el.text = ""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
    else:
        p0.add_run().text = text


def set_run(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def clear_slide_shapes(slide):
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def add_blank_content_slide(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    clear_slide_shapes(slide)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(180000), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    return slide


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_multiline(slide, left, top, width, height, lines, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=False, color=DARK)
    return box


def style_content_slide(slide, page_num, title, eyebrow, body_lines, footer=None):
    add_textbox(slide, Emu(220000), Emu(6540000), Emu(900000), Emu(600000), str(page_num), size=22, bold=True, color=GRAY)
    add_textbox(slide, Emu(2050000), Emu(1400000), Emu(21000000), Emu(1000000), title, size=26, bold=True, color=DARK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(2060000), Emu(2500000), Emu(12000000), Emu(70000))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    add_textbox(slide, Emu(2050000), Emu(2650000), Emu(18000000), Emu(500000), eyebrow, size=14, color=GRAY)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(2050000), Emu(3400000), Emu(20000000), Emu(8800000))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.fill.background()
    add_multiline(slide, Emu(2400000), Emu(3700000), Emu(19000000), Emu(8000000), body_lines, size=16)
    if footer:
        add_textbox(slide, Emu(2050000), Emu(12400000), Emu(20000000), Emu(600000), footer, size=11, color=GRAY)


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in new_items:
        sldIdLst.append(el)


def renumber_pages(prs):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if (sh.left or 0) < Emu(1200000) and (sh.width or 0) < Emu(1500000) and (sh.top or 0) > Emu(5000000):
                for t_el in sh.element.findall(f".//{NS}t"):
                    t_el.text = ""
                p0 = sh.text_frame.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = str(i)
                else:
                    p0.add_run().text = str(i)


def set_title_like(slide, text):
    cands = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if (sh.top or 0) < Emu(3500000) and (sh.width or 0) > Emu(8000000):
            cands.append((-(sh.width or 0), sh))
    if not cands:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() and (sh.width or 0) > Emu(5000000):
                cands.append((-(sh.width or 0), sh))
    if cands:
        cands.sort()
        replace_shape_text(cands[0][1], text)
        return True
    return False


def set_eyebrow_like(slide, text):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "Аспекты оказания" in t or "Тема 1" in t:
            replace_shape_text(sh, text)
            return True
    return False


def insert_content_before_thanks(prs, title, eyebrow, body_lines, footer=None):
    n0 = len(prs.slides)
    thanks_idx = n0 - 1
    # assume last is thanks if «БЛАГОДАР» in text
    last_texts = []
    for sh in prs.slides[thanks_idx].shapes:
        if sh.has_text_frame:
            last_texts.append(sh.text_frame.text)
    if not any("БЛАГОДАР" in t.upper() for t in last_texts):
        thanks_idx = None

    slide = add_blank_content_slide(prs)
    style_content_slide(slide, 0, title, eyebrow, body_lines, footer)
    new_i = n0
    if thanks_idx is None:
        return
    # order: 0..thanks-1, new, thanks
    order = list(range(thanks_idx)) + [new_i] + [thanks_idx]
    # remaining slides after thanks (shouldn't exist)
    for i in range(thanks_idx + 1, n0):
        order.append(i)
    reorder_slides(prs, order)


def insert_after_title(prs, title, eyebrow, body_lines, footer=None):
    n0 = len(prs.slides)
    slide = add_blank_content_slide(prs)
    style_content_slide(slide, 0, title, eyebrow, body_lines, footer)
    new_i = n0
    order = [0, new_i] + list(range(1, n0))
    reorder_slides(prs, order)


def update_existing_eyebrows(prs, eyebrow):
    for slide in prs.slides:
        set_eyebrow_like(slide, eyebrow)


# --- Content supplements from Minzdrav 2025 guide ---

ORG_5_COMPONENTS = [
    "По учебному пособию Минздрава (2025) система оказания первой помощи в РФ",
    "состоит из пяти основных компонентов (ранее в презентации указывали три):",
    "",
    "1. Организация и нормативно-правовое обеспечение.",
    "2. Обучение участников оказания первой помощи правилам и навыкам ее оказания.",
    "3. Оснащение участников средствами для ее оказания (аптечками, укладками, наборами, комплектами).",
    "4. Мотивирование на обучение и оказание первой помощи.",
    "5. Учет и анализ эффективности оказания первой помощи.",
]

ORG_220_CONSENT = [
    "По приказу Минздрава России от 03.05.2024 № 220н (с 01.09.2024):",
    "",
    "• Оказание первой помощи допускается, если отсутствует выраженный до начала",
    "  оказания первой помощи отказ гражданина или его законного представителя.",
    "• Первая помощь оказывается при условии отсутствия угрожающих жизни и здоровью",
    "  оказывающего ее лица факторов.",
    "• Первоочередность при двух и более пострадавших — исходя из тяжести состояния;",
    "  приоритет должен отдаваться детям (несовершеннолетним).",
    "• Разрешено применение автоматических наружных дефибрилляторов (АНД) при наличии.",
    "• Может использоваться помощь пострадавшему в принятии лекарственных препаратов,",
    "  назначенных врачом; допустимы подручные средства.",
]

KITS_ORDERS = [
    "Актуальные требования к комплектации (приказы Минздрава России от 24.05.2024):",
    "",
    "• № 260н — аптечка для оказания первой помощи пострадавшим в ДТП (автомобильная);",
    "• № 262н — аптечка для оказания работниками первой помощи пострадавшим;",
    "• № 261н — аптечка для организаций, осуществляющих образовательную деятельность",
    "  (введена впервые).",
    "",
    "Типовые компоненты: жгут; бинты; салфетки стерильные; лейкопластыри; устройство",
    "для ИВЛ «Рот–Устройство–Рот»; ножницы; перчатки медицинские; маска медицинская;",
    "покрывало спасательное изотермическое.",
    "",
    "Пополнять аптечку — по мере расхода и/или истечения срока годности.",
    "Замена компонентов автомобильной аптечки не допускается; для личного пользования",
    "водитель может дополнительно хранить свои лекарства/изделия.",
]

PRIORITY_LINES = [
    "Если пострадавших несколько, а участников оказания помощи недостаточно —",
    "определяют приоритетность:",
    "",
    "• в первую очередь — наиболее тяжело пострадавшим и несовершеннолетним детям;",
    "• для взрослых приоритетность определяется последовательностью мероприятий",
    "  Порядка оказания первой помощи (приказ № 220н);",
    "• в ряде случаев допустима самопомощь (например, прямое давление на рану),",
    "  пока оказывается помощь другому пострадавшему;",
    "• более опытный участник может координировать действия остальных,",
    "  направляя их к наиболее тяжелым пострадавшим.",
]

ORDER_220_STRUCTURE = [
    "Порядок оказания первой помощи (приказ Минздрава № 220н) включает:",
    "",
    "• общие организационные положения;",
    "• перечень из 9 состояний (приложение № 1) — в т.ч. острые психологические",
    "  реакции на стресс; судорожный приступ с потерей сознания; укусы/ужаливания;",
    "• перечень из 9 мероприятий и последовательность их проведения (приложение № 2).",
    "",
    "Важные изменения относительно прежнего приказа № 477н:",
    "• не проверяют пульс для оценки кровообращения;",
    "• из перечня мероприятий убраны пальцевое прижатие артерии и максимальное",
    "  сгибание конечности в суставе как обязательные техники для широкого обучения.",
]

CALL_EXTRA = [
    "Единый номер экстренных служб — 112 (также 101, 102, 103 и региональные номера).",
    "",
    "Поводы к вызову скорой медицинской помощи (в т.ч. не все входят в перечень",
    "состояний первой помощи, но требуют вызова СМП): нарушения сознания, дыхания,",
    "кровообращения; психические расстройства с опасностью для себя/окружающих;",
    "болевой синдром; травмы/отравления/ранения; термические и химические ожоги;",
    "кровотечения; роды, угроза прерывания беременности.",
    "",
    "Сообщить диспетчеру: место и суть происшествия; число пострадавших, повреждения,",
    "тяжесть; какая помощь оказывается.",
    "Трубку отключать после сообщения диспетчера о том, что вызов принят.",
    "Диспетчер/сотрудник СМП может дать команды по оказанию первой помощи.",
]

SAFETY_MASK = [
    "Дополнение по пособию Минздрава (2025):",
    "",
    "• В аптечке для оказания первой помощи работниками есть медицинские маски —",
    "  для снижения риска инфицирования оказывающего помощь.",
    "• Эти маски не используются для проведения искусственного дыхания",
    "  (для ИВЛ — отдельное устройство «Рот–Устройство–Рот»).",
    "• Перчатки медицинские — защита от контакта с кровью и другими биологическими",
    "  жидкостями пострадавшего.",
]


def process_file(spec):
    src = SRC / spec["src"]
    out = OUT / spec["out"]
    shutil.copy(src, out)
    prs = Presentation(str(out))

    # Title slide
    set_title_like(prs.slides[0], spec["title"])
    update_existing_eyebrows(prs, spec["eyebrow"])

    n = spec["num"]
    changes = [f"Присвоен порядковый номер {n} по списку темы 1.", f"Титул: «{spec['title']}»."]

    if n == 1:
        # Fix 3→5 components on org slide if present
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                t = sh.text_frame.text
                if "трех основных компонентов" in t or "трёх основных компонентов" in t:
                    new_t = t.replace("трех основных компонентов", "пяти основных компонентов").replace(
                        "трёх основных компонентов", "пяти основных компонентов"
                    )
                    # full replace with guide text is cleaner via new slide
                    replace_shape_text(sh, new_t)
                    changes.append("В тексте организации: «трех» → «пяти» компонентов (по методичке 2025).")
        insert_after_title(
            prs,
            "ПЯТЬ КОМПОНЕНТОВ СИСТЕМЫ ПЕРВОЙ ПОМОЩИ",
            spec["eyebrow"],
            ORG_5_COMPONENTS,
            "Дополнено по учебному пособию Минздрава, 2025.",
        )
        insert_content_before_thanks(
            prs,
            "ПОРЯДОК № 220н: СОГЛАСИЕ, ПРИОРИТЕТ, АНД",
            spec["eyebrow"],
            ORG_220_CONSENT,
            "Дополнено по приказу Минздрава № 220н / пособию 2025.",
        )
        changes.append("Добавлены слайды: 5 компонентов системы; согласие/приоритет/АНД по № 220н.")
        # Retitle notion slide contextually if exists
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == "ПОНЯТИЕ ПЕРВОЙ ПОМОЩИ":
                    # keep definition but note it's reference — definition still legally valid
                    pass

    elif n == 2:
        set_title_like(prs.slides[0], spec["title"])
        insert_after_title(
            prs,
            "ПРИКАЗЫ № 260н, 261н, 262н И СОСТАВ АПТЕЧЕК",
            spec["eyebrow"],
            KITS_ORDERS,
            "Дополнено по пособию Минздрава 2025 (ранее в слайдах не было № 261н и перечня компонентов).",
        )
        changes.append("Титул приведён к формулировке «аптечки, укладки, комплекты и наборы».")
        changes.append("Добавлен слайд: приказы 260н/261н/262н и типовые компоненты.")

    elif n == 3:
        insert_after_title(
            prs,
            "ПРИОРИТЕТНОСТЬ ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ",
            spec["eyebrow"],
            PRIORITY_LINES,
            "Дополнено по разделу «Приоритетность…» учебного пособия Минздрава 2025.",
        )
        changes.append("Добавлен слайд приоритетности (дети/тяжесть/последовательность/самопомощь/координация).")

    elif n == 4:
        insert_after_title(
            prs,
            "СТРУКТУРА ПОРЯДКА ОКАЗАНИЯ ПЕРВОЙ ПОМОЩИ (№ 220н)",
            spec["eyebrow"],
            ORDER_220_STRUCTURE,
            "Дополнено по пособию Минздрава 2025.",
        )
        changes.append("Добавлен слайд о структуре приказа № 220н и ключевых изменениях.")

    elif n == 5:
        set_title_like(prs.slides[0], spec["title"])
        insert_content_before_thanks(
            prs,
            "МАСКА И ПЕРЧАТКИ ИЗ АПТЕЧКИ",
            spec["eyebrow"],
            SAFETY_MASK,
            "Дополнено по пособию Минздрава 2025.",
        )
        changes.append("Титул сокращён под новую формулу (без акцента только на «личную безопасность»).")
        changes.append("Добавлен слайд про маску/перчатки из аптечки.")

    elif n == 6:
        insert_after_title(
            prs,
            "НОМЕР 112 И ПОВОДЫ К ВЫЗОВУ СМП",
            spec["eyebrow"],
            CALL_EXTRA,
            "Дополнено по пособию Минздрава 2025.",
        )
        changes.append("Добавлен слайд: 112/101–103, поводы к вызову, правила завершения разговора.")

    tmp = Path("/tmp") / spec["out"]
    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    renumber_pages(prs)
    prs.save(str(out))
    return changes, len(prs.slides)


def write_table(all_changes):
    wb = Workbook()
    ws = wb.active
    ws.title = "Нумерация и дополнения"
    ws.append(["№", "Новое имя файла", "Пункт списка темы 1", "Было (исходное имя)", "Что изменено / дополнено", "Слайдов"])
    src_names = {
        1: "1 Организационно-правовые аспекты…pptx",
        2: "1.2_Современные наборы…pptx",
        3: "1.3_Общая последовательность… / Порядок и приоритетность.pptx",
        4: "1.1_Понятие первая помощь.pptx",
        5: "1.4_Соблюдение правил личной безопасности.pptx",
        6: "1.5_Основные правила вызова скорой…pptx",
    }
    for spec, (changes, nslides) in zip(FILES, all_changes):
        ws.append([
            spec["num"],
            spec["out"],
            spec["list_item"],
            src_names[spec["num"]],
            " ".join(changes),
            nslides,
        ])
    ws2 = wb.create_sheet("Сверка с методичкой 2025")
    ws2.append(["Тема методички", "Было в презентациях", "Стало / дополнено"])
    rows = [
        ("Система ПП: 5 компонентов", "Указано «три основных компонента»", "Исправлено + отдельный слайд с пятью компонентами"),
        ("Приказы 260н, 261н, 262н; аптечка для образования", "Только визуалы авто/работникам", "Добавлен слайд с приказами и типовым составом; № 261н"),
        ("Приоритетность (дети, тяжесть, координация)", "Не было отдельного слайда", "Добавлен слайд в файл № 3"),
        ("Структура приказа № 220н и ключевые изменения", "Только перечень состояний/мероприятий", "Добавлен слайд в файл № 4"),
        ("Согласие на ПП, АНД, приоритет детям (п. 8 Порядка)", "Не раскрыто", "Добавлен слайд в файл № 1"),
        ("Вызов: 112, поводы a–и, не класть трубку до принятия вызова", "Кратко: что сказать диспетчеру", "Добавлен слайд в файл № 6"),
        ("Маска из аптечки ≠ для ИВЛ", "Не было", "Добавлен слайд в файл № 5"),
    ]
    for r in rows:
        ws2.append(list(r))
    for ws in wb.worksheets:
        for col in ws.columns:
            maxlen = max(len(str(c.value or "")) for c in col)
            ws.column_dimensions[col[0].column_letter].width = min(maxlen + 2, 60)
    wb.save(XLSX)
    shutil.copy(XLSX, OUT / XLSX.name)


def write_readme(all_changes):
    lines = [
        "# Тема 1 — нумерация по списку и дополнения по методичке Минздрава 2025",
        "",
        "Порядковые номера присвоены **по порядку тем темы 1** (первый пункт —",
        "«Организация оказания первой помощи в РФ. Нормативно-правовая база»).",
        "",
        "Список соответствует оглавлению темы 1 учебного пособия Минздрава «Первая помощь» (2025)",
        "и структуре Приложения № 2 к Порядку № 2464 (темы первой помощи).",
        "",
        "| № | Файл |",
        "|---|---|",
    ]
    for spec in FILES:
        lines.append(f"| {spec['num']} | `{spec['out']}` |")
    lines += ["", "## Дополнения по методичке", ""]
    for spec, (changes, nslides) in zip(FILES, all_changes):
        lines.append(f"### {spec['num']}. {spec['out']} ({nslides} слайдов)")
        for c in changes:
            lines.append(f"- {c}")
        lines.append("")
    README.write_text("\n".join(lines), encoding="utf-8")


def main():
    if OUT.exists():
        shutil.rmtree(OUT)
    OUT.mkdir(parents=True)

    all_changes = []
    for spec in FILES:
        print(f"Processing {spec['num']}: {spec['out']}")
        changes, nslides = process_file(spec)
        all_changes.append((changes, nslides))
        for c in changes:
            print("  -", c)

    write_table(all_changes)
    write_readme(all_changes)

    if ZIP.exists():
        ZIP.unlink()
    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in sorted(OUT.glob("*.pptx")):
            zf.write(p, p.name)
        zf.write(XLSX, XLSX.name)
        zf.write(README, README.name)
    print(f"\nZIP: {ZIP}")
    print(f"XLSX: {XLSX}")


if __name__ == "__main__":
    main()
