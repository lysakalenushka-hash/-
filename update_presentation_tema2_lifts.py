#!/usr/bin/env python3
"""Tema 2 lifts: general design — prune quiz, fill gaps from 2.docx."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________cc39.pptx"
LIST_TMPL_SRC = "/workspace/Презентация_Тема_1_Подъемники.pptx"
TMP = "/tmp/tema2_lifts_step1.pptx"
OUT = "/workspace/Презентация_Тема_2_Подъемники.pptx"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep all except quiz slide 3; drop duplicate purpose slide 5
KEEP = {0, 1, 2, 4, 6, 7, 8, 9}


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10271455), Emu(6492240), Emu(1200000), Emu(228600)
            )
            set_text(box, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    leftish = [s for s in content if (s.left or 0) < Emu(5500000)]
    hdr = leftish if len(leftish) >= 4 else content
    set_text(hdr[0], section, 11, True, RED)
    set_text(hdr[1], title, 24, True, DARK)
    set_text(hdr[2], intro, 12, False, GRAY)
    set_text(hdr[3], list_title, 11, True, RED)

    pairs, more_shape, note_shape = [], None, None
    for sh in content[4:]:
        if (sh.left or 0) >= Emu(5500000):
            if sh.top < Emu(6200000) and sh.text_frame.text.strip():
                set_text(sh, "", 12, False, DARK)
            continue
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    shown = items[:4]
    for i, item in enumerate(shown):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i + 1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(shown), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    if note_shape is None:
        for sh in content[4:]:
            if (
                sh.top > Emu(5600000)
                and sh.left < Emu(5500000)
                and sh not in [p[0] for p in pairs] + [p[1] for p in pairs if p[1]]
            ):
                note_shape = sh
                break
    if note_shape:
        if note:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            extra = items[4:]
            set_text(
                note_shape,
                ("· " + " · ".join(extra)) if extra else "",
                11,
                True,
                AMBER if extra else GRAY,
            )
    elif note:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, note, 12, True, AMBER)
    elif len(items) > 4:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, "· " + " · ".join(items[4:]), 11, True, AMBER)


def find_templates(prs):
    text_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if 5 <= n <= 8:
            titles = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(len(t) > 40 for t in titles):
                text_tpl = slide
                break
    if text_tpl is None:
        text_tpl = prs.slides[min(1, len(prs.slides) - 1)]

    list_tpl = None
    try:
        donor = Presentation(LIST_TMPL_SRC)
        for slide in donor.slides:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            ones = sum(1 for t in texts if t == "01")
            if ones == 1 and any(t == "02" for t in texts) and len(texts) >= 10:
                list_tpl = slide
                break
    except Exception:
        list_tpl = None
    if list_tpl is None:
        list_tpl = text_tpl
    return text_tpl, list_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def label_first(slide, label):
    try:
        shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], label, 11, True, RED)
    except Exception:
        pass


def enrich_kept(prs):
    # After prune KEEP order: 0,1,2,4,6,7,8,9 → indices 0..7
    shapes = sorted(
        [s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: (s.top, s.left),
    )
    if len(shapes) >= 4:
        set_text(shapes[0], "Общее устройство", 36, True, DARK)
        set_text(shapes[1], "подъёмников (вышек)", 36, True, RED)
        set_text(
            shapes[2],
            "Тема 2 · Конструкция, классификация, параметры, устойчивость, опоры и люлька",
            13,
            False,
            GRAY,
        )
        if len(shapes) >= 4:
            set_text(
                shapes[3],
                "Подъёмники · Вышки · Параметры · Аутригеры · Рабочее оборудование",
                12,
                False,
                GRAY,
            )

    fill_text_slide(
        prs.slides[1],
        "2 · Определение",
        "Подъёмники и вышки как грузоподъёмные машины",
        "Назначение класса машин",
        "Подъёмники и вышки — грузоподъёмные машины для вертикального или наклонного "
        "перемещения грузов и людей с одного уровня на другой в люльках и рабочих площадках, "
        "установленных на рабочем оборудовании. "
        "Ходовая часть — шасси автомобиля, спецшасси, гусеничное, пневмоколёсное или железнодорожное. "
        "Машины оснащают системами управления, сигнализацией, указателями, ограничителями и регистраторами.",
        note="Автомобильный подъёмник (вышка) — наиболее распространённый вид на шасси автомобиля.",
    )

    # slide 2 — figure with caption: light label only
    label_first(prs.slides[2], "2 · Схема устройства")

    # slide 3 (reach) — keep visual callout layout, only relabel
    label_first(prs.slides[3], "2 · Параметр «вылет»")

    label_first(prs.slides[4], "2 · Виды машин")

    # classification dump slide — short overview if enough text shapes
    try:
        shapes = sorted(text_shapes(prs.slides[5]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        if len(content) >= 4:
            fill_text_slide(
                prs.slides[5],
                "2 · Классификация",
                "Признаки классификации подъёмников (вышек)",
                "Назначение · конструкция колен · привод · степень поворота",
                "По назначению: общего назначения (строительно-монтажные работы) и специальные "
                "(пожарные, устройство линий связи и др.). "
                "По конструкции колен: одно-, двух-, трёхколенные и телескопические. "
                "По приводу: гидравлический, электрический, механический, электрогидравлический. "
                "По повороту: полноповоротный, неполноповоротный, неповоротный (вышка).",
                note="Детали по каждому признаку — на следующих слайдах.",
            )
        else:
            label_first(prs.slides[5], "2 · Классификация")
            if len(content) >= 2:
                set_text(content[1], "Признаки классификации подъёмников (вышек)", 22, True, DARK)
            if len(content) >= 3:
                set_text(
                    content[2],
                    "По назначению (общие / специальные); по конструкции колен; по приводу; "
                    "по степени поворота (полноповоротный, неполноповоротный, неповоротный).",
                    13,
                    False,
                    DARK,
                )
    except Exception:
        label_first(prs.slides[5], "2 · Классификация")

    try:
        fill_text_slide(
            prs.slides[6],
            "2 · Автоподъёмники",
            "Автомобильные подъёмники и автовышки",
            "Ведущее звено системы строительно-монтажных машин",
            "Автомобильные подъёмники и вышки обеспечивают комплексную механизацию строительства "
            "и высокие темпы работ. Ходовая часть — на шасси автомобиля или специального шасси "
            "автомобильного типа. Перед работой устанавливают выносные опоры для устойчивости.",
            note="Пример: автомобильный подъёмник ПСС 131.28Э на базе КАМАЗ.",
        )
    except Exception:
        label_first(prs.slides[6], "2 · Автоподъёмники")

    label_first(prs.slides[7], "2 · Пример: ПСС 131.28Э")


def append_missing(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    enrich_kept(prs)

    s = nl()
    fill_list_slide(
        s,
        "2 · Подъёмник и вышка",
        "Ключевое конструктивное отличие",
        "Тип рабочего оборудования определяет класс машины",
        "Сравнение",
        [
            "подъёмник — одно, два или более шарнирно сочленённых колена; люлька на оголовке верхнего колена",
            "вышка — мачта (телескопическая, ножничная и др.) с люлькой (платформой) наверху",
            "вышка предназначена только для вертикального подъёма людей",
            "оба типа имеют системы управления, сигнализацию, указатели и ограничители безопасности",
        ],
        note="Именно тип рабочего оборудования отличает подъёмник от вышки.",
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Конструкция колен",
        "Классификация по конструкции колен (стрелы)",
        "Рис. классификации: одно-, двух-, трёхколенный и телескопический",
        "Типы",
        [
            "одноколенный — одно колено (секция стрелы)",
            "двухколенный — два шарнирно соединённых колена",
            "трёхколенный — три колена для увеличения зоны обслуживания",
            "телескопический — выдвижные секции стрелы/мачты",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Привод и поворот",
        "Классификация по приводу и степени поворота",
        "Дополнительные признаки классификации",
        "Виды",
        [
            "привод: гидравлический, электрический, механический, электрогидравлический",
            "полноповоротный — вращение поворотной части с люлькой ≥ 360°",
            "неполноповоротный — ограниченный угол поворота в плане",
            "неповоротный — типичная схема вышки без вращения поворотной части",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Основные параметры",
        "Параметры и характеристики подъёмников (вышек)",
        "По параметрам определяют, какую работу можно выполнять",
        "Основные параметры",
        [
            "грузоподъёмность — наибольшая допускаемая масса груза",
            "высота подъёма H — наибольшее расстояние по вертикали от основания до пола люльки",
            "вылет — горизонтальное расстояние до рабочей площадки (зоны охвата)",
            "габариты, база, колея, опорный контур, глубина опускания, время подъёма люльки",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Устойчивость",
        "Сведения об устойчивости свободностоящих машин",
        "Устойчивость обеспечивается собственной силой тяжести",
        "Ключевые понятия",
        [
            "при работе машина опирается на колёса или выносные опоры; ребро опрокидывания — край опоры",
            "удерживающий момент МУ = G·б; опрокидывающий момент Могр = Q·а (плюс ветер и динамика)",
            "коэффициент собственной устойчивости — отношение удерживающего момента к опрокидывающему",
            "потеря устойчивости обычно результат нарушения правил эксплуатации",
        ],
        note="Подъёмники — машины повышенной опасности; необходим запас устойчивости.",
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Выносные опоры",
        "Аутригеры: назначение и требования",
        "Увеличение опорной поверхности и устойчивости при работе",
        "Правила",
        [
            "конструкции опор: поворотные, откидные, выдвижные",
            "усилие поднятия/выдвижения вручную — не более 200 Н; иначе — гидропривод или иной привод",
            "опоры устанавливают до начала работ; без опор работа на максимальном вылете недопустима",
            "опоры предохраняют от самопроизвольного выдвижения в транспортном положении",
        ],
        note="Выносные опоры обязательны к установке перед началом работы.",
    )

    s = nl()
    fill_list_slide(
        s,
        "2 · Рабочее оборудование",
        "Колена, механизмы и люлька",
        "Состав рабочего оборудования подъёмника",
        "Элементы",
        [
            "одно–три колена, шарнирно соединённых; механизмы подъёма и поворота колен",
            "люлька: пол (несущая площадка), ограждение по периметру, проём для входа/выхода",
            "раздвижная платформа — увеличение площади пола выдвижением вспомогательной секции",
            "поворотная / неповоротная люлька; электроизолированная — сопротивление изоляции ≥ 0,5 МОм",
        ],
        note="Электроизолированная платформа применяется при работах вблизи токоведущих частей.",
    )

    s = nl()
    fill_list_slide(
        s,
        "Итоги темы 2",
        "Общее устройство подъёмников (вышек)",
        "Что должен знать машинист",
        "Главные выводы",
        [
            "подъёмник — шарнирные колена; вышка — мачта для вертикального подъёма",
            "классификация: назначение, конструкция колен, привод, степень поворота",
            "параметры: грузоподъёмность, высота, вылет; устойчивость и аутригеры критичны для безопасности",
            "люлька: пол, ограждение, проём; возможны раздвижные, поворотные и электроизолированные платформы",
        ],
        note="Нарушение правил эксплуатации — типичная причина потери устойчивости.",
    )


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order length {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for item in new_items:
        sldIdLst.append(item)


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    prs.save(TMP)

    prs = Presentation(TMP)
    append_missing(prs)
    # kept 0..7: title, def, scheme, reach, kinds, class overview, auto, PSS example
    # appended 8 lift/tower, 9 knees, 10 drive/rotate, 11 params, 12 stability,
    # 13 outriggers, 14 cradle, 15 summary
    order = [
        0,  # title
        1,  # definition
        2,  # scheme figure
        8,  # lift vs tower
        4,  # kinds comparison
        6,  # auto lifts text
        7,  # PSS figure
        5,  # classification overview
        9,  # knees
        10,  # drive/rotate
        3,  # reach
        11,  # parameters
        12,  # stability
        13,  # outriggers
        14,  # working equipment
        15,  # summary
    ]
    reorder_slides(prs, order)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
