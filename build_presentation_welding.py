#!/usr/bin/env python3
"""Rebuild welding presentation with proper red Inter list layout."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_Сварочные_работы.pptx"
IMG = Path("/workspace/assets/welding_images")

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def is_bg(shape):
    """Full-bleed background rectangle — never put content here."""
    if shape.top is None or shape.left is None:
        return False
    return shape.top < Emu(50000) and (shape.width or 0) > Emu(10000000)


def is_accent_bar(shape):
    """Thin left red bar added for branding."""
    if shape.left is None or shape.width is None:
        return False
    return shape.left < Emu(20000) and shape.width < Emu(200000)


def is_footer(shape):
    return (
        shape.top is not None
        and shape.left is not None
        and shape.top > Emu(6200000)
        and shape.left > Emu(8000000)
    )


def content_shapes(slide):
    shapes = [
        s
        for s in slide.shapes
        if s.has_text_frame and not is_bg(s) and not is_footer(s) and not is_accent_bar(s)
    ]
    return sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def add_red_bar(slide):
    """Thin left accent bar — signature of the red template look."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(90000), Emu(6858000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    if bar.has_text_frame:
        bar.text_frame.clear()


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame or is_bg(sh) or is_accent_bar(sh):
                continue
            if is_footer(sh):
                set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                updated = True
                break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10271455), Emu(6492240), Emu(1200000), Emu(228600)
            )
            set_text(box, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = content_shapes(slide)
    for sh in shapes:
        set_text(sh, "", 12, False, DARK)
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 26, True, DARK)
    set_text(shapes[2], intro, 13, False, GRAY)
    set_text(shapes[3], body, 14, False, DARK)
    if len(shapes) > 4:
        set_text(shapes[4], note or "", 13, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None, start_num=1):
    """Red numbered list layout (01–04). items must be ≤4 — split outside."""
    assert len(items) <= 4, f"list slide max 4 items, got {len(items)}"
    shapes = content_shapes(slide)

    # Detect structure BEFORE clearing (numbers are short "01"/"02" text)
    hdr_candidates = []
    pairs, more_shape, note_shape = [], None, None
    for sh in shapes:
        t = sh.text_frame.text.strip()
        if (sh.left or 0) >= Emu(5500000):
            continue
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and (sh.left or 0) > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif (sh.top or 0) > Emu(5600000):
            note_shape = sh
        elif (sh.top or 0) < Emu(2800000):
            hdr_candidates.append(sh)

    hdr_candidates = sorted(hdr_candidates, key=lambda s: (s.top or 0, s.left or 0))
    if len(pairs) < 2:
        nums = sorted(
            [
                s
                for s in shapes
                if (s.left or 0) < Emu(900000) and Emu(2700000) < (s.top or 0) < Emu(5600000)
            ],
            key=lambda s: s.top or 0,
        )
        texts = sorted(
            [
                s
                for s in shapes
                if Emu(900000) < (s.left or 0) < Emu(5500000)
                and Emu(2700000) < (s.top or 0) < Emu(5600000)
            ],
            key=lambda s: s.top or 0,
        )
        pairs = [[n, texts[i] if i < len(texts) else None] for i, n in enumerate(nums[:4])]

    if len(hdr_candidates) < 4:
        hdr_candidates = sorted(
            [s for s in shapes if (s.left or 0) < Emu(5500000) and (s.top or 0) < Emu(2800000)],
            key=lambda s: (s.top or 0, s.left or 0),
        )

    # Wipe ALL content (kills leftover PS copy from template)
    for sh in shapes:
        set_text(sh, "", 12, False, DARK)
    for sh in slide.shapes:
        if sh.has_text_frame and not is_bg(sh) and not is_footer(sh) and (sh.left or 0) >= Emu(5500000):
            set_text(sh, "", 12, False, DARK)

    set_text(hdr_candidates[0], section, 11, True, RED)
    set_text(hdr_candidates[1], title, 24, True, DARK)
    set_text(hdr_candidates[2], intro, 13, False, GRAY)
    set_text(hdr_candidates[3], list_title, 12, True, RED)

    for i, item in enumerate(items):
        if i >= len(pairs):
            break
        num_sh, txt_sh = pairs[i]
        set_text(num_sh, f"{start_num + i:02d}", 16, True, RED)
        if txt_sh is not None:
            set_text(txt_sh, item, 13, False, DARK)
    for i in range(len(items), len(pairs)):
        set_text(pairs[i][0], "", 16, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 13, False, DARK)

    if more_shape:
        set_text(more_shape, "", 9, False, GRAY)

    if note:
        if note_shape:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            box = slide.shapes.add_textbox(
                Emu(731520), Emu(5850000), Emu(10000000), Emu(400000)
            )
            set_text(box, note, 12, True, AMBER)


def add_image(slide, path, side=True):
    path = Path(path)
    if not path.exists() or path.stat().st_size < 5000:
        return
    if side:
        # Right side, below title — leave room for list on left
        slide.shapes.add_picture(
            str(path), Emu(7200000), Emu(2200000), width=Emu(4500000), height=Emu(3600000)
        )
    else:
        slide.shapes.add_picture(
            str(path), Emu(731520), Emu(4000000), width=Emu(10700000), height=Emu(2000000)
        )


def chunks(items, n=4):
    for i in range(0, len(items), n):
        yield items[i : i + n]


def main():
    tpl = Presentation(TEMPLATE)
    slides = list(tpl.slides)
    text_elems = [deepcopy(s.element) for s in slides[0].shapes]
    list_elems = [deepcopy(s.element) for s in slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    def text_slide():
        return make_slide(prs, text_elems)

    def list_slide():
        return make_slide(prs, list_elems)

    # —— 1. Title ——
    s = text_slide()
    fill_text_slide(
        s,
        "ФНП · Сварочное производство",
        "Общие требования к производству сварочных работ на ОПО",
        "Персонал · аттестация · ПТД · организация · контроль · документация",
        "Производство сварочных работ — деятельность персонала сварочного производства "
        "с применением сварочных и родственных процессов, сварочных материалов и оборудования "
        "с соблюдением нормативных требований и проектной (конструкторской) документации.",
        note="На ОПО сварку выполняют только по аттестованным технологиям.",
    )
    add_image(s, IMG / "arc.jpg")

    # —— 2. Personnel ——
    s = list_slide()
    fill_list_slide(
        s,
        "Общие требования · Персонал",
        "Обязанности персонала сварочного производства",
        "Сварщики, операторы, специалисты и контролёры обеспечивают:",
        "Обязанности",
        [
            "техническую и технологическую подготовку и выполнение сварочных работ по ФНП и НД",
            "безопасную эксплуатацию, обслуживание и ремонт сварочного оборудования",
            "соблюдение технологий сварки",
            "контроль качества сварных соединений",
        ],
    )

    # —— 3. Certification personnel ——
    s = list_slide()
    fill_list_slide(
        s,
        "Аттестация · Персонал",
        "Квалификация и аттестация сварщиков и специалистов",
        "Допуск к сварочным работам на ОПО",
        "Требования",
        [
            "квалификация соответствует видам работ и применяемым технологиям сварки",
            "сварщики и специалисты аттестованы в установленном порядке",
            "допуск — только по положительным результатам аттестационных испытаний",
            "аттестация подтверждает способы сварки, виды конструкций, положения и материалы",
        ],
        note="Без положительных результатов аттестационных испытаний к работам не допускают.",
    )
    add_image(s, IMG / "ppe_weld.jpg")

    # —— 4. Technology readiness ——
    s = list_slide()
    fill_list_slide(
        s,
        "Аттестация · Технологии",
        "Проверка готовности к применению аттестованных технологий",
        "Организации и индивидуальные предприниматели",
        "Порядок",
        [
            "работы выполняют организации (ИП), прошедшие проверку готовности технологий на ОПО",
            "проверка подтверждает технические, организационные и квалификационные возможности",
            "выполняют контрольные сварные соединения (в т. ч. наплавку) для оценки характеристик",
            "положительные результаты оформляют документом с областью допуска",
        ],
    )

    # —— 5. Organization ——
    s = list_slide()
    fill_list_slide(
        s,
        "Организация работ",
        "Кто обеспечивает организацию и выполнение",
        "Ответственность руководства",
        "Ключевые правила",
        [
            "организацию обеспечивает руководитель организации (ИП) или уполномоченное лицо",
            "аттестационные процедуры — руководитель независимого аттестационного центра",
            "работы — по ПТД, разработанной специалистом сварочного производства",
            "ПТД — на основании проектной документации и требований НПА/НД РФ",
        ],
    )

    # —— 6–10. Full PTD (18 items → 5 slides × ≤4) ——
    ptd_items = [
        "способы сварки",
        "требования к квалификации, аттестации и допускным испытаниям сварщиков",
        "требования к сборке, прихваткам и приварке временных технологических креплений",
        "конструкция нестандартизированных сварных соединений (в т. ч. разной толщины)",
        "требования к хранению и подготовке сварочных материалов",
        "сочетания марок основных и сварочных материалов",
        "типоразмеры сварочных материалов (электрод, проволока, лента)",
        "используемое сварочное оборудование",
        "род и полярность сварочного тока",
        "типы выполняемых сварных соединений",
        "режимы сварки для конкретных сварных соединений",
        "необходимость, методы и режимы предварительного и сопутствующего подогрева",
        "пространственные положения при сварке",
        "порядок и последовательность выполнения шва (наплавки)",
        "способы защиты зоны сварки",
        "порядок и способы маркировки сварных соединений",
        "методы и объёмы НК и механических испытаний (при наличии требований НД)",
        "требования к исправлению дефектов и контролю после исправления",
    ]
    parts = list(chunks(ptd_items, 4))
    start = 1
    for idx, part in enumerate(parts, 1):
        s = list_slide()
        fill_list_slide(
            s,
            f"ПТД · Содержание ({idx}/{len(parts)})",
            "Производственно-технологическая документация по сварке",
            "В технологических инструкциях и картах сварки устанавливают:",
            "Содержание ПТД",
            part,
            note=(
                "Конструктивные элементы, режимы и контроль обеспечивают качество соединений."
                if idx == len(parts)
                else None
            ),
            start_num=start,
        )
        start += len(part)

    # —— Assembly (6 → 4+2) ——
    asm = [
        "способы подготовки поверхностей деталей под сварку",
        "приспособления и оборудование для сборки",
        "порядок и последовательность сборки; способы крепления деталей",
        "способы сварки, материалы и режимы прихваток и временных креплений",
        "размеры, количество и расположение прихваток",
        "методы контроля качества сборки",
    ]
    start = 1
    for idx, part in enumerate(chunks(asm, 4), 1):
        s = list_slide()
        fill_list_slide(
            s,
            f"Сборка · Требования ПТД ({idx}/2)",
            "Что должно быть в требованиях по сборке под сварку",
            "Состав требований к сборке деталей",
            "Параметры сборки",
            part,
            note=(
                "Стыковые кольцевые швы собирают с соосным позиционированием деталей."
                if idx == 2
                else None
            ),
            start_num=start,
        )
        if idx == 1:
            add_image(s, IMG / "pipe_weld.jpg")
        start += len(part)

    # —— Equipment ——
    s = list_slide()
    fill_list_slide(
        s,
        "Оборудование и материалы",
        "Сварочное оборудование и сварочные материалы",
        "Соответствие аттестованным технологиям",
        "Требования",
        [
            "оборудование и материалы соответствуют применяемым аттестованным технологиям",
            "сварочно-технологические характеристики обеспечивают качество соединений",
            "соответствие технологиям и нормам качества подтверждают в установленном порядке",
            "оборудование содержат в исправном состоянии по указаниям производителя",
        ],
    )
    add_image(s, IMG / "mig.jpg")

    # —— Supervisor ——
    s = list_slide()
    fill_list_slide(
        s,
        "Руководство сварочными работами",
        "Обязанности руководителя перед началом работ",
        "Лицо, осуществляющее руководство сварочными работами:",
        "Перед выполнением работ",
        [
            "проверить состав и квалификацию персонала, оборудование, материалы и технологию по ПТД",
            "ознакомить сварщиков с технологическими картами и изменениями — под подпись",
            "организовать проведение операционного контроля",
        ],
    )

    # —— Admission / workplace ——
    s = list_slide()
    fill_list_slide(
        s,
        "Допуск сварщика и рабочее место",
        "Допускные испытания и комплектация места сварки",
        "Перед допуском к сварке на объекте",
        "Правила",
        [
            "при первом допуске или после перерыва дольше нормы НД — допускные соединения",
            "место работ комплектуют исправным оборудованием, оснасткой и инструментом по ПТД",
            "место сварки защищают от осадков, влаги, сквозняков и иных воздействий",
            "условия выполнения работ соответствуют требованиям ПТД и НД",
        ],
    )
    add_image(s, IMG / "mma.jpg")

    # —— Control types ——
    s = list_slide()
    fill_list_slide(
        s,
        "Контроль · Виды",
        "Виды контроля при подготовке и выполнении сварочных работ",
        "Обязательные виды контроля",
        "Виды",
        [
            "входной — все партии свариваемых и сварочных материалов до применения",
            "операционный — подготовка кромок, сборка, прихватка, сварка, послесварочная обработка",
            "приёмочный — все выполненные сварные соединения (методы и объём — по ПТД)",
        ],
        note="При трещинах или недопустимых дефектах работы останавливают до устранения причин.",
    )

    # —— Incoming (6 → 4+2) ——
    incoming = [
        "документы о качестве продукции (идентификация, материалы, результаты испытаний)",
        "маркировка на каждом упаковочном месте: марка, сортамент, номер партии",
        "отсутствие повреждений упаковки и самих материалов",
        "покрытые электроды — номинальные размеры и состояние покрытия по сертификату",
        "проволока и лента — размеры, вид и состояние поверхности, маркировка",
        "флюс — цвет, однородность и гранулометрический состав",
    ]
    start = 1
    for idx, part in enumerate(chunks(incoming, 4), 1):
        s = list_slide()
        fill_list_slide(
            s,
            f"Контроль · Входной ({idx}/2)",
            "Входной контроль материалов",
            "Что проверяют до применения партий материалов",
            "Проверки",
            part,
            start_num=start,
        )
        if idx == 1:
            add_image(s, IMG / "electrodes.jpg")
        start += len(part)

    # —— Operational (9 → 4+4+1) ——
    operational = [
        "наличие маркировки",
        "размеры деталей и форму разделки кромок",
        "качество подготовленных под сварку поверхностей (если регламентировано ПТД)",
        "марки и типоразмеры сварочных материалов для прихваток",
        "надёжность фиксации и расположение деталей в сборочных приспособлениях",
        "чистоту и отсутствие повреждений кромок и прилегающих поверхностей",
        "размеры и расположение прихваток и швов приварки временных креплений",
        "зазоры, смещение кромок, перелом осей или плоскостей после прихваток",
        "размеры собранной под сварку конструкции",
    ]
    parts = list(chunks(operational, 4))
    start = 1
    for idx, part in enumerate(parts, 1):
        s = list_slide()
        fill_list_slide(
            s,
            f"Контроль · Операционный ({idx}/{len(parts)})",
            "Операционный контроль сборки и подготовки",
            "Что проверяют на операциях подготовки и сборки",
            "Проверки",
            part,
            start_num=start,
        )
        start += len(part)

    # —— During welding ——
    s = list_slide()
    fill_list_slide(
        s,
        "Контроль · В процессе сварки",
        "Минимальный контроль во время сварки",
        "В процессе сварки проводят, как минимум:",
        "Контроль",
        [
            "соответствие параметров режима сварки и технологических приёмов",
            "очерёдность выполнения сварных швов и участков наплавки",
            "отсутствие видимых дефектов",
            "иные параметры по технологическим (операционным) картам сварки",
        ],
        note="Устранение дефектов — по ПТД; после исправления — повторный контроль.",
    )
    add_image(s, IMG / "tig.jpg")

    # —— Marking / repair ——
    s = list_slide()
    fill_list_slide(
        s,
        "Маркировка и ремонт",
        "Маркировка швов и исправление дефектов",
        "После выполнения и при выявлении дефектов",
        "Правила",
        [
            "сварные соединения маркируют по требованиям ПТД",
            "маркировка — шифр клейма сварщика для однозначной идентификации",
            "недопустимые дефекты исправляют по ПТД и карте на ремонт",
            "число исправлений одного участка — не выше указанного в ПТД (НД)",
        ],
    )

    # —— Documentation duties ——
    s = list_slide()
    fill_list_slide(
        s,
        "Документация · Руководитель",
        "Что обязан обеспечить руководитель сварочных работ",
        "При производстве сварочных работ",
        "Обязанности",
        [
            "идентификацию применяемых сварочных материалов и оборудования",
            "выполнение соединений по технологическим (операционным) картам сварки",
            "регистрацию сведений о сварщиках, выполняющих соединения",
            "идентификацию мест швов в конструкции и мест исправления дефектов",
        ],
        note="Также регистрируют результаты контроля качества, включая контроль после ремонта.",
    )

    # —— Executive docs / summary ——
    s = list_slide()
    fill_list_slide(
        s,
        "Итоги темы",
        "Исполнительная документация и ключевые требования",
        "В процессе сварочных работ оформляют:",
        "Документы",
        [
            "журналы сварочных работ",
            "паспорта, акты и заключения по неразрушающему контролю",
            "протоколы испытаний сварных соединений",
            "иную исполнительную и эксплуатационную документацию по ПТД и НД",
        ],
        note="Сварка на ОПО — аттестованным персоналом, по ПТД и с полным циклом контроля.",
    )
    add_image(s, IMG / "mpi.jpg")

    for slide in prs.slides:
        add_red_bar(slide)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
