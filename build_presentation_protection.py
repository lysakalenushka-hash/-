#!/usr/bin/env python3
"""Build protective equipment presentation using uploaded template styling."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_Средства_защиты.pptx"
IMG = "/workspace/assets/protection_images"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

IMG_LEFT = Emu(6600000)
IMG_W = Emu(5200000)
TEXT_W = Emu(5800000)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, footer=None, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], body, 13, False, DARK)
    if footer and len(shapes) > 4:
        set_text(shapes[4], footer, 11, False, GRAY)
    elif len(shapes) > 4:
        set_text(shapes[4], "", 11, False, GRAY)
    if note:
        for sh in shapes:
            if sh.top > Emu(5800000) and sh.left < Emu(5000000):
                set_text(sh, note, 13, True, AMBER)
                break


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    for sh in shapes[5:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            set_text(pairs[i][1], item, 13, False, DARK)

    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        set_text(pairs[i][1], "", 14, False, DARK)

    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)

    if note:
        for sh in shapes:
            if sh.top > Emu(5900000) and sh.left > Emu(700000):
                set_text(sh, note, 13, True, AMBER)
                break


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3600000), width=Emu(10700000), height=Emu(2500000))
    else:
        slide.shapes.add_picture(path, IMG_LEFT, Emu(1200000), width=IMG_W, height=Emu(4800000))


def add_cards(slide, items, cols=3):
    """Grid of short requirement cards."""
    n = len(items)
    rows = (n + cols - 1) // cols
    card_w = Emu(3600000) if cols == 3 else Emu(5400000)
    card_h = Emu(1100000)
    gap_x = Emu(200000)
    gap_y = Emu(150000)
    start_left = Emu(548640)
    start_top = Emu(2100000)
    for i, text in enumerate(items):
        row, col = divmod(i, cols)
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RED
        box.line.width = Pt(1.0)
        tf = slide.shapes.add_textbox(
            left + Emu(100000), top + Emu(150000), card_w - Emu(200000), card_h - Emu(250000)
        ).text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        set_run(r, 12, False, DARK)


def add_category_cards(slide, items):
    """Six category cards in 2 rows of 3."""
    positions = [
        (Emu(548640), Emu(2100000)),
        (Emu(4400000), Emu(2100000)),
        (Emu(8250000), Emu(2100000)),
        (Emu(548640), Emu(4200000)),
        (Emu(4400000), Emu(4200000)),
        (Emu(8250000), Emu(4200000)),
    ]
    card_w = Emu(3600000)
    card_h = Emu(1700000)
    for (left, top), text in zip(positions, items):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        box.line.width = Pt(1.0)
        # accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(80000), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED
        bar.line.fill.background()
        tf = slide.shapes.add_textbox(
            left + Emu(180000), top + Emu(500000), card_w - Emu(280000), Emu(800000)
        ).text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        set_run(r, 14, True, DARK)


def main():
    tpl = Presentation(TEMPLATE)
    text_tpl = [deepcopy(s.element) for s in tpl.slides[0].shapes]
    list_tpl = [deepcopy(s.element) for s in tpl.slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # 1. Title
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Охрана труда",
        "Требования к средствам защиты при эксплуатации трубопроводов пара и горячей воды (ПиГВ)",
        "Коллективные и индивидуальные средства защиты",
        "Модуль охватывает требования к средствам защиты, классификацию коллективных "
        "и индивидуальных средств, оградительные, предохранительные и блокирующие "
        "устройства, специальную одежду и обувь, защиту головы и органов дыхания.",
        note="В производственных условиях требуются специальные приспособления и СИЗ.",
    )
    add_image(s, f"{IMG}/ppe_head_respiratory.png")

    # 2. Intro quote
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Средства защиты",
        "Зачем нужны средства защиты",
        "Если бы к нашим телам прилагались инструкции…",
        "«Можно эксплуатировать в большинстве ситуаций… "
        "но в производственных условиях — требуются специальные приспособления».\n\n"
        "Средства защиты снижают риск травм и профессиональных заболеваний "
        "при обслуживании трубопроводов пара и горячей воды.",
    )
    add_image(s, f"{IMG}/worker_special_clothing.png")

    # 3. Requirements
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Требования",
        "Требования к средствам защиты",
        "Средства защиты должны соответствовать следующим требованиям",
        "",
    )
    shapes = sorted(text_shapes(s), key=lambda sh: (sh.top, sh.left))
    if len(shapes) > 3:
        set_text(shapes[3], "", 13, False, DARK)
    add_cards(
        s,
        [
            "Не снижать производительность оборудования",
            "Предотвращать контакт",
            "Не создавать новых опасностей",
            "Допускать обслуживание без снятия ограждения",
            "Обеспечивать безопасность",
            "Не создавать помех",
            "Соответствовать требованиям технической эстетики",
            "Закрывать от падающих предметов",
        ],
        cols=3,
    )

    # 4. Collective vs Individual
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Средства защиты",
        "Коллективные и индивидуальные средства защиты",
        "Средства защиты подразделяются на две основные группы:",
        "Виды средств защиты",
        [
            "коллективные — защищают всех работников в зоне (ограждения, сигнализация, предохранительные устройства)",
            "индивидуальные — СИЗ конкретного работника (каска, очки, респиратор, спецодежда, спецобувь)",
        ],
        note="Оба вида применяются совместно при работах на ПиГВ.",
    )

    # 5. Collective categories
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Коллективные средства защиты",
        "Виды коллективных средств защиты",
        "Коллективные средства защиты включают:",
        "",
    )
    shapes = sorted(text_shapes(s), key=lambda sh: (sh.top, sh.left))
    if len(shapes) > 3:
        set_text(shapes[3], "", 13, False, DARK)
    add_category_cards(
        s,
        [
            "Сигнализирующие",
            "Оградительные",
            "Предохранительные",
            "Блокирующие",
            "Специальные устройства",
            "Системы дистанционного управления оборудованием",
        ],
    )

    # 6. Protective and safety devices
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Защитные устройства",
        "Защитные и предохранительные устройства",
        "Защитные и предохранительные устройства подразделяются на:",
        "Типы устройств",
        [
            "ограждения, исключающие доступ к движущимся или опасным частям механизмов",
            "ограждения, защищающие от высоких и низких температур и вредных факторов",
            "ограждения, уменьшающие разрушения при взрыве или предупреждающие их",
            "приспособления, обеспечивающие безопасность пуска и остановки механизмов",
            "устройства и приборы, предупреждающие об опасности",
            "приборы и устройства, устраняющие опасность",
        ],
    )
    add_image(s, f"{IMG}/safety_barriers.png")

    # 7. Fencing devices
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Оградительные средства",
        "Оградительные средства защиты",
        "Оградительные устройства:",
        "Виды ограждений",
        [
            "стационарные (несъёмные)",
            "подвижные (съёмные)",
            "переносные",
        ],
        note="Ограждения обязательны при аварийных работах на теплосетях.",
    )
    add_image(s, f"{IMG}/site_fencing.png")

    # 8. Safety valves
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Предохранительные средства",
        "Предохранительные средства защиты",
        "Предохранительные клапаны:",
        "Типы предохранительных клапанов",
        [
            "А. Рычажные",
            "Б. Пружинные",
            "В. Мембранные",
        ],
        note="Клапаны сбрасывают давление при превышении допустимого значения.",
    )
    add_image(s, f"{IMG}/safety_valves.png", bottom=True)

    # 9. Locking devices
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Блокирующие устройства",
        "Блокирующие устройства",
        "Блокировочные устройства:",
        "Виды блокировок",
        [
            "механические — фланец-лок, тросовые блокировки, кожухи на штурвалы",
            "предотвращают несанкционированное или случайное включение оборудования",
            "применяются при ремонте и обслуживании трубопроводов и арматуры",
        ],
    )
    add_image(s, f"{IMG}/lockout_devices.png", bottom=True)

    # 10. Protection means overview
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Средства защиты",
        "Средства защиты работников",
        "К средствам защиты относятся:",
        "Категории",
        [
            "специальная одежда и специальная обувь",
            "смывающие и обезвреживающие средства",
            "технические средства",
        ],
    )

    # 11. Special clothing
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "СИЗ",
        "Специальная одежда и специальная обувь",
        "Спецодежда и спецобувь защищают от следующих опасностей:",
        "Защищаемые опасности",
        [
            "механические опасности",
            "опасности, связанные с воздействием микроклимата и климатические опасности",
            "опасности, связанные с воздействием химического фактора",
        ],
    )
    add_image(s, f"{IMG}/worker_special_clothing.png")

    # 12. Thermal coveralls
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "СИЗ",
        "Специальные комбинезоны и специальная одежда",
        "При работах с термическим воздействием применяют:",
        "Термическая защита",
        [
            "опасности, связанные с воздействием термического фактора",
            "специальные комбинезоны с теплозащитными свойствами",
            "защитные шлемы со щитками, термостойкие перчатки",
        ],
        note="При работах в зоне пара и горячей воды термозащита обязательна.",
    )
    add_image(s, f"{IMG}/thermal_coveralls.png")

    # 13. Foot protection
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "СИЗ ног",
        "Средства индивидуальной защиты ног",
        "Спецобувь защищает от следующих опасностей:",
        "Опасности для ног",
        [
            "механические опасности (поскальзывание, защита носочной части 200 Дж)",
            "опасности, связанные с воздействием микроклимата и климатические опасности",
            "термические опасности",
            "опасности, связанные с воздействием химического фактора",
        ],
    )
    add_image(s, f"{IMG}/safety_footwear.png", bottom=True)

    # 14. Head and respiratory
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "СИЗ головы и дыхания",
        "Защита головы и органов дыхания",
        "К средствам защиты головы и дыхания относятся:",
        "Средства защиты",
        [
            "защитная каска",
            "защитный щиток",
            "защитные очки",
            "противошумные вкладыши и наушники",
            "респираторы и противогазы",
        ],
    )
    add_image(s, f"{IMG}/ppe_head_respiratory.png")

    # 15. Consequences of no protection
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Последствия",
        "Результаты отсутствия надлежащей защиты головы",
        "Отсутствие каски и защиты головы приводит к:",
        "Последствия травм головы",
        [
            "скальпированные раны мягких тканей головы",
            "травмы черепа, переломы",
            "повреждения мозга",
        ],
        note="Защитная каска — обязательный элемент СИЗ на производстве.",
    )

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
