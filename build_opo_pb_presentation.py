#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Красный стиль текста и оформления — все картинки и макет сохраняются."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOURCE = Path("Общие_положения_ПБ_на_ОПО_исходник.pptx")
OUTPUT = Path("Общие_положения_ПБ_на_ОПО.pptx")

# Палитра как в материалах Б.7.5 / котлы
RED = RGBColor(0xE3, 0x06, 0x13)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
CREAM = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)
FONT = "Inter"

# Исходные цвета текста в презентации
OLD_TITLE_RED = RGBColor(0xFF, 0x00, 0x00)
OLD_HEAD = RGBColor(0x2E, 0x3C, 0x4E)
OLD_BODY = RGBColor(0x38, 0x46, 0x53)
OLD_BLACK = RGBColor(0x00, 0x00, 0x00)

# Заливки карточек в исходнике
OLD_CARD_BLUE = RGBColor(0xD9, 0xED, 0xF2)
OLD_CARD_CREAM = RGBColor(0xFA, 0xF9, 0xF5)

TITLE_SIZES = {304800, 285750, 266700}
HEADER_SIZE = 152400
NUMBER_SIZES = {190500, 184150, 171450}


def rgb_eq(a: RGBColor, b: RGBColor) -> bool:
    return str(a) == str(b)


def get_run_rgb(run) -> RGBColor | None:
    try:
        if run.font.color.type == MSO_COLOR_TYPE.RGB:
            return run.font.color.rgb
    except AttributeError:
        pass
    return None


def style_run(run):
    size = int(run.font.size) if run.font.size else 0
    rgb = get_run_rgb(run)
    text = run.text.strip()

    run.font.name = FONT

    # Заголовки слайдов (крупный красный или тёмный в исходнике)
    if size in TITLE_SIZES or (size >= 260000 and run.font.bold):
        run.font.bold = True
        run.font.color.rgb = RED
        return

    # Номера блоков / пунктов
    if size in NUMBER_SIZES or (size == HEADER_SIZE and text.isdigit()):
        run.font.bold = True
        run.font.color.rgb = RED
        return

    # Подзаголовки карточек и секций
    if size == HEADER_SIZE or rgb_eq(rgb, OLD_HEAD) if rgb else False:
        run.font.bold = True
        run.font.color.rgb = RED
        return

    # Основной текст
    if size <= 127000 or rgb_eq(rgb, OLD_BODY) if rgb else True:
        run.font.color.rgb = GRAY
        return

    # Прочее
    if rgb_eq(rgb, OLD_BLACK) if rgb else False:
        run.font.color.rgb = DARK
        return

    run.font.color.rgb = GRAY


def style_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                style_run(run)


def get_shape_fill_rgb(shape) -> RGBColor | None:
    try:
        fill = shape.fill
        # MSO_FILL.SOLID = 1
        if fill.type != 1:
            return None
        fc = fill.fore_color
        if fc.type == MSO_COLOR_TYPE.RGB:
            return fc.rgb
    except (AttributeError, TypeError):
        pass
    return None


def style_shape_fill(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return
    rgb = get_shape_fill_rgb(shape)
    if rgb is None:
        return
    if rgb_eq(rgb, OLD_CARD_BLUE):
        shape.fill.solid()
        shape.fill.fore_color.rgb = CREAM
    elif rgb_eq(rgb, OLD_CARD_CREAM):
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG


def apply_red_style(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            style_shape_fill(shape)
            if shape.has_text_frame:
                style_text_frame(shape.text_frame)


def build():
    if not SOURCE.exists():
        raise FileNotFoundError(f"Нет исходника: {SOURCE}")
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(str(OUTPUT))
    apply_red_style(prs)
    prs.save(OUTPUT)

    # Проверка: картинки на месте
    src_pics = sum(
        1 for s in Presentation(str(SOURCE)).slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    out_pics = sum(
        1 for s in Presentation(str(OUTPUT)).slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    size_mb = OUTPUT.stat().st_size / (1024 * 1024)
    print(f"Created: {OUTPUT} ({size_mb:.1f} MB, {len(prs.slides)} slides, pictures {out_pics}/{src_pics})")
    if out_pics != src_pics:
        raise RuntimeError(f"Потеряны картинки: было {src_pics}, стало {out_pics}")


if __name__ == "__main__":
    build()
