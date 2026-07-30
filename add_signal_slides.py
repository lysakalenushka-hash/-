#!/usr/bin/env python3
"""Add sign-signaling illustration slides to PS presentation."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt, Inches

SRC = "/workspace/Презентация_ПС_обновленная.pptx"
OUT = "/workspace/Презентация_ПС_обновленная.pptx"
IMG_DIR = "/workspace/assets/ps_signals_hires"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
BG = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF3, 0xF4, 0xF6)

LABELS = [
    "01. Готовность подавать команду",
    "02. Остановка",
    "03. Замедление",
    "04. Подъём",
    "05. Опускание",
    "06. Указание направления",
    "07. Поднять колено (стрелу)",
    "08. Опустить колено (стрелу)",
    "09. Выдвинуть стрелу",
    "10. Втянуть стрелу",
]


def set_text(shape, text, size=11, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return dest


def insert_slide_copy(prs, template_index, insert_after_index):
    new_slide = duplicate_slide(prs, template_index)
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(insert_after_index + 1, new_id)
    return prs.slides[insert_after_index + 1]


def clear_content_shapes(slide, keep_bg=True):
    """Remove content shapes except background rect and footer connectors/page."""
    keep = []
    for sh in list(slide.shapes):
        # keep background
        if sh.name == "Rectangle 1":
            keep.append(sh)
            continue
        # remove everything else; rebuild
        sh.element.getparent().remove(sh.element)


def add_textbox(slide, left, top, width, height, text, size=11, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, size=size, bold=bold, color=color)
    return box


def build_signal_slide(prs, template_idx, insert_after, title, subtitle, items):
    """items: list of (img_path, label)"""
    slide = insert_slide_copy(prs, template_idx, insert_after)
    clear_content_shapes(slide)

    # header
    add_textbox(slide, Emu(548640), Emu(731520), Emu(9000000), Emu(228600),
                "Система сигнализации при выполнении работ", 11, True, RED)
    add_textbox(slide, Emu(548640), Emu(1005840), Emu(11094415), Emu(600000),
                title, 28, True, DARK)
    add_textbox(slide, Emu(548640), Emu(1650000), Emu(11094415), Emu(350000),
                subtitle, 11, False, GRAY)

    # grid 5 columns x 1 row for 5 items
    left0 = 548640
    top0 = 2100000
    gap_x = 160000
    card_w = 2080000
    card_h = 3900000

    for i, (img, label) in enumerate(items):
        left = left0 + i * (card_w + gap_x)
        # card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top0), Emu(card_w), Emu(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.fill.background()
        try:
            card.adjustments[0] = 0.08
        except Exception:
            pass

        # image — keep aspect ratio
        pic_left = left + 180000
        pic_top = top0 + 200000
        pic_w = card_w - 360000
        pic = slide.shapes.add_picture(img, Emu(pic_left), Emu(pic_top), width=Emu(pic_w))
        # center vertically in image area if short
        max_img_h = 2800000
        if pic.height > max_img_h:
            ratio = max_img_h / pic.height
            pic.height = max_img_h
            pic.width = int(pic.width * ratio)
            pic.left = left + (card_w - pic.width) // 2

        # label under image
        add_textbox(
            slide,
            Emu(left + 80000),
            Emu(top0 + 3100000),
            Emu(card_w - 160000),
            Emu(650000),
            label,
            10,
            True,
            DARK,
        )

    # footer line + page placeholder
    line = slide.shapes.add_connector(
        1, Emu(548640), Emu(6355080), Emu(11643055), Emu(6355080)
    )
    line.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    add_textbox(
        slide, Emu(10271455), Emu(6492240), Emu(1188720), Emu(228600), "00 / 00", 9, True, RED
    )
    return slide


def update_page_numbers(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def main():
    prs = Presentation(SRC)
    # Insert after signaling theory slide (index 6)
    items1 = [(f"{IMG_DIR}/signal_{i:02d}.png", LABELS[i-1]) for i in range(1, 6)]
    items2 = [(f"{IMG_DIR}/signal_{i:02d}.png", LABELS[i-1]) for i in range(6, 11)]

    build_signal_slide(
        prs, 6, 6,
        "Знаковая сигнализация (1/2)",
        "Рисунки 1–5. Сигналы, применяемые при работе подъемника (вышки) и крана",
        items1,
    )
    build_signal_slide(
        prs, 6, 7,
        "Знаковая сигнализация (2/2)",
        "Рисунки 6–10. Сигналы управления стрелой и указания направления",
        items2,
    )

    update_page_numbers(prs)
    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
