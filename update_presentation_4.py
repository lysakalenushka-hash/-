#!/usr/bin/env python3
from copy import deepcopy
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SRC = '/home/ubuntu/.cursor/projects/workspace/uploads/____________38a8.pptx'
OUT = '/workspace/Презентация_4_обновленная.pptx'
RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def set_text(shape, text, size, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = 'Inter'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    tf.word_wrap = True


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def find_prefix(slide, prefix):
    for sh in text_shapes(slide):
        if sh.text_frame.text.startswith(prefix):
            return sh
    return None


def find_exact(slide, txt):
    for sh in text_shapes(slide):
        if sh.text_frame.text.strip() == txt:
            return sh
    return None


def duplicate_slide(prs, index):
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in src.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), 'p:extLst')
    return dest


def insert_slide_copy(prs, template_index, insert_after):
    new = duplicate_slide(prs, template_index)
    ids = prs.slides._sldIdLst
    new_id = ids[-1]
    ids.remove(new_id)
    ids.insert(insert_after + 1, new_id)
    return prs.slides[insert_after + 1]


def fill_list(slide, section, title, subtitle, heading, items, note):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], subtitle, 11, False, GRAY)
    set_text(shapes[3], heading, 11, True, RED)
    nums = [s for s in shapes if s.text_frame.text.strip() in {'01','02','03','04','05','06'}]
    nums.sort(key=lambda s: s.top)
    bodies = []
    for sh in shapes:
        t = sh.text_frame.text.strip()
        if t.startswith(('а)', 'б)', 'в)', 'г)', 'д)', 'е)', 'Ж', 'Назначить', 'Условия', 'При', 'Световые', 'Обслуживание', 'ПС должна')) or t.startswith('… ещё'):
            bodies.append(sh)
    bodies.sort(key=lambda s: s.top)
    for i, item in enumerate(items):
        if i < len(nums):
            set_text(nums[i], f'{i+1:02d}', 14, True, RED)
        if i < len(bodies):
            set_text(bodies[i], item, 14, False, DARK)
    # fallback: after numbered boxes, body boxes are alternating after nums in vertical order
    if len(bodies) < len(items):
        extra_bodies = [s for s in shapes if s not in nums and s.text_frame.text not in {section,title,subtitle,heading}]
        extra_bodies.sort(key=lambda s: s.top)
        used=0
        for sh in extra_bodies:
            if sh.text_frame.text.strip().startswith(('Эксплуатация','Выявленные','Соблюдение')): continue
            if used >= len(items): break
            set_text(sh, items[used], 14, False, DARK)
            used += 1
    for sh in shapes:
        if sh.text_frame.text.strip().startswith(('Эксплуатация','Выявленные','Соблюдение')):
            set_text(sh, note, 14, True, AMBER)
            break


def fill_simple(slide, section, title, subtitle, body, footer):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], subtitle, 11, False, GRAY)
    set_text(shapes[3], body, 14, False, DARK)
    set_text(shapes[4], footer, 14, False, AMBER if 'запрещ' in footer.lower() or 'останов' in footer.lower() else GRAY)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            txt = sh.text_frame.text.strip()
            if ' / ' in txt:
                a,b = txt.split(' / ',1)
                if a.strip().isdigit() and b.strip().isdigit():
                    set_text(sh, f'{i:02d} / {total:02d}', 9, True, RED)

prs = Presentation(SRC)

# tweak existing slides
s3 = prs.slides[2]
h = find_exact(s3, '… ещё 2')
if h: set_text(h, '05  д) Не эксплуатировать ПС на неработоспособных рельсовых путях.  06  е) Не эксплуатировать ПС с нарушениями требований по установке.', 9, False, GRAY)

s6 = prs.slides[5]
shapes = sorted(text_shapes(s6), key=lambda s: (s.top, s.left))
set_text(shapes[1], 'Ключевые запреты и обязанности эксплуатирующей организации', 30, True, DARK)
set_text(shapes[2], 'Что должен обеспечить владелец ОПО с ПС', 11, False, GRAY)
set_text(shapes[3], '1. ОБЯЗАТЕЛЬНО назначить ответственных ИТР за производственный контроль, содержание ПС и безопасное производство работ.\n2. ОБЯЗАТЕЛЬНО организовать порядок допуска к самостоятельной работе, проверку знаний и выдачу удостоверений.\n3. ОБЯЗАТЕЛЬНО обеспечить ППР, ТК, схемы строповки, журналы, инструкции и контроль технологического процесса.\n4. ОБЯЗАТЕЛЬНО ограждать опасные зоны и исключать нахождение людей под грузом.\n5. ОБЯЗАТЕЛЬНО прекращать эксплуатацию ПС при угрозе аварии и при несоответствии технологическому процессу.\n6. ЗАПРЕЩАЕТСЯ допускать в работу ПС при неисправностях, отсутствии документов, ППР/ТК или неаттестованном персонале.', 14, False, DARK)
set_text(shapes[4], 'Соблюдение обязанностей эксплуатирующей организации — базовое условие безопасной эксплуатации ПС на ОПО.', 14, False, AMBER)

insert_pos = 5
new_slides = [
('list', 2, 'Требования к эксплуатации ПС', 'Ответственные лица и допуск персонала', 'Кого обязана назначить эксплуатирующая организация', 'Обязательные организационные меры', [
'Назначить внутренним распорядительным актом ответственных ИТР: за производственный контроль, за содержание ПС в работоспособном состоянии, за безопасное производство работ.',
'Установить порядок допуска к самостоятельной работе на ПС и контролировать его соблюдение.',
'Обеспечить выдачу удостоверений на право самостоятельной работы и периодическую проверку знаний персонала.',
'Назначить машинистов, крановщиков, стропальщиков, электромонтеров, слесарей, рабочих люльки и наладчиков.'
], 'Если хотя бы один обязательный ответственный не назначен — допуск ПС к работе запрещается.'),
('simple', 1, 'Организация безопасной эксплуатации ПС в составе ОПО', 'Производственный контроль и обслуживание', 'Что обязана организовать эксплуатирующая организация', 'Эксплуатирующая организация обязана обеспечить содержание ПС в работоспособном состоянии и безопасные условия их работы путем организации надлежащего надзора, обслуживания, технического освидетельствования и ремонта. Для этого устанавливаются порядок периодических осмотров, технического обслуживания и ремонтов, ведутся журналы, графики ППР, программы, схемы строповки и складирования, а персонал обеспечивается производственными инструкциями под подпись.', 'Если ПС невозможно привести в соответствие требованиям промышленной безопасности технологического процесса, его эксплуатация должна быть остановлена.'),
('list', 2, 'Организация безопасного производства работ', 'Что обязательно включают ППР и ТК', 'Раздел по безопасному производству работ с применением ПС', 'Обязательное содержание', [
'Условия совместной работы двух и более ПС, совместного подъема груза и применения координатной защиты.',
'Условия перемещения ПС с грузом, подачи грузов в проемы перекрытий и установки ПС над подземными коммуникациями.',
'Выписку из паспорта ПС о предельной силе ветра, требования к эксплуатации тары и порядок работы кранов с грейфером или магнитом.',
'Мероприятия при движении транспорта и пешеходов вблизи опасной зоны, а также требования пунктов 98–134 ФНП.'
], 'ППР и ТК должны быть выданы на участок до начала работ, а работники — ознакомлены с ними под подпись.'),
('list', 2, 'Организация безопасного производства работ', 'Совместная работа нескольких ПС и подача грузов в проемы', 'Ключевые безопасные расстояния и условия', 'Критически важные требования', [
'При совместной работе нескольких ПС расстояние по горизонтали между ними, стрелами и грузами должно быть не менее 5 м.',
'Для башенных кранов при пересечении зон обслуживания стрелы и противовесные консоли должны быть на разных уровнях; разность уровней — не менее 1 м.',
'При подаче грузов в проемы перекрытий расстояние между краем проема и грузом (или крюковой обоймой) должно быть не менее 0,5 м.',
'Если стропальщик вне видимости крановщика, требуется двусторонняя радио- или телефонная связь либо назначенные сигнальщики.'
], 'В местах подачи грузов в проемы обязательны световая сигнализация и знаки, запрещающие нахождение людей под грузом.'),
('list', 2, 'Оценка соответствия и проверка устройств безопасности', 'Что проверяют при техническом диагностировании ПС', 'Проверки работоспособности ограничителей, указателей и защит', 'Основные проверяемые устройства', [
'Световые и звуковые указатели, ограничитель грузоподъемности, ограничитель грузового момента, ограничители верхнего и нижнего положения.',
'Защиту от опасного приближения к ЛЭП, координатную защиту, блокировки, средства автоматической остановки и регистратор параметров.',
'Для подъемников — ловители, аварийные остановы, ограничители скорости, устройство аварийного опускания люльки, кнопки «Стоп».',
'Противоугонные захваты, тупиковые упоры, указатель угла наклона, сигнализатор предельной скорости ветра.'
], 'После положительных результатов проверки устройств безопасности выполняют статические и динамические испытания.'),
('list', 2, 'Запрет эксплуатации ПС', 'Нарушения, при которых ПС не допускается в работу', 'Основания для немедленного запрета эксплуатации', 'ПС должна быть остановлена, если', [
'Обслуживание ведется неаттестованным персоналом или не назначены обязательные ответственные ИТР.',
'Не проведены техническое освидетельствование, экспертиза промышленной безопасности либо не исполнены предписания надзора.',
'Выявлены неисправности: трещины и деформации металлоконструкций, неработоспособность заземления, ограничителей, тормозов, рельсового пути, тупиковых упоров.',
'Отсутствуют исправные грузозахватные приспособления, паспорт ПС, руководство по эксплуатации, ППР, ТК или наряды-допуски.'
], 'Продолжение работ при выявленных нарушениях ФНП, ППР, ТК и инструкций запрещено.'),
]

for entry in new_slides:
    if entry[0] == 'list':
        kind, template, section, title, subtitle, heading_or_body, items_or_footer, note = entry
    else:
        kind, template, section, title, subtitle, heading_or_body, items_or_footer = entry
        note = None
    slide = insert_slide_copy(prs, template, insert_pos)
    insert_pos += 1
    if kind == 'list':
        fill_list(slide, section, title, subtitle, heading_or_body, items_or_footer, note)
    else:
        fill_simple(slide, section, title, subtitle, heading_or_body, items_or_footer)

# update wear slides by title matching
for slide in prs.slides:
    if find_prefix(slide, 'Критерии выбраковки барабанов'):
        set_text(find_prefix(slide, 'Критерии выбраковки барабанов'), 'Критерии выбраковки ходовых колес и блоков', 30, True, DARK)
        if find_exact(slide, '0,5%'): set_text(find_exact(slide, '0,5%'), '0,5%', 32, True, RED)
        body = None
        for sh in text_shapes(slide):
            if 'Разность диаметров колёс' in sh.text or 'Бараба' in sh.text:
                body = sh; break
        if body: set_text(body, 'Разность диаметров колес, связанных кинематически, более 0,5 % — критерий выбраковки. Для блоков — износ ручья более 40 % от первоначального радиуса.', 14, False, DARK)
        src = None
        for sh in text_shapes(slide):
            if sh.text.startswith('Источник:'):
                src = sh; break
        if src: set_text(src, 'Источник: таблица браковки элементов ПС из 4.docx (ходовые колеса, блоки).', 11, False, GRAY)
    if find_prefix(slide, 'Критерий выбраковки крюков') and find_prefix(slide, 'Износ барабана'):
        set_text(find_prefix(slide, 'Критерий выбраковки крюков'), 'Критерий выбраковки барабанов', 30, True, DARK)
        b = find_prefix(slide, 'Дополнительный критерий')
        if b: set_text(b, 'Трещины любых размеров — безусловное основание для выбраковки барабана.', 14, False, DARK)
    if find_exact(slide, '… ещё 1'):
        set_text(find_exact(slide, '… ещё 1'), 'Ходовые колеса и блоки: трещины любых размеров, износ реборды колеса > 50 %, износ ручья блока > 40 % — основания для выбраковки.', 9, False, GRAY)

# final summaries: last two slides
slideA = prs.slides[-2]
shapes = sorted(text_shapes(slideA), key=lambda s: (s.top, s.left))
items = [
    'Эксплуатирующая организация обязана назначить ответственных ИТР, организовать производственный контроль и порядок допуска персонала к самостоятельной работе.',
    'ППР и ТК должны быть выданы на участок до начала работ; при работе нескольких ПС и при подаче грузов в проемы применяются специальные меры безопасности.',
    'При диагностировании обязательно проверяются ограничители, указатели, регистраторы, блокировки, защита от приближения к ЛЭП и иные устройства безопасности.',
    'ПС не допускается в работу при неисправностях, отсутствии документов, ППР/ТК, неаттестованном персонале и неисполненных предписаниях надзора.'
]
nums = [s for s in shapes if s.text_frame.text.strip() in {'01','02','03','04'}]
bodies = [s for s in shapes if s.text_frame.text.strip().startswith(('Периодичность','Браковка','Выявленные','Эксплуатирующая','ППР','При диагностировании','ПС не'))]
# safer explicit by positions after numbers
ordered_bodies=[]
for n in nums:
    below = [s for s in shapes if s.top == n.top and s.left > n.left]
    if below: ordered_bodies.append(sorted(below,key=lambda s:s.left)[0])
for sh,item in zip(ordered_bodies, items): set_text(sh,item,14,False,DARK)
for sh in shapes:
    if sh.text_frame.text.strip().startswith(('Выявленные','Нарушения')):
        set_text(sh,'Нарушения требований промышленной безопасности — основание для немедленного останова ПС.',14,True,AMBER)
        break

slideB = prs.slides[-1]
shapes = sorted(text_shapes(slideB), key=lambda s: (s.top, s.left))
items2 = [
    'Выбраковка элементов ПС выполняется по фактическим дефектам и установленным порогам износа: колеса, блоки, барабаны, крюки, тормозные шкивы и накладки.',
    'Критические пороги: износ зева крюка > 10 %, разность диаметров колес > 0,5 %, износ обода тормозного шкива > 25 %, тормозной накладки > 50 %.',
    'Результаты технического обслуживания, освидетельствования и проверки устройств безопасности фиксируются в паспорте ПС, журнале осмотров и актах.',
    'Если ПС невозможно привести в соответствие требованиям безопасной эксплуатации и технологического процесса, эксплуатация должна быть остановлена.'
]
nums = [s for s in shapes if s.text_frame.text.strip() in {'01','02','03','04'}]
ordered_bodies=[]
for n in nums:
    below = [s for s in shapes if s.top == n.top and s.left > n.left]
    if below: ordered_bodies.append(sorted(below,key=lambda s:s.left)[0])
for sh,item in zip(ordered_bodies, items2): set_text(sh,item,14,False,DARK)

update_pages(prs)
prs.save(OUT)
print('saved', OUT, 'slides', len(prs.slides))
