#!/usr/bin/env python3
"""
Тема 2 — 5 презентаций.

КЛЮЧЕВОЕ:
- Слайд НЕ является изображением.
- Весь текст — редактируемые текстовые блоки.
- Иллюстрации из исходника — отдельные объекты Picture (не фон).
- Текст исходных слайдов сохранён дословно.
- Новые слайды (пробелы) — Open Sans 46/26/32, RGB 221/225.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

OUT = Path("/workspace/tema2_bleeding")
ASSETS = OUT / "assets"
OUT.mkdir(parents=True, exist_ok=True)

SLIDE_W = Emu(24384000)
SLIDE_H = Emu(13716000)
M = Emu(700000)

# Общая палитра (сочетается с исходником + новые требования)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF7, 0xF7)
BG_BAR = RGBColor(221, 221, 221)
LINE = RGBColor(225, 225, 225)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED = RGBColor(0xCC, 0x00, 0x00)
ACCENT_BLUE = RGBColor(0x00, 0x55, 0xAA)
TABLE_HDR = RGBColor(0x44, 0x44, 0x44)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
CREAM = RGBColor(0xF5, 0xF0, 0xE1)  # как «облачко» определения в исходнике
FONT = "Open Sans"
SUB = "Оказание первой помощи при наружных кровотечениях и травмах"
COURSE = "Тема 2 · Оказание первой помощи при наружных кровотечениях"


def asset(name: str) -> Path:
    return ASSETS / name


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", FONT)


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def tbox(slide, l, t, w, h, text, *, size=32, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, *, size=28, alert=None, marker="•"):
    alert = alert or set()
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        font(r, size, bold=(i in alert), color=ACCENT_RED if i in alert else TEXT)
    return box


def pic_fit(slide, name, l, t, max_w, max_h):
    from PIL import Image as PILImage
    path = asset(name)
    if not path.exists():
        # placeholder
        rect(slide, l, t, max_w, max_h, WHITE, line=LINE)
        tbox(slide, l + Emu(100000), t + max_h // 2 - Emu(400000), max_w - Emu(200000),
             Emu(800000), f"[МЕСТО ДЛЯ ИЗОБРАЖЕНИЯ: {name}]",
             size=16, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def bg(slide):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)


def src_header(slide, title: str):
    """Заголовок в стиле исходника: крупный CAPS, линия, подзаголовок."""
    bg(slide)
    tbox(slide, M, Emu(280000), Emu(22000000), Emu(850000),
         title, size=30, bold=True, color=TEXT)
    rect(slide, M, Emu(1150000), Emu(16000000), Emu(35000), RGBColor(0x55, 0x55, 0x55))
    tbox(slide, M, Emu(1250000), Emu(22000000), Emu(450000),
         SUB, size=14, bold=False, color=MUTED)


def new_header(slide, title: str):
    """Заголовок новых слайдов — 46 pt."""
    bg(slide)
    tbox(slide, M, Emu(300000), Emu(23000000), Emu(1100000),
         title, size=46, bold=True, color=TEXT)
    rect(slide, M, Emu(1500000), Emu(23000000), Emu(20000), LINE)


# ───────── служебные новые слайды ─────────

def slide_title_src(prs):
    """Титул в стиле исходника — фигуры + редактируемый текст."""
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF2, 0xF2, 0xF2))
    rect(slide, Emu(17500000), 0, Emu(6884000), SLIDE_H, BG_BAR)
    rect(slide, 0, Emu(4500000), Emu(4500000), Emu(2800000), BG_BAR)
    tbox(slide, M, Emu(5000000), Emu(16500000), Emu(1600000),
         "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ НАРУЖНЫХ\nКРОВОТЕЧЕНИЯХ И ТРАВМАХ",
         size=34, bold=True, color=TEXT)
    rect(slide, M, Emu(6800000), Emu(14000000), Emu(35000), TEXT)
    return slide


def slide_thanks_src(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF2, 0xF2, 0xF2))
    rect(slide, Emu(17500000), 0, Emu(6884000), SLIDE_H, BG_BAR)
    tbox(slide, M, Emu(5500000), Emu(16000000), Emu(1200000),
         "БЛАГОДАРИМ ЗА ВНИМАНИЕ", size=40, bold=True, color=TEXT)
    rect(slide, M, Emu(6900000), Emu(10000000), Emu(35000), TEXT)
    return slide


def slide_toc(prs, items):
    slide = blank(prs)
    new_header(slide, "Содержание")
    y = Emu(1900000)
    for i, item in enumerate(items[:5], 1):
        oval(slide, M, y, Emu(800000), Emu(800000), ACCENT_BLUE)
        tbox(slide, M, y, Emu(800000), Emu(800000), str(i), size=26, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, M + Emu(1100000), y, Emu(20500000), Emu(800000), WHITE, line=LINE)
        tbox(slide, M + Emu(1400000), y, Emu(19500000), Emu(800000),
             item, size=28, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(1400000)
    return slide


def slide_new_list(prs, title, items, *, highlight=None, alert=None):
    slide = blank(prs)
    new_header(slide, title)
    h = Emu(8000000) if not highlight else Emu(6500000)
    bullets(slide, M, Emu(1900000), Emu(23000000), h, items[:6], size=32, alert=alert)
    if highlight:
        round_rect(slide, M, Emu(11000000), Emu(23000000), Emu(1600000), BG)
        rect(slide, M, Emu(11000000), Emu(120000), Emu(1600000), ACCENT_RED)
        tbox(slide, M + Emu(400000), Emu(11200000), Emu(22000000), Emu(1200000),
             highlight, size=26, bold=True, color=ACCENT_RED, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_new_scheme(prs, title, steps, *, key_idx=None):
    slide = blank(prs)
    new_header(slide, title)
    y = Emu(1900000)
    bh = Emu(1250000)
    for i, step in enumerate(steps):
        accent = key_idx is not None and i == key_idx
        fill = ACCENT_BLUE if accent else BG
        round_rect(slide, M, y, Emu(23000000), bh, fill, line=None if accent else LINE)
        oval(slide, M + Emu(300000), y + Emu(200000), Emu(800000), Emu(800000),
             WHITE if accent else ACCENT_BLUE)
        tbox(slide, M + Emu(300000), y + Emu(200000), Emu(800000), Emu(800000),
             str(i + 1), size=26, bold=True,
             color=ACCENT_BLUE if accent else WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tbox(slide, M + Emu(1400000), y + Emu(200000), Emu(20500000), Emu(800000),
             step, size=26, color=WHITE if accent else TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += bh + Emu(180000)
    return slide


def slide_new_table(prs, title, headers, rows, *, fs=22):
    slide = blank(prs)
    new_header(slide, title)
    n_rows, n_cols = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(
        n_rows, n_cols, M, Emu(1900000), Emu(23000000),
        min(Emu(8000000), Emu(1200000) * n_rows)
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                font(r, 24, True, WHITE)
    for ri, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
            cell.text = val
            alert = any(k in val for k in ("60", "30", "Алый", "Пульсир", "Критическ", "НЕ "))
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    font(r, fs, alert, ACCENT_RED if alert else TEXT)
    return slide


def slide_new_summary(prs, points, callout=None):
    slide = blank(prs)
    new_header(slide, "Главное запомнить")
    y = Emu(1900000)
    for pt in points[:4]:
        rect(slide, M, y, Emu(23000000), Emu(1600000), WHITE, line=LINE)
        rect(slide, M, y, Emu(120000), Emu(1600000), ACCENT_BLUE)
        tbox(slide, M + Emu(500000), y + Emu(350000), Emu(22000000), Emu(900000),
             pt, size=28, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Emu(1850000)
    if callout:
        tbox(slide, M, Emu(11500000), Emu(23000000), Emu(1000000),
             callout, size=26, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    return slide


def slide_new_compare(prs, title, left_t, left_items, right_t, right_items):
    slide = blank(prs)
    new_header(slide, title)
    rect(slide, M, Emu(1900000), Emu(11000000), Emu(9500000), WHITE, line=LINE)
    rect(slide, M + Emu(11800000), Emu(1900000), Emu(11000000), Emu(9500000), WHITE, line=LINE)
    rect(slide, M, Emu(1900000), Emu(11000000), Emu(1000000), ACCENT_BLUE)
    rect(slide, M + Emu(11800000), Emu(1900000), Emu(11000000), Emu(1000000), BG_BAR)
    tbox(slide, M + Emu(200000), Emu(2050000), Emu(10600000), Emu(700000),
         left_t, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tbox(slide, M + Emu(12000000), Emu(2050000), Emu(10600000), Emu(700000),
         right_t, size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    bullets(slide, M + Emu(400000), Emu(3200000), Emu(10000000), Emu(7800000), left_items, size=26)
    bullets(slide, M + Emu(12200000), Emu(3200000), Emu(10000000), Emu(7800000), right_items, size=26)
    return slide


# ───────── исходные слайды как РЕДАКТИРУЕМЫЙ контент + картинки ─────────

def src_bleed_definition(prs):
    """Стр.2 — как в исходнике: слева персонаж, облачко с определением, справа признаки."""
    slide = blank(prs)
    src_header(slide, "ПОНЯТИЕ «КРОВОТЕЧЕНИЕ»")
    # слева — иллюстрация (отдельный объект)
    pic_fit(slide, "def_man.png", M, Emu(2000000), Emu(7000000), Emu(10500000))
    # облачко с определением (редактируемый текст)
    round_rect(slide, M + Emu(6200000), Emu(2100000), Emu(7500000), Emu(4200000), CREAM)
    tbox(slide, M + Emu(6600000), Emu(2400000), Emu(6700000), Emu(3600000),
         "Под кровотечением понимают ситуацию, когда кровь по разным причинам "
         "покидает сосудистое русло, что приводит к острой кровопотере.",
         size=22, color=TEXT)
    # справа — признаки (редактируемый список)
    tbox(slide, M + Emu(14000000), Emu(2000000), Emu(9000000), Emu(700000),
         "Основные признаки острой кровопотери:", size=22, bold=True, color=TEXT)
    bullets(slide, M + Emu(14000000), Emu(2800000), Emu(9000000), Emu(9000000), [
        "резкая общая слабость;",
        "чувство жажды;",
        "головокружение;",
        "мелькание «мушек» перед глазами;",
        "обморок, чаще при попытке встать;",
        "бледная, влажная и холодная кожа;",
        "учащённое сердцебиение;",
        "частое дыхание.",
    ], size=20, marker="☐")
    return slide


def src_bleed_types(prs):
    """Стр.3 — три колонки: подпись (текст) + иллюстрация (объект)."""
    slide = blank(prs)
    src_header(slide, "ПРИЗНАКИ РАЗЛИЧНЫХ ВИДОВ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    cols = [
        ("bleed_arterial.png", "Артериальное кровотечение"),
        ("bleed_venous.png", "Венозное кровотечение"),
        ("bleed_capillary.png", "Капиллярное кровотечение"),
    ]
    x = M
    cw = Emu(7200000)
    for img, title in cols:
        tbox(slide, x, Emu(1900000), cw, Emu(800000), title, size=24, bold=True,
             color=TEXT, align=PP_ALIGN.CENTER)
        pic_fit(slide, img, x + Emu(200000), Emu(2900000), cw - Emu(400000), Emu(9000000))
        x += cw + Emu(300000)
    return slide


def src_overview(prs):
    """Стр.4 — как в исходнике: слева картинка, справа текст."""
    slide = blank(prs)
    src_header(slide, "ОБЗОРНЫЙ ОСМОТР ПОСТРАДАВШЕГО")
    # слева — иллюстрация (отдельный объект)
    pic_fit(slide, "overview_exam.png", M, Emu(2000000), Emu(12000000), Emu(10500000))
    # справа — редактируемый текст исходника
    tbox(slide, M + Emu(12800000), Emu(4500000), Emu(10000000), Emu(4500000),
         "Обзорный осмотр производится очень быстро, "
         "в течение 1-2 секунд, с головы до ног",
         size=28, bold=False, color=TEXT)
    return slide


def src_pressure_bandage(prs):
    """Стр.5 — две колонки: подпись + картинка."""
    slide = blank(prs)
    src_header(slide, "СПОСОБЫ ВРЕМЕННОЙ ОСТАНОВКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    for i, (img, title) in enumerate([
        ("direct_pressure.png", "Прямое давление на рану"),
        ("pressure_bandage.png", "Наложение давящей повязки"),
    ]):
        x = M + i * Emu(11500000)
        tbox(slide, x, Emu(1900000), Emu(11000000), Emu(700000), title, size=26, bold=True,
             color=TEXT, align=PP_ALIGN.CENTER)
        pic_fit(slide, img, x + Emu(300000), Emu(2800000), Emu(10400000), Emu(9500000))
    return slide


def src_artery_grid(prs, title_note: str, cells: list[tuple[str, str]]):
    """Пальцевое прижатие: сверху картинка, снизу текст (как в исходнике)."""
    slide = blank(prs)
    src_header(slide, "СПОСОБЫ ВРЕМЕННОЙ ОСТАНОВКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    tbox(slide, M, Emu(1700000), Emu(23000000), Emu(450000),
         "Пальцевое прижатие артерии", size=22, bold=True, color=MUTED)
    n = len(cells)
    gap = Emu(250000)
    cw = (Emu(23000000) - gap * (n - 1)) // n
    for i, (img, text) in enumerate(cells):
        x = M + i * (cw + gap)
        pic_fit(slide, img, x + Emu(50000), Emu(2300000), cw - Emu(100000), Emu(5200000))
        tbox(slide, x + Emu(100000), Emu(7800000), cw - Emu(200000), Emu(4500000),
             text, size=15, color=TEXT)
    return slide


def src_flexion(prs):
    slide = blank(prs)
    src_header(slide, "СПОСОБЫ ВРЕМЕННОЙ ОСТАНОВКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    tbox(slide, M, Emu(1750000), Emu(23000000), Emu(500000),
         "Максимальное сгибание конечности", size=24, bold=True, color=ACCENT_BLUE)
    cells = [
        ("flexion_arm.png",
         "Для остановки кровотечения из предплечья в локтевой сгиб вкладывают валик, "
         "конечность максимально сгибают в локтевом суставе и предплечье фиксируют к плечу "
         "в таком положении, например, ремнем"),
        ("flexion_leg.png",
         "При повреждении сосудов стопы, голени и подколенной ямки в последнюю вкладывают "
         "несколько бинтов или валик из ткани, после чего конечность сгибают в коленном "
         "суставе и фиксируют в этом положении бинтом"),
        ("flexion_thigh.png",
         "Для остановки кровотечения при травме бедра сверток из ткани или несколько бинтов "
         "вкладывают в область паховой складки, притягивают колено к груди и фиксируют "
         "руками или бинтом"),
    ]
    gap = Emu(250000)
    cw = (Emu(23000000) - 2 * gap) // 3
    for i, (img, text) in enumerate(cells):
        x = M + i * (cw + gap)
        rect(slide, x, Emu(2400000), cw, Emu(9800000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(2550000), cw - Emu(200000), Emu(4800000))
        tbox(slide, x + Emu(150000), Emu(7600000), cw - Emu(300000), Emu(4300000),
             text, size=15, color=TEXT)
    return slide


def src_tourniquet(prs):
    slide = blank(prs)
    src_header(slide, "СПОСОБЫ ВРЕМЕННОЙ ОСТАНОВКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    tbox(slide, M, Emu(1750000), Emu(23000000), Emu(500000),
         "Наложение кровоостанавливающего жгута", size=24, bold=True, color=ACCENT_BLUE)
    # важное ограничение — редактируемый текст, акцент
    tbox(slide, M, Emu(2350000), Emu(23000000), Emu(700000),
         "t max ⩽ 60 мин. в теплое время года · 30 мин. в холодное",
         size=26, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    for i, img in enumerate(["tourniquet_1.png", "tourniquet_2.png", "tourniquet_3.png"]):
        x = M + i * Emu(7700000)
        rect(slide, x, Emu(3300000), Emu(7400000), Emu(8800000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(150000), Emu(3500000), Emu(7100000), Emu(8400000))
    return slide


def src_improvised(prs):
    """Стр.13 — слева текст, справа картинка."""
    slide = blank(prs)
    src_header(slide, "СПОСОБЫ ВРЕМЕННОЙ ОСТАНОВКИ НАРУЖНОГО КРОВОТЕЧЕНИЯ")
    tbox(slide, M, Emu(2000000), Emu(11000000), Emu(1000000),
         "В качестве импровизированного жгута можно использовать подручные средства:",
         size=24, bold=True, color=TEXT)
    bullets(slide, M, Emu(3200000), Emu(11000000), Emu(6000000), [
        "тесьму,",
        "платок,",
        "галстук",
        "и другие подобные вещи.",
    ], size=28, marker="☐")
    pic_fit(slide, "improvised_tq.png", M + Emu(12000000), Emu(2000000),
            Emu(10000000), Emu(10000000))
    return slide


def src_nosebleed(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ НОСОВОМ КРОВОТЕЧЕНИИ")
    for i, img in enumerate(["nose_1.png", "nose_2.png", "nose_3.png", "nose_4.png"]):
        x = M + i * Emu(5750000)
        tbox(slide, x, Emu(1900000), Emu(5500000), Emu(600000),
             f"Шаг {i+1}", size=24, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        rect(slide, x, Emu(2600000), Emu(5500000), Emu(9500000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(2800000), Emu(5300000), Emu(9100000))
    return slide


def src_shock(prs):
    slide = blank(prs)
    src_header(slide, "ПОНЯТИЕ О ТРАВМАТИЧЕСКОМ ШОКЕ, ПРИЧИНЫ И ПРИЗНАКИ")
    round_rect(slide, M, Emu(2000000), Emu(10000000), Emu(3500000), CREAM)
    tbox(slide, M + Emu(400000), Emu(2300000), Emu(9200000), Emu(2800000),
         "Травматический шок – состояние, причинами развития которого являются "
         "тяжелые травмы и сильные кровотечения.",
         size=24, color=TEXT)
    tbox(slide, M + Emu(11000000), Emu(2000000), Emu(11000000), Emu(700000),
         "Признаками травматического шока являются:", size=22, bold=True, color=TEXT)
    bullets(slide, M + Emu(11000000), Emu(2800000), Emu(11000000), Emu(8000000), [
        "наличие тяжелой травмы и сильного кровотечения;",
        "нарушения дыхания и кровообращения (учащенное дыхание и сердцебиение);",
        "бледная холодная влажная кожа;",
        "возбуждение, сменяющееся апатией.",
    ], size=22, marker="☑")
    return slide


def src_shock_prev(prs):
    slide = blank(prs)
    src_header(slide, "МЕРОПРИЯТИЯ, ПРЕДУПРЕЖДАЮЩИЕ РАЗВИТИЕ ТРАВМАТИЧЕСКОГО ШОКА")
    tbox(slide, M, Emu(2000000), Emu(11000000), Emu(600000),
         "К этим мероприятиям относятся:", size=24, bold=True, color=TEXT)
    bullets(slide, M, Emu(2800000), Emu(11000000), Emu(8000000), [
        "остановка кровотечения;",
        "придание пострадавшему оптимального положения тела;",
        "иммобилизация травмированных конечностей;",
        "защита от переохлаждения (укутывание подручными средствами или покрывалом "
        "спасательным изотермическим).",
    ], size=24, marker="☑")
    rect(slide, M + Emu(12000000), Emu(2000000), Emu(10000000), Emu(10000000), WHITE, line=LINE)
    pic_fit(slide, "shock_prev.png", M + Emu(12200000), Emu(2200000),
            Emu(9600000), Emu(9600000))
    return slide


def src_exam(prs, part: int):
    slide = blank(prs)
    src_header(slide, "ПОДРОБНЫЙ ОСМОТР ПОСТРАДАВШЕГО")
    if part == 1:
        cells = [
            ("exam_head.png", "Шаг 1 Осмотр головы"),
            ("exam_neck.png", "Шаг 2 Осмотр шеи"),
            ("exam_chest.png", "Шаг 3 Осмотр груди"),
        ]
    else:
        cells = [
            ("exam_abdomen.png", "Шаг 4 Осмотр живота и таза"),
            ("exam_legs.png", "Шаг 5 Осмотр ног"),
            ("exam_arms.png", "Шаг 6 Осмотр рук"),
        ]
    gap = Emu(300000)
    cw = (Emu(23000000) - 2 * gap) // 3
    for i, (img, cap) in enumerate(cells):
        x = M + i * (cw + gap)
        tbox(slide, x, Emu(1900000), cw, Emu(700000), cap, size=22, bold=True,
             color=TEXT, align=PP_ALIGN.CENTER)
        rect(slide, x, Emu(2700000), cw, Emu(9500000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(2900000), cw - Emu(200000), Emu(9100000))
    return slide


def src_head(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ ГОЛОВЫ")
    cells = [
        ("head_trauma_1.png",
         "Придание пострадавшему с травмой головы и находящемуся без сознания, "
         "устойчивого бокового положения."),
        ("head_trauma_2.png",
         "Остановка кровотечения при ранении головы путем прямого давления на рану. "
         "Нельзя применять этот способ при открытых черепно-мозговых травмах."),
        ("head_trauma_3.png",
         "Остановка кровотечения при травме головы наложением давящей повязки."),
    ]
    gap = Emu(250000)
    cw = (Emu(23000000) - 2 * gap) // 3
    for i, (img, text) in enumerate(cells):
        x = M + i * (cw + gap)
        rect(slide, x, Emu(1900000), cw, Emu(10500000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(2100000), cw - Emu(200000), Emu(5500000))
        tbox(slide, x + Emu(150000), Emu(7900000), cw - Emu(300000), Emu(4200000),
             text, size=16, color=TEXT)
    return slide


def src_eye_nose(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМАХ ГЛАЗА И НОСА")
    rect(slide, M, Emu(1900000), Emu(11000000), Emu(10500000), WHITE, line=LINE)
    tbox(slide, M + Emu(300000), Emu(2100000), Emu(10400000), Emu(600000),
         "При травме глаза:", size=24, bold=True, color=TEXT)
    tbox(slide, M + Emu(300000), Emu(2800000), Emu(10400000), Emu(1200000),
         "Наложить стерильную повязку на оба глаза.", size=22, color=TEXT)
    pic_fit(slide, "eye_injury.png", M + Emu(500000), Emu(4300000), Emu(10000000), Emu(7500000))

    rect(slide, M + Emu(11800000), Emu(1900000), Emu(11000000), Emu(10500000), WHITE, line=LINE)
    tbox(slide, M + Emu(12100000), Emu(2100000), Emu(10400000), Emu(600000),
         "При травме носа:", size=24, bold=True, color=TEXT)
    tbox(slide, M + Emu(12100000), Emu(2800000), Emu(10400000), Emu(1800000),
         "Усадить пострадавшего со слегка наклоненной вперед головой, "
         "зажать нос, приложить холод на переносицу.", size=22, color=TEXT)
    pic_fit(slide, "nose_injury.png", M + Emu(12300000), Emu(5000000),
            Emu(10000000), Emu(6800000))
    return slide


def src_steps_images(prs, title: str, cells: list[tuple[str, str]]):
    slide = blank(prs)
    src_header(slide, title)
    n = len(cells)
    gap = Emu(250000)
    cw = (Emu(23000000) - gap * (n - 1)) // n
    for i, (img, cap) in enumerate(cells):
        x = M + i * (cw + gap)
        tbox(slide, x, Emu(1850000), cw, Emu(900000), cap, size=18, bold=True,
             color=TEXT, align=PP_ALIGN.CENTER)
        rect(slide, x, Emu(2900000), cw, Emu(9400000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(3100000), cw - Emu(200000), Emu(9000000))
    return slide


def src_limbs_text(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ КОНЕЧНОСТЕЙ")
    cells = [
        ("limb_1.png", "Шаг 1 Остановить кровотечение"),
        ("limb_2.png",
         "Шаг 2 Иммобилизация с обездвиживанием двух соседних суставов с помощью шин, "
         "наложенных поверх одежды.\nВ качестве иммобилизующего средства можно использовать "
         "шины или плоские узкие предметы: палки, доски, линейки, прутья, фанеру, картон и др."),
        ("limb_3.png",
         "Шаг 3 При открытых переломах нельзя прикладывать шину к местам, "
         "где выступают наружу костные отломки."),
    ]
    gap = Emu(250000)
    cw = (Emu(23000000) - 2 * gap) // 3
    for i, (img, text) in enumerate(cells):
        x = M + i * (cw + gap)
        rect(slide, x, Emu(1900000), cw, Emu(10500000), WHITE, line=LINE)
        pic_fit(slide, img, x + Emu(100000), Emu(2100000), cw - Emu(200000), Emu(5000000))
        tbox(slide, x + Emu(150000), Emu(7400000), cw - Emu(300000), Emu(4700000),
             text, size=15, color=TEXT)
    return slide


def src_immobilization(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ КОНЕЧНОСТЕЙ")
    tbox(slide, M, Emu(2200000), Emu(11000000), Emu(5000000),
         "Иммобилизация – это создание неподвижности поврежденной части тела "
         "с помощью подручных средств, готовых транспортных шин или используя "
         "здоровые части тела пострадавшего (аутоиммобилизация).",
         size=26, color=TEXT)
    rect(slide, M + Emu(12000000), Emu(2000000), Emu(10000000), Emu(10000000), WHITE, line=LINE)
    pic_fit(slide, "immobilization.png", M + Emu(12200000), Emu(2200000),
            Emu(9600000), Emu(9600000))
    return slide


def src_spine(prs):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ ПОЗВОНОЧНИКА")
    rect(slide, M, Emu(1900000), Emu(11000000), Emu(10500000), WHITE, line=LINE)
    pic_fit(slide, "spine_1.png", M + Emu(200000), Emu(2100000), Emu(10600000), Emu(7000000))
    tbox(slide, M + Emu(300000), Emu(9400000), Emu(10400000), Emu(2500000),
         "При перемещении пострадавший должен находиться на ровной, жесткой, "
         "горизонтальной поверхности", size=20, color=TEXT)

    rect(slide, M + Emu(11800000), Emu(1900000), Emu(11000000), Emu(10500000), WHITE, line=LINE)
    pic_fit(slide, "spine_2.png", M + Emu(12000000), Emu(2100000), Emu(10600000), Emu(7000000))
    tbox(slide, M + Emu(12100000), Emu(9400000), Emu(10400000), Emu(2500000),
         "Перемещение или перекладывание пострадавшего следует осуществлять с помощью "
         "нескольких человек, особое внимание следует уделить фиксации шейного отдела "
         "позвоночника", size=20, color=TEXT)
    return slide


def src_abdomen(prs, open_: bool):
    slide = blank(prs)
    src_header(slide, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ ЖИВОТА И ТАЗА")
    if not open_:
        tbox(slide, M, Emu(1850000), Emu(23000000), Emu(500000),
             "Закрытая травма живота", size=24, bold=True, color=ACCENT_BLUE)
        steps = [
            "Шаг 1 Вызвать скорую помощь",
            "Шаг 2 Положить холод на живот",
            "Шаг 3 Придать положение на спине с валиком под полусогнутыми "
            "разведенными в стороны ногами",
        ]
        img = "abdomen_closed.png"
    else:
        tbox(slide, M, Emu(1850000), Emu(23000000), Emu(500000),
             "Открытая травма живота", size=24, bold=True, color=ACCENT_BLUE)
        steps = [
            "Шаг 1 Поверх валиков наложить повязку, не прижимая выпавшие органы",
            "Шаг 2 Положить холод на живот",
            "Шаг 3 Вызвать скорую помощь",
        ]
        img = "abdomen_open.png"
    bullets(slide, M, Emu(2500000), Emu(23000000), Emu(2500000), steps, size=24)
    rect(slide, M, Emu(5200000), Emu(23000000), Emu(7000000), WHITE, line=LINE)
    pic_fit(slide, img, M + Emu(200000), Emu(5400000), Emu(22600000), Emu(6600000))
    return slide


# ═══════════════════ СБОРКА 5 ПРЕЗЕНТАЦИЙ ═══════════════════

def build_pres1():
    prs = new_prs()
    slide_title_src(prs)
    slide_toc(prs, [
        "Определение кровотечения и признаки кровопотери",
        "Цель и порядок обзорного осмотра",
        "Обзорный и подробный осмотр: различия",
        "Когда обзорный осмотр критически важен",
    ])
    src_bleed_definition(prs)
    src_overview(prs)
    # новые
    slide_new_list(prs, "Что такое обзорный осмотр и зачем он нужен", [
        "Быстрая визуальная оценка пострадавшего с головы до ног",
        "Цель — выявить продолжающееся наружное кровотечение",
        "Выполняется сразу после оценки обстановки и безопасности",
        "При кровотечении — немедленно начать временную остановку",
        "По Приказу Минздрава № 220н — до подробного осмотра",
    ], highlight="Приоритет: сначала угрожающее жизни кровотечение")
    slide_new_compare(
        prs, "Обзорный и подробный осмотр — различия",
        "Обзорный осмотр",
        ["Цель: найти наружное кровотечение", "Время: около 1–2 секунд",
         "Объём: быстрый взгляд с головы до ног", "Когда: сразу после оценки обстановки"],
        "Подробный осмотр",
        ["Цель: выявить травмы и другие угрозы", "Время: несколько минут",
         "Объём: голова → шея → грудь → … → руки", "Когда: после остановки кровотечения"],
    )
    slide_new_list(prs, "Когда обзорный осмотр критически важен", [
        "ДТП — раны под одеждой, кровотечение из конечностей",
        "Падение с высоты — множественные повреждения",
        "Производственные травмы — риск артериального кровотечения",
        "ЧС с несколькими пострадавшими — быстрая сортировка",
    ])
    slide_new_summary(prs, [
        "Кровотечение — выход крови из сосудов с риском кровопотери",
        "Обзорный осмотр — быстрый поиск кровотечения (1–2 с)",
        "Сначала останавливаем кровь, затем подробный осмотр",
        "Особенно важен при ДТП и падении с высоты",
    ], "Время — критический фактор!")
    slide_thanks_src(prs)
    p = OUT / "Кровотечение_и_обзорный_осмотр.pptx"
    prs.save(p)
    return p


def build_pres2():
    prs = new_prs()
    slide_title_src(prs)
    slide_toc(prs, [
        "Виды наружного кровотечения",
        "Сравнительная характеристика",
        "Признаки кровопотери и объём",
        "Скрытое (внутреннее) кровотечение",
    ])
    src_bleed_definition(prs)
    src_bleed_types(prs)
    slide_new_table(prs, "Сравнительная таблица видов кровотечения",
                    ["Признак", "Артериальное", "Венозное", "Капиллярное"],
                    [
                        ["Цвет крови", "Алый, яркий", "Тёмный, вишнёвый", "Красный"],
                        ["Характер", "Пульсирующей струёй", "Равномерной струёй", "Сочится"],
                        ["Скорость", "Очень быстрая", "Умеренная / быстрая", "Медленная"],
                        ["Опасность", "Критическая — минуты", "Высокая", "Обычно низкая"],
                    ])
    slide_new_table(prs, "Оценка объёма кровопотери",
                    ["Объём", "Доля ОЦК", "Состояние"],
                    [
                        ["До ~500 мл", "≈ 10%", "Слабость, жажда"],
                        ["500–1000 мл", "≈ 10–20%", "Бледность, тахикардия"],
                        ["1000–1500 мл", "≈ 20–30%", "Обмороки, холодный пот"],
                        [">1500–2000 мл", "> 30%", "Шок, угроза жизни"],
                    ])
    slide_new_list(prs, "Признаки скрытого (внутреннего) кровотечения", [
        "Снаружи крови может не быть видно",
        "Бледность, холодный пот, нарастающая слабость, жажда",
        "Слабый учащённый пульс, частое дыхание",
        "Боль и напряжение живота",
        "Действия: СМП, положение, холод, не кормить и не поить",
    ], highlight="Внутреннее кровотечение на месте не останавливают — нужна СМП")
    slide_new_summary(prs, [
        "Артериальное — алая пульсирующая струя",
        "Венозное — тёмная струя; капиллярное — сочение",
        "Потеря >30% ОЦК угрожает жизни",
        "При подозрении на внутреннее — срочный вызов СМП",
    ])
    slide_thanks_src(prs)
    p = OUT / "Признаки_наружного_кровотечения.pptx"
    prs.save(p)
    return p


def build_pres3():
    prs = new_prs()
    slide_title_src(prs)
    slide_toc(prs, [
        "Прямое давление и давящая повязка",
        "Пальцевое прижатие и максимальное сгибание",
        "Жгут и подручные средства",
        "Алгоритм, записка, ошибки и ограничения",
    ])
    src_pressure_bandage(prs)
    src_artery_grid(prs, "сонная", [
        ("carotid_point.png",
         "Общая сонная артерия прижимается на передней поверхности шеи снаружи от гортани "
         "на стороне повреждения."),
        ("carotid_4fingers.png",
         "Давление в указанную точку может осуществляться четырьмя пальцами одновременно "
         "по направлению к позвоночнику, при этом сонная артерия придавливается к нему."),
        ("carotid_thumb.png",
         "Другим вариантом пальцевого прижатия сонной артерии является давление в ту же точку "
         "большим пальцем по направлению к позвоночнику."),
    ])
    src_artery_grid(prs, "подключичная", [
        ("subclavian_1.png",
         "Подключичная артерия прижимается в ямке над ключицей к первому ребру"),
        ("subclavian_2.png",
         "Осуществлять давление в точку прижатия подключичной артерии можно с помощью "
         "четырех выпрямленных пальцев"),
        ("subclavian_3.png",
         "Другим способом пальцевого прижатия подключичной артерии является давление "
         "согнутыми пальцами"),
    ])
    src_artery_grid(prs, "плечевая", [
        ("brachial_1.png",
         "Плечевая артерия прижимается к плечевой кости с внутренней стороны между бицепсом "
         "и трицепсом в средней трети плеча"),
        ("brachial_2.png",
         "Давление на точку прижатия осуществляется с помощью четырех пальцев кисти, "
         "обхватывающей плечо пострадавшего сверху или снизу"),
    ])
    src_artery_grid(prs, "подмышечная", [
        ("axillary_1.png",
         "Подмышечная артерия прижимается к плечевой кости в подмышечной впадине при "
         "кровотечении из раны плеча ниже плечевого сустава."),
        ("axillary_2.png",
         "Давление производится прямыми, жестко зафиксированными пальцами с достаточной силой "
         "в направлении плечевого сустава. Область плечевого сустава придерживать другой рукой."),
    ])
    src_artery_grid(prs, "бедренная", [
        ("femoral_1.png",
         "Бедренная артерия прижимается ниже паховой складки при кровотечении из ран "
         "в области бедра."),
        ("femoral_2.png",
         "Давление выполняется кулаком, зафиксированным второй рукой, весом тела участника "
         "оказания первой помощи"),
    ])
    src_flexion(prs)
    src_tourniquet(prs)
    src_improvised(prs)

    slide_new_scheme(prs, "Пошаговый алгоритм наложения жгута", [
        "Показания: давление и повязка неэффективны / отрыв",
        "Место: выше раны, на плечо или бедро",
        "На одежду / подкладку; первый тур — максимальное натяжение",
        "Кровь остановилась; пульс дистальнее не определяется",
        "Записка со временем; жгут должен быть виден",
        "Лимит: 60 мин (тепло) / 30 мин (холод)",
    ], key_idx=5)
    slide_new_list(prs, "Записка под жгут — что указать", [
        "Точное время наложения: ЧЧ:ММ",
        "Фамилия / кто наложил — по возможности",
        "Кратко — место происшествия",
        "Вложить под жгут или прикрепить на видном месте",
        "Сообщить время бригаде СМП",
    ], highlight="Без записки легко превысить допустимое время жгута")
    slide_new_list(prs, "Ошибки при наложении жгута", [
        "Слабое натяжение — усиление кровотечения",
        "Слишком далеко от раны",
        "На голую кожу без подкладки",
        "Скрытие жгута одеждой",
        "Нет записки со временем",
        "Жгут без показаний",
    ])
    slide_new_list(prs, "Когда НЕЛЬЗЯ накладывать жгут", [
        "Кровотечение останавливается давлением или повязкой",
        "Капиллярное и умеренное венозное кровотечение",
        "Раны шеи, груди, живота, головы",
        "Нет конечности проксимальнее раны",
    ], alert={2}, highlight="На шею, грудь, живот и голову жгут не накладывают")
    slide_new_list(prs, "Гемостатические средства", [
        "Гемостатические салфетки / бинты из аптечки",
        "Плотно вложить в рану / на рану и сильно придавить",
        "Затем — давящая повязка",
        "Там, где жгут невозможен, или как усиление давления",
    ])
    slide_new_summary(prs, [
        "Приоритет: давление → повязка → жгут",
        "Пальцевое прижатие и сгибание — временные приёмы",
        "Жгут: на одежду, с запиской; ≤60 / ≤30 мин",
        "Гемостатики усиливают давление там, где уместно",
    ], "t max ≤ 60 мин (тепло) · ≤ 30 мин (холод)")
    slide_thanks_src(prs)
    p = OUT / "Способы_временной_остановки_кровотечения.pptx"
    prs.save(p)
    return p


def build_pres4():
    prs = new_prs()
    slide_title_src(prs)
    slide_toc(prs, [
        "Полный алгоритм действий",
        "Приоритетность способов остановки",
        "Профилактика травматического шока",
        "Подробный осмотр и контроль",
    ])
    slide_new_scheme(prs, "Полный алгоритм (чек-лист)", [
        "Оценка обстановки → безопасность (СИЗ)",
        "Обзорный осмотр — поиск кровотечения",
        "Временная остановка кровотечения",
        "Признаки жизни: сознание → дыхание",
        "СЛР при отсутствии признаков жизни",
        "Вызов скорой медицинской помощи",
        "Подробный осмотр и помощь при травмах",
        "Положение, тепло, контроль, передача СМП",
    ], key_idx=2)
    slide_new_list(prs, "Оценка состояния пострадавшего", [
        "При массивном кровотечении сначала остановите кровь",
        "Сознание: обратитесь, осторожно потрясите за плечи",
        "Дыхание: смотрю–слушаю–ощущаю не более 10 секунд",
        "Сначала то, что угрожает жизни сейчас",
    ])
    slide_new_scheme(prs, "Приоритетность способов остановки", [
        "Прямое давление на рану",
        "Давящая повязка",
        "Жгут — по показаниям",
    ], key_idx=0)
    slide_new_list(prs, "Приоритет при множественных травмах", [
        "Угрожающее жизни кровотечение",
        "Дыхательные пути и дыхание / СЛР",
        "Остальные кровотечения и раны",
        "Иммобилизация и положение",
    ])
    src_shock(prs)
    src_shock_prev(prs)
    src_exam(prs, 1)
    src_exam(prs, 2)
    slide_new_list(prs, "Действия после остановки кровотечения", [
        "Контроль повязки: при промокании усилить, не снимая",
        "Сознание, дыхание, цвет кожи",
        "Жгут: время, видимость, записка",
        "Согреть, обеспечить покой",
        "Передать бригаде: что сделано и время жгута",
    ])
    slide_new_list(prs, "Когда и как ослаблять жгут", [
        "Планово ослаблять на месте не рекомендуется",
        "Если превышен лимит 60 / 30 мин и СМП нет:",
        "Подготовить давление / повязку",
        "Медленно ослабить, оценивая кровь",
        "При струе — сразу затянуть и обновить время",
    ], alert={1}, highlight="Лимит жгута: ≤ 60 мин (тепло) · ≤ 30 мин (холод)")
    slide_new_summary(prs, [
        "Безопасность → осмотр → кровь → признаки жизни → СМП",
        "Способы: давление → повязка → жгут",
        "Профилактика шока: кровь, положение, тепло",
        "После остановки — контроль до передачи СМП",
    ], "Сначала угрожающее жизни кровотечение!")
    slide_thanks_src(prs)
    p = OUT / "Последовательность_остановки_кровотечения.pptx"
    prs.save(p)
    return p


def build_pres5():
    prs = new_prs()
    slide_title_src(prs)
    slide_toc(prs, [
        "Голова, глаза, нос, носовое кровотечение",
        "Шея, грудь, инородный предмет",
        "Живот, таз, конечности",
        "Позвоночник и травматическая ампутация",
    ])
    src_head(prs)
    src_eye_nose(prs)
    src_nosebleed(prs)
    src_steps_images(prs, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ ШЕИ", [
        ("neck_1.png", "Шаг 1 Остановка кровотечения"),
        ("neck_2.png", "Шаг 2 Фиксация шеи при перемещении"),
        ("neck_3.png", "Шаг 3 Фиксация шейного отдела"),
    ])
    slide_new_list(prs, "Кровотечение при травме шеи: артерия или вена", [
        "Артериальное (сонная): алая кровь, пульсирующая струя",
        "Венозное (яремные): тёмная обильная струя; риск воздушной эмболии",
        "Прямое давление на рану (не пережимать обе сонные артерии)",
        "Давящая повязка через плечо — не тугая круговая на шее",
        "Жгут на шею НЕ накладывают",
    ], alert={4}, highlight="Не сдавливайте дыхательные пути повязкой")
    src_steps_images(prs, "ОКАЗАНИЕ ПЕРВОЙ ПОМОЩИ ПРИ ТРАВМЕ ГРУДИ", [
        ("chest_1.png", "Шаг 1 Придать полусидячее положение"),
        ("chest_2.png", "Шаг 2 Наложить герметизирующую (окклюзионную) повязку"),
        ("chest_3.png", "Шаг 3 Обложить инородный предмет салфетками или бинтами, "
                        "наложив поверх них давящую повязку"),
    ])
    slide_new_scheme(prs, "Техника окклюзионной повязки", [
        "Придать полусидячее положение",
        "Воздухонепроницаемый материал на рану",
        "Зафиксировать пластырем",
        "При сквозном — закрыть вход и выход",
        "Контролировать дыхание",
    ], key_idx=1)
    slide_new_list(prs, "Инородный предмет в ране", [
        "НЕ ИЗВЛЕКАТЬ — предмет может тампонировать сосуд",
        "Обложить салфетками / бинтами валиками вокруг",
        "Давящая повязка поверх валиков с фиксацией",
        "Не продавливать предмет вглубь",
        "Вызвать СМП",
    ], highlight="Правило: зафиксировать — не извлекать!")
    src_abdomen(prs, False)
    src_abdomen(prs, True)
    src_limbs_text(prs)
    src_immobilization(prs)
    slide_new_list(prs, "Иммобилизация при переломе с кровотечением", [
        "Сначала остановите кровотечение",
        "Затем иммобилизация двух соседних суставов",
        "Шины — поверх одежды",
        "При открытом переломе не класть шину на выступающие отломки",
        "Костные отломки не вправлять",
    ])
    slide_new_list(prs, "Первая помощь при отрыве конечности", [
        "Немедленно жгут выше места отрыва (плечо / бедро)",
        "Давящая повязка на культю",
        "Записка со временем наложения жгута",
        "Сегмент: ткань → пакет → пакет со холодом (не мочить)",
        "Передать сегмент вместе с пострадавшим",
    ], alert={0}, highlight="Жгут при отрыве показан. t max ≤ 60 мин / 30 мин")
    src_spine(prs)
    slide_new_list(prs, "Правила перемещения при травме позвоночника", [
        "Перемещать только при угрозе жизни или для эвакуации",
        "Поверхность: ровная, жёсткая, горизонтальная",
        "Нужно несколько человек (обычно 3–5)",
        "Один постоянно удерживает голову и шейный отдел",
        "Перекладывание одним движением, без скручивания",
    ])
    slide_new_summary(prs, [
        "Голова / нос: повязки; при потере сознания — боковое положение",
        "Шея и грудь: без кругового жгута; окклюзионная повязка",
        "Живот: органы не вправлять; конечности — кровь, затем шина",
        "Позвоночник — жёсткая поверхность и фиксация шеи",
    ], "Инородный предмет — зафиксировать, не извлекать!")
    slide_thanks_src(prs)
    p = OUT / "Кровотечение_при_ранениях_областей_тела.pptx"
    prs.save(p)
    return p


def main():
    paths = []
    for fn in (build_pres1, build_pres2, build_pres3, build_pres4, build_pres5):
        p = fn()
        prs = Presentation(str(p))
        full = 0
        pics = texts = 0
        for s in prs.slides:
            for sh in s.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pics += 1
                    if sh.width >= SLIDE_W * 0.95 and sh.height >= SLIDE_H * 0.95:
                        full += 1
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts += 1
        if full:
            raise SystemExit(f"FAIL full-slide image in {p.name}")
        print(f"OK: {p.name} · {len(prs.slides)} slides · pics={pics} texts={texts} · "
              f"{p.stat().st_size // 1024} KB")
        paths.append(p)

    (OUT / "README.md").write_text(
        "# Тема 2. Наружные кровотечения\n\n"
        "Все слайды — **редактируемый текст**. Иллюстрации из исходника — "
        "**отдельные объекты**, не фон слайда.\n"
        "Новые слайды (пробелы) — Open Sans 46/26/32.\n",
        encoding="utf-8",
    )
    print("Done.")


if __name__ == "__main__":
    main()
