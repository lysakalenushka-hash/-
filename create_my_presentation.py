#!/usr/bin/env python3
"""Создаёт презентацию «Моя презентация» из 3 слайдов (python-pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

# Целевой путь (Windows)
OUTPUT = Path(
    r"C:\Users\Lukina Marina\Desktop\Cursor__pptx"
    r"\Office-MCP-Assistant-main\Office-PowerPoint-MCP-Server-main"
    r"\workspace\my_presentation.pptx"
)

# Widescreen 16:9
SLIDE_W = Emu(24384000)
SLIDE_H = Emu(13716000)

DARK = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT = RGBColor(0x00, 0x82, 0xB9)
BG = RGBColor(0xF7, 0xF7, 0xF7)
MUTED = RGBColor(0x6F, 0x6F, 0x6F)
BAR = RGBColor(0xD9, 0xD9, 0xD9)


def add_rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, *, size=24, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # --- Слайд 1: титульный ---
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_rect(s1, Emu(23000000), 0, Emu(1384000), SLIDE_H, BAR)
    add_text(
        s1, Emu(900000), Emu(4800000), Emu(20000000), Emu(1200000),
        "МОЯ ПРЕЗЕНТАЦИЯ", size=40, bold=True,
    )
    add_rect(s1, Emu(900000), Emu(6200000), Emu(9000000), Emu(25000), ACCENT)
    add_text(
        s1, Emu(900000), Emu(6500000), Emu(20000000), Emu(500000),
        "Титульный слайд", size=16, color=MUTED,
    )

    # --- Слайд 2: содержание ---
    s2 = prs.slides.add_slide(blank)
    add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_text(
        s2, Emu(900000), Emu(600000), Emu(20000000), Emu(800000),
        "СОДЕРЖАНИЕ", size=28, bold=True,
    )
    add_rect(s2, Emu(900000), Emu(1450000), Emu(16000000), Emu(25000), BAR)

    items = [
        "Слайд 1 — титульный",
        "Слайд 2 — содержание",
        "Слайд 3 — спасибо за внимание",
    ]
    box = s2.shapes.add_textbox(Emu(900000), Emu(2000000), Emu(20000000), Emu(6000000))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.space_after = Pt(12)
        run = paragraph.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(20)
        run.font.color.rgb = DARK
        run.font.name = "Arial"

    # --- Слайд 3: спасибо за внимание ---
    s3 = prs.slides.add_slide(blank)
    add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_rect(s3, Emu(23000000), 0, Emu(1384000), SLIDE_H, BAR)
    add_text(
        s3, Emu(900000), Emu(5200000), Emu(20000000), Emu(1000000),
        "СПАСИБО ЗА ВНИМАНИЕ", size=32, bold=True,
    )
    add_rect(s3, Emu(900000), Emu(6400000), Emu(10000000), Emu(25000), DARK)
    add_text(
        s3, Emu(900000), Emu(6700000), Emu(20000000), Emu(500000),
        "Моя презентация", size=14, color=MUTED,
    )

    return prs


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUTPUT))
    print(f"Сохранено: {OUTPUT}")
    print(f"Слайдов: {len(prs.slides)}")


if __name__ == "__main__":
    main()
