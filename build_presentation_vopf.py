#!/usr/bin/env python3
"""Build VOPF presentation using uploaded template styling."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_ВОПФ_трубопроводы.pptx"
IMG = "/workspace/assets/vopf_images"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

IMG_LEFT = Emu(6600000)
IMG_W = Emu(5200000)
TEXT_W = Emu(5800000)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, footer=None, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], body, 14, False, DARK)
    if footer and len(shapes) > 4:
        set_text(shapes[4], footer, 11, False, GRAY)
    elif len(shapes) > 4:
        set_text(shapes[4], "", 11, False, GRAY)
    if note:
        for sh in shapes:
            if sh.top > Emu(5800000) and sh.left < Emu(5000000):
                set_text(sh, note, 14, True, AMBER)
                break


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 30, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    for sh in shapes[5:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            set_text(pairs[i][1], item, 14, False, DARK)

    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        set_text(pairs[i][1], "", 14, False, DARK)

    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)

    if note:
        for sh in shapes:
            if sh.top > Emu(5900000) and sh.left > Emu(700000):
                set_text(sh, note, 14, True, AMBER)
                break


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3600000), width=Emu(10700000), height=Emu(2500000))
    else:
        slide.shapes.add_picture(path, IMG_LEFT, Emu(1200000), width=IMG_W, height=Emu(4800000))


def add_highlight(slide, parts):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(548640), Emu(2400000), TEXT_W, Emu(650000))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    box.line.color.rgb = RED
    box.line.width = Pt(1.5)
    tf = slide.shapes.add_textbox(Emu(640000), Emu(2520000), TEXT_W - Emu(100000), Emu(500000)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        set_run(r, 14, bold, color)


def add_definition_block(slide, items):
    """Add term-definition pairs as highlighted blocks."""
    top = Emu(2200000)
    for term, definition in items:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(548640), top, TEXT_W, Emu(900000))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        box.line.width = Pt(0.75)
        tf = slide.shapes.add_textbox(Emu(640000), top + Emu(80000), TEXT_W - Emu(100000), Emu(750000)).text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{term} — "
        set_run(r1, 13, True, RED)
        r2 = p.add_run()
        r2.text = definition
        set_run(r2, 13, False, GRAY)
        top += Emu(1000000)


def main():
    tpl = Presentation(TEMPLATE)
    text_tpl = [deepcopy(s.element) for s in tpl.slides[0].shapes]
    list_tpl = [deepcopy(s.element) for s in tpl.slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # 1. Title
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Охрана труда",
        "Вредные и опасные производственные факторы (ВОПФ) при работе с трубопроводами пара и горячей воды",
        "Промышленная безопасность и охрана труда",
        "Модуль охватывает классификацию вредных и опасных производственных факторов, "
        "термины и определения, источники опасности и виды опасных зон при обслуживании "
        "трубопроводов пара и горячей воды.",
        note="Все требования основаны на нормах охраны труда и промышленной безопасности.",
    )
    add_image(s, f"{IMG}/worker_industrial_bw.png")

    # 2. Definitions
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Основные понятия",
        "Опасный и вредный производственный фактор",
        "Определения по охране труда",
        "Опасный производственный фактор — производственный фактор, воздействие которого "
        "на работника может привести к его травме.\n\n"
        "Вредный производственный фактор — производственный фактор, воздействие которого "
        "на работника может привести к его профессиональному заболеванию.",
    )
    add_image(s, f"{IMG}/worker_industrial_bw.png")

    # 3. Factor groups
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Классификация",
        "Укрупнённые группы ВОПФ",
        "Вредные и опасные производственные факторы можно разделить на группы:",
        "Основные группы факторов",
        [
            "физические — шум, вибрация, давление, температура, электрический ток",
            "химические — пыль, токсичные и ядовитые газы и жидкости",
            "биологические — микроорганизмы и макроорганизмы",
            "психофизиологические — физические и нервно-психические перегрузки",
        ],
    )

    # 4. Factor diagram (text)
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Классификация",
        "Физические, химические, биологические и психофизиологические факторы",
        "При эксплуатации трубопроводов пара и горячей воды наиболее значимы:",
        "Категории факторов",
        [
            "физические: вибрация, давление, электрический ток, температура, шум, излучения",
            "химические: пыль, токсичные и ядовитые газы и жидкости",
            "биологические: микроорганизмы, макроорганизмы",
            "психофизиологические: физические и нервно-психические перегрузки",
            "работа на высоте, недостаточная освещённость, стеснённые условия",
        ],
        note="Физические факторы — основная группа рисков при работе с трубопроводами.",
    )

    # 5. Production factors list (part 1)
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "ВОПФ при эксплуатации",
        "Вредные и опасные производственные факторы (1/2)",
        "На работника при обслуживании трубопроводов пара и горячей воды могут воздействовать:",
        "Опасные и вредные факторы",
        [
            "движущиеся машины и механизмы, подвижные части производственного оборудования",
            "повышенное напряжение в электрической цепи",
            "повышенная или пониженная температура воздуха рабочей зоны",
            "повышенная температура поверхностей оборудования, материалов",
            "повышенное давление воды, пара в трубопроводах",
            "повышенный уровень шума на рабочем месте",
            "повышенная запылённость и загазованность воздуха",
            "движущиеся транспортные средства, грузоподъёмные механизмы",
        ],
    )
    add_image(s, f"{IMG}/industrial_steam_sidebar.png")

    # 6. Production factors list (part 2)
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "ВОПФ при эксплуатации",
        "Вредные и опасные производственные факторы (2/2)",
        "Дополнительные факторы при выполнении работ:",
        "Дополнительные факторы риска",
        [
            "расположение рабочего места на значительной высоте (глубине) относительно поверхности земли",
            "возможность концентрации вредных и опасных веществ в воздухе рабочей зоны (пропан, метан, углекислый газ и др.)",
            "недостаточная освещённость рабочей зоны",
            "стеснённые условия работы (в камерах, колодцах)",
        ],
    )

    # 7. Terms and definitions
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Термины и определения",
        "Фактор риска, опасность, источник опасности",
        "Ключевые понятия промышленной безопасности",
        "",
    )
    add_definition_block(s, [
        (
            "Фактор риска",
            "характеристика производственной среды и (или) трудового процесса (источника опасности), "
            "которая при воздействии на организм работника может привести к утрате здоровья и (или) травмированию.",
        ),
        (
            "Опасность",
            "объект, ситуация или действие, которые способны нанести вред человеку в виде травмы "
            "или ухудшения здоровья, или их сочетания.",
        ),
        (
            "Источник опасности",
            "объекты или деятельность, которые являются причиной возникновения рисков.",
        ),
    ])
    add_image(s, f"{IMG}/steam_gauge_leak.png")

    # 8. Source of danger concept
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Источники опасности",
        "Источник опасности → Опасность",
        "Связь между источником опасности и последствиями",
        "Источник опасности — объект или деятельность, вызывающая риск. "
        "При воздействии источника опасности на работника возникает опасность — "
        "ситуация, способная нанести вред в виде травмы или ухудшения здоровья.\n\n"
        "Пример: разгерметизация трубопровода (источник) → ожоги, травмы (опасность).",
        note="Идентификация источников опасности — основа оценки профессиональных рисков.",
    )
    add_image(s, f"{IMG}/hazard_pipe_leak.png")

    # 9. Hazard sources types
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Источники опасности",
        "Классификация источников опасности",
        "При эксплуатации трубопроводов выделяют следующие источники:",
        "Типы источников опасности",
        [
            "источники давления — пневматические и гидравлические системы",
            "термические — трубопроводы с высокой температурой поверхности",
            "электрические — электрооборудование, электродвигатели, кабельные линии",
            "химические — химводоподготовка, агрессивные среды",
        ],
    )

    # 10. Hazardous zones overview
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Опасные зоны",
        "Виды опасных зон на производстве",
        "Зоны, в которых расположены источники опасности",
        "При обслуживании трубопроводов пара и горячей воды выделяют зоны "
        "пневматических и гидравлических, термических, электрических и химических "
        "источников опасности. Каждая зона требует соблюдения специальных мер безопасности.",
    )

    # 11. Pneumatic/hydraulic zone
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Опасные зоны",
        "Зона пневматических и гидравлических источников опасности",
        "Трубопроводы, находящиеся под давлением",
        "Опасности при эксплуатации трубопроводов под давлением:\n"
        "• получение ожогов под воздействием высоких температур\n"
        "• получение иных травм вследствие ограниченной видимости из-за пара",
    )
    add_image(s, f"{IMG}/zone_pneumatic_pipes.png", bottom=True)

    # 12. Thermal zone
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Опасные зоны",
        "Зона термических источников опасности",
        "Трубопроводы с высокой температурой поверхности",
        "Опасности при эксплуатации трубопроводов с высокой температурой поверхности:\n"
        "• получение ожогов от прикосновений\n"
        "• получение тепловых ударов вследствие высокой температуры производственного помещения",
    )
    add_image(s, f"{IMG}/zone_thermal_pipes.png", bottom=True)

    # 13. Electrical zone
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Опасные зоны",
        "Зона электрических источников опасности",
        "Электрооборудование насосных и котельных установок",
        "Зона электрических источников опасности включает электродвигатели, "
        "насосное оборудование, распределительные устройства. "
        "Основные опасности: поражение электрическим током, дуговой разряд, "
        "возгорание при коротком замыкании.",
    )
    add_image(s, f"{IMG}/zone_electrical_motors.png", bottom=True)

    # 14. Chemical zone
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Опасные зоны",
        "Зона химических источников опасности",
        "Химводоподготовка и обработка теплоносителя",
        "Зона химических источников опасности — химводоподготовка. "
        "Опасности: воздействие химических реагентов, отравление парами, "
        "ожоги кислотами и щелочами при обслуживании систем водоподготовки.",
    )
    add_image(s, f"{IMG}/zone_chemical_treatment.png", bottom=True)

    # 15. Pneumatic causes
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Пневмогидравлические опасности",
        "Зона пневматических и гидравлических источников опасности",
        "Причины разгерметизации или разрушения систем повышенного давления:",
        "Причины аварий",
        [
            "внешние механические воздействия",
            "старение систем и износ элементов",
            "коррозия металла трубопроводов",
            "ошибки проектирования и монтажа",
            "превышение допустимого рабочего давления",
            "нарушение правил эксплуатации и ремонта",
        ],
    )
    add_image(s, f"{IMG}/zone_field_pipes.png", bottom=True)

    # 16. Steam work hazards
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "Меры безопасности",
        "Работы в зонах повышенного давления и температуры",
        "Требования при обслуживании трубопроводов",
        "При выполнении работ в зонах пневматических, гидравлических и термических "
        "источников опасности необходимо:\n"
        "• использовать СИЗ (перчатки, каска, защитные очки)\n"
        "• не прикасаться к горячим поверхностям без защиты\n"
        "• ограждать место работы при появлении пара\n"
        "• соблюдать инструкции по охране труда и наряд-допуск",
        note="ОСТОРОЖНО! ГОРЯЧАЯ ПОВЕРХНОСТЬ. ЗАПРЕЩАЕТСЯ ПРИКАСАТЬСЯ — ОПАСНО!",
    )
    add_image(s, f"{IMG}/zone_steam_worker.png", bottom=True)

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
