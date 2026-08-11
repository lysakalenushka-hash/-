#!/usr/bin/env python3
"""Тема 3 первой помощи: 6 презентаций по п. 3.1–3.6 (стиль темы 1 + пособие Минздрава)."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from docx.parts.image import ImagePart
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

GUIDE = Path("/workspace/first_aid_uploads_src/uchebnoe_posobie_pp_d.docx")
IMG_DIR = Path("/tmp/guide_images_tema3")
OUT = Path("/workspace/first_aid_tema3")
ZIP = Path("/workspace/Тема3_презентации_3.1-3.6.zip")

# Стиль темы 1 (после restyle): PFDinDisplayPro / Open Sans, серый текст, красная линия
TITLE_FONT = "PFDinDisplayPro-Regular"
EYEBROW_FONT = "Open Sans"
BODY_FONT = "Open Sans"
TITLE_COLOR = RGBColor(0x40, 0x40, 0x40)
EYEBROW_COLOR = RGBColor(0x53, 0x53, 0x53)
BODY_COLOR = RGBColor(0x40, 0x40, 0x40)
RED = RGBColor(0xE3, 0x06, 0x13)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF7, 0xF8)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
NS_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def extract_images_and_fig_map() -> dict[int, Path]:
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    d = Document(str(GUIDE))
    img_by_rid: dict[str, Path] = {}
    n = 0
    for rel_id, rel in d.part.rels.items():
        if isinstance(rel.target_part, ImagePart):
            n += 1
            ct = rel.target_part.content_type
            ext = {"image/png": "png", "image/jpeg": "jpg", "image/jpg": "jpg"}.get(ct, "bin")
            path = IMG_DIR / f"img_{n:03d}.{ext}"
            path.write_bytes(rel.target_part.blob)
            img_by_rid[rel_id] = path

    events: list[tuple[str, object]] = []
    for child in d.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            texts = [t.text for t in child.findall(f".//{NS_W}t") if t.text]
            text = "".join(texts).strip()
            if text:
                events.append(("text", text))
            for b in child.findall(f".//{NS_A}blip"):
                rid = b.get(f"{NS_R}embed")
                if rid in img_by_rid:
                    events.append(("img", img_by_rid[rid]))
        elif tag == "tbl":
            for b in child.findall(f".//{NS_A}blip"):
                rid = b.get(f"{NS_R}embed")
                if rid in img_by_rid:
                    events.append(("img", img_by_rid[rid]))

    fig_map: dict[int, Path] = {}
    for i, (kind, val) in enumerate(events):
        if kind != "text":
            continue
        m = re.search(r"Рис(?:унок|\.)\s*(\d+)", str(val), re.I)
        if not m:
            continue
        num = int(m.group(1))
        for j in range(i - 1, max(-1, i - 6), -1):
            if events[j][0] == "img":
                fig_map[num] = events[j][1]  # type: ignore
                break
        if num not in fig_map:
            for j in range(i + 1, min(len(events), i + 6)):
                if events[j][0] == "img":
                    fig_map[num] = events[j][1]  # type: ignore
                    break
    return fig_map


def set_run(run, *, size=18, bold=False, color=BODY_COLOR, font=BODY_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def clear_shapes(slide):
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def add_blank(prs) -> object:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    clear_shapes(slide)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return slide


def add_header(slide, title: str, eyebrow: str, num: int | None = None):
    # Title
    box = slide.shapes.add_textbox(Emu(900000), Emu(480000), Emu(10400000), Emu(900000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=28, bold=True, color=TITLE_COLOR, font=TITLE_FONT)

    # Red underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(900000), Emu(1350000), Emu(7800000), Emu(36000))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    # Eyebrow
    eb = slide.shapes.add_textbox(Emu(900000), Emu(1420000), Emu(9000000), Emu(360000))
    p = eb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = eyebrow
    set_run(run, size=12, color=EYEBROW_COLOR, font=EYEBROW_FONT)

    if num is not None:
        nb = slide.shapes.add_textbox(Emu(20000), Emu(3000000), Emu(500000), Emu(500000))
        p = nb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(num)
        set_run(run, size=18, bold=True, color=RED, font=TITLE_FONT)


def add_bullets(slide, lines: list[str], *, left=Emu(900000), top=Emu(1900000), width=Emu(10400000), height=Emu(4300000), size=16):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.fill.background()
    try:
        box.adjustments[0] = 0.05
    except Exception:
        pass

    tb = slide.shapes.add_textbox(left + Emu(200000), top + Emu(160000), width - Emu(400000), height - Emu(320000))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=BODY_COLOR, font=BODY_FONT)


def add_picture_right(slide, path: Path | None, caption: str = ""):
    if path is None or not path.exists():
        return
    # Right column picture area
    left, top, width, height = Emu(7200000), Emu(1900000), Emu(4300000), Emu(4000000)
    slide.shapes.add_picture(str(path), left, top, width=width)
    if caption:
        cap = slide.shapes.add_textbox(left, Emu(6000000), width, Emu(400000))
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        set_run(run, size=11, color=EYEBROW_COLOR, font=EYEBROW_FONT)


def add_title_slide(prs, title: str, subtitle: str, eyebrow: str):
    slide = add_blank(prs)
    # accent bar left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(160000), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    eb = slide.shapes.add_textbox(Emu(900000), Emu(1800000), Emu(10000000), Emu(400000))
    p = eb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = eyebrow
    set_run(run, size=14, color=EYEBROW_COLOR, font=EYEBROW_FONT)

    tb = slide.shapes.add_textbox(Emu(900000), Emu(2300000), Emu(10400000), Emu(1800000))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=36, bold=True, color=TITLE_COLOR, font=TITLE_FONT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(900000), Emu(4300000), Emu(4000000), Emu(36000))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    sb = slide.shapes.add_textbox(Emu(900000), Emu(4500000), Emu(10000000), Emu(800000))
    p = sb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    set_run(run, size=16, color=BODY_COLOR, font=BODY_FONT)

    src = slide.shapes.add_textbox(Emu(900000), Emu(6200000), Emu(10000000), Emu(300000))
    p = src.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "По учебному пособию Минздрава России по первой помощи (2025) · Прил. № 2 к ПП РФ № 2464"
    set_run(run, size=11, color=EYEBROW_COLOR, font=EYEBROW_FONT)
    return slide


def add_content_slide(prs, title, eyebrow, lines, *, fig=None, fig_map=None, caption="", num=None, half=False):
    slide = add_blank(prs)
    add_header(slide, title, eyebrow, num=num)
    if half and fig and fig_map and fig_map.get(fig):
        add_bullets(
            slide,
            lines,
            left=Emu(700000),
            top=Emu(1900000),
            width=Emu(6200000),
            height=Emu(4300000),
            size=15,
        )
        add_picture_right(slide, fig_map.get(fig), caption=caption or f"Рис. {fig} (пособие Минздрава)")
    else:
        add_bullets(slide, lines, size=16)
        if fig and fig_map and fig_map.get(fig):
            # small image bottom-right overlay if full-width text not used — skip to avoid clutter
            pass
    return slide


def add_image_grid_slide(prs, title, eyebrow, items: list[tuple[int, str]], fig_map, num=None):
    """items: list of (fig_num, label) up to 3."""
    slide = add_blank(prs)
    add_header(slide, title, eyebrow, num=num)
    n = len(items)
    if n == 0:
        return slide
    gap = Emu(200000)
    usable = Emu(10400000)
    w = (usable - gap * (n - 1)) // n
    left0 = Emu(900000)
    top = Emu(1950000)
    for i, (fig, label) in enumerate(items):
        left = left0 + i * (w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Emu(4200000))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.fill.background()
        try:
            card.adjustments[0] = 0.04
        except Exception:
            pass
        path = fig_map.get(fig)
        if path and path.exists():
            slide.shapes.add_picture(str(path), left + Emu(120000), top + Emu(120000), width=w - Emu(240000))
        lb = slide.shapes.add_textbox(left + Emu(120000), top + Emu(3400000), w - Emu(240000), Emu(700000))
        tf = lb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_run(run, size=12, bold=True, color=BODY_COLOR, font=BODY_FONT)
    return slide


def build_all(fig_map: dict[int, Path]):
    OUT.mkdir(parents=True, exist_ok=True)
    decks = []

    # ---------- 3.1 ----------
    spec = {
        "out": "1_Определение_признаков_жизни.pptx",
        "title": "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ",
        "eyebrow": "Тема 3 · п. 3.1  Определение признаков жизни",
        "list_item": "Определение признаков жизни.",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        spec["title"],
        "Простейшие способы проверки сознания и дыхания. Ориентир для начала СЛР.",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "ЗАЧЕМ ПРОВЕРЯТЬ ПРИЗНАКИ ЖИЗНИ",
        spec["eyebrow"],
        [
            "Внезапная смерть (остановка дыхания и кровообращения) может быть вызвана заболеваниями",
            "(инфаркт, нарушения ритма) или внешним воздействием (травма, электротравма, утопление и др.).",
            "",
            "Вне зависимости от причины исчезновения признаков жизни сердечно-легочная реанимация",
            "проводится по единому алгоритму.",
            "",
            "Решение о начале СЛР принимается по отсутствию сознания и нормального дыхания",
            "(проверка пульса на магистральных артериях для этого не рекомендуется — недостаточная точность).",
        ],
        num=2,
    )
    add_content_slide(
        prs,
        "ПРОВЕРКА СОЗНАНИЯ",
        spec["eyebrow"],
        [
            "• Аккуратно потормошите пострадавшего за плечи.",
            "• Громко спросите: «Что с вами? Нужна ли вам помощь?»",
            "• Человек без сознания не ответит и не отреагирует.",
            "",
            "Нельзя: давить на болевые точки, хлопать по щекам и т.п.",
        ],
        fig=31,
        fig_map=fig_map,
        half=True,
        caption="Рис. 31. Проверка сознания",
        num=3,
    )
    add_content_slide(
        prs,
        "ПРОВЕРКА ДЫХАНИЯ",
        spec["eyebrow"],
        [
            "1. Восстановите проходимость дыхательных путей: рука на лбу, двумя пальцами подбородок —",
            "   запрокинуть голову и поднять подбородок.",
            "2. Наклонитесь щекой и ухом ко рту и носу на 10 секунд:",
            "   увидеть движения груди, услышать дыхание, почувствовать выдох на щеке.",
            "",
            "Нет нормального дыхания / агональное (редкое, ненормальное) → вызвать СМП и начать СЛР.",
        ],
        fig=32,
        fig_map=fig_map,
        half=True,
        caption="Рис. 32. Проверка дыхания",
        num=4,
    )
    add_content_slide(
        prs,
        "КРАТКИЙ АЛГОРИТМ",
        spec["eyebrow"],
        [
            "1. Безопасность места происшествия.",
            "2. Проверка сознания (тормошение за плечи + громкий вопрос).",
            "3. При отсутствии сознания — открытие дыхательных путей.",
            "4. Проверка нормального дыхания ≤ 10 секунд.",
            "5. Нет нормального дыхания → вызов СМП (112 / 103) и СЛР.",
            "",
            "Ориентир: нет сознания + нет нормального дыхания = нет признаков жизни для целей СЛР.",
        ],
        num=5,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # ---------- 3.2 ----------
    spec = {
        "out": "2_Восстановление_и_поддержание_проходимости_дыхательных_путей.pptx",
        "title": "ВОССТАНОВЛЕНИЕ И ПОДДЕРЖАНИЕ ПРОХОДИМОСТИ ДЫХАТЕЛЬНЫХ ПУТЕЙ",
        "eyebrow": "Тема 3 · п. 3.2  Проходимость дыхательных путей",
        "list_item": "Восстановление и поддержание проходимости дыхательных путей.",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        "ПРОХОДИМОСТЬ ДЫХАТЕЛЬНЫХ ПУТЕЙ",
        "Открытие дыхательных путей и устойчивое боковое положение.",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "ВОССТАНОВЛЕНИЕ ПРОХОДИМОСТИ",
        spec["eyebrow"],
        [
            "При отсутствии сознания перед проверкой дыхания и при СЛР:",
            "• одну руку — на лоб пострадавшего;",
            "• двумя пальцами другой руки — за подбородок;",
            "• запрокинуть голову и поднять подбородок.",
            "",
            "Не рекомендуется: валики под шею/спину, обязательный спинальный щит",
            "для поддержания проходимости дыхательных путей при СЛР.",
        ],
        num=2,
    )
    add_content_slide(
        prs,
        "КОГДА НУЖНО УСТОЙЧИВОЕ БОКОВОЕ ПОЛОЖЕНИЕ",
        spec["eyebrow"],
        [
            "Пострадавший без сознания, но с сохранённым дыханием, лёжа на спине,",
            "рискует перекрытием дыхательных путей языком, рвотой, кровью.",
            "",
            "Возможные ситуации: обморок, травма головы, отравление, перегревание и др.",
            "",
            "Наиболее эффективный способ поддержания проходимости — устойчивое боковое положение.",
            "Если его нельзя придать (травма и др.) — удерживать запрокинутую голову с подъёмом подбородка.",
        ],
        num=3,
    )
    add_image_grid_slide(
        prs,
        "УСТОЙЧИВОЕ БОКОВОЕ ПОЛОЖЕНИЕ: ШАГИ 1–3",
        spec["eyebrow"],
        [
            (42, "Шаг 1. Ближняя рука под 90°"),
            (43, "Шаг 2. Дальняя рука к щеке"),
            (44, "Шаг 3. Поворот на себя"),
        ],
        fig_map,
        num=4,
    )
    add_content_slide(
        prs,
        "ШАГ 4 И НАБЛЮДЕНИЕ",
        spec["eyebrow"],
        [
            "Шаг 4: после поворота набок подтянуть верхнюю ногу ближе к животу,",
            "слегка запрокинуть голову, проверить дыхание у рта и носа.",
            "",
            "Далее: наблюдать до приезда СМП, регулярно оценивать дыхание;",
            "каждые 30 минут поворачивать на другой бок.",
        ],
        fig=46,
        fig_map=fig_map,
        half=True,
        caption="Рис. 46. Итоговое положение",
        num=5,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # ---------- 3.3 ----------
    spec = {
        "out": "3_Техника_проведения_сердечно-легочной_реанимации.pptx",
        "title": "ТЕХНИКА ПРОВЕДЕНИЯ СЕРДЕЧНО-ЛЕГОЧНОЙ РЕАНИМАЦИИ",
        "eyebrow": "Тема 3 · п. 3.3  Техника СЛР",
        "list_item": "Техника проведения сердечно-легочной реанимации.",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        "СЕРДЕЧНО-ЛЕГОЧНАЯ РЕАНИМАЦИЯ",
        "Последовательность и техника по пособию Минздрава (2025).",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "СТАРТ АЛГОРИТМА",
        spec["eyebrow"],
        [
            "1. Оценка безопасности → устранение/минимизация угроз.",
            "2. Проверка сознания.",
            "3. Открытие дыхательных путей + проверка дыхания ≤ 10 с.",
            "4. Нет нормального дыхания → громко позвать конкретного помощника вызвать СМП",
            "   (или вызвать самому по громкой связи).",
            "",
            "При вызове сообщить: место и что произошло; число пострадавших и что с ними;",
            "какая помощь оказывается. Не отключаться, пока диспетчер не подтвердит приём.",
            "Телефоны: 112 (также 101, 102, 103 или региональные).",
        ],
        fig=33,
        fig_map=fig_map,
        half=True,
        caption="Рис. 33–34. Вызов помощи",
        num=2,
    )
    add_content_slide(
        prs,
        "ДАВЛЕНИЕ НА ГРУДИНУ",
        spec["eyebrow"],
        [
            "Пострадавший на спине на твёрдой ровной поверхности.",
            "• Основание ладони — на центр груди; вторая рука сверху; кисти в замок.",
            "• Руки прямые в локтях; плечи над пострадавшим; давление вертикально вниз.",
            "• Глубина примерно 5–6 см у взрослых; полное расправление груди после каждого надавливания.",
            "• Частота 100–120 в минуту.",
            "",
            "После 30 надавливаний — 2 вдоха искусственного дыхания.",
        ],
        fig=35,
        fig_map=fig_map,
        half=True,
        caption="Рис. 35. Положение рук",
        num=3,
    )
    add_content_slide(
        prs,
        "ИСКУССТВЕННОЕ ДЫХАНИЕ",
        spec["eyebrow"],
        [
            "Метод «рот ко рту»: открыть дыхательные пути, зажать нос, 2 последовательных вдоха.",
            "Признак эффективности — начало подъёма груди. На 2 вдоха — не более 10 секунд;",
            "не более двух попыток за цикл.",
            "",
            "Рекомендуется устройство «Рот–Устройство–Рот» из аптечки.",
            "При повреждении губ — метод «рот к носу».",
            "Если вдохи невозможны — только компрессии.",
            "После ИВЛ — прополоскать рот оказывающему помощь.",
            "",
            "Цикл: 30 надавливаний : 2 вдоха — до прекращения СЛР.",
        ],
        fig=37,
        fig_map=fig_map,
        half=True,
        caption="Рис. 37. Искусственное дыхание",
        num=4,
    )
    add_content_slide(
        prs,
        "ПРЕКРАЩЕНИЕ СЛР И ТИПИЧНЫЕ ОШИБКИ",
        spec["eyebrow"],
        [
            "Продолжать до прибытия СМП/спецслужб и их распоряжения; до появления признаков жизни;",
            "можно прекратить при угрозе для оказывающего помощь.",
            "При усталости — привлечь помощника (смена примерно каждые 2 минуты).",
            "Можно не начинать при явных признаках нежизнеспособности.",
            "",
            "Ошибки: неверная последовательность; глубина/частота/точка давления;",
            "неполное открытие дыхательных путей; паузы > 10 с; соотношение не 30:2.",
            "Осложнение: перелом рёбер — чаще при неверной точке/избыточной силе.",
        ],
        num=5,
    )
    add_content_slide(
        prs,
        "ОСОБЕННОСТИ: ДЕТИ И УТОПЛЕНИЕ",
        spec["eyebrow"],
        [
            "Дети: та же последовательность и соотношение 30:2; глубина ~⅓ переднезаднего размера",
            "(≈4 см до 1 года, ≈5 см старше). До 1 года — двумя пальцами; старше — одной/двумя руками.",
            "После определения отсутствия признаков жизни у ребёнка эффективнее сразу 5 вдохов,",
            "затем 30:2 (как в пособии).",
            "",
            "Утопление: после извлечения из воды при отсутствии признаков жизни — 5 вдохов,",
            "затем 30 надавливаний : 2 вдоха.",
        ],
        fig=38,
        fig_map=fig_map,
        half=True,
        caption="Рис. 38–39. СЛР у детей",
        num=6,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # ---------- 3.4 ----------
    spec = {
        "out": "4_Использование_АНД.pptx",
        "title": "ИСПОЛЬЗОВАНИЕ АВТОМАТИЧЕСКОГО НАРУЖНОГО ДЕФИБРИЛЛЯТОРА",
        "eyebrow": "Тема 3 · п. 3.4  АНД (при наличии)",
        "list_item": "Использование автоматического наружного дефибриллятора (при его наличии).",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        "АВТОМАТИЧЕСКИЙ НАРУЖНЫЙ ДЕФИБРИЛЛЯТОР (АНД)",
        "Дополнение к СЛР: анализ ритма и разряд по команде устройства.",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "ЧТО ТАКОЕ АНД",
        spec["eyebrow"],
        [
            "АНД автоматически анализирует ритм сердца и при необходимости наносит разряд,",
            "способствующий восстановлению нормальной работы сердца.",
            "",
            "Применяется в дополнение к СЛР и значительно повышает шансы на успех.",
            "Разрешён приказом Минздрава № 220н (при наличии).",
        ],
        num=2,
    )
    add_content_slide(
        prs,
        "ПОРЯДОК ПРИМЕНЕНИЯ",
        spec["eyebrow"],
        [
            "1. При утрате признаков жизни одновременно с вызовом СМП доставить АНД на место.",
            "2. Включить АНД и следовать голосовым/экранным подсказкам.",
            "3. Электроды: один — справа под ключицей (правее грудины, не на грудину!);",
            "   второй — на левую половину груди.",
            "4. При втором участнике — компрессии не прерывать во время наложения электродов.",
            "5. Обильный волосяной покров — сбрить (станок часто в комплекте).",
            "   Не накладывать на имплантируемый кардиостимулятор/дефибриллятор.",
        ],
        fig=40,
        fig_map=fig_map,
        half=True,
        caption="Рис. 40. АНД",
        num=3,
    )
    add_content_slide(
        prs,
        "РАЗРЯД И ДАЛЬНЕЙШИЕ ДЕЙСТВИЯ",
        spec["eyebrow"],
        [
            "АНД анализирует ритм и сообщает, нужен ли разряд.",
            "Перед разрядом убедиться, что никто (включая вас) не касается пострадавшего → «Разряд».",
            "",
            "После разряда сразу продолжить СЛР, пока АНД снова не попросит паузу для анализа.",
            "Если разряд не рекомендован — немедленно возобновить СЛР и следовать командам АНД.",
            "",
            "Дети 1–8 лет: детские электроды/режим; иначе стандартные на грудь и спину + детские настройки.",
            "Утопление: извлечь из воды, осушить грудь перед электродами.",
        ],
        fig=41,
        fig_map=fig_map,
        half=True,
        caption="Рис. 41. Разряд АНД",
        num=4,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # ---------- 3.5 ----------
    spec = {
        "out": "5_Инородное_тело_в_дыхательных_путях.pptx",
        "title": "ИНОРОДНОЕ ТЕЛО В ВЕРХНИХ ДЫХАТЕЛЬНЫХ ПУТЯХ",
        "eyebrow": "Тема 3 · п. 3.5  Нарушение проходимости инородным телом",
        "list_item": "Первая помощь при нарушении проходимости верхних дыхательных путей, вызванном инородным телом.",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        "ИНОРОДНОЕ ТЕЛО В ДЫХАТЕЛЬНЫХ ПУТЯХ",
        "Частичная и полная закупорка. Взрослые, тучные, беременные, дети.",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "ПРИЗНАКИ И СТЕПЕНЬ НАРУШЕНИЯ",
        spec["eyebrow"],
        [
            "Типичная поза: рука у горла + попытки кашлять.",
            "Спросите: «Вы подавились?»",
            "",
            "Частичное нарушение: отвечает, может кашлять → предложить покашлять.",
            "Полное нарушение: не говорит, не дышит (или крайне затруднённо, шумно),",
            "хватает горло, может кивать → удаление инородного тела.",
        ],
        fig=47,
        fig_map=fig_map,
        half=True,
        caption="Рис. 47. Частичная закупорка — кашель",
        num=2,
    )
    add_content_slide(
        prs,
        "ПОЛНАЯ ЗАКУПОРКА У ВЗРОСЛОГО",
        spec["eyebrow"],
        [
            "1. Встать сбоку и немного сзади.",
            "2. Придерживая одной рукой, наклонить вперёд.",
            "3. 5 резких ударов основанием ладони между лопатками; после каждого — проверка.",
            "4. Если не помогло: сзади обхватить на уровне верхней части живота;",
            "   кулак над пупком большим пальцем к себе; другой рукой обхватить кулак;",
            "   надавливания внутрь и кверху до 5 раз.",
            "5. Чередовать 5 ударов по спине и 5 надавливаний на живот.",
            "При потере сознания — начать СЛР; следить за появлением инородного тела во рту.",
        ],
        fig=48,
        fig_map=fig_map,
        half=True,
        caption="Рис. 48–50. Удары и надавливания",
        num=3,
    )
    add_image_grid_slide(
        prs,
        "ОСОБЫЕ СЛУЧАИ",
        spec["eyebrow"],
        [
            (52, "Тучные / беременные: давление на нижнюю часть груди"),
            (53, "Самопомощь: упор животом в спинку стула"),
            (54, "Ребёнок до 1 года: удары по спине на предплечье"),
        ],
        fig_map,
        num=4,
    )
    add_content_slide(
        prs,
        "ДЕТИ И ВАЖНЫЕ ЗАПРЕТЫ",
        spec["eyebrow"],
        [
            "До 1 года: беззвучный плач, возбуждение, синюшность. Удалить видимое тело изо рта.",
            "Далее: на предплечье головой вниз — 5 ударов по спине; затем на спине — 5 надавливаний",
            "двумя пальцами в центр груди; чередовать; осматривать полость рта.",
            "Старше 1 года: удары по спине и надавливания над пупком с дозированием усилия.",
            "При потере сознания — СЛР.",
            "",
            "Важно! Не переворачивать ребёнка вниз головой, держа за ноги, и не трясти.",
        ],
        fig=55,
        fig_map=fig_map,
        half=True,
        caption="Рис. 55. Надавливания у младенца",
        num=5,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # ---------- 3.6 ----------
    spec = {
        "out": "6_Иные_угрожающие_нарушения_дыхания.pptx",
        "title": "ИНЫЕ УГРОЖАЮЩИЕ ЖИЗНИ И ЗДОРОВЬЮ НАРУШЕНИЯ ДЫХАНИЯ",
        "eyebrow": "Тема 3 · п. 3.6  Иные нарушения дыхания",
        "list_item": "Первая помощь при иных угрожающих жизни и здоровью нарушениях дыхания.",
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        "ИНЫЕ НАРУШЕНИЯ ДЫХАНИЯ",
        "Не инородное тело: одышка при острых и хронических заболеваниях.",
        spec["eyebrow"],
    )
    add_content_slide(
        prs,
        "КАК ПРОЯВЛЯЕТСЯ",
        spec["eyebrow"],
        [
            "Острые и хронические заболевания могут вызвать угрожающие нарушения дыхания.",
            "Типичные проявления:",
            "• одышка (затруднённое и учащённое дыхание);",
            "• бледность и синюшность кожи (особенно лица);",
            "• общее плохое самочувствие.",
            "",
            "Это не алгоритм удаления инородного тела и не СЛР при отсутствии признаков жизни —",
            "но состояние требует первой помощи и вызова СМП.",
        ],
        num=2,
    )
    add_content_slide(
        prs,
        "ЧТО ДЕЛАТЬ",
        spec["eyebrow"],
        [
            "1. Помочь пострадавшему принять лекарственные препараты, назначенные ранее лечащим врачом",
            "   (приказ № 220н допускает такую помощь).",
            "2. Придать и поддерживать оптимальное положение тела.",
            "3. Вызвать скорую медицинскую помощь.",
            "4. Контролировать состояние и оказывать психологическую поддержку",
            "   до передачи бригаде СМП.",
            "",
            "Если сознание и дыхание исчезли — перейти к алгоритму СЛР (п. 3.1–3.3).",
        ],
        num=3,
    )
    add_content_slide(
        prs,
        "СВЯЗЬ С ДРУГИМИ ПУНКТАМИ ТЕМЫ 3",
        spec["eyebrow"],
        [
            "• Нет сознания / нет нормального дыхания → СЛР ± АНД.",
            "• Есть дыхание, нет сознания → проходимость дыхательных путей (боковое положение).",
            "• Инородное тело → удаление по п. 3.5.",
            "• Иные нарушения дыхания при сохранённом сознании → положение, лекарства по назначению,",
            "  вызов СМП, контроль и поддержка.",
        ],
        num=4,
    )
    path = OUT / spec["out"]
    prs.save(path)
    decks.append((spec, path))

    # README
    readme = OUT / "README.md"
    lines = [
        "# Тема 3. Оказание первой помощи при отсутствии сознания, остановке дыхания и кровообращения",
        "",
        "Шесть презентаций по пунктам Приложения № 2 к ПП РФ № 2464 (ред. с № 805).",
        "Стиль — как у презентаций темы 1 (PFDinDisplayPro / Open Sans, красный акцент).",
        "Содержание и рисунки — по учебному пособию Минздрава по первой помощи (2025).",
        "",
        "| № | Файл | Пункт Прил. 2 |",
        "|---|------|---------------|",
    ]
    for i, (spec, path) in enumerate(decks, 1):
        lines.append(f"| {i} | `{path.name}` | {spec['list_item']} |")
    readme.write_text("\n".join(lines) + "\n", encoding="utf-8")

    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as zf:
        for _, path in decks:
            zf.write(path, arcname=path.name)
        zf.write(readme, arcname=readme.name)

    print("Saved", len(decks), "decks →", OUT)
    print("ZIP", ZIP, ZIP.stat().st_size)


def main():
    fig_map = extract_images_and_fig_map()
    print("figures mapped:", len(fig_map))
    build_all(fig_map)


if __name__ == "__main__":
    main()
