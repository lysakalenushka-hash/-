#!/usr/bin/env python3
"""Обновить исходные презентации 2.1/2.3/2.4 под № 805 + пособие, сохранив стиль."""

from __future__ import annotations

import copy
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

SRC = Path("/workspace/first_aid_tema3_src_update")
OUT = Path("/workspace/first_aid_tema3_updated")
ZIP = Path("/workspace/Тема3_обновление_2.1_2.3_2.4.zip")

EYEBROW_THEME = "Тема 3. Оказание первой помощи при отсутствии сознания, остановке дыхания и кровообращения"
EYEBROW_31 = "Тема 3 · п. 3.1  Определение признаков жизни  ·  ред. № 805 с 01.09.2026"
EYEBROW_EXCL = "Доп. материал · исключено из Прил. № 2 с 01.09.2026 (ПП РФ № 805)"

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def first_run_style(shape):
    """Return (font_name, size, bold, color_rgb_or_None) from first non-empty run."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text:
                continue
            color = None
            try:
                if run.font.color is not None and run.font.color.rgb is not None:
                    color = run.font.color.rgb
            except Exception:
                pass
            return run.font.name, run.font.size, run.font.bold, color
    return "Open Sans", Pt(28), False, RGBColor(0x40, 0x40, 0x40)


def set_shape_plain(shape, text: str, *, keep_style=True):
    """Replace all text in shape with plain text, preserving first-run style."""
    name, size, bold, color = first_run_style(shape) if keep_style else ("Open Sans", Pt(28), False, RGBColor(0x40, 0x40, 0x40))
    tf = shape.text_frame
    # clear all paragraphs except first
    for p in list(tf.paragraphs)[1:]:
        p_el = p._p
        p_el.getparent().remove(p_el)
    p0 = tf.paragraphs[0]
    # clear runs
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    # also clear direct text
    for t_el in p0._p.findall(f".//{A_NS}t"):
        t_el.text = ""
    run = p0.add_run()
    run.text = text
    if name:
        run.font.name = name
    if size:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def set_shape_multiline(shape, lines: list[tuple[str, bool]], *, font="Open Sans", size=Pt(28), color=RGBColor(0x53, 0x53, 0x53)):
    """lines: list of (text, bold)."""
    # try inherit size/font from existing
    name0, size0, _, color0 = first_run_style(shape)
    font = name0 or font
    size = size0 or size
    color = color0 or color
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    # clear first
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for t_el in p0._p.findall(f".//{A_NS}t"):
        t_el.text = ""

    for i, (text, bold) in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        if i > 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = size
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def find_shape(slide, name_substr=None, *, has_text=None, top=None, left=None):
    for sh in slide.shapes:
        if name_substr and name_substr not in (sh.name or ""):
            continue
        if has_text is not None and sh.has_text_frame:
            joined = " ".join(p.text for p in sh.text_frame.paragraphs)
            if has_text in joined:
                return sh
        elif has_text is None:
            return sh
    return None


def update_21():
    src = SRC / "src_2.1.pptx"
    out = OUT / "3.1_Определение_признаков_жизни.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))
    slides = list(range(len(list(prs.slides._sldIdLst))))

    # Slide 1 title
    s = prs.slides[0]
    for sh in s.shapes:
        if sh.has_text_frame and "ПРИЗНАКИ ЖИЗНИ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ")

    # Slide 2 content
    s = prs.slides[1]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "ОСНОВНЫЕ ПРИЗНАКИ" in t and sh.name.startswith("TextBox"):
            set_shape_plain(sh, "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, EYEBROW_31)
        elif "К основным признакам жизни" in t:
            set_shape_multiline(
                sh,
                [
                    ("По новой редакции Прил. № 2 (№ 805): теория — «Определение признаков жизни».", False),
                    ("", False),
                    ("Признаки жизни для решения о СЛР: наличие сознания и нормального дыхания.", True),
                    ("Проверка пульса на магистральных артериях для решения о СЛР не рекомендуется (недостаточная точность).", False),
                    ("", False),
                    ("Внезапная смерть может быть вызвана заболеваниями (инфаркт, нарушения ритма) или внешним воздействием (травма, электротравма, утопление и др.). Алгоритм СЛР единый.", False),
                ],
            )

    # Slide 3 section divider — causes still in methodicheka as intro, but not separate App2 item
    s = prs.slides[2]
    for sh in s.shapes:
        if sh.has_text_frame and "ПРИЧИНЫ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "КОНТЕКСТ: ПРИЧИНЫ ОСТАНОВКИ ДЫХАНИЯ И КРОВООБРАЩЕНИЯ")

    # Slide 4
    s = prs.slides[3]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "ОСНОВНЫЕ ПРИЗНАКИ" in t and "TextBox" in sh.name:
            set_shape_plain(sh, "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, EYEBROW_31)
        elif "Внезапная смерть" in t:
            set_shape_multiline(
                sh,
                [
                    ("Заболевания: инфаркт миокарда, нарушения ритма сердца и др.", False),
                    ("Внешние воздействия: травма, поражение электрическим током, утопление и др.", False),
                    ("", False),
                    ("Вне зависимости от причины — сердечно-легочная реанимация проводится по единому алгоритму (пособие Минздрава).", True),
                ],
            )

    # Slide 5 section — remove "кровообращения" from title per new App2 focus
    s = prs.slides[4]
    for sh in s.shapes:
        if sh.has_text_frame and "СПОСОБЫ ПРОВЕРКИ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "СПОСОБЫ ПРОВЕРКИ СОЗНАНИЯ И ДЫХАНИЯ")

    # Slide 6 algorithm — refine labels + eyebrow
    s = prs.slides[5]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "АЛГОРИТМ" in t:
            set_shape_plain(sh, "ОПРЕДЕЛЕНИЕ ПРИЗНАКОВ ЖИЗНИ: ПОРЯДОК ПРОВЕРКИ")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, EYEBROW_31)
        elif "Шаг 1" in t:
            set_shape_plain(sh, "Шаг 1. Сознание: тормошение за плечи + «Что с вами? Нужна ли вам помощь?»")
        elif "Шаг 2" in t:
            set_shape_plain(sh, "Шаг 2. Открытие дыхательных путей: запрокинуть голову, поднять подбородок")
        elif "Шаг 3" in t:
            set_shape_plain(sh, "Шаг 3. Дыхание ≤ 10 с (видеть, слышать, чувствовать). Нет / агональное → СМП + СЛР")

    prs.save(out)
    print("saved", out)
    return out


def update_23():
    src = SRC / "src_2.3.pptx"
    out = OUT / "доп_исключено_Ошибки_осложнения_и_прекращение_СЛР.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    # Slide 1
    s = prs.slides[0]
    for sh in s.shapes:
        if sh.has_text_frame and "ОШИБКИ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "ОШИБКИ И ОСЛОЖНЕНИЯ СЛР (ДОП. МАТЕРИАЛ)")

    # Slide 2 — full methodicheka error list + exclusion note in eyebrow
    s = prs.slides[1]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "ОШИБКИ И ОСЛОЖНЕНИЯ" in t and "TextBox" in sh.name:
            set_shape_plain(sh, "ОШИБКИ И ОСЛОЖНЕНИЯ ПРИ СЛР")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, EYEBROW_EXCL)
        elif "К основным ошибкам" in t:
            set_shape_multiline(
                sh,
                [
                    ("С 01.09.2026 отдельный пункт программы исключён (Прил. № 2, № 805). Ниже — по пособию Минздрава.", True),
                    ("", False),
                    ("Основные ошибки:", True),
                    ("• нарушение последовательности мероприятий СЛР;", False),
                    ("• неправильная техника давления на грудину (точка, глубина, частота, неполное расправление);", False),
                    ("• неправильная техника искусственного дыхания;", False),
                    ("• неправильное соотношение надавливаний и вдохов (нужно 30 : 2);", False),
                    ("• частота надавливаний < 100 или > 120 в минуту;", False),
                    ("• паузы между циклами надавливаний более 10 секунд.", False),
                    ("", False),
                    ("Частое осложнение — перелом рёбер (неверная точка, избыточная сила, хрупкость костей).", False),
                ],
                size=Pt(22),
            )

    # Slide 3 divider
    s = prs.slides[2]
    for sh in s.shapes:
        if sh.has_text_frame and "ПРЕКРАЩЕНИЮ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "ПРЕКРАЩЕНИЕ СЛР (ПО ПОСОБИЮ)")

    # Slide 4 — align with methodicheka exactly
    s = prs.slides[3]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "ПОКАЗАНИЯ К ПРЕКРАЩЕНИЮ" in t:
            set_shape_plain(sh, "КОГДА ПРЕКРАТИТЬ / НЕ НАЧИНАТЬ СЛР")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, EYEBROW_EXCL)
        elif "Реанимационные мероприятия продолжаются" in t:
            set_shape_multiline(
                sh,
                [
                    ("Продолжать до прибытия СМП/спецслужб и их распоряжения о прекращении;", False),
                    ("до появления явных признаков жизни (дыхание, кашель, движения);", False),
                    ("можно прекратить при угрозе для оказывающего помощь.", False),
                ],
                size=Pt(24),
            )
        elif "В случае длительного проведения" in t:
            set_shape_multiline(
                sh,
                [
                    ("При усталости — привлечь помощника (рекомендуется смена около каждых 2 минут).", False),
                    ("", False),
                    ("Можно не начинать: явные признаки нежизнеспособности (разложение, травма, несовместимая с жизнью) или исход длительного неизлечимого заболевания в терминальной стадии.", False),
                ],
                size=Pt(24),
            )

    # Slide 5 divider — after CPR when life signs return → airway (still relevant via 3.2)
    s = prs.slides[4]
    for sh in s.shapes:
        if sh.has_text_frame and "ПОСЛЕ ПРЕКРАЩЕНИЯ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "ЕСЛИ ПОЯВИЛИСЬ ПРИЗНАКИ ЖИЗНИ")

    # Slides 6-7: only "Шаг N" labels — add methodicheka-aligned captions under steps
    # Stable side position is the key post-ROSC action in methodicheka theme 3
    s = prs.slides[5]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "МЕРОПРИЯТИЯ, ВЫПОЛНЯЕМЫЕ ПОСЛЕ" in t:
            set_shape_plain(sh, "ПОДДЕРЖАНИЕ ПРОХОДИМОСТИ ПОСЛЕ СЛР")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, "Связь с п. 3.2 · устойчивое боковое положение (пособие)")
        elif t.strip() == "Шаг 1":
            set_shape_plain(sh, "Шаг 1. Ближняя рука под 90°")
        elif t.strip() == "Шаг 2":
            set_shape_plain(sh, "Шаг 2. Дальняя рука к щеке")
        elif t.strip() == "Шаг 3":
            set_shape_plain(sh, "Шаг 3. Поворот набок")

    s = prs.slides[6]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "МЕРОПРИЯТИЯ, ВЫПОЛНЯЕМЫЕ ПОСЛЕ" in t:
            set_shape_plain(sh, "ПОДДЕРЖАНИЕ ПРОХОДИМОСТИ ПОСЛЕ СЛР")
        elif "Оказание первой помощи при отсутствии" in t:
            set_shape_plain(sh, "Связь с п. 3.2 · устойчивое боковое положение (пособие)")
        elif t.strip() == "Шаг 4":
            set_shape_plain(sh, "Шаг 4. Нога к животу, проверить дыхание")
        elif t.strip() == "Шаг 5":
            set_shape_plain(sh, "Шаг 5. Наблюдение; каждые 30 мин — другой бок")

    prs.save(out)
    print("saved", out)
    return out


def update_24():
    src = SRC / "src_2.4.pptx"
    out = OUT / "доп_исключено_Особенности_СЛР_у_детей.pptx"
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    # Slide 1
    s = prs.slides[0]
    for sh in s.shapes:
        if sh.has_text_frame and "ОСОБЕННОСТИ РЕАНИМАЦИИ" in "".join(p.text for p in sh.text_frame.paragraphs):
            set_shape_plain(sh, "ОСОБЕННОСТИ СЛР У ДЕТЕЙ (ДОП. МАТЕРИАЛ)")

    # Slide 2 — currently almost no body text; add a content box with methodicheka text
    s = prs.slides[1]
    title_sh = None
    eyebrow_sh = None
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = "".join(p.text for p in sh.text_frame.paragraphs)
        if "ОСОБЕННОСТИ СЛР" in t:
            title_sh = sh
            set_shape_plain(sh, "ОСОБЕННОСТИ СЛР У ДЕТЕЙ")
        elif "Оказание первой помощи при отсутствии" in t:
            eyebrow_sh = sh
            set_shape_plain(sh, EYEBROW_EXCL)

    # Add text box with methodicheka content (left/center above pictures area)
    # Pictures are around T=4000000+. Put text band between eyebrow and pictures.
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(1800000),
        Emu(3600000),
        Emu(20800000),
        Emu(2200000),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF8)
    box.line.fill.background()
    try:
        box.adjustments[0] = 0.04
    except Exception:
        pass

    tb = s.shapes.add_textbox(Emu(2000000), Emu(3720000), Emu(20400000), Emu(2000000))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        ("С 01.09.2026 отдельный пункт «Особенности реанимации у детей» исключён из Прил. № 2 (№ 805). Текст — по пособию Минздрава.", True),
        ("• Та же последовательность и соотношение 30 надавливаний : 2 вдоха, что у взрослых.", False),
        ("• Глубина ≈ ⅓ переднезаднего размера груди (~4 см до 1 года; ~5 см старше).", False),
        ("• До 1 года — давление двумя пальцами; старше — одной или двумя руками.", False),
        ("• После отсутствия признаков жизни у ребёнка эффективнее сначала 5 вдохов, затем 30:2.", False),
        ("• При утоплении после извлечения из воды — также 5 вдохов, затем 30:2.", False),
    ]
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.name = "Open Sans"
        run.font.size = Pt(22)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # Move pictures slightly down if overlapping? They start ~4.0M, box ends ~5.8M — overlap.
    # Shrink text box and place higher, move pictures down.
    box.top = Emu(3450000)
    box.height = Emu(1600000)
    tb.top = Emu(3520000)
    tb.height = Emu(1450000)
    for sh in s.shapes:
        if sh.shape_type is not None and sh.shape_type == 13:  # PICTURE
            if sh.top < Emu(5500000):
                sh.top = Emu(5200000)

    prs.save(out)
    print("saved", out)
    return out


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    paths = [update_21(), update_23(), update_24()]
    readme = OUT / "README.md"
    readme.write_text(
        """# Обновление исходных презентаций темы 3 (бывш. 2.1 / 2.3 / 2.4)

Учтены изменения **ПП РФ № 805** к Приложению № 2 № 2464 (с **01.09.2026**) и учебное пособие Минздрава по первой помощи (2025).
Стиль исходников сохранён (PFDinDisplayPro / Open Sans / Montserrat, исходный размер слайдов).

| Файл | Статус по № 805 | Что сделано |
|------|-----------------|-------------|
| `3.1_Определение_признаков_жизни.pptx` | **Обязательный** п. 3.1 | Переименован акцент с «основные признаки…» на «определение признаков жизни»; акцент на сознание + нормальное дыхание; проверка пульса не для решения о СЛР; шаги проверки по пособию |
| `доп_исключено_Ошибки_осложнения_и_прекращение_СЛР.pptx` | **Исключено** как отдельный пункт | Помечено как доп. материал; список ошибок расширен по пособию; прекращение СЛР уточнено; «после СЛР» связано с боковым положением (п. 3.2) |
| `доп_исключено_Особенности_СЛР_у_детей.pptx` | **Исключено** как отдельный пункт | Помечено как доп. материал; добавлен текст по пособию (глубина, пальцы/руки, 5 вдохов, утопление) |

""",
        encoding="utf-8",
    )
    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in paths:
            zf.write(p, arcname=p.name)
        zf.write(readme, arcname=readme.name)
    print("ZIP", ZIP, ZIP.stat().st_size)


if __name__ == "__main__":
    main()
