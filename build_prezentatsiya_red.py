#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Красный стиль Б.7.5 для презентация.pptx — макет и картинки сохраняются."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOURCE = Path("/home/ubuntu/.cursor/projects/workspace/uploads/____________d6e2.pptx")
OUTPUT = Path("презентация.pptx")

RED = RGBColor(0xE3, 0x06, 0x13)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
CREAM = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
FONT = "Inter"

# Исходная палитра (синий корпоративный стиль)
OLD_ACCENT = RGBColor(0x1B, 0x3A, 0x5B)
OLD_SUB = RGBColor(0x5B, 0x68, 0x78)
OLD_BODY = RGBColor(0x1A, 0x1A, 0x1A)
OLD_WARN = RGBColor(0xA3, 0x72, 0x00)
OLD_OK = RGBColor(0x0E, 0x7C, 0x5C)
OLD_CARD = RGBColor(0xFA, 0xEF, 0xD4)
OLD_CARD_GREEN = RGBColor(0xE6, 0xF4, 0xEE)
OLD_LINE = RGBColor(0xC7, 0xCF, 0xDA)

TITLE_SIZES = {508000, 381000, 254000}
TAG_SIZE = 139700
SUBTITLE_SIZE = 215900
HEADER_SIZE = 152400
BODY_SIZE = 190500
PAGE_SIZE = 114300


def rgb_eq(a: RGBColor | None, b: RGBColor) -> bool:
    return a is not None and str(a) == str(b)


def get_run_rgb(run) -> RGBColor | None:
    try:
        if run.font.color.type == MSO_COLOR_TYPE.RGB:
            return run.font.color.rgb
    except AttributeError:
        pass
    return None


def style_run(run):
    size = int(run.font.size) if run.font.size else 0
    rgb = get_run_rgb(run)
    text = run.text.strip()

    run.font.name = FONT

    if rgb_eq(rgb, OLD_ACCENT) or rgb_eq(rgb, OLD_WARN):
        run.font.color.rgb = RED
        return

    if rgb_eq(rgb, OLD_OK):
        run.font.color.rgb = RED
        return

    if rgb_eq(rgb, OLD_SUB):
        run.font.color.rgb = GRAY
        return

    if size in TITLE_SIZES and run.font.bold:
        run.font.color.rgb = RED
        return

    if size == TAG_SIZE or size == PAGE_SIZE:
        run.font.bold = True
        run.font.color.rgb = RED
        return

    if size == HEADER_SIZE and (run.font.bold or text.isdigit() or len(text) <= 4):
        run.font.bold = True
        run.font.color.rgb = RED
        return

    if size == SUBTITLE_SIZE:
        run.font.color.rgb = GRAY
        return

    if size == BODY_SIZE or (size <= BODY_SIZE and rgb_eq(rgb, OLD_BODY)):
        if run.font.bold:
            run.font.color.rgb = RED
        else:
            run.font.color.rgb = GRAY
        return

    if rgb_eq(rgb, OLD_BODY):
        run.font.color.rgb = DARK if run.font.bold and size >= HEADER_SIZE else GRAY
        return

    run.font.color.rgb = GRAY


def style_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                style_run(run)


def get_shape_fill_rgb(shape) -> RGBColor | None:
    try:
        fill = shape.fill
        if fill.type != 1:
            return None
        fc = fill.fore_color
        if fc.type == MSO_COLOR_TYPE.RGB:
            return fc.rgb
    except (AttributeError, TypeError):
        pass
    return None


def style_shape_fill(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return
    rgb = get_shape_fill_rgb(shape)
    if rgb is None:
        return
    if rgb_eq(rgb, OLD_ACCENT):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
    elif rgb_eq(rgb, OLD_CARD):
        shape.fill.solid()
        shape.fill.fore_color.rgb = CREAM
    elif rgb_eq(rgb, OLD_CARD_GREEN):
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN_BG


def style_line(shape):
    if shape.shape_type != MSO_SHAPE_TYPE.LINE:
        return
    try:
        if shape.line.color.type != MSO_COLOR_TYPE.RGB:
            return
        rgb = shape.line.color.rgb
        if rgb_eq(rgb, OLD_ACCENT):
            shape.line.color.rgb = RED
        elif rgb_eq(rgb, OLD_LINE):
            shape.line.color.rgb = BORDER
    except (AttributeError, TypeError):
        pass


def apply_red_style(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            style_line(shape)
            style_shape_fill(shape)
            if shape.has_text_frame:
                style_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        style_text_frame(cell.text_frame)


def build():
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(str(OUTPUT))
    apply_red_style(prs)
    prs.save(OUTPUT)

    src_pics = sum(
        1 for s in Presentation(str(SOURCE)).slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    out_pics = sum(
        1 for s in Presentation(str(OUTPUT)).slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides, pictures {out_pics}/{src_pics})")
    if out_pics != src_pics:
        raise RuntimeError(f"Потеряны картинки: было {src_pics}, стало {out_pics}")


if __name__ == "__main__":
    build()
