#!/usr/bin/env python3
"""Update Tema 4 OT presentation: prune quizzes, add missing from 4.docx, ~20 slides."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________4669.pptx"
TMP = "/tmp/tema4_step1.pptx"
OUT = "/workspace/Презентация_Тема_4_Охрана_труда.pptx"
IMG = Path("/workspace/assets/oteu_images")

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep 0-12; drop duplicate tools (13) and quizzes (14-16)
KEEP = set(range(13))


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(Emu(10271455), Emu(6492240), Emu(1200000), Emu(228600))
            set_text(box, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], list_title, 11, True, RED)

    pairs, more_shape, note_shape = [], None, None
    for sh in content[4:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    shown = items[:4]
    for i, item in enumerate(shown):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(shown), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    # Always clear leftover cloned notes
    if note_shape is None:
        for sh in content[4:]:
            if sh.top > Emu(5600000) and sh.left < Emu(9000000) and sh not in [p[0] for p in pairs] + [p[1] for p in pairs if p[1]]:
                note_shape = sh
                break
    if note_shape:
        if note:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            extra = items[4:]
            set_text(
                note_shape,
                ("· " + " · ".join(extra)) if extra else "",
                11,
                True,
                AMBER if extra else GRAY,
            )


def add_image(slide, path, bottom=False):
    path = Path(path)
    if not path.exists():
        return
    if bottom:
        slide.shapes.add_picture(
            str(path), Emu(731520), Emu(4000000), width=Emu(10700000), height=Emu(2000000)
        )
    else:
        slide.shapes.add_picture(
            str(path), Emu(6600000), Emu(1200000), width=Emu(5000000), height=Emu(4500000)
        )


def find_templates(prs):
    text_tpl = list_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if text_tpl is None and 5 <= n <= 8:
            # Prefer content slides, not title (title has short stacked words)
            titles = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(len(t) > 40 for t in titles):
                text_tpl = slide
        if list_tpl is None and n >= 10:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(t in {"01", "02", "1", "2"} for t in texts):
                list_tpl = slide
        if text_tpl and list_tpl:
            break
    if text_tpl is None:
        text_tpl = prs.slides[1]
    if list_tpl is None:
        for slide in prs.slides:
            n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
            if n >= 8:
                list_tpl = slide
                break
    return text_tpl, list_tpl or text_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def enrich_kept(prs):
    """Refresh sparse early slides from 4.docx without changing layout family."""
    fill_text_slide(
        prs.slides[1],
        "4.1 · Требования к работникам",
        "Допуск работников к работам в электроустановках",
        "Приказ Минтруда России от 15.12.2020 № 903н",
        "Работники проходят обучение безопасным методам работ и первой помощи "
        "(электротехнический персонал — также освобождение от тока). "
        "Группа III — только с 18 лет. Специальные работы: на высоте; под напряжением; "
        "испытания повышенным напряжением; при наведённом напряжении > 25 В. "
        "Специалисты по ОТ потребителей — группа IV, субъектов электроэнергетики — группа V. "
        "Удостоверение постоянно при работнике; право на спецработы фиксируется в нём.",
        note="Группа I — неэлектротехническому персоналу; присваивает работник с группой ≥ III.",
    )
    fill_text_slide(
        prs.slides[2],
        "4.2 · Оперативное обслуживание",
        "Требования к персоналу при оперативном обслуживании",
        "Единоличное обслуживание и осмотры",
        "Оперативное обслуживание — оперативный и оперативно-ремонтный персонал "
        "(административно-технический — по ОРД). "
        "Единолично: выше 1000 В — группа ≥ IV (старший смены), до 1000 В — ≥ III. "
        "Единоличный осмотр: выше 1000 В — группа V, до 1000 В — IV (по ОРД). "
        "Во время осмотра любая работа запрещена. "
        "При замыкании на землю 3–35 кВ — не ближе 4 м в ЗРУ и 8 м в ОРУ.",
        note="При НС напряжение снимается немедленно, без предварительного разрешения.",
    )
    fill_text_slide(
        prs.slides[3],
        "4.3 · Работы в действующих ЭУ",
        "Три допустимых порядка производства работ",
        "Наряд-допуск · распоряжение · перечень текущей эксплуатации",
        "Самовольные работы и расширение задания запрещены. "
        "Работы по другому наряду согласовываются с выдавшим первый наряд "
        "(запись «Согласовано» на лицевой стороне). "
        "Капремонт > 1000 В, работы без снятия напряжения > 1000 В, ремонт ВЛ — не менее 2 человек. "
        "Под напряжением до 1000 В: диэлектрические галоши / подставка / ковёр; "
        "изолированный инструмент и перчатки; оградить другие ТВЧ.",
        note="Не работать в коротких рукавах и металлическими ножовками, напильниками, метрами.",
    )
    # Slide 4 covers breaks/end of work
    try:
        shapes = sorted(text_shapes(prs.slides[4]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.4 · Перерыв, перевод, окончание", 11, True, RED)
    except Exception:
        pass
    fill_text_slide(
        prs.slides[8],
        "4.7 · Текущая эксплуатация",
        "Перечень работ в порядке текущей эксплуатации",
        "Небольшие работы в течение рабочей смены",
        "Перечень утверждает руководитель организации (обособленного подразделения). "
        "Учитывают: необходимость и возможность безопасного выполнения силами оперативного "
        "или оперативно-ремонтного персонала; подготовку РМ; квалификацию персонала; "
        "условия в служебных помещениях, складах, мастерских. "
        "Работы по перечню выполняют оперативный/оперативно-ремонтный персонал "
        "на закреплённом оборудовании в течение одной смены.",
        note="Оформление наряда-допуска или распоряжения для работ по перечню не требуется.",
    )
    try:
        shapes = sorted(text_shapes(prs.slides[9]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.8 · Подготовка РМ и допуск", 11, True, RED)
    except Exception:
        pass
    try:
        shapes = sorted(text_shapes(prs.slides[10]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.9 · Аккумуляторные батареи", 11, True, RED)
    except Exception:
        pass
    try:
        shapes = sorted(text_shapes(prs.slides[11]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.10 · Воздушные линии", 11, True, RED)
    except Exception:
        pass
    try:
        shapes = sorted(text_shapes(prs.slides[12]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.11 · Электроинструмент", 11, True, RED)
    except Exception:
        pass
    try:
        shapes = sorted(text_shapes(prs.slides[5]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.5 · Наряд-допуск", 11, True, RED)
        shapes = sorted(text_shapes(prs.slides[6]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.5 · Наряд-допуск", 11, True, RED)
        shapes = sorted(text_shapes(prs.slides[7]), key=lambda s: (s.top, s.left))
        content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
        set_text(content[0], "4.6 · Работы по распоряжению", 11, True, RED)
    except Exception:
        pass


def append_missing(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    enrich_kept(prs)

    # Light title refresh — keep stacked brand layout of source
    s0 = prs.slides[0]
    shapes = sorted(
        [s for s in s0.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: (s.top, s.left),
    )
    if len(shapes) >= 5:
        set_text(shapes[0], "Охрана труда", 40, True, DARK)
        set_text(shapes[1], "при эксплуатации", 40, True, DARK)
        set_text(shapes[2], "электроустановок", 40, True, RED)
        set_text(
            shapes[3],
            "Тема 4 · Приказ Минтруда России от 15.12.2020 № 903н (ПОТЭУ)",
            14,
            False,
            GRAY,
        )
        set_text(
            shapes[4],
            "Работники · обслуживание · наряд / распоряжение / перечень · допуск · АБ · ВЛ · инструмент · командированные · СМО",
            12,
            False,
            GRAY,
        )

    s = nl()
    fill_list_slide(
        s,
        "4.2 · Ключи и предохранители",
        "Предохранители и ключи от электроустановок",
        "Осмотры, оперативное обслуживание и технологическое управление",
        "Основные правила",
        [
            "предохранители — при снятом напряжении; без нагрузки — допускается под напряжением",
            "под напряжением и нагрузкой: ТН, пробковые и предохранители вторичных систем",
            "выше 1000 В — клещи/штанга, перчатки, защита лица; до 1000 В — клещи, перчатки, защита лица",
            "ключи пронумерованы, в запираемом ящике; выдача и возврат — в журнале",
            "запасной комплект обязателен; передача по смене — в оперативном журнале",
        ],
        note="Ключи выдают имеющим право осмотра, допускающему, руководителю/производителю работ.",
    )

    s = nl()
    fill_list_slide(
        s,
        "4.4 · Ответственные лица",
        "Работники, ответственные за безопасное ведение работ",
        "Организационные мероприятия по обеспечению безопасности",
        "Роли и ответственность",
        [
            "мероприятия: наряд/распоряжение/перечень; разрешение на подготовку РМ и допуск; допуск; надзор; перерыв/перевод/окончание",
            "выдающий наряд — достаточность мер безопасности, квалификация и состав бригады",
            "допускающий — правильность допуска и соответствие подготовленного РМ наряду",
            "производитель работ — безопасное проведение работы, СИЗ, инструмент, надзор за бригадой",
            "наблюдающий — надзор за членами бригады в части безопасности",
        ],
        note="Ответственный руководитель работ назначается при работах по наряду в случаях, установленных Правилами.",
    )

    s = nl()
    fill_list_slide(
        s,
        "4.8 · Целевые инструктажи",
        "Инструктажи при первичном допуске бригады",
        "Без целевых инструктажей допуск к работе не разрешается",
        "Кто проводит",
        [
            "по наряду: выдающий → ответственному/производителю; допускающий → бригаде; производитель → членам",
            "по распоряжению: отдающий → производителю/исполнителю; допускающий → бригаде; производитель → членам",
            "допускается инструктаж по телефону работником, выдающим наряд или распоряжение",
            "нового члена бригады инструктирует производитель работ или наблюдающий",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "4.11 · Электроинструмент",
        "Запреты при работе с ручным электроинструментом",
        "Дополнение к проверкам перед началом работы",
        "Запрещается",
        [
            "передавать электроинструмент другим работникам",
            "разбирать и производить ремонт самостоятельно",
            "держаться за провод электроинструмента, касаться вращающихся частей",
            "работать с приставных лестниц; вносить трансформаторы/преобразователи в котлы и резервуары",
            "заземление вторичной обмотки разделительного трансформатора не допускается",
        ],
    )

    s = nl()
    fill_list_slide(
        s,
        "4.12 · Командированный персонал",
        "Организация работ командированного персонала",
        "Приказ Минтруда № 903н",
        "Ключевые требования",
        [
            "командирующая организация отвечает за группы ЭБ и подготовку персонала",
            "принимающая сторона — вводный и первичный инструктажи, предоставление РМ",
            "право выдачи нарядов командированным — по решению принимающей организации",
            "работают в составе бригады по наряду принимающей стороны после инструктажей",
        ],
        note="Удостоверение с группой ЭБ не заменяет инструктажи на месте командировки.",
    )

    s = nl()
    fill_list_slide(
        s,
        "4.13 · Персонал СМО",
        "Допуск строительно-монтажных организаций",
        "Работы на территории действующей электроустановки и в охранной зоне",
        "Акт-допуск определяет",
        [
            "границы зоны работ СМО на территории действующей электроустановки",
            "место и вид ограждений, исключающих ошибочное проникновение",
            "места входа/выхода и въезда/выезда в зону работ",
            "персонал СМО проходит вводный и первичный инструктажи принимающей стороны",
        ],
        note="СМО не начинает работы без акта-допуска и согласования с владельцем ЭУ.",
    )

    s = nl()
    fill_list_slide(
        s,
        "Итоги темы 4",
        "Ключевые требования охраны труда в электроустановках",
        "Что должен знать каждый работник",
        "Главные выводы",
        [
            "работы — только по наряду, распоряжению или перечню текущей эксплуатации",
            "организационные мероприятия и назначение ответственных обязательны",
            "допуск — только после целевых инструктажей",
            "особые правила: АБ, ВЛ, электроинструмент, командированные, СМО",
            "самовольные работы и расширение задания запрещены",
        ],
        note="Незнание Правил не освобождает от ответственности.",
    )


def reorder_slides(prs, order):
    """Reorder slides by current indices. order is a list of old indices."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order length {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for item in new_items:
        sldIdLst.append(item)


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    prs.save(TMP)

    prs = Presentation(TMP)
    append_missing(prs)
    # After append: 0-12 kept, 13 keys, 14 roles, 15 briefings, 16 bans, 17 seconded, 18 SMO, 19 summary
    # Place keys after 4.2 ops, roles before breaks, briefings after admission, bans after tool checks
    order = [
        0,   # title
        1,   # 4.1
        2,   # 4.2 ops
        13,  # keys / fuses
        3,   # 4.3
        14,  # 4.4 roles
        4,   # breaks / end
        5,   # naryad 1
        6,   # naryad 2
        7,   # rasporjazhenie
        8,   # current ops
        9,   # admission
        15,  # briefings
        10,  # AB
        11,  # VL
        12,  # tool checks
        16,  # tool bans
        17,  # seconded
        18,  # SMO
        19,  # summary
    ]
    reorder_slides(prs, order)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
