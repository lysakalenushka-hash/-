#!/usr/bin/env python3
"""Build pipeline operation requirements presentation from 5.docx."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

TEMPLATE = "/home/ubuntu/.cursor/projects/workspace/uploads/____________54c7.pptx"
OUT = "/workspace/Презентация_Эксплуатация_трубопроводов.pptx"
IMG = "/workspace/assets/operation_images"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

IMG_LEFT = Emu(6600000)
IMG_W = Emu(5200000)


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def clear_slides(prs):
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def make_slide(prs, shape_elements):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    for element in shape_elements:
        slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
    return slide


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        for sh in text_shapes(slide):
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], body, 13, False, DARK)
    if len(shapes) > 4:
        set_text(shapes[4], "", 11, False, GRAY)
    if note:
        for sh in shapes:
            if sh.top > Emu(5800000) and sh.left < Emu(5000000):
                set_text(sh, note, 13, True, AMBER)
                break


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    set_text(shapes[0], section, 11, True, RED)
    set_text(shapes[1], title, 28, True, DARK)
    set_text(shapes[2], intro, 12, False, GRAY)
    set_text(shapes[3], list_title, 11, True, RED)

    pairs = []
    more_shape = None
    for sh in shapes[5:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None:
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh

    for i, item in enumerate(items[:4]):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i+1:02d}", 14, True, RED)
            set_text(pairs[i][1], item, 13, False, DARK)

    for i in range(len(items[:4]), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        set_text(pairs[i][1], "", 14, False, DARK)

    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)

    if note:
        for sh in shapes:
            if sh.top > Emu(5900000) and sh.left > Emu(700000):
                set_text(sh, note, 13, True, AMBER)
                break


def add_image(slide, path, bottom=False):
    if bottom:
        slide.shapes.add_picture(path, Emu(731520), Emu(3600000), width=Emu(10700000), height=Emu(2500000))
    else:
        slide.shapes.add_picture(path, IMG_LEFT, Emu(1200000), width=IMG_W, height=Emu(4800000))


def main():
    tpl = Presentation(TEMPLATE)
    text_tpl = [deepcopy(s.element) for s in tpl.slides[0].shapes]
    list_tpl = [deepcopy(s.element) for s in tpl.slides[2].shapes]

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # 1. Title
    s = make_slide(prs, text_tpl)
    fill_text_slide(
        s,
        "ФНП · Оборудование под давлением",
        "Требования к эксплуатации трубопроводов пара и горячей воды на ОПО",
        "Производственные инструкции, контроль, ввод в эксплуатацию, прокладка",
        "Модуль охватывает требования к производственным инструкциям, контролю "
        "параметров при эксплуатации, проверке манометров и предохранительных "
        "клапанов, ремонту, вводу в эксплуатацию и прокладке трубопроводов "
        "пара и горячей воды на опасных производственных объектах.",
        note="Требования основаны на ФНП при использовании оборудования под избыточным давлением.",
    )
    add_image(s, f"{IMG}/pipeline_system_intro.png")

    # 2. Production instructions content
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Эксплуатация",
        "Производственная инструкция по эксплуатации",
        "На рабочих местах должен быть комплект инструкций и схем трубопроводов. В инструкции регламентируют:",
        "Содержание инструкции",
        [
            "состав схемы трубопровода и входящее оборудование",
            "обязанности персонала по наблюдению и контролю за работой",
            "порядок проверки КИП, арматуры, предохранительных устройств и автоматики",
            "подготовка к пуску, пуск, остановка и меры безопасности при ремонте",
            "случаи немедленной остановки и порядок действий при аварии",
            "порядок ведения сменного (оперативного) журнала",
        ],
        note="Описание операций должно содержать последовательность действий с номерами устройств по схеме.",
    )

    # 3. What to control during operation
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Эксплуатация",
        "Контроль при эксплуатации трубопроводов и арматуры",
        "В соответствии с производственными инструкциями контролируют:",
        "Параметры контроля",
        [
            "тепловые перемещения трубопроводов по индикаторам (реперам)",
            "отсутствие защемлений и повышенной вибрации",
            "плотность предохранительных устройств, арматуры и фланцев",
            "температурный режим металла при пусках и остановах",
            "затяжку пружин подвесок и опор — не реже 1 раза в 2 года",
            "герметичность сальников и положение регулирующей арматуры",
        ],
    )
    add_image(s, f"{IMG}/industrial_steam_sidebar.png")

    # 4. Manometers
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "КИПиА",
        "Манометры: класс точности и проверка",
        "Требования к манометрам при эксплуатации трубопроводов",
        "Класс точности и сроки проверки",
        [
            "до 2,5 МПа — класс точности не ниже 2,5",
            "от 2,5 до 14 МПа — класс точности не ниже 1,5",
            "более 14 МПа — класс точности не ниже 1,0",
            "проверка: до 1,4 МПа — раз в смену; 1,4–4,0 МПа — раз в сутки",
            "свыше 4 МПа и ТЭС — по инструкции главного инженера",
            "поверка манометров — не реже 1 раза в 12 месяцев",
        ],
        note="На шкале — красная черта разрешённого рабочего давления; перед манометром — трёхходовой кран.",
    )
    add_image(s, f"{IMG}/kipia_instruments.png", bottom=True)

    # 5. Safety valves and reducing devices
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Предохранительные устройства",
        "Предохранительные клапаны и редуцирующие устройства",
        "Требования безопасности при эксплуатации",
        "Основные требования",
        [
            "исправность клапанов проверяют кратковременным подрывом или на стенде",
            "отводящие трубопроводы обязательны; запорная арматура на них запрещена",
            "дренажи отводов — без запорных устройств, защита от замерзания",
            "при меньшем давлении трубопровода — редуцирующее устройство с манометром и ПУ",
            "редуцирующие устройства — автоматическое регулирование давления",
            "РОУ — автоматическое регулирование давления и температуры",
        ],
    )
    add_image(s, f"{IMG}/safety_valves.png", bottom=True)

    # 6. Repair
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Ремонт",
        "Ремонтные работы и журналы",
        "Порядок организации ремонта трубопроводов",
        "Требования к ремонту",
        [
            "ведётся ремонтный журнал (бумажный или электронный с идентификацией лица)",
            "сведения о ремонте, вызывающем внеочередное освидетельствование — в паспорт",
            "до ремонта трубопровод отделяют заглушками или отсоединяют от действующего оборудования",
            "ремонт, установка и снятие заглушек — только по наряду-допуску",
        ],
        note="Ответственный за исправное состояние и безопасную эксплуатацию подписывает записи в журнале.",
    )

    # 7. Commissioning
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Ввод в эксплуатацию",
        "Ввод и пуск трубопроводов в работу",
        "Решение о вводе принимает руководитель эксплуатирующей организации",
        "Проверка готовности",
        [
            "проверка документации изготовителя, монтажа и технического освидетельствования",
            "исправность арматуры, КИП, приборов безопасности и технологических защит",
            "наличие обученного персонала и аттестованных специалистов",
            "наличие должностных и производственных инструкций",
            "оформление Акта готовности и распорядительного документа о вводе",
            "перед пуском — табличка с номером, параметрами и сроками НО/ГИ",
        ],
        note="Пуск и остановка — по письменному распоряжению ответственного лица.",
    )
    add_image(s, f"{IMG}/pipeline_system_intro.png")

    # 8. Laying requirements
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Прокладка",
        "Прокладка трубопроводов пара и горячей воды",
        "Требования к размещению на ОПО",
        "Основные нормы прокладки",
        [
            "горизонтальные участки: уклон ≥ 0,004; тепловые сети ≥ 0,002",
            "полупроходные каналы: высота ≥ 1,5 м, проход ≥ 600 мм",
            "проходные тоннели: высота ≥ 2 м, проход ≥ 0,7 м",
            "люки в проходных каналах — не реже чем через 300 м (совместно — 50 м)",
            "камеры подземных трубопроводов — не менее двух люков",
            "антикоррозионная, тепловая и гидроизоляционная защита обязательна",
        ],
    )

    # 9. Valves and drainage
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Арматура и дренаж",
        "Запорная арматура, дренажи и конденсатоотвод",
        "Требования к оснащению трубопроводов",
        "Ключевые требования",
        [
            "запорная арматура — на выводах от источника тепла и секционирующих участках",
            "дренажи для слива после ГИ и воздушники в верхних точках",
            "непрерывный отвод конденсата через конденсатоотводчики",
            "пусковой дренаж: через 400–500 м (попутный) / 200–300 м (встречный уклон)",
            "отключаемые участки паропроводов — штуцер с запорным устройством",
            "при давлении > 2,2 МПа — запорное и регулирующее устройства последовательно",
        ],
    )
    add_image(s, f"{IMG}/condensate_drainage.png", bottom=True)

    # 10. Thermal insulation
    s = make_slide(prs, list_tpl)
    fill_list_slide(
        s,
        "Теплоизоляция",
        "Тепловая изоляция трубопроводов",
        "Требования к изоляции и покрытию",
        "Основные требования",
        [
            "температура на поверхности изоляции при +25 °C — не более 55 °C",
            "изоляция фланцев, арматуры и контрольных участков — съёмная",
            "на открытом воздухе и у маслопроводов — металлическое или иное покрытие",
            "трубопроводы у кабельных линий — металлическое покрытие",
            "при температуре среды ниже окружающей — гидро- и теплоизоляция",
            "материалы изоляции не должны вызывать коррозию металла",
        ],
        note="Опознавательная окраска и маркировка — по проекту и схеме трубопровода.",
    )
    add_image(s, f"{IMG}/thermal_insulation_cutaway.png", bottom=True)

    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
