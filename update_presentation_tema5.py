#!/usr/bin/env python3
"""Update Tema 5: accidents investigation + emergency schedules from 5.docx."""

from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

SRC = "/home/ubuntu/.cursor/projects/workspace/uploads/____________dec2.pptx"
TMP = "/tmp/tema5_step1.pptx"
OUT = "/workspace/Презентация_Тема_5_Расследование_аварий.pptx"

RED = RGBColor(0xE3, 0x06, 0x13)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Keep 0-8; drop quizzes 9-12
KEEP = set(range(9))


def set_run(run, size_pt, bold=False, color=DARK):
    run.font.name = "Inter"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size_pt=13, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size_pt, bold=bold, color=color)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = sldIdLst[idx]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def update_pages(prs):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        updated = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is not None and sh.left is not None:
                if sh.top > Emu(6200000) and sh.left > Emu(8000000):
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
            t = sh.text_frame.text.strip()
            if " / " in t:
                left, right = [x.strip() for x in t.split(" / ", 1)]
                if left.isdigit() and right.isdigit():
                    set_text(sh, f"{i:02d} / {total:02d}", 9, True, RED)
                    updated = True
                    break
        if not updated:
            box = slide.shapes.add_textbox(
                Emu(10271455), Emu(6492240), Emu(1200000), Emu(228600)
            )
            set_text(box, f"{i:02d} / {total:02d}", 9, True, RED)


def fill_text_slide(slide, section, title, intro, body, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], body, 13, False, DARK)
    if len(content) > 4 and content[4].top < Emu(6200000):
        set_text(content[4], note or "", 12, True, AMBER if note else GRAY)


def fill_list_slide(slide, section, title, intro, list_title, items, note=None):
    shapes = sorted(text_shapes(slide), key=lambda s: (s.top, s.left))
    content = [s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))]
    set_text(content[0], section, 11, True, RED)
    set_text(content[1], title, 24, True, DARK)
    set_text(content[2], intro, 12, False, GRAY)
    set_text(content[3], list_title, 11, True, RED)

    pairs, more_shape, note_shape = [], None, None
    for sh in content[4:]:
        t = sh.text_frame.text.strip()
        if t in {f"{i:02d}" for i in range(1, 10)} or t in {str(i) for i in range(1, 10)}:
            pairs.append([sh, None])
        elif pairs and pairs[-1][1] is None and sh.left > Emu(700000):
            pairs[-1][1] = sh
        elif t.startswith("…"):
            more_shape = sh
        elif sh.top > Emu(5600000) and sh.left < Emu(9000000):
            note_shape = sh

    shown = items[:4]
    for i, item in enumerate(shown):
        if i < len(pairs):
            set_text(pairs[i][0], f"{i + 1:02d}", 14, True, RED)
            if pairs[i][1] is not None:
                set_text(pairs[i][1], item, 12, False, DARK)
    for i in range(len(shown), len(pairs)):
        set_text(pairs[i][0], "", 14, True, RED)
        if pairs[i][1] is not None:
            set_text(pairs[i][1], "", 12, False, DARK)
    if more_shape:
        rest = max(0, len(items) - 4)
        set_text(more_shape, f"… ещё {rest}" if rest else "", 9, False, GRAY)
    if note_shape is None:
        for sh in content[4:]:
            if (
                sh.top > Emu(5600000)
                and sh.left < Emu(9000000)
                and sh not in [p[0] for p in pairs] + [p[1] for p in pairs if p[1]]
            ):
                note_shape = sh
                break
    if note_shape:
        if note:
            set_text(note_shape, note, 12, True, AMBER)
        else:
            extra = items[4:]
            set_text(
                note_shape,
                ("· " + " · ".join(extra)) if extra else "",
                11,
                True,
                AMBER if extra else GRAY,
            )
    elif note:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, note, 12, True, AMBER)
    elif len(items) > 4:
        box = slide.shapes.add_textbox(Emu(731520), Emu(5850000), Emu(10000000), Emu(400000))
        set_text(box, "· " + " · ".join(items[4:]), 11, True, AMBER)


def find_templates(prs):
    text_tpl = list_tpl = None
    for slide in prs.slides:
        n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        if text_tpl is None and 5 <= n <= 8:
            titles = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(len(t) > 40 for t in titles):
                text_tpl = slide
        if list_tpl is None and n >= 10:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            if any(t in {"01", "02", "1", "2"} for t in texts):
                list_tpl = slide
        if text_tpl and list_tpl:
            break
    if text_tpl is None:
        text_tpl = prs.slides[min(1, len(prs.slides) - 1)]
    if list_tpl is None:
        list_tpl = text_tpl
        for slide in prs.slides:
            n = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
            if n >= 8:
                list_tpl = slide
                break
    return text_tpl, list_tpl


def clone_slide(prs, elems):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(blank.shapes):
        shape.element.getparent().remove(shape.element)
    for el in elems:
        blank.shapes._spTree.insert_element_before(deepcopy(el), "p:extLst")
    return blank


def enrich_kept(prs):
    shapes = sorted(
        [s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: (s.top, s.left),
    )
    if len(shapes) >= 4:
        set_text(shapes[0], "Аварии в электроэнергетике:", 32, True, DARK)
        set_text(shapes[1], "расследование и аварийные ограничения", 28, True, RED)
        set_text(
            shapes[2],
            "Тема 5 · Расследование причин аварий · графики аварийного ограничения · противоаварийная автоматика",
            13,
            False,
            GRAY,
        )
        set_text(
            shapes[3],
            "Электроэнергетика · НПА · Диспетчерское управление",
            12,
            False,
            GRAY,
        )

    fill_text_slide(
        prs.slides[1],
        "5.1 · Определение",
        "Что такое авария в электроэнергетике",
        "Нормативное определение",
        "Авария — технологические нарушения на объекте электроэнергетики и (или) "
        "энергопринимающей установке, приведшие к разрушению или повреждению зданий, "
        "сооружений и (или) технических устройств, неконтролируемому взрыву и (или) "
        "выбросу опасных веществ; отклонению от установленного режима работы; "
        "полному или частичному ограничению режима потребления; возникновению или "
        "угрозе возникновения аварийного электроэнергетического режима.",
        note="Расследованию и учёту подлежат аварии на всех объектах на территории РФ.",
    )

    # Label section brows on dense kept slides
    labels = {
        2: "5.1 · Термины",
        3: "5.1 · Сфера применения",
        4: "5.1 · Кто расследует",
        5: "5.1 · Надзор (1/2)",
        6: "5.1 · Надзор (2/2)",
        7: "5.1 · Пороги надзора",
        8: "5.1 · Собственник (1/2)",
    }
    for idx, label in labels.items():
        try:
            shapes = sorted(text_shapes(prs.slides[idx]), key=lambda s: (s.top, s.left))
            content = [
                s for s in shapes if not (s.top > Emu(6200000) and s.left > Emu(8000000))
            ]
            set_text(content[0], label, 11, True, RED)
        except Exception:
            pass


def append_missing(prs):
    text_tpl, list_tpl = find_templates(prs)
    text_elems = [deepcopy(s.element) for s in text_tpl.shapes]
    list_elems = [deepcopy(s.element) for s in list_tpl.shapes]

    def nt():
        return clone_slide(prs, text_elems)

    def nl():
        return clone_slide(prs, list_elems)

    enrich_kept(prs)

    s = nl()
    fill_list_slide(
        s,
        "5.1 · Надзор (продолжение)",
        "Дополнительные критерии федерального надзора",
        "пп. з–л (в дополнение к а–ж)",
        "Критерии",
        [
            "з) отключение сетей ≥110 кВ / генерации ≥100 МВт на ≥2 объектах с обесточением ≥200 тыс. чел. / ≥500 МВт нагрузки",
            "и) нарушения противоаварийной/режимной автоматики с отключением объекта ≥110 кВ или генерации ≥100 МВт",
            "к) нарушение в сетях с отклонением частоты на шинах РУ АЭС ≥110 кВ за пределы, угрожающие безопасности",
            "л) потеря диспетчерской связи и дистанционного управления / телеметрии / управляющих воздействий ПА",
        ],
        note="Системные и межобъектные события — компетенция федерального энергетического надзора.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.1 · Собственник (2/2)",
        "Аварии, расследуемые собственником / эксплуатирующей организацией",
        "пп. г–и (в дополнение к а–в)",
        "Критерии",
        [
            "г) потеря управляемости объекта ≥1 ч (СН, оперативный ток, сжатый воздух, масло и др.)",
            "д) неправильные действия защитных устройств и (или) систем автоматики",
            "е) вывод электрооборудования системы электропитания АЭС действием РЗА от повышения напряжения",
            "ж–и) превышение лимитов выбросов ×5; пожар как причина/следствие; повреждения сетей <6 кВ",
        ],
        note="Локальные события объекта — зона собственника, если не достигнуты пороги надзора.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.1 · Порядок расследования",
        "Сроки уведомления, решения и работы комиссии",
        "Действия после аварии",
        "Ключевые сроки",
        [
            "собственник незамедлительно уведомляет орган федерального энергетического надзора",
            "решение о расследовании надзором — не позднее 24 часов с момента получения информации",
            "собственник в срок ≤24 ч создаёт комиссию при расследовании «своих» аварий",
            "срок расследования — ≤20 календарных дней; продление — не более чем на 45 дней",
        ],
        note="Председатель комиссии надзора — должностное лицо органа федерального энергетического надзора.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.1 · Комиссия и действия",
        "Состав комиссии и обязательные действия при расследовании",
        "Выявление причин аварии",
        "Что делают",
        [
            "в комиссию при необходимости включают представителей Минэнерго, Ростехнадзора смежных сфер, сетей, генерации, СО, крупных потребителей (>50 МВт)",
            "сохраняют обстановку; изымают регистрограммы и записи переговоров; фиксируют положение защит",
            "осмотр, фото/видео, схема места; опрос очевидцев и оперативного персонала",
            "оценка действий персонала, проверка норм, проекта, области применения, СИЗ и техдокументации",
        ],
        note="Действия комиссии оформляются протоколом и подписываются председателем.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.1 · Акт и учёт",
        "Оформление результатов и систематизация информации",
        "Акт о расследовании причин аварии",
        "Требования",
        [
            "акт составляется в 2 экземплярах; особое мнение члена комиссии прилагается",
            "копии акта с приложениями — в 3-дневный срок субъектам и потребителям, чьи интересы затронуты",
            "материалы хранятся надзором / собственником ≥5 лет; формируется отдельное дело с описью",
            "ежемесячный сводный отчёт собственника; электронные копии актов — в базу данных об авариях",
        ],
        note="Материалы используют при планировании режимов и разработке противоаварийных мероприятий.",
    )

    s = nt()
    fill_text_slide(
        s,
        "5.2 · Аварийные ограничения",
        "Графики аварийного ограничения и противоаварийная автоматика",
        "Правила разработки и применения графиков",
        "Правила определяют: порядок разработки графиков аварийного ограничения режима "
        "потребления электрической энергии (мощности); порядок их применения "
        "(введение диспетчерским центром); правила использования противоаварийной автоматики. "
        "Сетевые организации и владельцы сетей/генерации формируют перечни потребителей, "
        "в отношении которых может осуществляться аварийное ограничение.",
        note="Графики разрабатываются ежегодно: 1 октября — 30 сентября следующего года.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.2 · Виды графиков",
        "Виды графиков аварийного ограничения и основания введения",
        "Утверждаются первичными получателями команд",
        "Состав и основания",
        [
            "три вида: ограничение энергии; ограничение мощности; временное отключение потребления",
            "задание диспетчерского центра на разработку — ежегодно до 15 июня",
            "основания: дефицит с частотой <49,8 Гц; перегрузка сечений/линий; повреждение сетей; повреждение АСУ/связи/ПА",
            "вводятся при невозможности предотвратить угрозу иными мерами; без согласования с потребителем",
        ],
        note="Автоматический перевод отключаемой нагрузки на другие центры питания не допускается.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.2 · Введение графиков",
        "Диспетчерские команды и получатели ограничений",
        "Порядок исполнения",
        "Ключевые правила",
        [
            "вводит диспетчерский центр командой/распоряжением; запись в оперативном журнале с основанием",
            "первичный получатель распределяет объёмы и передаёт вторичным / потребителям",
            "потребитель самостоятельно выполняет технические мероприятия; при невыполнении — принудительное ограничение сетями",
            "графики ограничения энергии/мощности — с 00:00 следующих суток, уведомление потребителя ≤14:00 текущих",
        ],
        note="Временное отключение — если ограничение нельзя ввести вовремя или потребитель не исполнил распоряжение.",
    )

    s = nl()
    fill_list_slide(
        s,
        "5.2 · Противоаварийная автоматика",
        "Использование ПА, действующей на отключение нагрузки",
        "Предотвращение и ликвидация аварийных режимов",
        "Требования",
        [
            "диспетчерский центр определяет необходимость ПА, тип, факторы запуска и объёмы воздействий",
            "задания на изменение уставок/алгоритмов — в сроки задания; АЧР/ЧАПВ — ≤5 месяцев; иначе — по согласованию",
            "сведения о настройке АЧР: до 1 сентября и до 20 февраля; внеочередные замеры — за 10 рабочих дней",
            "минимально необходимый уровень потребления сохраняется; АРИП — по категории надёжности потребителя",
        ],
        note="При срабатывании ПА перевод отключаемой нагрузки на другие центры питания не допускается.",
    )

    s = nl()
    fill_list_slide(
        s,
        "Итоги темы 5",
        "Расследование аварий и аварийные ограничения",
        "Что должен знать работник",
        "Главные выводы",
        [
            "авария — технологическое нарушение с последствиями для объекта/режима/потребления",
            "крупные и системные события расследует надзор; локальные — собственник/эксплуатирующая организация",
            "сроки: уведомление немедленное; решение ≤24 ч; расследование ≤20 дней (+≤45)",
            "графики аварийного ограничения и ПА вводятся диспетчерским центром без согласования с потребителем",
        ],
        note="Невыполнение диспетчерских команд об аварийных ограничениях недопустимо.",
    )


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if len(order) != len(items):
        raise ValueError(f"order length {len(order)} != slides {len(items)}")
    new_items = [items[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for item in new_items:
        sldIdLst.append(item)


def main():
    prs = Presentation(SRC)
    delete = set(range(len(prs.slides))) - KEEP
    delete_slides(prs, delete)
    prs.save(TMP)

    prs = Presentation(TMP)
    append_missing(prs)
    # kept 0-8; appended 9 nadzor+, 10 owner+, 11 procedure, 12 commission,
    # 13 act, 14 5.2 intro, 15 schedules, 16 apply, 17 PA, 18 summary
    order = [
        0, 1, 2, 3, 4, 5, 6,  # title … nadzor 2/2
        9,  # nadzor з–л
        7,  # thresholds
        8,  # owner 1/2
        10,  # owner 2/2
        11, 12, 13,  # procedure, commission, act
        14, 15, 16, 17,  # 5.2 block
        18,  # summary
    ]
    reorder_slides(prs, order)
    update_pages(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
